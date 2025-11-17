#!/usr/bin/env python3
"""
마지막 "감사합니다" 슬라이드만 분석하는 스크립트
"""

import sys
from pathlib import Path
from pptx import Presentation

def analyze_thanks_slide(ppt_path):
    """마지막 감사합니다 슬라이드만 분석합니다."""
    
    if not Path(ppt_path).exists():
        print(f"❌ 파일이 존재하지 않습니다: {ppt_path}")
        return
        
    try:
        prs = Presentation(ppt_path)
        print(f"📊 PPT 분석: {Path(ppt_path).name}")
        print(f"총 슬라이드 수: {len(prs.slides)}")
        
        # 마지막 슬라이드 분석
        last_slide = prs.slides[-1]
        slide_idx = len(prs.slides)
        
        print(f"\n🎯 마지막 슬라이드 (슬라이드 {slide_idx}) 분석:")
        
        # 텍스트 내용 수집
        text_contents = []
        
        for shape_idx, shape in enumerate(last_slide.shapes):
            # 텍스트 박스/플레이스홀더 텍스트
            if hasattr(shape, 'text_frame') and shape.text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    text_contents.append({
                        'shape_idx': shape_idx,
                        'shape_type': str(shape.shape_type),
                        'shape_name': getattr(shape, 'name', 'Unknown'),
                        'text': text
                    })
        
        print(f"  📦 Shape 개수: {len(last_slide.shapes)}")
        for i, shape in enumerate(last_slide.shapes):
            print(f"    - Shape {i}: {shape.shape_type} ({getattr(shape, 'name', 'Unknown')})")
        
        print(f"\n  📝 텍스트 내용:")
        for text_item in text_contents:
            print(f"    - Shape {text_item['shape_idx']} ({text_item['shape_name']}): '{text_item['text']}'")
        
        # 감사합니다 슬라이드 검증
        if len(text_contents) == 1 and text_contents[0]['text'] == '감사합니다':
            print(f"\n  ✅ 완벽! 제목만 '감사합니다'로 표시됨")
        elif len(text_contents) == 2:
            title_text = text_contents[0]['text'] if text_contents else ""
            subtitle_text = text_contents[1]['text'] if len(text_contents) > 1 else ""
            
            if title_text == '감사합니다' and not subtitle_text:
                print(f"\n  ✅ 좋음! 제목 '감사합니다', 부제목 빈 문자열")
            elif title_text == '감사합니다' and subtitle_text:
                print(f"\n  ❌ 문제! 부제목에 내용이 있음: '{subtitle_text}'")
            else:
                print(f"\n  ❌ 예상과 다른 구조: 제목='{title_text}', 부제목='{subtitle_text}'")
        else:
            print(f"\n  ❌ 예상과 다른 텍스트 개수: {len(text_contents)}개")
            
    except Exception as e:
        print(f"❌ 분석 중 오류 발생: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    # 새로 생성된 테스트 PPT 파일 분석
    new_ppt_path = "/home/admin/wkms-aws/uploads/test_fixed_insulin_pump.pptx"
    
    print("🔍 마지막 '감사합니다' 슬라이드 수정 확인")
    analyze_thanks_slide(new_ppt_path)
