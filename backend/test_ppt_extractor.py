import json
from pptx import Presentation
from pptx.util import Emu

# 분석할 PPTX 파일 경로
PPTX_FILE_PATH = 'uploads/templates/제품소개서 샘플.pptx'

# 슬라이드 크기를 기반으로 위치를 해석하기 위한 기준값 (pt 단위)
SLIDE_WIDTH = 720  # 표준 슬라이드 가로 크기
SLIDE_HEIGHT = 540 # 표준 슬라이드 세로 크기

def get_shape_details(shape):
    """도형(Shape) 객체의 상세 정보를 딕셔너리로 반환합니다."""
    details = {
        "type": shape.shape_type.name, # 도형 타입 (예: TEXT_BOX, PICTURE)
        "position": {
            "left": Emu(shape.left).pt, # pt 단위로 변환
            "top": Emu(shape.top).pt,
        },
        "size": {
            "width": Emu(shape.width).pt,
            "height": Emu(shape.height).pt,
        },
        "content": None
    }

    # 도형이 텍스트 프레임을 가지고 있는지 확인
    if shape.has_text_frame:
        # 텍스트가 비어있지 않은 경우에만 내용 추가
        if shape.text.strip():
            details["content"] = shape.text
            # 추가적으로 폰트 정보 등도 추출 가능
            # 예: 첫 번째 문단의 첫 번째 run의 폰트 크기
            try:
                font = shape.text_frame.paragraphs[0].runs[0].font
                details["font"] = {
                    "name": font.name,
                    "size_pt": font.size.pt if font.size else None,
                    "bold": font.bold,
                    "italic": font.italic
                }
            except IndexError:
                # 텍스트는 있지만 run이 없는 경우 (거의 없음)
                details["font"] = None


    # 도형이 이미지인 경우
    elif hasattr(shape, 'image'):
        details["content"] = f"Image: {shape.image.filename or 'embedded_image'}"

    return details


def analyze_presentation(file_path):
    """PPTX 파일을 분석하여 구조를 JSON 형식으로 반환합니다."""
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"Error opening presentation file: {e}")
        return None

    presentation_data = {
        "fileName": file_path,
        "totalPages": len(prs.slides),
        "slides": []
    }

    for i, slide in enumerate(prs.slides):
        slide_details = {
            "pageNumber": i + 1,
            # 슬라이드 레이아웃의 이름을 가져옵니다.
            "layout": slide.slide_layout.name,
            "elements": []
        }
        for shape in slide.shapes:
            shape_info = get_shape_details(shape)
            slide_details["elements"].append(shape_info)
        
        presentation_data["slides"].append(slide_details)
        
    return presentation_data

def interpret_position(shape_info, slide_width, slide_height):
    """도형의 좌표와 크기를 바탕으로 'center', 'top-left' 등 추상적인 위치를 반환"""
    pos = shape_info['position']
    size = shape_info['size']
    
    # 도형의 중심 좌표 계산
    center_x = pos['left'] + size['width'] / 2
    center_y = pos['top'] + size['height'] / 2
    
    # 위치 구분을 위한 기준점 설정
    left_boundary = slide_width * 0.3
    right_boundary = slide_width * 0.7
    top_boundary = slide_height * 0.3
    bottom_boundary = slide_height * 0.7
    
    # 수직 위치 판단
    if center_y < top_boundary:
        vertical = "top"
    elif center_y > bottom_boundary:
        vertical = "bottom"
    else:
        vertical = "middle"
    
    # 수평 위치 판단
    if center_x < left_boundary:
        horizontal = "left"
    elif center_x > right_boundary:
        horizontal = "right"
    else:
        horizontal = "center"
    
    # 조합하여 최종 위치 반환
    if vertical == "middle" and horizontal == "center":
        return "center"
    elif vertical == "top" and horizontal == "center":
        return "top-center"
    elif vertical == "bottom" and horizontal == "center":
        return "bottom-center"
    else:
        return f"{vertical}-{horizontal}"

def interpret_element_type(shape_info):
    """도형 타입과 내용을 바탕으로 UI 친화적인 요소 타입을 결정"""
    shape_type = shape_info.get('type', '')
    content = shape_info.get('content', '')
    
    # 기본적으로 텍스트가 있으면 textbox
    if content and content.strip():
        # 제목처럼 보이는 텍스트 (큰 폰트, 짧은 내용)
        font_info = shape_info.get('font', {})
        font_size = font_info.get('size_pt', 12) if font_info else 12
        
        if font_size and font_size >= 24:
            return "title"
        elif len(content.strip()) < 50 and '\n' not in content:
            return "heading"
        elif '\n' in content and content.count('\n') >= 2:
            return "paragraph"
        else:
            return "textbox"
    
    # 이미지 타입
    elif shape_type == 'PICTURE':
        return "image"
    
    # 표 타입
    elif shape_type == 'TABLE':
        return "table"
    
    # 도형 타입
    elif shape_type in ['AUTO_SHAPE', 'GROUP']:
        return "shape"
    
    # 선 타입
    elif shape_type == 'LINE':
        return "line"
    
    return "unknown"

def interpret_style(shape_info):
    """폰트 정보를 바탕으로 스타일 객체를 생성"""
    style = {}
    font_info = shape_info.get('font')
    
    if font_info:
        if font_info.get('size_pt'):
            style['fontSize'] = f"{int(font_info['size_pt'])}pt"
        
        if font_info.get('bold'):
            style['fontWeight'] = "bold"
        elif font_info.get('bold') == False:
            style['fontWeight'] = "normal"
        
        if font_info.get('italic'):
            style['fontStyle'] = "italic"
        elif font_info.get('italic') == False:
            style['fontStyle'] = "normal"
            
        if font_info.get('name'):
            style['fontFamily'] = font_info['name']
    
    return style

def interpret_slide_layout(slide_data):
    """슬라이드의 요소들을 분석하여 더 의미있는 레이아웃 이름을 결정"""
    elements = slide_data.get('elements', [])
    layout_name = slide_data.get('layout', '')
    
    # 제목 슬라이드 감지
    title_elements = [e for e in elements if e.get('content') and len(e['content'].strip()) < 100]
    if len(title_elements) <= 3 and any('제품' in str(e.get('content', '')) or 'TITLE' in str(e.get('content', '')) for e in title_elements):
        return "제목 슬라이드"
    
    # 목차 슬라이드 감지
    if '목차' in str([e.get('content', '') for e in elements]) or layout_name == "목차":
        return "목차"
    
    # 내용 슬라이드 패턴 감지
    text_elements = [e for e in elements if e.get('content') and e['content'].strip()]
    if len(text_elements) >= 3:
        return f"내용 슬라이드 ({len(text_elements)}개 요소)"
    
    return layout_name or "일반 슬라이드"

def create_interpreted_json(raw_data):
    """1단계에서 추출된 원시 데이터를 2단계에서 해석하여 UI 친화적인 JSON으로 변환"""
    interpreted_data = {
        "presentationTitle": "스마트 인슐린 펌프 제품소개서",  # 파일명에서 추출
        "extractionMethod": "2-stage hybrid approach (extraction + interpretation)",
        "totalPages": raw_data['totalPages'],
        "slideSize": {
            "width": f"{SLIDE_WIDTH}pt",
            "height": f"{SLIDE_HEIGHT}pt"
        },
        "slides": []
    }
    
    for slide_data in raw_data['slides']:
        # 슬라이드 레이아웃 해석
        interpreted_layout = interpret_slide_layout(slide_data)
        
        new_slide = {
            "pageNumber": slide_data['pageNumber'],
            "layout": interpreted_layout,
            "originalLayout": slide_data.get('layout', ''),
            "elements": []
        }
        
        for shape_info in slide_data['elements']:
            # 내용이 없는 요소는 건너뛰기 (단, 이미지나 도형은 포함)
            if not shape_info.get('content') and shape_info.get('type') not in ['PICTURE', 'AUTO_SHAPE', 'GROUP', 'LINE', 'TABLE']:
                continue
            
            new_element = {
                "type": interpret_element_type(shape_info),
                "originalType": shape_info.get('type', 'UNKNOWN'),
                "content": shape_info.get('content', ''),
                "position": interpret_position(shape_info, SLIDE_WIDTH, SLIDE_HEIGHT),
                "rawPosition": shape_info.get('position', {}),
                "size": shape_info.get('size', {}),
                "style": interpret_style(shape_info)
            }
            new_slide['elements'].append(new_element)
        
        interpreted_data['slides'].append(new_slide)
    
    return interpreted_data

# --- 메인 실행 부분 ---
if __name__ == "__main__":
    print("🔍 PPT 분석을 시작합니다...")
    print("=" * 50)
    
    # === 1단계: 데이터 추출 (Extraction) ===
    print("📤 1단계: 원시 데이터 추출 중...")
    raw_data = analyze_presentation(PPTX_FILE_PATH)

    if not raw_data:
        print("❌ PPT 파일 분석 실패")
        exit(1)
    
    # 원시 데이터 저장
    raw_output_filename = "pptx_raw_extraction.json"
    with open(raw_output_filename, 'w', encoding='utf-8') as f:
        json.dump(raw_data, f, indent=2, ensure_ascii=False)
    print(f"✅ 원시 데이터 저장: '{raw_output_filename}'")
    
    # === 2단계: 데이터 해석 (Interpretation) ===
    print("🧠 2단계: 데이터 해석 및 변환 중...")
    interpreted_data = create_interpreted_json(raw_data)
    
    # 해석된 데이터 저장
    interpreted_output_filename = "pptx_interpreted_analysis.json"
    with open(interpreted_output_filename, 'w', encoding='utf-8') as f:
        json.dump(interpreted_data, f, indent=2, ensure_ascii=False)
    print(f"✅ 해석된 데이터 저장: '{interpreted_output_filename}'")
    
    # === 결과 요약 출력 ===
    print("\n" + "=" * 50)
    print("📊 분석 결과 요약:")
    print(f"   총 슬라이드 수: {interpreted_data['totalPages']}개")
    
    for slide in interpreted_data['slides']:
        elements_count = len(slide['elements'])
        print(f"   슬라이드 {slide['pageNumber']}: {slide['layout']} ({elements_count}개 요소)")
    
    print(f"\n🎯 생성된 파일:")
    print(f"   1. {raw_output_filename} - 추출된 원시 데이터")
    print(f"   2. {interpreted_output_filename} - 해석된 UI 친화적 데이터")
    print("\n✅ 2단계 하이브리드 분석이 완료되었습니다!")
