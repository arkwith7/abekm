import json
import logging
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 로깅 설정
logging.basicConfig(level=logging.INFO, format='%(message)s')
logger = logging.getLogger(__name__)

def analyze_ppt_compliance(pptx_path, metadata_path):
    """
    생성된 PPT가 템플릿 메타데이터를 준수했는지 검증하는 스크립트
    """
    logger.info(f"🔍 PPT 품질 검증 시작: {pptx_path}")
    
    # 1. 메타데이터 로드
    try:
        with open(metadata_path, 'r', encoding='utf-8') as f:
            metadata = json.load(f)
        logger.info(f"✅ 메타데이터 로드 완료: {len(metadata.get('slide_templates', []))}개 슬라이드 정의됨")
    except Exception as e:
        logger.error(f"❌ 메타데이터 로드 실패: {e}")
        return

    # 2. PPT 로드
    try:
        prs = Presentation(pptx_path)
        logger.info(f"✅ PPT 로드 완료: {len(prs.slides)}개 슬라이드")
    except Exception as e:
        logger.error(f"❌ PPT 로드 실패: {e}")
        return

    # 3. 슬라이드별 비교 분석
    total_issues = 0
    
    # 메타데이터 구조 확인 (slides vs slide_templates)
    meta_slides = metadata.get('slides', [])
    if not meta_slides:
        meta_slides = metadata.get('slide_templates', [])
        
    logger.info(f"✅ 메타데이터 로드 완료: {len(meta_slides)}개 슬라이드 정의됨")
    
    for i, slide in enumerate(prs.slides):
        slide_idx = i + 1
        logger.info(f"\n--- [Slide {slide_idx}] 분석 ---")
        
        # 메타데이터에서 해당 슬라이드 정보 찾기
        # 'index' 또는 'slide_number' 키 사용
        slide_meta = next((s for s in meta_slides if s.get('index', s.get('slide_number')) == slide_idx), None)
        
        if not slide_meta:
            logger.warning(f"⚠️ 메타데이터에 정의되지 않은 슬라이드입니다.")
            continue
            
        # layout_boxes 또는 shapes 키 사용
        meta_boxes = slide_meta.get('shapes', slide_meta.get('layout_boxes', []))
        logger.info(f"📋 메타데이터 정의: {len(meta_boxes)}개 박스")
        
        # 실제 슬라이드의 텍스트 박스/도형 분석
        real_shapes = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                real_shapes.append(shape)
        
        logger.info(f"📊 실제 텍스트 요소: {len(real_shapes)}개")
        
        # Issue 1: 텍스트 박스 수 불일치 (자의적 생성/삭제 의심)
        # 주의: 메타데이터의 모든 박스가 텍스트용은 아닐 수 있으므로 단순 개수 비교는 참고용
        # 메타데이터에서 텍스트가 있는 요소만 카운트
        meta_text_count = 0
        for box in meta_boxes:
            # 'text' 필드가 있고 내용이 있거나, 'type'이 TEXT_BOX인 경우
            has_text = False
            if box.get('type') == 'TEXT_BOX':
                has_text = True
            elif box.get('text') and isinstance(box.get('text'), dict) and box['text'].get('raw'):
                has_text = True
            
            if has_text:
                meta_text_count += 1
                
        if len(real_shapes) > len(meta_boxes):
            logger.warning(f"🚨 [Issue] 텍스트 요소 과다: 메타({len(meta_boxes)}) < 실제({len(real_shapes)})")
            logger.warning(f"   -> AI가 자의적으로 텍스트 박스를 추가했거나, 텍스트가 없어야 할 도형에 텍스트를 넣었을 가능성 있음")
            total_issues += 1
            
        # Issue 2: 위치/크기 기반 매칭 및 검증
        # 메타데이터 좌표는 픽셀(px) 또는 인치(inch) 단위일 수 있음
        # pptx 라이브러리는 EMU 단위 사용 (1 inch = 914400 EMU, 1 px = 9525 EMU approx)
        
        matched_count = 0
        unexpected_text_shapes = []
        
        # 메타데이터 단위 추정 (width가 20 이하면 inch, 100 이상이면 px)
        is_px_unit = False
        if meta_boxes and meta_boxes[0].get('width_px'):
             is_px_unit = True
        elif meta_boxes and meta_boxes[0].get('width', 0) > 50:
             is_px_unit = True
             
        for shape in real_shapes:
            # 좌표 변환
            if is_px_unit:
                # EMU -> PX (96 DPI 기준: 1 px = 9525 EMU)
                left_val = shape.left / 9525
                top_val = shape.top / 9525
                width_val = shape.width / 9525
                height_val = shape.height / 9525
                tolerance = 20 # 20px 오차 허용
            else:
                # EMU -> Inch
                left_val = shape.left / 914400
                top_val = shape.top / 914400
                width_val = shape.width / 914400
                height_val = shape.height / 914400
                tolerance = 0.5 # 0.5 inch 오차 허용
            
            # 메타데이터 박스와 매칭 (위치 기반 근사 매칭)
            is_matched = False
            for box in meta_boxes:
                # 키 이름 호환성 (x/left, y/top, width, height)
                box_x = box.get('left_px', box.get('x', box.get('left', 0)))
                box_y = box.get('top_px', box.get('y', box.get('top', 0)))
                
                if (abs(box_x - left_val) < tolerance and 
                    abs(box_y - top_val) < tolerance):
                    is_matched = True
                    matched_count += 1
                    break
            
            if not is_matched:
                unexpected_text_shapes.append(shape)

        if unexpected_text_shapes:
            logger.warning(f"🚨 [Issue] 메타데이터에 없는 위치의 텍스트 요소 발견 ({len(unexpected_text_shapes)}개)")
            for s in unexpected_text_shapes:
                text_preview = s.text[:20].replace('\n', ' ')
                pos_str = f"{s.left/9525:.1f}px, {s.top/9525:.1f}px" if is_px_unit else f"{s.left/914400:.2f}in, {s.top/914400:.2f}in"
                logger.warning(f"   - 텍스트: '{text_preview}...' (위치: {pos_str})")
            total_issues += 1
            
    logger.info(f"\n=== 🏁 분석 완료: 총 {total_issues}개의 잠재적 문제 발견 ===")

if __name__ == "__main__":
    from pathlib import Path
    backend_root = Path(__file__).parent.parent
    pptx_file = str(backend_root / "uploads" / "자동차 산업의 특허분석 방법론.pptx")
    metadata_file = str(backend_root / "uploads" / "templates" / "users" / "8" / "metadata" / "제품소개서_샘플_metadata.json")
    
    analyze_ppt_compliance(pptx_file, metadata_file)
