"""
Template Application Tool for Enhanced PPT Generation

템플릿 PPTX 파일에 DeckSpec 내용을 적용하는 도구

Author: Presentation System
Created: 2025-01-20
Phase: 2.1
"""

import logging
from pathlib import Path
from typing import Any, ClassVar, Dict, List, Optional, Type

from langchain_core.tools import BaseTool
from pydantic import BaseModel, Field

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from app.services.presentation.ppt_models import DeckSpec, SlideSpec
from app.core.config import settings

logger = logging.getLogger(__name__)


class TemplateApplicationInput(BaseModel):
    """Input schema for TemplateApplicationTool"""
    deck_spec: Dict[str, Any] = Field(..., description="DeckSpec dictionary")
    template_path: str = Field(..., description="템플릿 PPTX 파일 경로")
    text_box_mappings: Optional[List[Dict[str, Any]]] = Field(
        default=None,
        description="텍스트박스 매핑 정보 (UI에서 제공)"
    )
    file_basename: Optional[str] = Field(default=None, description="출력 파일명")


class TemplateApplicationTool(BaseTool):
    """
    템플릿 PPTX 파일에 구조화된 콘텐츠(DeckSpec)를 적용하는 도구
    
    기능:
    - 템플릿 파일 로드 및 분석
    - 템플릿 레이아웃에 DeckSpec 매핑
    - 플레이스홀더 자동 감지 및 채우기
    - 사용자 정의 텍스트박스 매핑 지원
    - 콘텐츠 적용 후 PPTX 파일 생성
    
    입력:
    - deck_spec: SlideSpec 리스트 포함
    - template_path: 템플릿 파일 경로
    - text_box_mappings: 선택적 매핑 정보
    
    출력:
    {
        "success": True,
        "file_path": "/path/to/output.pptx",
        "slides_processed": 5,
        "template_used": "business_template.pptx"
    }
    """
    
    name: str = "template_application_tool"
    description: str = (
        "템플릿 PPTX 파일에 구조화된 프레젠테이션 콘텐츠를 적용합니다. "
        "플레이스홀더 자동 감지 및 사용자 정의 매핑을 지원합니다."
    )
    args_schema: Type[BaseModel] = TemplateApplicationInput
    
    # 클래스 변수
    upload_dir: Path = settings.resolved_upload_dir
    
    def _run(self, *args, **kwargs):
        """Synchronous wrapper for async _arun."""
        import asyncio
        try:
            loop = asyncio.get_event_loop()
        except RuntimeError:
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
        return loop.run_until_complete(self._arun(*args, **kwargs))

    async def _arun(
        self,
        deck_spec: Dict[str, Any],
        template_path: str,
        text_box_mappings: Optional[List[Dict[str, Any]]] = None,
        file_basename: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        템플릿에 DeckSpec 적용 (비동기)
        
        Args:
            deck_spec: DeckSpec dictionary
            template_path: 템플릿 파일 경로
            text_box_mappings: 선택적 매핑 정보
            file_basename: 출력 파일명
        
        Returns:
            Dict with success status and file path
        """
        try:
            # Ensure upload directory exists
            self.upload_dir.mkdir(parents=True, exist_ok=True)
            
            # Parse DeckSpec
            spec = DeckSpec(**deck_spec)
            logger.info(f"🎨 [TemplateApp] 시작: template={template_path}")
            logger.info(f"📊 슬라이드 수: {len(spec.slides)}개")
            
            # 템플릿 경로 검증
            template_file = Path(template_path)
            if not template_file.exists():
                raise FileNotFoundError(f"템플릿 파일이 존재하지 않습니다: {template_path}")
            
            # 템플릿 로드
            prs = Presentation(str(template_file))
            logger.info(f"✅ 템플릿 로드 완료: {len(prs.slides)}개 슬라이드")
            
            # 매핑 방식 결정
            if text_box_mappings:
                # 사용자 정의 매핑 사용
                logger.info(f"🔧 사용자 정의 매핑 사용: {len(text_box_mappings)}개")
                prs = self._apply_custom_mappings(prs, spec, text_box_mappings)
            else:
                # 자동 플레이스홀더 매핑
                logger.info("🤖 자동 플레이스홀더 매핑 사용")
                prs = self._apply_automatic_mapping(prs, spec)
            
            # 파일 저장
            from datetime import datetime
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            basename = file_basename or f"presentation_template_{timestamp}"
            if not basename.endswith('.pptx'):
                basename += '.pptx'
            
            output_path = self.upload_dir / basename
            prs.save(str(output_path))
            
            logger.info(f"✅ [TemplateApp] 완료: {output_path}")
            
            return {
                "success": True,
                "file_path": str(output_path),
                "filename": basename,
                "slides_processed": len(spec.slides),
                "template_used": template_file.name,
                "mapping_mode": "custom" if text_box_mappings else "automatic"
            }
            
        except Exception as e:
            logger.error(f"❌ [TemplateApp] 실패: {e}", exc_info=True)
            return {
                "success": False,
                "error": str(e),
                "file_path": None,
            }

    def _apply_custom_mappings(
        self,
        prs: Presentation,
        spec: DeckSpec,
        mappings: List[Dict[str, Any]]
    ) -> Presentation:
        """
        사용자 정의 매핑 적용
        
        매핑 형식:
        {
            "slideIndex": 0,
            "shapeIndex": 2,
            "contentType": "title" | "bullets" | "key_message",
            "placeholder": "{{title}}" | "{{content}}"
        }
        """
        try:
            logger.info(f"🔧 사용자 매핑 적용 시작: {len(mappings)}개")
            
            # 매핑을 슬라이드별로 그룹화
            mappings_by_slide = {}
            for mapping in mappings:
                slide_idx = mapping.get('slideIndex', 0)
                if slide_idx not in mappings_by_slide:
                    mappings_by_slide[slide_idx] = []
                mappings_by_slide[slide_idx].append(mapping)
            
            # 각 슬라이드에 매핑 적용
            for slide_idx, slide_mappings in mappings_by_slide.items():
                if slide_idx < len(spec.slides) and slide_idx < len(prs.slides):
                    slide = prs.slides[slide_idx]
                    slide_spec = spec.slides[slide_idx]
                    
                    self._apply_mappings_to_slide(slide, slide_spec, slide_mappings)
                    logger.info(f"  ✅ 슬라이드 {slide_idx}: {len(slide_mappings)}개 매핑 적용")
            
            return prs
            
        except Exception as e:
            logger.error(f"사용자 매핑 적용 실패: {e}")
            return prs

    def _apply_mappings_to_slide(
        self,
        slide,
        slide_spec: SlideSpec,
        mappings: List[Dict[str, Any]]
    ):
        """개별 슬라이드에 매핑 적용"""
        try:
            for mapping in mappings:
                shape_idx = mapping.get('shapeIndex')
                content_type = mapping.get('contentType', 'text')
                
                if shape_idx is None or shape_idx >= len(slide.shapes):
                    continue
                
                shape = slide.shapes[shape_idx]
                
                # 콘텐츠 타입에 따라 처리
                if content_type == 'title':
                    self._set_shape_text(shape, slide_spec.title)
                elif content_type == 'key_message':
                    self._set_shape_text(shape, slide_spec.key_message)
                elif content_type == 'bullets':
                    self._set_shape_bullets(shape, slide_spec.bullets)
                elif content_type == 'text':
                    # 일반 텍스트 (placeholder에 따라)
                    placeholder = mapping.get('placeholder', '')
                    if '{{title}}' in placeholder:
                        self._set_shape_text(shape, slide_spec.title)
                    elif '{{content}}' in placeholder:
                        self._set_shape_bullets(shape, slide_spec.bullets)
                    
        except Exception as e:
            logger.warning(f"슬라이드 매핑 적용 중 오류: {e}")

    def _apply_automatic_mapping(
        self,
        prs: Presentation,
        spec: DeckSpec
    ) -> Presentation:
        """
        자동 플레이스홀더 매핑
        
        템플릿의 플레이스홀더를 자동으로 감지하여 DeckSpec 내용 채우기
        """
        try:
            logger.info("🤖 자동 매핑 시작")
            
            # 템플릿 슬라이드 수와 spec 슬라이드 수 비교
            template_slide_count = len(prs.slides)
            spec_slide_count = len(spec.slides)
            
            logger.info(f"  템플릿: {template_slide_count}개, Spec: {spec_slide_count}개")
            
            # 슬라이드 추가가 필요한 경우
            if spec_slide_count > template_slide_count:
                # 마지막 슬라이드 레이아웃을 복제하여 추가
                last_layout = prs.slides[template_slide_count - 1].slide_layout
                for _ in range(spec_slide_count - template_slide_count):
                    prs.slides.add_slide(last_layout)
                logger.info(f"  ➕ {spec_slide_count - template_slide_count}개 슬라이드 추가")
            
            # 각 슬라이드에 콘텐츠 적용
            for i, slide_spec in enumerate(spec.slides):
                if i < len(prs.slides):
                    slide = prs.slides[i]
                    self._fill_slide_placeholders(slide, slide_spec)
                    logger.info(f"  ✅ 슬라이드 {i+1}: '{slide_spec.title}' 적용")
            
            return prs
            
        except Exception as e:
            logger.error(f"자동 매핑 실패: {e}")
            return prs

    def _fill_slide_placeholders(self, slide, slide_spec: SlideSpec):
        """슬라이드의 플레이스홀더를 자동으로 채우기"""
        try:
            title_filled = False
            content_filled = False
            
            # 플레이스홀더 검색 및 채우기
            for shape in slide.shapes:
                if not hasattr(shape, 'placeholder_format'):
                    continue
                
                ph_type = shape.placeholder_format.type
                
                # 제목 플레이스홀더 (TITLE = 1, CENTER_TITLE = 3)
                if ph_type in [1, 3] and not title_filled:
                    self._set_shape_text(shape, slide_spec.title)
                    title_filled = True
                    logger.debug(f"    제목 적용: {shape.placeholder_format.idx}")
                
                # 본문 플레이스홀더 (BODY = 2, OBJECT = 7)
                elif ph_type in [2, 7] and not content_filled:
                    if slide_spec.key_message:
                        # 키 메시지 + 불릿
                        text = slide_spec.key_message
                        if slide_spec.bullets:
                            text += "\n\n" + "\n".join(f"• {b}" for b in slide_spec.bullets[:6])
                        self._set_shape_text(shape, text)
                    else:
                        # 불릿만
                        self._set_shape_bullets(shape, slide_spec.bullets)
                    content_filled = True
                    logger.debug(f"    본문 적용: {shape.placeholder_format.idx}")
            
            # 플레이스홀더가 없는 경우 텍스트박스 검색
            if not title_filled or not content_filled:
                self._fill_text_boxes(slide, slide_spec, title_filled, content_filled)
                
        except Exception as e:
            logger.warning(f"플레이스홀더 채우기 실패: {e}")

    def _fill_text_boxes(
        self,
        slide,
        slide_spec: SlideSpec,
        title_filled: bool,
        content_filled: bool
    ):
        """플레이스홀더가 없는 경우 일반 텍스트박스에 채우기"""
        try:
            text_boxes = []
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and not hasattr(shape, 'placeholder_format'):
                    # 위치 기반으로 제목/본문 구분 (상단 30% = 제목)
                    slide_height = slide.height if hasattr(slide, 'height') else Inches(7.5)
                    is_title_area = shape.top < (slide_height * 0.3)
                    text_boxes.append((shape, is_title_area))
            
            # 제목 영역 텍스트박스에 제목 채우기
            if not title_filled:
                for shape, is_title in text_boxes:
                    if is_title:
                        self._set_shape_text(shape, slide_spec.title)
                        title_filled = True
                        break
            
            # 본문 영역 텍스트박스에 본문 채우기
            if not content_filled:
                for shape, is_title in text_boxes:
                    if not is_title:
                        self._set_shape_bullets(shape, slide_spec.bullets)
                        content_filled = True
                        break
                        
        except Exception as e:
            logger.warning(f"텍스트박스 채우기 실패: {e}")

    def _set_shape_text(self, shape, text: str):
        """Shape에 텍스트 설정"""
        try:
            if not text or not hasattr(shape, 'text_frame'):
                return
            
            text_frame = shape.text_frame
            text_frame.clear()  # 기존 내용 제거
            text_frame.text = text
            
            # 기본 스타일 적용
            if text_frame.paragraphs:
                para = text_frame.paragraphs[0]
                para.font.size = Pt(18)
                para.font.name = '맑은 고딕'
                
        except Exception as e:
            logger.warning(f"텍스트 설정 실패: {e}")

    def _set_shape_bullets(self, shape, bullets: List[str]):
        """Shape에 불릿 포인트 설정"""
        try:
            if not bullets or not hasattr(shape, 'text_frame'):
                return
            
            text_frame = shape.text_frame
            text_frame.clear()
            
            # 불릿 아이콘
            bullet_icons = ["🔹", "🔸", "💎", "⭐", "🎯", "📌"]
            
            for i, bullet in enumerate(bullets[:8]):  # 최대 8개
                if i == 0:
                    para = text_frame.paragraphs[0]
                else:
                    para = text_frame.add_paragraph()
                
                para.text = f"{bullet_icons[i % len(bullet_icons)]} {bullet}"
                para.level = 0
                para.font.size = Pt(14)
                para.font.name = '맑은 고딕'
                para.space_after = Pt(6)
                
        except Exception as e:
            logger.warning(f"불릿 설정 실패: {e}")


# 전역 인스턴스
template_application_tool = TemplateApplicationTool()
