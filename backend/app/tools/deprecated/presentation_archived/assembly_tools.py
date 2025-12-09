"""Assembly Tools for Presentation Agent."""
from __future__ import annotations

import os
import re
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

from loguru import logger
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

try:
    from langchain_core.tools import BaseTool
except ImportError:
    from langchain.tools import BaseTool

from app.core.config import settings
from app.services.presentation.ppt_models import DeckSpec, SlideSpec
from app.services.presentation.enhanced_object_processor import EnhancedPPTObjectProcessor
from app.services.presentation.ppt_template_manager import template_manager
from app.tools.presentation.archived.visualization_tools import DiagramBuilderTool
from app.tools.presentation.archived.design_tools import TemplateSelectorTool, LayoutOptimizerTool


class SlideBuilderTool(BaseTool):
    """Tool for building individual slides with content."""
    
    name: str = "slide_builder"
    description: str = "Add content (text, images) to a slide."
    
    diagram_builder: DiagramBuilderTool = DiagramBuilderTool()

    def _run(self, slide: Any, spec: SlideSpec, colors: Dict[str, Any], palette: List[RGBColor], 
             template_style: str, include_charts: bool) -> None:
        """Build a single slide."""
        
        # 1. Add Text Content
        self._add_text_content(slide, spec, colors)
        
        # 2. Add Visual Suggestion (Image)
        if spec.visual_suggestion:
            self._maybe_add_image(slide, spec.visual_suggestion)
            
        # 3. Add Diagram/Chart
        if include_charts and spec.diagram and spec.diagram.type != 'none':
            self.diagram_builder._run(slide, spec.diagram, colors, palette)
            
        # 4. Add Speaker Notes
        if spec.speaker_notes:
            try:
                notes = slide.part.notes_slide.notes_text_frame
                notes.text = spec.speaker_notes[:1500]
            except Exception:
                pass

    async def _arun(self, *args, **kwargs):
        return self._run(*args, **kwargs)

    def _add_text_content(self, slide: Any, spec: SlideSpec, colors: Dict[str, Any]) -> None:
        tf = self._find_text_frame(slide)
        
        if tf:
            try:
                tf.clear()
                tf.word_wrap = True
            except Exception:
                pass
            
            agenda_mode = (spec.style and spec.style.get('agenda')) or (spec.title in ['목차','Agenda','Contents'] and not spec.key_message)
            
            if agenda_mode:
                self._add_agenda_content(tf, spec, colors)
            else:
                self._add_standard_content(tf, spec, colors)
        else:
            logger.warning(f"No text frame found for slide '{spec.title}'")

    def _find_text_frame(self, slide: Any) -> Any:
        # Try placeholder 1
        try:
            if len(slide.placeholders) > 1:
                ph = slide.placeholders[1]
                if getattr(ph, 'has_text_frame', False) and hasattr(ph, 'text_frame'):
                    return ph.text_frame
        except Exception:
            pass
            
        # Fallback: find first non-title shape
        try:
            for sh in slide.shapes:
                if getattr(sh, 'has_text_frame', False):
                    try:
                        if sh is slide.shapes.title:
                            continue
                    except Exception:
                        pass
                    return sh.text_frame
        except Exception:
            pass
        return None

    def _add_agenda_content(self, tf: Any, spec: SlideSpec, colors: Dict[str, Any]) -> None:
        for i, b in enumerate(spec.bullets):
            txt = b.strip()
            if len(txt) > 60: txt = txt[:60] + '…'
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = txt
            p.level = 0
            p.font.size = Pt(20)
            p.font.color.rgb = colors['text']

    def _add_standard_content(self, tf: Any, spec: SlideSpec, colors: Dict[str, Any]) -> None:
        # Key Message
        if spec.key_message and spec.key_message.strip():
            p = tf.paragraphs[0]
            p.text = spec.key_message.strip()
            p.level = 0
            try:
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = colors['text']
            except Exception:
                pass
        
        # Bullets
        for i, bullet in enumerate(spec.bullets):
            if bullet and bullet.strip():
                txt = bullet.strip()
                if len(txt) > 80: txt = txt[:80] + '…'
                
                if spec.key_message and spec.key_message.strip():
                    p = tf.add_paragraph()
                else:
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                
                p.text = f"• {txt}"
                p.level = 1
                try:
                    p.font.size = Pt(18)
                    p.font.color.rgb = colors['text']
                except Exception:
                    pass

    def _maybe_add_image(self, slide: Any, suggestion: str) -> None:
        if not (os.environ.get('PPT_IMAGE_FETCH') == '1' and os.environ.get('UNSPLASH_ACCESS_KEY')):
            return
        try:
            import requests
            query = suggestion.split()[0][:40] if suggestion else 'technology'
            resp = requests.get(
                'https://api.unsplash.com/photos/random', 
                params={'query': query, 'content_filter':'high','orientation':'landscape'}, 
                headers={'Authorization': f"Client-ID {os.environ['UNSPLASH_ACCESS_KEY']}"}, 
                timeout=4
            )
            if resp.status_code != 200: return
            image_url = resp.json().get('urls', {}).get('small')
            if not image_url: return
            
            img = requests.get(image_url, timeout=4).content
            with tempfile.NamedTemporaryFile(delete=False, suffix='.jpg') as tmp:
                tmp.write(img)
                path = tmp.name
            
            slide.shapes.add_picture(path, Inches(8.2), Inches(3.0), width=Inches(4.5))
        except Exception as e:
            logger.debug(f"Image fetch failed: {e}")


class SlideAssemblerTool(BaseTool):
    """Tool for assembling the final presentation file."""
    
    name: str = "slide_assembler"
    description: str = "Assemble slides into a PPTX file."
    
    template_selector: TemplateSelectorTool = TemplateSelectorTool()
    layout_optimizer: LayoutOptimizerTool = LayoutOptimizerTool()
    slide_builder: SlideBuilderTool = SlideBuilderTool()
    object_processor: EnhancedPPTObjectProcessor = EnhancedPPTObjectProcessor()

    def _run(self, spec: DeckSpec, file_basename: Optional[str] = None,
             template_style: str = "business", include_charts: bool = True,
             custom_template_path: Optional[str] = None,
             user_template_id: Optional[str] = None,
             text_box_mappings: Optional[List[Dict[str, Any]]] = None,
             content_segments: Optional[List[Dict[str, Any]]] = None) -> str:
        """Build the presentation."""
        
        # 1. Custom Template Build
        if custom_template_path and Path(custom_template_path).exists():
            return self._build_from_template(spec, Path(custom_template_path), file_basename, text_box_mappings, content_segments)
        
        # 2. Legacy/Standard Build
        return self._build_standard_pptx(spec, file_basename, template_style, include_charts)

    async def _arun(self, *args, **kwargs):
        return self._run(*args, **kwargs)

    def _build_from_template(self, spec: DeckSpec, template_path: Path, file_basename: Optional[str],
                           text_box_mappings: Optional[List[Dict[str, Any]]],
                           content_segments: Optional[List[Dict[str, Any]]]) -> str:
        try:
            prs = Presentation(str(template_path))
            
            if text_box_mappings:
                prs = self.object_processor.apply_object_mappings(prs, text_box_mappings, content_segments)
            else:
                template_spec = template_manager.analyze_template(template_path)
                if template_spec:
                    adapted_spec = template_manager.map_deck_to_template(spec, template_spec)
                    prs = template_manager.build_from_template(adapted_spec, template_path)
            
            return self._save_presentation(prs, spec.topic, file_basename)
            
        except Exception as e:
            logger.error(f"Template build failed: {e}")
            return self._build_standard_pptx(spec, file_basename, "business", True)

    def _build_standard_pptx(self, spec: DeckSpec, file_basename: Optional[str],
                           template_style: str, include_charts: bool) -> str:
        
        # Load Template & Optimize Layout
        prs = self.template_selector._run(template_style)
        design = self.layout_optimizer._run(template_style, spec)
        colors = design["colors"]
        palette = design["palette"]

        # Build Slides
        for idx, slide_spec in enumerate(spec.slides[:spec.max_slides]):
            # Select Layout
            if idx == 0:
                layout = prs.slide_layouts[0]
            else:
                if slide_spec.diagram and slide_spec.diagram.type == 'chart' and slide_spec.layout not in ['two-content', 'title-only']:
                    slide_spec.layout = 'two-content'
                layout_map = {"title-only":5, "title-and-content":1, "two-content":3, "section-header":2, "blank":6}
                layout = prs.slide_layouts[layout_map.get(slide_spec.layout, 1)]
            
            slide = prs.slides.add_slide(layout)
            
            # Apply Background (Modern/Playful)
            if template_style in ['modern', 'playful']:
                self._apply_background(slide, colors)

            # Set Title
            if slide.shapes.title:
                slide.shapes.title.text = slide_spec.title or spec.topic
                ft = slide.shapes.title.text_frame.paragraphs[0].font
                ft.color.rgb = colors['primary']
                ft.size = Pt(42 if idx == 0 else 26)
                ft.bold = True

            # Build Content
            if idx > 0: # Skip content for title slide (handled differently usually, or simple title)
                 self.slide_builder._run(slide, slide_spec, colors, palette, template_style, include_charts)
            
            # Add Footer
            if idx >= 1:
                self._add_footer(slide, spec.topic, idx, spec.max_slides, colors)

        return self._save_presentation(prs, spec.topic, file_basename)

    def _apply_background(self, slide: Any, colors: Dict[str, Any]) -> None:
        try:
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
            bg.fill.solid()
            bg.fill.fore_color.rgb = colors['background']
            bg.line.fill.background()
        except Exception:
            pass

    def _add_footer(self, slide: Any, topic: str, idx: int, total: int, colors: Dict[str, Any]) -> None:
        try:
            margin_x = Inches(0.6)
            footer_text_top = Inches(6.6)
            width_main = Inches(9.5)
            
            box = slide.shapes.add_textbox(margin_x, footer_text_top, width_main, Inches(0.25))
            tf = box.text_frame; tf.clear(); p = tf.paragraphs[0]
            p.text = f"{topic} | {datetime.now().strftime('%Y-%m-%d')}"
            p.font.size = Pt(8); p.font.color.rgb = colors['text']
            
            page_left = Inches(11.2)
            pbox = slide.shapes.add_textbox(page_left, footer_text_top, Inches(1.5), Inches(0.25))
            ptf = pbox.text_frame; ptf.clear(); p2 = ptf.paragraphs[0]
            p2.text = f"{idx+1}/{total}"
            p2.font.size = Pt(8); p2.font.color.rgb = colors['text']; p2.alignment = PP_ALIGN.RIGHT
        except Exception:
            pass

    def _save_presentation(self, prs: Presentation, topic: str, file_basename: Optional[str]) -> str:
        upload_dir = settings.resolved_upload_dir
        upload_dir.mkdir(parents=True, exist_ok=True)
        
        def _sanitize_name(text: str) -> str:
            cleaned = re.sub(r"\[\[PPT_OPTS:.*?\]\]", "", text)
            cleaned = re.sub(r"[\n\r\t]", " ", cleaned).strip() or "presentation"
            cleaned = re.sub(r"[^0-9A-Za-z가-힣 _-]", "_", cleaned)
            cleaned = re.sub(r"\s+", "_", cleaned)[:40]
            return cleaned or "presentation"
            
        base_name = _sanitize_name(topic)
        if file_basename:
            fb_clean = _sanitize_name(file_basename)
            fname = fb_clean if not re.search(r"(해주세요|해줘|\?$)", fb_clean) else f"enhanced_presentation_{base_name}"
        else:
            fname = f"enhanced_presentation_{base_name}"
            
        out_path = upload_dir / f"{fname}.pptx"
        prs.save(str(out_path))
        logger.info(f"PPTX Saved: {out_path}")
        return str(out_path)