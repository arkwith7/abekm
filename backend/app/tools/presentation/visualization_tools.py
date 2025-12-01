"""Visualization Tools for Presentation Agent."""
from __future__ import annotations

from typing import Any, Dict, List, Optional

from loguru import logger
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

try:
    from langchain_core.tools import BaseTool
except ImportError:
    from langchain.tools import BaseTool

from app.services.presentation.ppt_models import ChartData, DiagramData


class ChartGeneratorTool(BaseTool):
    """Tool for generating charts in PowerPoint slides."""
    
    name: str = "chart_generator"
    description: str = "Generate charts (column, bar, line, pie) on a slide."

    def _run(self, slide: Any, chart_data: ChartData, colors: Dict[str, Any], palette: List[RGBColor]) -> None:
        """Add a chart to the slide."""
        if not chart_data.categories or not chart_data.series:
            return

        data = CategoryChartData()
        data.categories = chart_data.categories
        for s in chart_data.series:
            data.add_series(s.get('name', 'Series'), s.get('values', []))

        chart_type_map = {
            'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'bar': XL_CHART_TYPE.BAR_CLUSTERED,
            'line': XL_CHART_TYPE.LINE,
            'pie': XL_CHART_TYPE.PIE
        }
        chart_type = chart_type_map.get(chart_data.type, XL_CHART_TYPE.COLUMN_CLUSTERED)

        # Default position and size
        x, y, cx, cy = Inches(4), Inches(2), Inches(5), Inches(4)
        
        try:
            chart = slide.shapes.add_chart(chart_type, x, y, cx, cy, data).chart
            
            # Apply colors
            try:
                for i, s in enumerate(chart.series):
                    f = s.format.fill
                    f.solid()
                    f.fore_color.rgb = palette[i % len(palette)]
            except Exception:
                pass

            # Set title
            if hasattr(chart, 'chart_title') and chart_data.title:
                chart.chart_title.text_frame.text = chart_data.title
                
        except Exception as e:
            logger.error(f"Failed to add chart: {e}")

    async def _arun(self, *args, **kwargs):
        return self._run(*args, **kwargs)


class DiagramBuilderTool(BaseTool):
    """Tool for building diagrams (tables, flows) in PowerPoint slides."""
    
    name: str = "diagram_builder"
    description: str = "Build diagrams like tables and flowcharts on a slide."
    
    chart_generator: ChartGeneratorTool = ChartGeneratorTool()

    def _run(self, slide: Any, diagram: DiagramData, colors: Dict[str, Any], palette: List[RGBColor]) -> None:
        """Add a diagram to the slide."""
        try:
            if diagram.type == 'chart' and diagram.chart:
                self.chart_generator._run(slide, diagram.chart, colors, palette)
            elif diagram.type == 'table' and diagram.data:
                self._add_table(slide, diagram.data, colors)
            elif diagram.type == 'flow' and diagram.data:
                self._add_flow(slide, diagram.data, colors)
        except Exception as e:
            logger.warning(f"Failed to add diagram: {e}")

    async def _arun(self, *args, **kwargs):
        return self._run(*args, **kwargs)

    def _add_table(self, slide: Any, table_data: Dict[str, Any], colors: Dict[str, Any]) -> None:
        headers = table_data.get('headers', [])
        rows = table_data.get('rows', [])
        if not headers or not rows:
            return

        x, y, cx, cy = Inches(1), Inches(3), Inches(8), Inches(3)
        table = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, cx, cy).table

        # Header
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = str(h)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = colors.get('secondary', RGBColor(100, 100, 100))

        # Rows
        for r_idx, row in enumerate(rows):
            for c_idx, val in enumerate(row):
                if c_idx < len(headers):
                    table.cell(r_idx + 1, c_idx).text = str(val)

    def _add_flow(self, slide: Any, flow: Dict[str, Any], colors: Dict[str, Any]) -> None:
        steps = flow.get('steps', [])
        if not steps:
            return

        box_w, box_h = Inches(1.5), Inches(0.8)
        start_x, start_y, spacing = Inches(1), Inches(3), Inches(2)

        for i, step in enumerate(steps[:5]):
            x = start_x + i * spacing
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, start_y, box_w, box_h)
            
            shape.fill.solid()
            shape.fill.fore_color.rgb = colors.get('primary', RGBColor(0, 0, 255))
            shape.line.color.rgb = colors.get('text', RGBColor(0, 0, 0))
            
            tf = shape.text_frame
            tf.text = str(step)
            p = tf.paragraphs[0]
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER

            # Connector
            if i < len(steps) - 1 and i < 4:
                slide.shapes.add_connector(
                    1, 
                    x + box_w + Inches(0.1), start_y + box_h / 2, 
                    x + box_w + Inches(0.4), start_y + box_h / 2
                )
