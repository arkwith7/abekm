"""Style analysis tool for PPTX references."""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any, Dict, List, Optional, Type

from loguru import logger
from pptx import Presentation
from pptx.dml.color import RGBColor
from pydantic import BaseModel, Field

try:  # pragma: no cover - optional dependency for LangChain compatibility
    from langchain_core.tools import BaseTool  # type: ignore
except ImportError:  # pragma: no cover
    from langchain.tools import BaseTool  # type: ignore


class StyleAnalysisInput(BaseModel):
    template_path: str = Field(..., description="분석할 PPTX 파일 경로")


class StyleAnalysisTool(BaseTool):
    """Analyzes a PPTX file and extracts styling metadata."""

    name: str = "style_analysis_tool"
    description: str = (
        "Analyzes an existing PPTX file to extract color palette, font families, "
        "and representative slide layouts so we can reuse the style in newly "
        "generated presentations."
    )
    args_schema: Type[BaseModel] = StyleAnalysisInput

    def _run(self, template_path: str) -> Dict[str, Any]:  # pragma: no cover - sync fallback
        import asyncio

        return asyncio.run(self._arun(template_path=template_path))

    async def _arun(self, template_path: str) -> Dict[str, Any]:
        try:
            ppt_path = Path(template_path)
            if not ppt_path.exists():
                raise FileNotFoundError(f"템플릿을 찾을 수 없습니다: {ppt_path}")

            presentation = Presentation(str(ppt_path))
            style_metadata = {
                "color_palette": self._extract_colors(presentation),
                "fonts": self._extract_fonts(presentation),
                "layouts": self._extract_layouts(presentation),
            }

            logger.info(
                "🎨 [StyleAnalysis] 완료: colors=%d, fonts=%d, layouts=%d",
                len(style_metadata["color_palette"]),
                len(style_metadata["fonts"]),
                len(style_metadata["layouts"]),
            )

            return {
                "success": True,
                "style_metadata": style_metadata,
                "color_palette": style_metadata["color_palette"],
                "fonts": style_metadata["fonts"],
                "layouts": style_metadata["layouts"],
            }
        except Exception as exc:  # pragma: no cover - defensive
            logger.error("❌ [StyleAnalysis] 실패: %s", exc, exc_info=True)
            return {"success": False, "error": str(exc)}

    def _extract_colors(self, presentation: Presentation) -> List[str]:
        colors: List[str] = []
        try:
            theme = presentation.slide_master.theme
            theme_colors = getattr(theme, "theme_colors", None)
            if theme_colors:
                for color in theme_colors:
                    rgb = getattr(color, "rgb", None)
                    if isinstance(rgb, RGBColor):
                        colors.append(self._rgb_to_hex(rgb))
        except Exception:
            pass
        # Fallback: inspect shapes in first few slides
        if not colors:
            for slide in presentation.slides[:3]:
                for shape in slide.shapes:
                    fill = getattr(shape, "fill", None)
                    if fill and fill.type == 1 and getattr(fill, "fore_color", None):
                        rgb = getattr(fill.fore_color, "rgb", None)
                        if isinstance(rgb, RGBColor):
                            colors.append(self._rgb_to_hex(rgb))
                    line = getattr(shape, "line", None)
                    if line and getattr(line, "color", None):
                        rgb = getattr(line.color, "rgb", None)
                        if isinstance(rgb, RGBColor):
                            colors.append(self._rgb_to_hex(rgb))
        return list(dict.fromkeys(colors))[:10]

    def _extract_fonts(self, presentation: Presentation) -> List[str]:
        fonts: List[str] = []
        for slide in presentation.slides[:5]:
            for shape in slide.shapes:
                text_frame = getattr(shape, "text_frame", None)
                if not text_frame:
                    continue
                for paragraph in text_frame.paragraphs:
                    font = getattr(paragraph.font, "name", None)
                    if font:
                        fonts.append(font)
        return list(dict.fromkeys(fonts))[:5]

    def _extract_layouts(self, presentation: Presentation) -> List[Dict[str, Any]]:
        layouts: List[Dict[str, Any]] = []
        for layout in presentation.slide_layouts:
            shapes_summary = []
            for shape in layout.shapes:
                shapes_summary.append(
                    {
                        "name": getattr(shape, "name", ""),
                        "has_text_frame": bool(getattr(shape, "text_frame", None)),
                        "has_table": bool(getattr(shape, "table", None)),
                    }
                )
            layouts.append(
                {
                    "name": layout.name,
                    "shape_count": len(layout.shapes),
                    "shapes": shapes_summary[:5],
                }
            )
        return layouts[:10]

    @staticmethod
    def _rgb_to_hex(rgb: RGBColor) -> str:
        return f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


style_analysis_tool = StyleAnalysisTool()

__all__ = ["StyleAnalysisTool", "style_analysis_tool"]
