"""PPT Quality Validator Tool - Validates generated PPT against source content."""

from __future__ import annotations

import re
from collections import Counter
from difflib import SequenceMatcher
from pathlib import Path
from typing import Any, Dict, List, Optional, Set, Tuple, Type

from loguru import logger
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pydantic import BaseModel, Field

try:
    from langchain_core.tools import BaseTool
except ImportError:
    from langchain_core.tools import BaseTool


class PPTQualityValidatorInput(BaseModel):
    """Input schema for PPTQualityValidatorTool."""

    pptx_path: str = Field(..., description="생성된 PPTX 파일 경로")
    source_content: str = Field(..., description="원본 RAG 응답 또는 컨텍스트 텍스트")
    source_outline: Optional[Dict] = Field(
        default=None, description="원본 아웃라인 (DeckSpec)"
    )


class SlideAnalysis(BaseModel):
    """Analysis result for a single slide."""

    slide_index: int
    title: str
    text_content: str
    bullet_count: int
    has_chart: bool
    has_table: bool
    has_image: bool
    has_diagram: bool
    shape_count: int
    word_count: int


class ContentCoverage(BaseModel):
    """Content coverage analysis."""

    total_keywords: int
    matched_keywords: int
    coverage_rate: float
    missing_keywords: List[str]
    matched_keywords_list: List[str]


class StructureAnalysis(BaseModel):
    """Structure comparison analysis."""

    source_sections: int
    ppt_slides: int
    section_mapping: List[Dict[str, Any]]
    structure_match_rate: float


class VisualizationAnalysis(BaseModel):
    """Visualization appropriateness analysis."""

    total_data_opportunities: int
    visualized_count: int
    visualization_rate: float
    suggestions: List[Dict[str, str]]


class QualityReport(BaseModel):
    """Complete quality validation report."""

    overall_score: float
    content_coverage: ContentCoverage
    structure_analysis: StructureAnalysis
    visualization_analysis: VisualizationAnalysis
    slide_details: List[SlideAnalysis]
    recommendations: List[str]


class PPTQualityValidatorTool(BaseTool):
    """
    Validates generated PPT quality by comparing with source content.

    This tool performs:
    1. Content Coverage: Checks if key information from source is in PPT
    2. Structure Analysis: Validates section-to-slide mapping
    3. Visualization Check: Evaluates if data/processes are properly visualized
    4. Quality Score: Generates overall quality assessment
    """

    name: str = "ppt_quality_validator_tool"
    description: str = (
        "Validates PPT quality by extracting content from generated PPTX and "
        "comparing with source content. Returns coverage rate, missing items, "
        "and improvement recommendations."
    )
    args_schema: Type[BaseModel] = PPTQualityValidatorInput

    def _run(
        self,
        pptx_path: Optional[str] = None,
        source_content: Optional[str] = None,
        source_outline: Optional[Dict] = None,
        file_path: Optional[str] = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """Synchronous run."""
        import asyncio

        return asyncio.run(
            self._arun(pptx_path=pptx_path, source_content=source_content, 
                      source_outline=source_outline, file_path=file_path, **kwargs)
        )

    async def _arun(
        self,
        pptx_path: Optional[str] = None,
        source_content: Optional[str] = None,
        source_outline: Optional[Dict] = None,
        # 호환성을 위한 별칭
        file_path: Optional[str] = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """
        Validate PPT quality asynchronously.

        Args:
            pptx_path: Path to generated PPTX file
            file_path: Alias for pptx_path (for compatibility)
            source_content: Original RAG response or context text (optional for basic validation)
            source_outline: Optional original outline specification

        Returns:
            Quality report with scores and recommendations
        """
        # 파라미터 호환성 처리
        actual_path = pptx_path or file_path
        if not actual_path:
            return {"success": False, "error": "pptx_path 또는 file_path가 필요합니다"}
        
        # source_content가 없으면 기본 검증만 수행
        if not source_content:
            source_content = ""
            logger.info("⚠️ [PPTValidator] source_content 없음 - 기본 검증만 수행")
        
        logger.info(f"🔍 [PPTValidator] 검증 시작: {actual_path}")

        try:
            # 1. Load and analyze PPT
            ppt_path = Path(actual_path)
            if not ppt_path.exists():
                raise FileNotFoundError(f"PPTX 파일을 찾을 수 없습니다: {actual_path}")

            presentation = Presentation(str(ppt_path))
            slide_analyses = self._analyze_slides(presentation)

            # 2. Extract all text from PPT
            ppt_text = self._extract_all_text(presentation)

            # 3. Content coverage analysis
            content_coverage = self._analyze_content_coverage(
                source_content, ppt_text
            )

            # 4. Structure analysis
            structure_analysis = self._analyze_structure(
                source_content, slide_analyses, source_outline
            )

            # 5. Visualization analysis
            visualization_analysis = self._analyze_visualization(
                source_content, slide_analyses
            )

            # 6. Calculate overall score
            overall_score = self._calculate_overall_score(
                content_coverage, structure_analysis, visualization_analysis
            )

            # 7. Generate recommendations
            recommendations = self._generate_recommendations(
                content_coverage,
                structure_analysis,
                visualization_analysis,
                slide_analyses,
            )

            report = QualityReport(
                overall_score=overall_score,
                content_coverage=content_coverage,
                structure_analysis=structure_analysis,
                visualization_analysis=visualization_analysis,
                slide_details=slide_analyses,
                recommendations=recommendations,
            )

            logger.info(
                f"✅ [PPTValidator] 완료: 점수={overall_score:.1f}%, "
                f"내용커버리지={content_coverage.coverage_rate:.1f}%"
            )

            return {
                "success": True,
                "report": report.model_dump(),
                "overall_score": overall_score,
                "summary": self._generate_summary(report),
            }

        except Exception as e:
            logger.error(f"❌ [PPTValidator] 실패: {e}", exc_info=True)
            return {
                "success": False,
                "error": str(e),
            }

    def _analyze_slides(self, presentation: Presentation) -> List[SlideAnalysis]:
        """Analyze each slide in the presentation."""
        analyses = []

        for idx, slide in enumerate(presentation.slides):
            title = ""
            text_parts = []
            bullet_count = 0
            has_chart = False
            has_table = False
            has_image = False
            has_diagram = False
            shape_count = len(slide.shapes)

            for shape in slide.shapes:
                # Check shape types
                if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    has_chart = True
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    has_table = True
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    has_image = True
                elif shape.shape_type in (
                    MSO_SHAPE_TYPE.GROUP,
                    MSO_SHAPE_TYPE.FREEFORM,
                ):
                    has_diagram = True

                # Extract text
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            text_parts.append(text)
                            # Count bullets (paragraphs with level > 0 or list markers)
                            if para.level > 0 or text.startswith(("-", "•", "▪")):
                                bullet_count += 1

                # Get title
                if shape.is_placeholder:
                    ph = shape.placeholder_format
                    if ph.type == 1:  # Title placeholder
                        if hasattr(shape, "text"):
                            title = shape.text.strip()

            text_content = "\n".join(text_parts)
            word_count = len(text_content.split())

            analyses.append(
                SlideAnalysis(
                    slide_index=idx,
                    title=title,
                    text_content=text_content,
                    bullet_count=bullet_count,
                    has_chart=has_chart,
                    has_table=has_table,
                    has_image=has_image,
                    has_diagram=has_diagram,
                    shape_count=shape_count,
                    word_count=word_count,
                )
            )

        return analyses

    def _extract_all_text(self, presentation: Presentation) -> str:
        """Extract all text content from presentation."""
        all_text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text.append(shape.text)
                if hasattr(shape, "table"):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            all_text.append(cell.text)
        return "\n".join(all_text)

    def _analyze_content_coverage(
        self, source_content: str, ppt_text: str
    ) -> ContentCoverage:
        """Analyze how much of source content is covered in PPT."""
        # Extract keywords from source
        source_keywords = self._extract_keywords(source_content)
        ppt_keywords = self._extract_keywords(ppt_text)

        # Normalize for comparison
        source_set = {self._normalize_keyword(k) for k in source_keywords}
        ppt_set = {self._normalize_keyword(k) for k in ppt_keywords}

        matched = source_set & ppt_set
        missing = source_set - ppt_set

        coverage_rate = (len(matched) / len(source_set) * 100) if source_set else 100.0

        return ContentCoverage(
            total_keywords=len(source_set),
            matched_keywords=len(matched),
            coverage_rate=coverage_rate,
            missing_keywords=list(missing)[:20],  # Top 20 missing
            matched_keywords_list=list(matched)[:20],
        )

    def _extract_keywords(self, text: str) -> List[str]:
        """Extract important keywords from text."""
        # Remove markdown formatting
        text = re.sub(r"[#*_`\[\]()]", " ", text)
        text = re.sub(r"\s+", " ", text)

        # Split into words
        words = text.split()

        # Filter: length >= 2, not pure numbers, not common stopwords
        stopwords = {
            "의", "를", "을", "이", "가", "은", "는", "에", "와", "과", "도", "로",
            "으로", "에서", "부터", "까지", "라고", "하는", "하고", "하여", "하면",
            "the", "a", "an", "is", "are", "was", "were", "be", "been", "being",
            "have", "has", "had", "do", "does", "did", "will", "would", "could",
            "should", "may", "might", "must", "shall", "can", "need", "dare",
            "and", "or", "but", "if", "then", "else", "when", "where", "why",
            "how", "what", "which", "who", "whom", "this", "that", "these",
            "those", "am", "is", "are", "was", "were", "been", "being",
            "for", "from", "with", "about", "into", "through", "during",
            "before", "after", "above", "below", "to", "of", "in", "on", "at",
        }

        keywords = []
        for word in words:
            word = word.strip(".,;:!?\"'()[]{}").lower()
            if len(word) >= 2 and not word.isdigit() and word not in stopwords:
                keywords.append(word)

        # Return unique keywords by frequency (top 50)
        counter = Counter(keywords)
        return [kw for kw, _ in counter.most_common(50)]

    def _normalize_keyword(self, keyword: str) -> str:
        """Normalize keyword for comparison."""
        return keyword.lower().strip()

    def _analyze_structure(
        self,
        source_content: str,
        slide_analyses: List[SlideAnalysis],
        source_outline: Optional[Dict],
    ) -> StructureAnalysis:
        """Analyze structure mapping between source and PPT."""
        # Extract sections from source (markdown headers)
        source_sections = self._extract_sections(source_content)
        ppt_slides = len(slide_analyses)

        # Map sections to slides
        section_mapping = []
        for i, section in enumerate(source_sections):
            # Find best matching slide
            best_match = None
            best_score = 0.0

            for slide in slide_analyses:
                # Compare section title with slide title and content
                title_sim = SequenceMatcher(
                    None, section["title"].lower(), slide.title.lower()
                ).ratio()
                content_sim = SequenceMatcher(
                    None,
                    section["content"][:200].lower(),
                    slide.text_content[:200].lower(),
                ).ratio()
                score = title_sim * 0.6 + content_sim * 0.4

                if score > best_score:
                    best_score = score
                    best_match = slide.slide_index

            section_mapping.append(
                {
                    "section_index": i,
                    "section_title": section["title"],
                    "matched_slide": best_match,
                    "match_score": round(best_score, 2),
                }
            )

        # Calculate structure match rate
        matched_sections = sum(
            1 for m in section_mapping if m["match_score"] > 0.3
        )
        match_rate = (
            (matched_sections / len(source_sections) * 100)
            if source_sections
            else 100.0
        )

        return StructureAnalysis(
            source_sections=len(source_sections),
            ppt_slides=ppt_slides,
            section_mapping=section_mapping,
            structure_match_rate=match_rate,
        )

    def _extract_sections(self, text: str) -> List[Dict[str, str]]:
        """Extract sections from markdown text."""
        sections = []
        # Match ## or ### headers
        pattern = r"^(#{2,3})\s+(.+?)$"
        lines = text.split("\n")

        current_section = None
        current_content = []

        for line in lines:
            match = re.match(pattern, line)
            if match:
                # Save previous section
                if current_section:
                    sections.append(
                        {
                            "title": current_section,
                            "content": "\n".join(current_content).strip(),
                        }
                    )
                current_section = match.group(2).strip()
                current_content = []
            elif current_section:
                current_content.append(line)

        # Save last section
        if current_section:
            sections.append(
                {
                    "title": current_section,
                    "content": "\n".join(current_content).strip(),
                }
            )

        return sections

    def _analyze_visualization(
        self, source_content: str, slide_analyses: List[SlideAnalysis]
    ) -> VisualizationAnalysis:
        """Analyze if data/processes are properly visualized."""
        # Detect opportunities for visualization in source
        opportunities = self._detect_visualization_opportunities(source_content)

        # Count actual visualizations in PPT
        visualized_count = sum(
            1
            for s in slide_analyses
            if s.has_chart or s.has_table or s.has_diagram or s.has_image
        )

        # Generate suggestions
        suggestions = []
        for opp in opportunities:
            # Check if corresponding slide has visualization
            has_viz = False
            for slide in slide_analyses:
                if self._content_matches_slide(opp["content"], slide):
                    if slide.has_chart or slide.has_table or slide.has_diagram:
                        has_viz = True
                        break

            if not has_viz:
                suggestions.append(
                    {
                        "type": opp["type"],
                        "content_hint": opp["content"][:100],
                        "recommendation": opp["recommendation"],
                    }
                )

        viz_rate = (
            (visualized_count / len(opportunities) * 100)
            if opportunities
            else 100.0
        )

        return VisualizationAnalysis(
            total_data_opportunities=len(opportunities),
            visualized_count=visualized_count,
            visualization_rate=min(viz_rate, 100.0),
            suggestions=suggestions[:5],  # Top 5 suggestions
        )

    def _detect_visualization_opportunities(
        self, text: str
    ) -> List[Dict[str, str]]:
        """Detect content that should be visualized."""
        opportunities = []

        # Pattern 1: Numerical data (percentages, statistics)
        number_pattern = r"(\d+(?:\.\d+)?%|\d+(?:,\d{3})*(?:\.\d+)?[억만천]?)"
        number_matches = re.findall(number_pattern, text)
        if len(number_matches) >= 3:
            opportunities.append(
                {
                    "type": "chart",
                    "content": f"숫자 데이터 {len(number_matches)}개 발견",
                    "recommendation": "막대 차트 또는 파이 차트로 시각화 권장",
                }
            )

        # Pattern 2: Process/steps (1. 2. 3. or First, Second, Third)
        step_pattern = r"(^|\n)\s*(\d+\.|\(?[①②③④⑤]\)?|첫째|둘째|셋째|Step\s*\d+)"
        step_matches = re.findall(step_pattern, text, re.IGNORECASE)
        if len(step_matches) >= 2:
            opportunities.append(
                {
                    "type": "flowchart",
                    "content": f"단계별 프로세스 {len(step_matches)}개 발견",
                    "recommendation": "플로우차트 또는 타임라인으로 시각화 권장",
                }
            )

        # Pattern 3: Comparisons (vs, 비교, 대비)
        compare_pattern = r"(vs\.?|비교|대비|반면|그러나|한편)"
        if re.search(compare_pattern, text, re.IGNORECASE):
            opportunities.append(
                {
                    "type": "comparison",
                    "content": "비교 내용 발견",
                    "recommendation": "비교 테이블 또는 대비 다이어그램 권장",
                }
            )

        # Pattern 4: Lists (bullet points, numbered items)
        list_pattern = r"(^|\n)\s*[-•▪]\s+"
        list_matches = re.findall(list_pattern, text)
        if len(list_matches) >= 4:
            opportunities.append(
                {
                    "type": "list_visual",
                    "content": f"리스트 항목 {len(list_matches)}개 발견",
                    "recommendation": "아이콘이 있는 리스트 또는 그리드로 시각화 권장",
                }
            )

        # Pattern 5: Hierarchy/Categories
        hierarchy_pattern = r"(분류|종류|유형|카테고리|계층|상위|하위)"
        if re.search(hierarchy_pattern, text):
            opportunities.append(
                {
                    "type": "hierarchy",
                    "content": "계층/분류 구조 발견",
                    "recommendation": "트리 다이어그램 또는 조직도 권장",
                }
            )

        return opportunities

    def _content_matches_slide(
        self, content_hint: str, slide: SlideAnalysis
    ) -> bool:
        """Check if content hint matches slide content."""
        combined = f"{slide.title} {slide.text_content}".lower()
        keywords = content_hint.lower().split()[:5]
        return any(kw in combined for kw in keywords if len(kw) > 2)

    def _calculate_overall_score(
        self,
        content: ContentCoverage,
        structure: StructureAnalysis,
        visualization: VisualizationAnalysis,
    ) -> float:
        """Calculate overall quality score (0-100)."""
        # Weights: Content 50%, Structure 30%, Visualization 20%
        content_score = content.coverage_rate * 0.5
        structure_score = structure.structure_match_rate * 0.3
        viz_score = visualization.visualization_rate * 0.2

        return round(content_score + structure_score + viz_score, 1)

    def _generate_recommendations(
        self,
        content: ContentCoverage,
        structure: StructureAnalysis,
        visualization: VisualizationAnalysis,
        slides: List[SlideAnalysis],
    ) -> List[str]:
        """Generate improvement recommendations."""
        recommendations = []

        # Content recommendations
        if content.coverage_rate < 70:
            recommendations.append(
                f"⚠️ 내용 커버리지가 {content.coverage_rate:.1f}%로 낮습니다. "
                f"누락된 키워드: {', '.join(content.missing_keywords[:5])}"
            )

        # Structure recommendations
        if structure.structure_match_rate < 70:
            recommendations.append(
                f"⚠️ 구조 매칭률이 {structure.structure_match_rate:.1f}%입니다. "
                "원본 섹션 구조를 더 잘 반영해 주세요."
            )

        # Low section mapping scores
        low_matches = [
            m for m in structure.section_mapping if m["match_score"] < 0.3
        ]
        if low_matches:
            sections = ", ".join([m["section_title"][:20] for m in low_matches[:3]])
            recommendations.append(
                f"📌 다음 섹션이 PPT에 잘 반영되지 않았습니다: {sections}"
            )

        # Visualization recommendations
        if visualization.visualization_rate < 50:
            recommendations.append(
                "📊 시각화 요소가 부족합니다. 차트, 다이어그램, 표를 추가하세요."
            )

        for suggestion in visualization.suggestions[:2]:
            recommendations.append(
                f"💡 {suggestion['type']}: {suggestion['recommendation']}"
            )

        # Slide-specific recommendations
        empty_slides = [s for s in slides if s.word_count < 20]
        if empty_slides:
            recommendations.append(
                f"⚠️ {len(empty_slides)}개 슬라이드의 내용이 너무 적습니다."
            )

        text_heavy = [s for s in slides if s.word_count > 150 and s.bullet_count < 3]
        if text_heavy:
            recommendations.append(
                f"📝 {len(text_heavy)}개 슬라이드의 텍스트가 너무 많습니다. "
                "글머리 기호로 분리하세요."
            )

        return recommendations[:10]

    def _generate_summary(self, report: QualityReport) -> str:
        """Generate human-readable summary."""
        score = report.overall_score
        if score >= 80:
            grade = "우수 ✅"
        elif score >= 60:
            grade = "양호 ⚠️"
        else:
            grade = "개선 필요 ❌"

        summary = f"""
## PPT 품질 검증 결과

**종합 점수**: {score:.1f}점 ({grade})

### 세부 분석
- **내용 커버리지**: {report.content_coverage.coverage_rate:.1f}% ({report.content_coverage.matched_keywords}/{report.content_coverage.total_keywords} 키워드)
- **구조 일치도**: {report.structure_analysis.structure_match_rate:.1f}% ({report.structure_analysis.source_sections}개 섹션 → {report.structure_analysis.ppt_slides}개 슬라이드)
- **시각화 적용률**: {report.visualization_analysis.visualization_rate:.1f}% ({report.visualization_analysis.visualized_count}/{report.visualization_analysis.total_data_opportunities} 기회)

### 주요 권장 사항
"""
        for i, rec in enumerate(report.recommendations[:5], 1):
            summary += f"{i}. {rec}\n"

        if report.content_coverage.missing_keywords:
            summary += f"\n### 누락된 주요 키워드\n"
            summary += ", ".join(report.content_coverage.missing_keywords[:10])

        return summary.strip()


# Singleton instance
ppt_quality_validator_tool = PPTQualityValidatorTool()

__all__ = ["PPTQualityValidatorTool", "ppt_quality_validator_tool", "QualityReport"]
