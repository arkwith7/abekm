"""Template PPT Comparator Tool - Compares generated PPT with template for quality validation."""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any, Dict, List, Optional, Type

from loguru import logger
from pptx import Presentation
from pydantic import BaseModel, Field

try:
    from langchain_core.tools import BaseTool
except ImportError:
    from langchain_core.tools import BaseTool


class TemplatePPTComparatorInput(BaseModel):
    """Input schema for TemplatePPTComparatorTool."""

    generated_pptx_path: str = Field(..., description="생성된 PPTX 파일 경로")
    template_pptx_path: str = Field(..., description="템플릿 PPTX 파일 경로")
    template_metadata_path: Optional[str] = Field(
        default=None, description="템플릿 메타데이터 JSON 파일 경로"
    )
    expected_content: Optional[Dict] = Field(
        default=None, description="기대되는 콘텐츠 정보 (outline)"
    )


class SlideComparison(BaseModel):
    """Single slide comparison result."""

    slide_index: int
    layout_match: bool
    shape_count_match: bool
    text_replaced_count: int
    text_unchanged_count: int
    table_issues: List[str]
    image_issues: List[str]
    issues: List[str]
    warnings: List[str]


class TableIssue(BaseModel):
    """Table-specific issue."""

    slide_index: int
    table_name: str
    issue_type: str  # "unchanged", "partial", "missing_data"
    description: str
    template_sample: str
    generated_sample: str


class TextIssue(BaseModel):
    """Text-specific issue."""

    slide_index: int
    shape_name: str
    issue_type: str  # "unchanged", "placeholder"
    template_text: str
    generated_text: str


class ComparisonReport(BaseModel):
    """Complete comparison report."""

    overall_quality_score: float  # 0-100
    total_slides: int
    slides_with_issues: int
    critical_issues: List[str]
    warnings: List[str]
    slide_comparisons: List[SlideComparison]
    table_issues: List[TableIssue]
    text_issues: List[TextIssue]
    recommendations: List[str]
    passed: bool


class TemplatePPTComparatorTool(BaseTool):
    """
    Compares generated PPT with template to validate content replacement.

    This tool performs:
    1. Structure Validation: Ensures slide count, layout, and shape structure match
    2. Content Replacement Check: Verifies all template content was replaced
    3. Table Data Validation: Checks if table data was properly replaced
    4. Text Placeholder Check: Identifies unchanged template text
    5. Quality Scoring: Generates pass/fail assessment with detailed issues

    Returns actionable recommendations for the AI agent to fix issues.
    """

    name: str = "template_ppt_comparator_tool"
    description: str = (
        "Compares generated PPT with template PPT to validate content replacement quality. "
        "Detects unchanged template data (especially tables), placeholder text not replaced, "
        "and structural mismatches. Returns detailed issue report with recommendations."
    )
    args_schema: Type[BaseModel] = TemplatePPTComparatorInput

    def _run(
        self,
        generated_pptx_path: str,
        template_pptx_path: str,
        template_metadata_path: Optional[str] = None,
        expected_content: Optional[Dict] = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """Synchronous run."""
        import asyncio

        return asyncio.run(
            self._arun(
                generated_pptx_path=generated_pptx_path,
                template_pptx_path=template_pptx_path,
                template_metadata_path=template_metadata_path,
                expected_content=expected_content,
                **kwargs,
            )
        )

    async def _arun(
        self,
        generated_pptx_path: str,
        template_pptx_path: str,
        template_metadata_path: Optional[str] = None,
        expected_content: Optional[Dict] = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """
        Compare generated PPT with template asynchronously.

        Args:
            generated_pptx_path: Path to generated PPTX file
            template_pptx_path: Path to template PPTX file
            template_metadata_path: Optional path to template metadata JSON
            expected_content: Optional expected content structure

        Returns:
            Comparison report with pass/fail status and recommendations
        """
        logger.info(
            f"🔍 [TemplatePPTComparator] 비교 시작\n"
            f"  생성 파일: {generated_pptx_path}\n"
            f"  템플릿: {template_pptx_path}"
        )

        try:
            # 1. Load presentations
            gen_path = Path(generated_pptx_path)
            tmp_path = Path(template_pptx_path)

            if not gen_path.exists():
                raise FileNotFoundError(f"생성 파일 없음: {generated_pptx_path}")
            if not tmp_path.exists():
                raise FileNotFoundError(f"템플릿 파일 없음: {template_pptx_path}")

            generated = Presentation(str(gen_path))
            template = Presentation(str(tmp_path))

            # 2. Load metadata if provided
            metadata = None
            if template_metadata_path:
                meta_path = Path(template_metadata_path)
                if meta_path.exists():
                    with open(meta_path, "r", encoding="utf-8") as f:
                        metadata = json.load(f)
                    logger.info(f"✅ 메타데이터 로드 완료")

            # 3. Compare slides
            slide_comparisons = []
            table_issues = []
            text_issues = []

            for idx in range(min(len(generated.slides), len(template.slides))):
                comparison = self._compare_slide(
                    generated.slides[idx],
                    template.slides[idx],
                    idx,
                    metadata,
                )
                slide_comparisons.append(comparison)

                # Collect detailed issues
                slide_table_issues, slide_text_issues = self._extract_detailed_issues(
                    generated.slides[idx], template.slides[idx], idx
                )
                table_issues.extend(slide_table_issues)
                text_issues.extend(slide_text_issues)

            # 4. Identify critical issues
            critical_issues = self._identify_critical_issues(
                generated, template, slide_comparisons, table_issues, text_issues
            )

            # 5. Generate warnings
            warnings = self._generate_warnings(slide_comparisons, text_issues)

            # 6. Calculate quality score
            quality_score = self._calculate_quality_score(
                slide_comparisons, table_issues, text_issues
            )

            # 7. Generate recommendations
            recommendations = self._generate_recommendations(
                critical_issues, warnings, table_issues, text_issues
            )

            # 8. Determine pass/fail
            passed = len(critical_issues) == 0 and quality_score >= 70.0

            report = ComparisonReport(
                overall_quality_score=quality_score,
                total_slides=len(generated.slides),
                slides_with_issues=len(
                    [s for s in slide_comparisons if len(s.issues) > 0]
                ),
                critical_issues=critical_issues,
                warnings=warnings,
                slide_comparisons=slide_comparisons,
                table_issues=table_issues,
                text_issues=text_issues,
                recommendations=recommendations,
                passed=passed,
            )

            logger.info(
                f"✅ [TemplatePPTComparator] 비교 완료\n"
                f"  품질 점수: {quality_score:.1f}/100\n"
                f"  결과: {'✅ PASS' if passed else '❌ FAIL'}\n"
                f"  치명적 문제: {len(critical_issues)}개\n"
                f"  경고: {len(warnings)}개"
            )

            return {
                "success": True,
                "report": report.model_dump(),
                "passed": passed,
                "quality_score": quality_score,
                "critical_issues_count": len(critical_issues),
                "warnings_count": len(warnings),
                "summary": self._generate_summary(report),
            }

        except Exception as e:
            logger.error(f"❌ [TemplatePPTComparator] 오류 발생: {str(e)}")
            return {
                "success": False,
                "error": str(e),
                "error_type": type(e).__name__,
            }

    def _compare_slide(
        self,
        gen_slide,
        tmp_slide,
        idx: int,
        metadata: Optional[Dict],
    ) -> SlideComparison:
        """Compare a single slide."""
        issues = []
        warnings = []
        table_issues = []
        image_issues = []

        # 1. Layout check
        layout_match = gen_slide.slide_layout.name == tmp_slide.slide_layout.name
        if not layout_match:
            issues.append(
                f"레이아웃 불일치: '{tmp_slide.slide_layout.name}' → '{gen_slide.slide_layout.name}'"
            )

        # 2. Shape count check
        shape_count_match = len(gen_slide.shapes) == len(tmp_slide.shapes)
        if not shape_count_match:
            issues.append(
                f"Shape 개수 불일치: {len(tmp_slide.shapes)} → {len(gen_slide.shapes)}"
            )

        # 3. Text replacement check
        text_replaced = 0
        text_unchanged = 0

        for gen_shape, tmp_shape in zip(gen_slide.shapes, tmp_slide.shapes):
            if hasattr(gen_shape, "text_frame") and gen_shape.text_frame:
                gen_text = gen_shape.text_frame.text.strip()
                tmp_text = tmp_shape.text_frame.text.strip()

                if gen_text and tmp_text:
                    if gen_text != tmp_text:
                        text_replaced += 1
                    else:
                        text_unchanged += 1
                        # 긴 텍스트가 동일하면 경고
                        if len(gen_text) > 15 and not self._is_template_metadata(
                            gen_text
                        ):
                            warnings.append(
                                f"[{gen_shape.name}] 템플릿 텍스트 미교체: '{gen_text[:50]}...'"
                            )

            # 4. Table check
            if hasattr(gen_shape, "table") and hasattr(tmp_shape, "table"):
                table_identical = self._compare_tables(
                    gen_shape.table, tmp_shape.table
                )
                if table_identical:
                    table_issues.append(
                        f"[{gen_shape.name}] 테이블 데이터가 템플릿과 동일"
                    )
                    issues.append(
                        f"테이블 데이터 미교체: {gen_shape.name}"
                    )

        return SlideComparison(
            slide_index=idx + 1,
            layout_match=layout_match,
            shape_count_match=shape_count_match,
            text_replaced_count=text_replaced,
            text_unchanged_count=text_unchanged,
            table_issues=table_issues,
            image_issues=image_issues,
            issues=issues,
            warnings=warnings,
        )

    def _compare_tables(self, gen_table, tmp_table) -> bool:
        """Check if two tables are identical."""
        if len(gen_table.rows) != len(tmp_table.rows):
            return False
        if len(gen_table.columns) != len(tmp_table.columns):
            return False

        # Compare first 3 rows (enough to detect template data)
        for r_idx in range(min(3, len(gen_table.rows))):
            for c_idx in range(len(gen_table.columns)):
                gen_cell = gen_table.cell(r_idx, c_idx).text.strip()
                tmp_cell = tmp_table.cell(r_idx, c_idx).text.strip()

                if gen_cell != tmp_cell:
                    return False

        return True

    def _is_template_metadata(self, text: str) -> bool:
        """Check if text is template metadata that should remain unchanged."""
        metadata_patterns = [
            "Company or Team Name",
            "Logo",
            "THANK YOU",
            "FOR WATCHING",
        ]
        return any(pattern.lower() in text.lower() for pattern in metadata_patterns)

    def _extract_detailed_issues(
        self, gen_slide, tmp_slide, idx: int
    ) -> tuple[List[TableIssue], List[TextIssue]]:
        """Extract detailed issues for tables and text."""
        table_issues = []
        text_issues = []

        for gen_shape, tmp_shape in zip(gen_slide.shapes, tmp_slide.shapes):
            # Table issues
            if hasattr(gen_shape, "table") and hasattr(tmp_shape, "table"):
                gen_table = gen_shape.table
                tmp_table = tmp_shape.table

                if self._compare_tables(gen_table, tmp_table):
                    # Extract sample data
                    tmp_sample = self._get_table_sample(tmp_table)
                    gen_sample = self._get_table_sample(gen_table)

                    table_issues.append(
                        TableIssue(
                            slide_index=idx + 1,
                            table_name=gen_shape.name,
                            issue_type="unchanged",
                            description="테이블 데이터가 템플릿과 완전히 동일합니다",
                            template_sample=tmp_sample,
                            generated_sample=gen_sample,
                        )
                    )

            # Text issues
            if hasattr(gen_shape, "text_frame") and gen_shape.text_frame:
                gen_text = gen_shape.text_frame.text.strip()
                tmp_text = tmp_shape.text_frame.text.strip()

                if (
                    gen_text
                    and tmp_text
                    and gen_text == tmp_text
                    and len(gen_text) > 15
                    and not self._is_template_metadata(gen_text)
                ):
                    text_issues.append(
                        TextIssue(
                            slide_index=idx + 1,
                            shape_name=gen_shape.name,
                            issue_type="unchanged",
                            template_text=tmp_text[:100],
                            generated_text=gen_text[:100],
                        )
                    )

        return table_issues, text_issues

    def _get_table_sample(self, table) -> str:
        """Get sample data from table (first 2 rows)."""
        samples = []
        for r_idx in range(min(2, len(table.rows))):
            row_data = []
            for c_idx in range(len(table.columns)):
                cell_text = table.cell(r_idx, c_idx).text.strip()
                row_data.append(cell_text[:20] if cell_text else "(empty)")
            samples.append(" | ".join(row_data))
        return "\n".join(samples)

    def _identify_critical_issues(
        self,
        generated,
        template,
        slide_comparisons: List[SlideComparison],
        table_issues: List[TableIssue],
        text_issues: List[TextIssue],
    ) -> List[str]:
        """Identify critical issues that require fixing."""
        critical = []

        # 1. Slide count mismatch
        if len(generated.slides) != len(template.slides):
            critical.append(
                f"슬라이드 수 불일치: 템플릿({len(template.slides)}) vs 생성({len(generated.slides)})"
            )

        # 2. Table data not replaced
        if table_issues:
            critical.append(
                f"테이블 데이터 미교체: {len(table_issues)}개 테이블에서 템플릿 데이터 그대로 유지됨"
            )

        # 3. Layout mismatches
        layout_mismatches = [s for s in slide_comparisons if not s.layout_match]
        if layout_mismatches:
            critical.append(
                f"레이아웃 불일치: {len(layout_mismatches)}개 슬라이드"
            )

        # 4. Shape count mismatches
        shape_mismatches = [s for s in slide_comparisons if not s.shape_count_match]
        if shape_mismatches:
            critical.append(
                f"Shape 구조 변경: {len(shape_mismatches)}개 슬라이드에서 Shape 개수 불일치"
            )

        return critical

    def _generate_warnings(
        self,
        slide_comparisons: List[SlideComparison],
        text_issues: List[TextIssue],
    ) -> List[str]:
        """Generate warning messages."""
        warnings = []

        # 1. Unchanged text
        if text_issues:
            warnings.append(
                f"텍스트 미교체: {len(text_issues)}개 텍스트가 템플릿과 동일"
            )

        # 2. Low replacement rate
        total_replaced = sum(s.text_replaced_count for s in slide_comparisons)
        total_unchanged = sum(s.text_unchanged_count for s in slide_comparisons)

        if total_replaced + total_unchanged > 0:
            replacement_rate = (
                total_replaced / (total_replaced + total_unchanged)
            ) * 100
            if replacement_rate < 50:
                warnings.append(
                    f"낮은 텍스트 교체율: {replacement_rate:.1f}% (기대: 70% 이상)"
                )

        return warnings

    def _calculate_quality_score(
        self,
        slide_comparisons: List[SlideComparison],
        table_issues: List[TableIssue],
        text_issues: List[TextIssue],
    ) -> float:
        """Calculate overall quality score (0-100)."""
        score = 100.0

        # Deduct points for issues
        for slide in slide_comparisons:
            if not slide.layout_match:
                score -= 5.0
            if not slide.shape_count_match:
                score -= 5.0
            score -= len(slide.issues) * 3.0

        # Table issues are critical
        score -= len(table_issues) * 15.0

        # Text issues are warnings
        score -= len(text_issues) * 2.0

        return max(0.0, score)

    def _generate_recommendations(
        self,
        critical_issues: List[str],
        warnings: List[str],
        table_issues: List[TableIssue],
        text_issues: List[TextIssue],
    ) -> List[str]:
        """Generate actionable recommendations for the AI agent."""
        recommendations = []

        if critical_issues:
            recommendations.append(
                "🔴 치명적 문제 발견 - 즉시 수정 필요:"
            )
            for issue in critical_issues:
                recommendations.append(f"  • {issue}")

        if table_issues:
            recommendations.append(
                "\n📊 테이블 데이터 수정 방법:"
            )
            recommendations.append(
                "  1. 프론트엔드에서 tableData를 metadata에 포함하여 전송"
            )
            recommendations.append(
                "  2. 백엔드 enhanced_object_processor.py의 테이블 처리 로직 확인"
            )
            recommendations.append(
                "  3. 데이터가 없을 경우 테이블 셀을 빈 문자열로 초기화"
            )

            for issue in table_issues[:3]:  # Show first 3
                recommendations.append(
                    f"  • 슬라이드 {issue.slide_index} [{issue.table_name}]: {issue.description}"
                )

        if text_issues:
            recommendations.append(
                "\n📝 텍스트 교체 필요:"
            )
            for issue in text_issues[:5]:  # Show first 5
                recommendations.append(
                    f"  • 슬라이드 {issue.slide_index} [{issue.shape_name}]: '{issue.template_text[:50]}...'"
                )

        if not critical_issues and not warnings:
            recommendations.append(
                "✅ 품질 검증 통과 - 문제 없음"
            )

        return recommendations

    def _generate_summary(self, report: ComparisonReport) -> str:
        """Generate human-readable summary."""
        if report.passed:
            return (
                f"✅ 품질 검증 통과 (점수: {report.overall_quality_score:.1f}/100)\n"
                f"모든 콘텐츠가 정상적으로 교체되었습니다."
            )
        else:
            return (
                f"❌ 품질 검증 실패 (점수: {report.overall_quality_score:.1f}/100)\n"
                f"치명적 문제: {len(report.critical_issues)}개\n"
                f"경고: {len(report.warnings)}개\n"
                f"수정이 필요합니다."
            )


# Singleton instance
template_ppt_comparator_tool = TemplatePPTComparatorTool()
