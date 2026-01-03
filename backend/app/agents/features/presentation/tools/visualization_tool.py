"""
Visualization Tool for PPT Generation

슬라이드 내용 기반으로 시각화 요소(차트, 표, 프로세스 다이어그램) 추가

Author: Presentation System
Created: 2025-01-20
"""

import re
import logging
from typing import Any, ClassVar, Dict, List, Optional, Type

from langchain_core.tools import BaseTool
from pydantic import BaseModel, Field

from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

logger = logging.getLogger(__name__)


class VisualizationToolInput(BaseModel):
    """Input schema for VisualizationTool"""
    slide_title: str = Field(..., description="슬라이드 제목")
    key_message: str = Field(default="", description="키 메시지 (선택)")
    bullets: List[str] = Field(default_factory=list, description="불릿 포인트 리스트")
    slide_index: int = Field(default=0, description="슬라이드 인덱스 (로깅용)")


class VisualizationTool(BaseTool):
    """
    슬라이드 내용을 분석하여 시각화 요소를 감지하고 힌트를 반환하는 도구
    
    기능:
    - 차트 필요성 감지 (증가/감소/비율 키워드 + 숫자 데이터)
    - 표 필요성 감지 (항목/구분/사양 키워드 + 콜론 패턴)
    - 프로세스 다이어그램 필요성 감지 (단계/과정/절차 키워드)
    - 시각화 타입 및 배치 정보 반환
    
    반환값:
    {
        "success": True,
        "chart": bool,
        "table": bool,
        "process": bool,
        "chart_type": "pie" | "line" | "column" | None,
        "table_from_bullets": bool,
        "process_from_bullets": bool,
        "numeric_score": int,
        "recommendations": List[str]
    }
    """
    
    name: str = "visualization_tool"
    description: str = (
        "슬라이드 내용을 분석하여 시각화 요소(차트, 표, 프로세스 다이어그램)의 "
        "필요성을 감지하고 시각화 힌트를 반환합니다. "
        "입력: slide_title, key_message, bullets, slide_index"
    )
    args_schema: Type[BaseModel] = VisualizationToolInput
    return_direct: bool = False
    
    # 색상 팔레트 (ClassVar로 정의하여 Pydantic 필드 검증 우회)
    colors: ClassVar[Dict[str, RGBColor]] = {
        "primary": RGBColor(0, 70, 150),       # 진한 파란색
        "secondary": RGBColor(240, 247, 255),  # 연한 파란색
        "accent": RGBColor(255, 153, 0),       # 주황색
        "text": RGBColor(30, 30, 30),          # 거의 검은색
    }

    def _run(
        self,
        slide_title: str,
        key_message: str = "",
        bullets: Optional[List[str]] = None,
        slide_index: int = 0,
    ) -> Dict[str, Any]:
        """시각화 힌트 감지 (동기 버전)"""
        try:
            bullets = bullets or []
            hints = self._detect_visualization_hints(slide_title, key_message, bullets)
            
            # 권장사항 생성
            recommendations = self._generate_recommendations(hints)
            hints["recommendations"] = recommendations
            hints["success"] = True
            
            logger.info(f"🎨 슬라이드 {slide_index} 시각화 힌트: {hints}")
            return hints
            
        except Exception as e:
            logger.error(f"시각화 힌트 감지 실패: {e}", exc_info=True)
            return {
                "success": False,
                "error": str(e),
                "chart": False,
                "table": False,
                "process": False,
                "recommendations": [],
            }

    async def _arun(
        self,
        slide_title: str,
        key_message: str = "",
        bullets: Optional[List[str]] = None,
        slide_index: int = 0,
    ) -> Dict[str, Any]:
        """시각화 힌트 감지 (비동기 버전)"""
        return self._run(slide_title, key_message, bullets, slide_index)

    def _detect_visualization_hints(
        self,
        slide_title: str,
        key_message: str,
        bullets: List[str]
    ) -> Dict[str, Any]:
        """슬라이드 내용에서 시각화 힌트를 감지합니다."""
        hints = {
            "chart": False,
            "table": False,
            "diagram": False,
            "process": False,
            "comparison": False,
            "chart_type": None,
            "chart_data": None,
            "numeric_score": 0,
            "table_from_bullets": False,
            "process_from_bullets": False,
        }
        
        # 모든 텍스트 결합
        all_text = f"{slide_title} {key_message} {' '.join(bullets)}"
        all_text_lower = all_text.lower()
        
        # 차트 관련 키워드
        chart_keywords = [
            "증가", "감소", "성장", "하락", "비율", "퍼센트", "%", "추이", "변화",
            "비교", "대비", "점유율", "시장", "매출", "수익", "통계", "데이터"
        ]
        
        # 표 관련 키워드  
        table_keywords = [
            "항목", "구분", "분류", "목록", "리스트", "사양", "스펙", "기능",
            "가격", "요금", "비용", "계획", "일정", "단계"
        ]
        
        # 프로세스/다이어그램 키워드
        process_keywords = [
            "단계", "과정", "절차", "순서", "흐름", "프로세스", "워크플로",
            "다음", "이후", "진행", "구조", "조직도", "관계"
        ]
        
        # 숫자 신호 집계 (차트 신뢰도 향상)
        num_count = 0
        percent_count = 0
        for bullet in bullets:
            s = (bullet or "").strip()
            num_count += len(re.findall(r"\b\d+(?:[\.,]\d+)?\b", s))
            percent_count += s.count("%")
        hints["numeric_score"] = num_count + percent_count

        # 차트 감지: 키워드 + 숫자 신호가 있는 경우에만 활성화
        if any(keyword in all_text_lower for keyword in chart_keywords) and hints["numeric_score"] >= 2:
            hints["chart"] = True
            # 간단한 차트 타입 결정
            if any(word in all_text_lower for word in ["증가", "감소", "성장", "추이"]):
                hints["chart_type"] = "line"
            elif any(word in all_text_lower for word in ["비율", "점유율", "퍼센트", "%"]):
                hints["chart_type"] = "pie"
            else:
                hints["chart_type"] = "column"
        
        # 표 감지
        if any(keyword in all_text_lower for keyword in table_keywords):
            hints["table"] = True
            hints["table_from_bullets"] = True
        else:
            # 불릿 포인트에 ":" 또는 " - " 패턴이 2개 이상 존재하면 표로 간주
            colon_style_count = 0
            for bullet in bullets:
                s = (bullet or "").strip()
                if not s:
                    continue
                if ":" in s or " - " in s or "|" in s:
                    colon_style_count += 1
            if colon_style_count >= 2:
                hints["table"] = True
                hints["table_from_bullets"] = True
        
        # 프로세스 다이어그램 감지
        if any(keyword in all_text_lower for keyword in process_keywords) or any(
            (bullet or "").strip().startswith(("1.", "2.", "-", "•", "*")) for bullet in bullets
        ):
            hints["process"] = True
            hints["process_from_bullets"] = True
            
        # 비교 구조 감지
        if any(word in all_text_lower for word in ["vs", "대비", "비교", "차이"]):
            hints["comparison"] = True
            
        logger.debug(f"🎨 시각화 힌트 감지 결과: {hints}")
        return hints

    def _generate_recommendations(self, hints: Dict[str, Any]) -> List[str]:
        """감지된 힌트를 바탕으로 시각화 권장사항 생성"""
        recommendations = []
        
        if hints["chart"]:
            chart_type = hints.get("chart_type", "column")
            recommendations.append(f"📊 {chart_type} 차트 추가 권장 (숫자 데이터 {hints['numeric_score']}개)")
        
        if hints["table"]:
            source = "불릿 포인트" if hints.get("table_from_bullets") else "내용"
            recommendations.append(f"📋 표 추가 권장 ({source}에서 생성)")
        
        if hints["process"]:
            recommendations.append("🔄 프로세스 다이어그램 추가 권장")
        
        if hints["comparison"]:
            recommendations.append("⚖️ 비교 차트 또는 표 추가 권장")
        
        if not recommendations:
            recommendations.append("💬 텍스트 중심 슬라이드 (시각화 없음)")
        
        return recommendations

    @staticmethod
    def create_sample_chart(
        slide,
        chart_type: str,
        title: str,
        x: Optional[float] = None,
        y: Optional[float] = None,
        cx: Optional[float] = None,
        cy: Optional[float] = None
    ):
        """
        샘플 차트를 생성합니다.
        
        Args:
            slide: pptx slide 객체
            chart_type: "pie", "line", "column" 중 하나
            title: 차트 제목
            x, y, cx, cy: 차트 위치 및 크기 (Inches 단위)
        
        Returns:
            chart 객체 또는 None (실패시)
        """
        try:
            # 차트 데이터 준비
            chart_data = CategoryChartData()
            
            if chart_type == "pie":
                # 파이 차트
                chart_data.categories = ['제품 A', '제품 B', '제품 C', '기타']
                chart_data.add_series('시장 점유율', (40, 30, 20, 10))
            elif chart_type == "line":
                # 선 차트 (추이)
                chart_data.categories = ['1Q', '2Q', '3Q', '4Q']
                chart_data.add_series('매출 증가율', (15, 25, 35, 45))
            else:
                # 기본: 컬럼 차트
                chart_data.categories = ['현재', '목표', '예상']
                chart_data.add_series('성과 지표', (75, 100, 95))
            
            # 차트 타입 결정
            if chart_type == "pie":
                chart_type_enum = XL_CHART_TYPE.PIE
            elif chart_type == "line":
                chart_type_enum = XL_CHART_TYPE.LINE
            else:
                chart_type_enum = XL_CHART_TYPE.COLUMN_CLUSTERED
            
            # 차트 위치 및 크기 (기본값)
            if x is None or y is None or cx is None or cy is None:
                x, y, cx, cy = Inches(1), Inches(3), Inches(8), Inches(4)
            
            chart = slide.shapes.add_chart(chart_type_enum, x, y, cx, cy, chart_data).chart
            
            # 차트 제목
            if chart.has_title:
                chart.chart_title.text_frame.text = title or "데이터 시각화"
                
            logger.info(f"✅ {chart_type} 차트 생성 완료")
            return chart
            
        except Exception as e:
            logger.warning(f"⚠️ 차트 생성 실패: {e}")
            return None

    @staticmethod
    def create_simple_table(
        slide,
        title: str,
        bullets: List[str],
        x: Optional[float] = None,
        y: Optional[float] = None,
        cx: Optional[float] = None,
        cy: Optional[float] = None,
        colors: Optional[Dict[str, RGBColor]] = None
    ):
        """
        간단한 표를 생성합니다.
        
        Args:
            slide: pptx slide 객체
            title: 표 제목 (헤더에 사용)
            bullets: 불릿 포인트 리스트 (표 데이터로 변환)
            x, y, cx, cy: 표 위치 및 크기 (Inches 단위)
            colors: 색상 팔레트 (선택)
        
        Returns:
            table 객체 또는 None (실패시)
        """
        try:
            # 불릿 포인트를 표 형태로 변환
            rows = min(len(bullets) + 1, 6)  # 최대 5개 데이터 행 + 헤더
            cols = 2
            
            # 표 위치 및 크기 (기본값)
            if x is None or y is None or cx is None or cy is None:
                x, y, cx, cy = Inches(1), Inches(2.5), Inches(8), Inches(4)
            
            # 색상 팔레트 (기본값)
            if colors is None:
                colors = {"primary": RGBColor(0, 70, 150)}
            
            # 표 생성
            table = slide.shapes.add_table(rows, cols, x, y, cx, cy).table
            
            # 헤더 설정
            table.cell(0, 0).text = "항목"
            table.cell(0, 1).text = "내용"
            
            # 헤더 스타일링
            for col in range(cols):
                cell = table.cell(0, col)
                cell.fill.solid()
                cell.fill.fore_color.rgb = colors["primary"]
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].font.bold = True
            
            # 데이터 행 채우기
            for i, bullet in enumerate(bullets[:rows-1]):
                if i + 1 < rows:
                    # 간단한 파싱: "항목: 내용" 형태 감지
                    if ":" in bullet:
                        parts = bullet.split(":", 1)
                        table.cell(i + 1, 0).text = parts[0].strip()
                        table.cell(i + 1, 1).text = parts[1].strip()
                    else:
                        table.cell(i + 1, 0).text = f"항목 {i + 1}"
                        table.cell(i + 1, 1).text = bullet

            # 본문 셀 기본 스타일 (가독성)
            for r in range(1, rows):
                for c in range(cols):
                    try:
                        p = table.cell(r, c).text_frame.paragraphs[0]
                        p.font.size = Pt(12)
                        p.font.name = '맑은 고딕'
                        p.font.color.rgb = RGBColor(30, 30, 30)
                    except Exception:
                        pass
            
            logger.info(f"✅ 표 생성 완료: {rows}x{cols}")
            return table
            
        except Exception as e:
            logger.warning(f"⚠️ 표 생성 실패: {e}")
            return None

    @staticmethod
    def create_process_diagram(
        slide,
        title: str,
        bullets: List[str],
        y: Optional[float] = None,
        colors: Optional[Dict[str, RGBColor]] = None
    ):
        """
        프로세스 다이어그램을 생성합니다.
        
        Args:
            slide: pptx slide 객체
            title: 다이어그램 제목 (로깅용)
            bullets: 단계 설명 리스트
            y: 다이어그램 세로 위치 (Inches 단위)
            colors: 색상 팔레트 (선택)
        
        Returns:
            True (성공) 또는 False (실패)
        """
        try:
            # 색상 팔레트 (기본값)
            if colors is None:
                colors = {
                    "primary": RGBColor(0, 70, 150),
                    "secondary": RGBColor(240, 247, 255),
                }
            
            # 단계별 박스 생성
            step_count = min(len(bullets), 5)  # 최대 5단계
            box_width = Inches(1.5)
            box_height = Inches(0.8)
            spacing = Inches(0.3)
            
            # 중앙 정렬을 위한 시작 위치 계산
            total_width = (box_width * step_count) + (spacing * (step_count - 1))
            start_x = (Inches(10) - total_width) / 2
            box_y = y if y is not None else Inches(3.5)
            
            for i in range(step_count):
                x = start_x + (box_width + spacing) * i
                
                # 단계 박스 생성
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, x, box_y, box_width, box_height
                )
                
                # 박스 스타일링
                shape.fill.solid()
                shape.fill.fore_color.rgb = colors["secondary"]
                shape.line.color.rgb = colors["primary"]
                
                # 텍스트 추가
                text_frame = shape.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                p.text = f"단계 {i + 1}"
                p.font.bold = True
                p.font.size = Pt(12)
                p.alignment = PP_ALIGN.CENTER
                
                # 단계 내용 추가
                if i < len(bullets):
                    p2 = text_frame.add_paragraph()
                    p2.text = bullets[i][:20] + "..." if len(bullets[i]) > 20 else bullets[i]
                    p2.font.size = Pt(10)
                    p2.alignment = PP_ALIGN.CENTER
                
                # 화살표 추가 (마지막 단계 제외)
                if i < step_count - 1:
                    arrow_x = x + box_width + spacing/4
                    arrow_y = box_y + box_height/2 - Inches(0.1)
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW, arrow_x, arrow_y, spacing/2, Inches(0.2)
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = colors["primary"]
            
            logger.info(f"✅ 프로세스 다이어그램 생성 완료: {step_count}단계")
            return True
            
        except Exception as e:
            logger.warning(f"⚠️ 프로세스 다이어그램 생성 실패: {e}")
            return False


# 전역 인스턴스
visualization_tool = VisualizationTool()
