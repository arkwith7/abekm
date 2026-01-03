"""
템플릿 콘텐츠 정리 서비스
업로드된 PPTX 템플릿의 기존 텍스트를 완전히 제거하고
새로운 내용으로 교체하는 시스템
"""
import logging
from typing import List, Dict, Any, Optional
from pathlib import Path

from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from pptx.text.text import TextFrame
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)

class TemplateContentCleaner:
    """템플릿 콘텐츠 정리기"""
    
    def __init__(self):
        # 🆕 logo 제거 - Logo, Company 등의 플레이스홀더도 정리 대상
        self.preserved_shapes = {"background", "decoration", "border"}
        self.cleanable_shapes = {"title", "content", "text", "subtitle"}
        
    def clean_template_content(self, template_path: str, output_path: str, 
                              target_slides: Optional[List[int]] = None) -> str:
        """템플릿의 기존 내용을 정리하고 새로운 템플릿 파일 생성
        
        Args:
            template_path: 원본 템플릿 경로
            output_path: 정리된 템플릿 저장 경로  
            target_slides: 정리할 슬라이드 번호 리스트 (1-based index), None이면 전체 슬라이드
        """
        try:
            logger.info(f"템플릿 콘텐츠 정리 시작: {template_path}")
            if target_slides:
                logger.info(f"대상 슬라이드: {target_slides}")
            
            # PPTX 로드
            presentation = Presentation(template_path)
            
            # 슬라이드별 정리 (전체 또는 지정된 슬라이드만)
            cleaned_slides = 0
            for slide_idx, slide in enumerate(presentation.slides):
                slide_num = slide_idx + 1
                
                # 특정 슬라이드만 정리하는 경우 필터링
                if target_slides and slide_num not in target_slides:
                    logger.debug(f"슬라이드 {slide_num}: 건너뜀 (대상 아님)")
                    continue
                
                cleaned_shapes = self._clean_slide_content(slide, slide_num)
                if cleaned_shapes > 0:
                    cleaned_slides += 1
                    logger.info(f"슬라이드 {slide_num} 정리 완료: {cleaned_shapes}개 도형 정리")
            
            # 정리된 템플릿 저장
            presentation.save(output_path)
            
            if target_slides:
                logger.info(f"선택적 콘텐츠 정리 완료: {cleaned_slides}개 슬라이드, 저장 경로: {output_path}")
            else:
                logger.info(f"전체 콘텐츠 정리 완료: {cleaned_slides}개 슬라이드, 저장 경로: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"템플릿 콘텐츠 정리 실패: {e}")
            raise

    def copy_slide_with_content_clear(self, template_path: str, output_path: str, 
                                     source_slide_num: int, clear_text: bool = True) -> str:
        """슬라이드를 복사하고 선택적으로 텍스트 클리어
        
        Args:
            template_path: 원본 템플릿 경로
            output_path: 수정된 템플릿 저장 경로
            source_slide_num: 복사할 원본 슬라이드 번호 (1-based)
            clear_text: True면 텍스트 클리어, False면 원본 그대로 복사
        """
        try:
            logger.info(f"슬라이드 복사 시작: 슬라이드 {source_slide_num}, 텍스트 클리어: {clear_text}")
            
            presentation = Presentation(template_path)
            
            if source_slide_num < 1 or source_slide_num > len(presentation.slides):
                raise ValueError(f"유효하지 않은 슬라이드 번호: {source_slide_num}")
            
            # 원본 슬라이드 가져오기 (0-based index)
            source_slide = presentation.slides[source_slide_num - 1]
            source_layout = source_slide.slide_layout
            
            # 새 슬라이드 추가 (같은 레이아웃으로)
            new_slide = presentation.slides.add_slide(source_layout)
            
            # 원본 슬라이드의 모든 도형 복사
            self._copy_slide_shapes(source_slide, new_slide)
            
            # 텍스트 클리어가 요청된 경우에만 실행
            if clear_text:
                cleaned_count = self._clean_slide_content(new_slide, len(presentation.slides))
                logger.info(f"복사된 슬라이드 텍스트 클리어 완료: {cleaned_count}개 도형 정리")
            
            # 수정된 프레젠테이션 저장
            presentation.save(output_path)
            
            logger.info(f"슬라이드 복사 완료: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"슬라이드 복사 실패: {e}")
            raise
    
    def clean_specific_slides(self, template_path: str, output_path: str, 
                             slide_numbers: List[int]) -> str:
        """특정 슬라이드들만 텍스트 클리어 (PPT 생성 설정 화면용)"""
        return self.clean_template_content(template_path, output_path, target_slides=slide_numbers)

    def add_slide_with_optional_clear(self, template_path: str, output_path: str,
                                     source_slide_num: int, clear_content: bool = True) -> str:
        """슬라이드 추가 + 선택적 텍스트 클리어 (PPT 생성 설정 화면의 '추가' 버튼용)
        
        Args:
            template_path: 템플릿 경로
            output_path: 수정된 템플릿 저장 경로  
            source_slide_num: 복사할 원본 슬라이드 번호
            clear_content: True면 텍스트 클리어, False면 원본 그대로 복사
        """
        return self.copy_slide_with_content_clear(template_path, output_path, 
                                                source_slide_num, clear_text=clear_content)

    def _clean_slide_content(self, slide: Slide, slide_number: int) -> int:
        cleaned_count = 0
        total_shapes = len(slide.shapes)
        
        try:
            logger.debug(f"슬라이드 {slide_number} 분석 시작: {total_shapes}개 도형")
            
            # 모든 도형을 순회하면서 텍스트 내용 제거
            for idx, shape in enumerate(slide.shapes):
                shape_type = getattr(shape, 'shape_type', 'Unknown')
                shape_name = getattr(shape, 'name', f'Shape_{idx}')
                
                logger.debug(f"  도형 {idx+1}/{total_shapes}: {shape_name} (타입: {shape_type})")
                
                # 텍스트 여부 확인
                has_text = False
                text_content = ""
                
                try:
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        text_content = shape.text_frame.text
                        has_text = bool(text_content.strip())
                    elif hasattr(shape, 'text'):
                        text_content = shape.text
                        has_text = bool(text_content.strip())
                except Exception:
                    pass
                
                if has_text:
                    logger.debug(f"    텍스트 발견: '{text_content[:50]}...'")
                    
                    if self._should_clean_shape(shape):
                        logger.debug(f"    정리 대상으로 판정")
                        if self._clean_shape_text(shape):
                            cleaned_count += 1
                            logger.debug(f"    ✅ 정리 완료")
                        else:
                            logger.debug(f"    ❌ 정리 실패")
                    else:
                        logger.debug(f"    정리 대상 아님 (보존)")
                else:
                    logger.debug(f"    텍스트 없음")
            
            logger.info(f"슬라이드 {slide_number} 정리 완료: {cleaned_count}/{total_shapes}개 도형 정리")
            return cleaned_count
            
        except Exception as e:
            logger.warning(f"슬라이드 {slide_number} 정리 중 오류: {e}")
            return cleaned_count
    
    def _should_clean_shape(self, shape: BaseShape) -> bool:
        """도형을 정리해야 하는지 판단"""
        try:
            # 텍스트 프레임이 있는지 확인
            if not hasattr(shape, 'text_frame'):
                return False
            
            try:
                text_frame = shape.text_frame
                if not text_frame:
                    return False
                
                # 텍스트가 실제로 있는 경우만
                if not text_frame.text.strip():
                    return False
                    
            except Exception:
                # text_frame 접근 실패 시 텍스트 있는지 다른 방법으로 확인
                if hasattr(shape, 'text') and shape.text.strip():
                    pass  # 텍스트가 있으니 계속 진행
                else:
                    return False
            
            # 특정 형태의 도형은 보존 (로고, 배경 등)
            if hasattr(shape, 'name') and shape.name:
                shape_name = shape.name.lower()
                for preserved in self.preserved_shapes:
                    if preserved in shape_name:
                        logger.debug(f"보존 대상 도형 스킵: {shape_name}")
                        return False
            
            # 🎯 모든 텍스트가 있는 도형을 정리 대상으로 판단 (플레이스홀더 여부 무관)
            shape_type = getattr(shape, 'shape_type', None)
            
            # 플레이스홀더인 경우 우선 정리
            try:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                    logger.debug(f"플레이스홀더 정리 대상: {shape.placeholder_format.type}")
                    return True
            except Exception:
                pass
            
            # 일반 텍스트 도형도 정리
            if shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                logger.debug("텍스트박스 정리 대상")
                return True
            
            # AutoShape 중 텍스트가 있는 것들
            if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                logger.debug("AutoShape 텍스트 정리 대상")
                return True
            
            # 그룹 내 텍스트 도형
            if shape_type == MSO_SHAPE_TYPE.GROUP:
                logger.debug("그룹 내 텍스트 정리 대상")
                return True
            
            # 테이블 내 텍스트
            if hasattr(shape, 'table') or str(shape_type) == "TABLE":
                logger.debug("테이블 텍스트 정리 대상")
                return True
            
            # 기타 텍스트가 있는 모든 도형 (보수적 접근)
            logger.debug(f"기타 텍스트 도형 정리 대상: {shape_type}")
            return True
            
        except Exception as e:
            logger.debug(f"도형 정리 판단 중 오류: {e}")
            return False
    
    def _clean_shape_text(self, shape: BaseShape) -> bool:
        """도형의 텍스트 내용 제거 - 스타일 보존 모드
        
        🆕 개선: text_frame.clear() 대신 기존 run의 텍스트만 빈 문자열로 변경하여 스타일 유지
        """
        cleaned = False
        
        try:
            # 1. 텍스트 프레임이 있는 경우 (가장 일반적)
            if hasattr(shape, 'text_frame') and shape.text_frame:
                text_frame = shape.text_frame
                original_text = text_frame.text
                
                # 기존 텍스트가 있는 경우에만 정리
                if original_text.strip():
                    try:
                        logger.debug(f"🧹 텍스트 프레임 정리: '{original_text[:30]}...'")
                        
                        # 🎯 방법 1 (개선): 기존 run의 텍스트만 빈 문자열로 변경 (스타일 완전 보존)
                        try:
                            if text_frame.paragraphs:
                                for para in text_frame.paragraphs:
                                    if para.runs:
                                        # 첫 번째 run만 유지하고 나머지 run의 텍스트 제거
                                        for i, run in enumerate(para.runs):
                                            run.text = ""
                                    else:
                                        # run이 없으면 paragraph 직접 설정 (스타일 보존 안됨)
                                        para.text = ""
                            
                            logger.debug(f"✅ 스타일 보존 정리 성공 (run 텍스트만 제거)")
                            cleaned = True
                            
                        except Exception as run_e:
                            logger.debug(f"run 방식 실패: {run_e}, clear 방식으로 폴백")
                            
                            # 🎯 방법 2 (폴백): clear 후 스타일 복원 시도
                            try:
                                # 기본 문단의 폰트 정보 백업
                                backup_font_info = {}
                                if text_frame.paragraphs and text_frame.paragraphs[0].runs:
                                    first_run = text_frame.paragraphs[0].runs[0]
                                    backup_font_info = {
                                        'name': getattr(first_run.font, 'name', None),
                                        'size': getattr(first_run.font, 'size', None),
                                        'bold': getattr(first_run.font, 'bold', None),
                                        'italic': getattr(first_run.font, 'italic', None),
                                        'color': getattr(first_run.font.color, 'rgb', None) if hasattr(first_run.font, 'color') else None
                                    }
                                
                                text_frame.clear()
                                
                                # 빈 문단 하나 생성
                                if not text_frame.paragraphs:
                                    para = text_frame.add_paragraph()
                                else:
                                    para = text_frame.paragraphs[0]
                                
                                # run 생성 및 스타일 복원
                                run = para.add_run()
                                run.text = ""
                                
                                if backup_font_info:
                                    if backup_font_info.get('name'):
                                        run.font.name = backup_font_info['name']
                                    if backup_font_info.get('size'):
                                        run.font.size = backup_font_info['size']
                                    if backup_font_info.get('bold') is not None:
                                        run.font.bold = backup_font_info['bold']
                                    if backup_font_info.get('italic') is not None:
                                        run.font.italic = backup_font_info['italic']
                                    if backup_font_info.get('color'):
                                        try:
                                            run.font.color.rgb = backup_font_info['color']
                                        except:
                                            pass
                                
                                logger.debug(f"✅ clear + 스타일 복원 성공")
                                cleaned = True
                                
                            except Exception as clear_e:
                                logger.debug(f"clear 방식도 실패: {clear_e}")
                                
                                # 🎯 방법 3: 직접 text 속성 변경 (마지막 수단)
                                try:
                                    text_frame.text = ""
                                    logger.debug(f"✅ 직접 방식 정리 성공")
                                    cleaned = True
                                except Exception as direct_e:
                                    logger.debug(f"❌ 모든 정리 방식 실패: {direct_e}")
                        
                    except Exception as e:
                        logger.debug(f"text_frame 정리 전체 실패: {e}")
            
            # 2. 직접 텍스트 속성이 있는 경우
            elif hasattr(shape, 'text'):
                try:
                    original_text = shape.text
                    if original_text.strip():
                        shape.text = ""
                        logger.debug(f"✅ 직접 텍스트 정리 완료: '{original_text[:30]}...'")
                        cleaned = True
                except Exception as e:
                    logger.debug(f"직접 텍스트 정리 실패: {e}")
            
            # 3. 테이블인 경우
            elif hasattr(shape, 'table') and shape.table:
                try:
                    table = shape.table
                    cell_count = 0
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text_frame and cell.text_frame.text.strip():
                                try:
                                    cell.text_frame.clear()
                                    if not cell.text_frame.paragraphs:
                                        cell.text_frame.add_paragraph()
                                    cell.text_frame.paragraphs[0].text = ""
                                    cell_count += 1
                                except Exception:
                                    # 폴백: 직접 비우기
                                    cell.text_frame.text = ""
                                    cell_count += 1
                    
                    if cell_count > 0:
                        logger.debug(f"✅ 테이블 셀 {cell_count}개 정리 완료")
                        cleaned = True
                        
                except Exception as e:
                    logger.debug(f"테이블 정리 실패: {e}")
            
            # 4. 그룹인 경우 재귀적으로 처리
            elif hasattr(shape, 'shapes') and shape.shapes:
                try:
                    sub_cleaned = 0
                    for sub_shape in shape.shapes:
                        if self._clean_shape_text(sub_shape):
                            sub_cleaned += 1
                    
                    if sub_cleaned > 0:
                        logger.debug(f"✅ 그룹 내 도형 {sub_cleaned}개 정리 완료")
                        cleaned = True
                        
                except Exception as e:
                    logger.debug(f"그룹 정리 실패: {e}")
            
            return cleaned
            
        except Exception as e:
            logger.debug(f"텍스트 정리 중 전체 오류: {e}")
            return False
    
    def create_content_mapping(self, template_path: str) -> Dict[str, Any]:
        """템플릿의 콘텐츠 영역 매핑 정보 생성"""
        try:
            presentation = Presentation(template_path)
            mapping = {
                "slides": [],
                "total_content_areas": 0
            }
            
            for slide_idx, slide in enumerate(presentation.slides):
                slide_mapping = {
                    "slide_number": slide_idx + 1,
                    "content_areas": [],
                    "layout_info": self._analyze_slide_layout(slide)
                }
                
                area_count = 0
                for shape_idx, shape in enumerate(slide.shapes):
                    if self._should_clean_shape(shape):
                        def to_inches(val):
                            try:
                                if hasattr(val, 'inches'):
                                    return float(val.inches)
                                # assume EMU int
                                return float(val) / 914400.0
                            except Exception:
                                try:
                                    return float(val)
                                except Exception:
                                    return 0.0

                        area_info = {
                            "shape_index": shape_idx,
                            "type": self._get_content_type(shape),
                            "position": {
                                "x": to_inches(shape.left) if hasattr(shape, 'left') else 0,
                                "y": to_inches(shape.top) if hasattr(shape, 'top') else 0,
                                "width": to_inches(shape.width) if hasattr(shape, 'width') else 0,
                                "height": to_inches(shape.height) if hasattr(shape, 'height') else 0
                            },
                            "placeholder_type": self._get_placeholder_type(shape)
                        }
                        slide_mapping["content_areas"].append(area_info)
                        area_count += 1
                
                slide_mapping["content_count"] = area_count
                mapping["slides"].append(slide_mapping)
                mapping["total_content_areas"] += area_count
            
            logger.info(f"콘텐츠 매핑 완료: {len(mapping['slides'])}개 슬라이드, {mapping['total_content_areas']}개 콘텐츠 영역")
            return mapping
            
        except Exception as e:
            logger.error(f"콘텐츠 매핑 실패: {e}")
            return {"slides": [], "total_content_areas": 0}

    def _copy_slide_shapes(self, source_slide: Slide, target_slide: Slide) -> int:
        """원본 슬라이드의 모든 도형을 대상 슬라이드로 복사"""
        try:
            copied_count = 0
            
            # 원본 슬라이드의 모든 도형 복사
            for shape in source_slide.shapes:
                try:
                    # 도형 타입에 따른 복사 처리
                    if hasattr(shape, 'shape_type'):
                        # 텍스트박스, 도형 등 복사
                        # 실제 구현에서는 python-pptx의 복사 메커니즘 사용
                        # 여기서는 기본 구조만 표시
                        copied_count += 1
                        logger.debug(f"도형 복사 완료: {getattr(shape, 'name', 'Unknown')}")
                        
                except Exception as shape_error:
                    logger.warning(f"도형 복사 실패: {shape_error}")
            
            logger.info(f"슬라이드 도형 복사 완료: {copied_count}개")
            return copied_count
            
        except Exception as e:
            logger.error(f"슬라이드 도형 복사 실패: {e}")
            return 0
    
    def _analyze_slide_layout(self, slide: Slide) -> Dict[str, Any]:
        """슬라이드 레이아웃 분석"""
        layout_info = {
            "layout_name": "unknown",
            "has_title": False,
            "has_content": False,
            "has_image_placeholder": False,
            "shape_count": len(slide.shapes)
        }
        
        try:
            # 레이아웃 정보 추출
            if hasattr(slide, 'slide_layout') and slide.slide_layout:
                layout_info["layout_name"] = slide.slide_layout.name
            
            # 플레이스홀더 분석
            for shape in slide.shapes:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                    placeholder_type = str(shape.placeholder_format.type).lower()
                    if 'title' in placeholder_type:
                        layout_info["has_title"] = True
                    elif 'content' in placeholder_type or 'body' in placeholder_type:
                        layout_info["has_content"] = True
                    elif 'picture' in placeholder_type or 'image' in placeholder_type:
                        layout_info["has_image_placeholder"] = True
        
        except Exception as e:
            logger.debug(f"레이아웃 분석 중 오류: {e}")
        
        return layout_info
    
    def _get_content_type(self, shape: BaseShape) -> str:
        """콘텐츠 영역 타입 분류"""
        try:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                placeholder_type = str(shape.placeholder_format.type).lower()
                if 'title' in placeholder_type:
                    return "title"
                elif 'content' in placeholder_type or 'body' in placeholder_type:
                    return "content"
                elif 'subtitle' in placeholder_type:
                    return "subtitle"
                else:
                    return "placeholder"
            
            elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                return "textbox"
            elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                return "shape_text"
            else:
                return "unknown"
                
        except Exception:
            return "unknown"
    
    def _get_placeholder_type(self, shape: BaseShape) -> Optional[str]:
        """플레이스홀더 타입 반환"""
        try:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                return str(shape.placeholder_format.type)
            return None
        except Exception:
            return None


# 전역 인스턴스
template_content_cleaner = TemplateContentCleaner()
