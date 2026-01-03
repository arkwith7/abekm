"""
PPT 슬라이드 썸네일 생성 서비스
사용자 업로드 템플릿의 각 페이지를 썸네일 이미지로 변환
LibreOffice를 사용해 실제 슬라이드를 이미지로 렌더링
"""
import os
import logging
import subprocess
import tempfile
import shutil
from typing import Optional, List, Dict, Any
from pathlib import Path
import hashlib
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
import io
import base64

logger = logging.getLogger(__name__)

class ThumbnailGenerator:
    """슬라이드 썸네일 생성기"""
    
    def __init__(self):
        # 썸네일 캐시 디렉토리 - uploads/templates 디렉토리 사용
        backend_root = Path(__file__).parent.parent.parent.parent
        upload_dir = Path(os.getenv('UPLOAD_DIR', str(backend_root / "uploads")))
        self.thumbnail_cache_dir = upload_dir / "templates" / "thumbnails"
        self.thumbnail_cache_dir.mkdir(parents=True, exist_ok=True)
        
        # 썸네일 설정
        self.thumbnail_width = 640
        self.thumbnail_height = 480
        self.cache_duration = 86400 * 7  # 7일
        
        # LibreOffice 경로 확인
        self.libreoffice_path = self._find_libreoffice()
    
    def _find_libreoffice(self) -> Optional[str]:
        """LibreOffice 실행 파일 찾기"""
        for path in ['/usr/bin/libreoffice', '/usr/bin/soffice', 
                     '/usr/local/bin/libreoffice', '/usr/local/bin/soffice']:
            if os.path.exists(path):
                return path
        
        # which 명령으로 찾기
        try:
            result = subprocess.run(['which', 'libreoffice'], capture_output=True, text=True)
            if result.returncode == 0:
                return result.stdout.strip()
        except Exception:
            pass
        
        return None
    
    def generate_template_thumbnails(self, pptx_file_path: str, template_id: str) -> List[Dict[str, Any]]:
        """템플릿의 모든 슬라이드 썸네일 생성"""
        try:
            logger.info(f"썸네일 생성 시작: {pptx_file_path}")
            
            # 캐시 키 생성 (파일 내용 기반)
            cache_key = self._generate_cache_key(pptx_file_path, template_id)
            
            # 캐시 확인
            cached_thumbnails = self._get_cached_thumbnails(cache_key)
            if cached_thumbnails:
                logger.info(f"캐시된 썸네일 사용: {template_id}")
                return cached_thumbnails
            
            # 새로 생성
            thumbnails = self._create_thumbnails(pptx_file_path, template_id, cache_key)
            
            # 캐시 저장
            self._save_thumbnail_cache(cache_key, thumbnails)
            
            logger.info(f"썸네일 생성 완료: {len(thumbnails)}개 슬라이드")
            return thumbnails
            
        except Exception as e:
            logger.error(f"썸네일 생성 실패: {e}")
            return self._create_fallback_thumbnails(pptx_file_path, template_id)
    
    def _create_thumbnails(self, pptx_file_path: str, template_id: str, cache_key: str) -> List[Dict[str, Any]]:
        """LibreOffice를 사용해 실제 슬라이드를 이미지로 변환"""
        thumbnails = []
        
        try:
            # LibreOffice로 슬라이드를 PNG로 내보내기
            png_files = self._convert_pptx_to_images(pptx_file_path, cache_key)
            
            if png_files:
                # 슬라이드 정보 추출을 위해 PPT 로드
                prs = Presentation(pptx_file_path)
                
                for slide_idx, png_path in enumerate(png_files):
                    try:
                        # 슬라이드 정보 추출
                        if slide_idx < len(prs.slides):
                            slide_info = self._extract_slide_info(prs.slides[slide_idx], slide_idx)
                        else:
                            slide_info = {"title": f"슬라이드 {slide_idx + 1}", "layout_name": "", 
                                         "description": "", "shape_count": 0, "text_shapes": 0,
                                         "has_images": False, "has_charts": False}
                        
                        # 이미지를 썸네일 크기로 리사이즈
                        thumbnail_data = self._resize_and_encode(png_path, cache_key, slide_idx)
                        
                        thumbnail = {
                            "slide_index": slide_idx,
                            "slide_title": slide_info["title"],
                            "layout_name": slide_info["layout_name"],
                            "description": slide_info["description"],
                            "thumbnail_data": thumbnail_data,
                            "thumbnail_url": f"/api/v1/agent/presentation/templates/{template_id}/thumbnails/{slide_idx}",
                            "shape_count": slide_info.get("shape_count", 0),
                            "text_shapes": slide_info.get("text_shapes", 0),
                            "has_images": slide_info.get("has_images", False),
                            "has_charts": slide_info.get("has_charts", False)
                        }
                        
                        thumbnails.append(thumbnail)
                        
                    except Exception as e:
                        logger.warning(f"슬라이드 {slide_idx} 처리 실패: {e}")
                        thumbnails.append(self._create_fallback_slide_thumbnail(slide_idx, template_id))
                
                return thumbnails
            
        except Exception as e:
            logger.error(f"LibreOffice 변환 실패: {e}")
        
        # LibreOffice 실패 시 기본 썸네일 생성
        return self._create_fallback_thumbnails(pptx_file_path, template_id)
    
    def _convert_pptx_to_images(self, pptx_file_path: str, cache_key: str) -> List[Path]:
        """LibreOffice를 사용해 PPTX를 PNG 이미지로 변환 (PDF 경유)"""
        if not self.libreoffice_path:
            logger.warning("LibreOffice를 찾을 수 없습니다")
            return []
        
        try:
            # 임시 디렉토리 생성
            temp_dir = tempfile.mkdtemp(prefix="ppt_thumb_")
            
            try:
                # PDF로 먼저 변환 (모든 슬라이드 포함)
                return self._convert_via_pdf(pptx_file_path, temp_dir, cache_key)
                
            finally:
                # 임시 디렉토리 정리
                shutil.rmtree(temp_dir, ignore_errors=True)
                
        except subprocess.TimeoutExpired:
            logger.error("LibreOffice 변환 시간 초과")
        except Exception as e:
            logger.error(f"이미지 변환 실패: {e}")
        
        return []
    
    def _convert_via_pdf(self, pptx_file_path: str, temp_dir: str, cache_key: str) -> List[Path]:
        """PDF 중간 변환을 통한 이미지 생성"""
        try:
            # PDF로 먼저 변환
            cmd = [
                self.libreoffice_path,
                '--headless',
                '--invisible', 
                '--convert-to', 'pdf',
                '--outdir', temp_dir,
                pptx_file_path
            ]
            
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=120,
                env={**os.environ, 'HOME': temp_dir}
            )
            
            if result.returncode != 0:
                logger.warning(f"PDF 변환 실패: {result.stderr}")
                return []
            
            # PDF 파일 찾기
            pdf_file = Path(temp_dir) / f"{Path(pptx_file_path).stem}.pdf"
            if not pdf_file.exists():
                return []
            
            # PDF를 이미지로 변환 (pdf2image 또는 pdftoppm 사용)
            try:
                from pdf2image import convert_from_path
                images = convert_from_path(str(pdf_file), dpi=150)
                
                result_files = []
                for idx, img in enumerate(images):
                    dest = self.thumbnail_cache_dir / f"{cache_key}_{idx}.png"
                    img.save(str(dest), 'PNG')
                    result_files.append(dest)
                
                return result_files
                
            except ImportError:
                logger.warning("pdf2image 모듈 없음, pdftoppm 시도")
                # pdftoppm 사용
                return self._convert_pdf_with_pdftoppm(pdf_file, temp_dir, cache_key)
                
        except Exception as e:
            logger.error(f"PDF 변환 실패: {e}")
            return []
    
    def _convert_pdf_with_pdftoppm(self, pdf_file: Path, temp_dir: str, cache_key: str) -> List[Path]:
        """pdftoppm을 사용해 PDF를 이미지로 변환"""
        try:
            output_prefix = Path(temp_dir) / "slide"
            cmd = ['pdftoppm', '-png', '-r', '150', str(pdf_file), str(output_prefix)]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            
            if result.returncode != 0:
                logger.warning(f"pdftoppm 실패: {result.stderr}")
                return []
            
            # 생성된 파일 찾기
            png_files = sorted(Path(temp_dir).glob("slide*.png"))
            
            result_files = []
            for idx, png_file in enumerate(png_files):
                dest = self.thumbnail_cache_dir / f"{cache_key}_{idx}.png"
                shutil.copy(png_file, dest)
                result_files.append(dest)
            
            return result_files
            
        except Exception as e:
            logger.error(f"pdftoppm 변환 실패: {e}")
            return []
    
    def _resize_and_encode(self, image_path: Path, cache_key: str, slide_idx: int) -> str:
        """이미지를 썸네일 크기로 리사이즈하고 Base64 인코딩"""
        try:
            img = Image.open(image_path)
            
            # 비율 유지하면서 리사이즈
            img.thumbnail((self.thumbnail_width, self.thumbnail_height), Image.Resampling.LANCZOS)
            
            # 최종 썸네일 저장
            thumb_path = self.thumbnail_cache_dir / f"{cache_key}_{slide_idx}.png"
            img.save(thumb_path, 'PNG')
            
            # Base64 인코딩
            buffer = io.BytesIO()
            img.save(buffer, format='PNG')
            buffer.seek(0)
            
            return base64.b64encode(buffer.getvalue()).decode('utf-8')
            
        except Exception as e:
            logger.error(f"이미지 리사이즈 실패: {e}")
            return self._create_fallback_image_data(slide_idx)
    
    def _extract_slide_info(self, slide, slide_idx: int) -> Dict[str, Any]:
        """슬라이드 정보 추출"""
        info = {
            "title": f"슬라이드 {slide_idx + 1}",
            "layout_name": "사용자 정의",
            "description": "",
            "shape_count": len(slide.shapes),
            "text_shapes": 0,
            "has_images": False,
            "has_charts": False
        }
        
        # 제목 추출 시도
        title_texts = []
        for shape in slide.shapes:
            try:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text = shape.text_frame.text.strip()
                    if text and len(text) < 100:  # 제목으로 보이는 짧은 텍스트
                        title_texts.append(text)
                        info["text_shapes"] += 1
                
                # 이미지 확인
                if shape.shape_type == 13:  # PICTURE
                    info["has_images"] = True
                
                # 차트 확인
                if shape.shape_type == 3:  # CHART
                    info["has_charts"] = True
                    
            except Exception:
                continue
        
        # 가장 적절한 제목 선택
        if title_texts:
            # 가장 짧고 의미있는 텍스트를 제목으로 선택
            info["title"] = min(title_texts, key=len)
            if len(info["title"]) > 30:
                info["title"] = info["title"][:27] + "..."
        
        # 설명 생성
        features = []
        if info["text_shapes"] > 0:
            features.append(f"텍스트 {info['text_shapes']}개")
        if info["has_images"]:
            features.append("이미지")
        if info["has_charts"]:
            features.append("차트")
        
        info["description"] = ", ".join(features) if features else "기본 레이아웃"
        
        return info
    
    def _create_fallback_image_data(self, slide_idx: int) -> str:
        """기본 썸네일 이미지 데이터 생성"""
        img = Image.new('RGB', (self.thumbnail_width, self.thumbnail_height), color='#F5F5F5')
        draw = ImageDraw.Draw(img)
        
        # 테두리
        draw.rectangle([0, 0, self.thumbnail_width-1, self.thumbnail_height-1], outline='#DDDDDD', width=2)
        
        # 텍스트
        text = f"슬라이드 {slide_idx + 1}"
        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 16)
        except:
            font = ImageFont.load_default()
        
        bbox = draw.textbbox((0, 0), text, font=font)
        text_width = bbox[2] - bbox[0]
        text_height = bbox[3] - bbox[1]
        
        x = (self.thumbnail_width - text_width) // 2
        y = (self.thumbnail_height - text_height) // 2
        
        draw.text((x, y), text, fill='#666666', font=font)
        
        buffer = io.BytesIO()
        img.save(buffer, format='PNG')
        buffer.seek(0)
        
        return base64.b64encode(buffer.getvalue()).decode('utf-8')
    
    def _create_fallback_thumbnails(self, pptx_file_path: str, template_id: str) -> List[Dict[str, Any]]:
        """기본 썸네일 목록 생성"""
        try:
            prs = Presentation(pptx_file_path)
            thumbnails = []
            
            for slide_idx in range(len(prs.slides)):
                thumbnail = self._create_fallback_slide_thumbnail(slide_idx, template_id)
                thumbnails.append(thumbnail)
            
            return thumbnails
            
        except Exception as e:
            logger.error(f"기본 썸네일 생성 실패: {e}")
            return []
    
    def _create_fallback_slide_thumbnail(self, slide_idx: int, template_id: str) -> Dict[str, Any]:
        """기본 슬라이드 썸네일 생성"""
        return {
            "slide_index": slide_idx,
            "slide_title": f"슬라이드 {slide_idx + 1}",
            "layout_name": "기본 레이아웃",
            "description": "썸네일 생성 중...",
            "thumbnail_data": self._create_fallback_image_data(slide_idx),
            "thumbnail_url": f"/api/v1/agent/presentation/templates/{template_id}/thumbnails/{slide_idx}",
            "shape_count": 0,
            "text_shapes": 0,
            "has_images": False,
            "has_charts": False
        }
    
    def _generate_cache_key(self, pptx_file_path: str, template_id: str) -> str:
        """캐시 키 생성"""
        try:
            with open(pptx_file_path, 'rb') as f:
                file_hash = hashlib.md5(f.read()).hexdigest()[:8]
            return f"{template_id}_{file_hash}"
        except Exception:
            return f"{template_id}_{hash(pptx_file_path) % 100000000}"
    
    def _get_cached_thumbnails(self, cache_key: str) -> Optional[List[Dict[str, Any]]]:
        """캐시된 썸네일 조회"""
        # 구현 시 JSON 파일이나 데이터베이스에서 읽기
        return None
    
    def _save_thumbnail_cache(self, cache_key: str, thumbnails: List[Dict[str, Any]]):
        """썸네일 캐시 저장"""
        # 구현 시 JSON 파일이나 데이터베이스에 저장
        pass
    
    def get_slide_thumbnail(self, template_id: str, slide_index: int, user_id: Optional[str] = None) -> Optional[bytes]:
        """특정 슬라이드 썸네일 이미지 반환"""
        try:
            # 먼저 정확한 template_id로 캐시에서 찾기
            cache_files = list(self.thumbnail_cache_dir.glob(f"{template_id}_*_{slide_index}.png"))
            if cache_files:
                with open(cache_files[0], 'rb') as f:
                    return f.read()
            
            # 캐시에 없으면 템플릿 경로 찾기
            template_path = None
            
            # 1. user_template_manager에서 찾기 (사용자 ID가 있는 경우)
            if user_id:
                try:
                    from app.agents.features.presentation.services.user_template_manager import user_template_manager
                    template_path = user_template_manager.get_template_path(user_id, template_id)
                except Exception:
                    pass
            
            # 2. 기존 ppt_template_manager에서 찾기 (fallback)
            if not template_path:
                try:
                    from app.agents.features.presentation.services.ppt_template_manager import template_manager
                    template_path = template_manager.get_template_path(template_id)
                except Exception:
                    pass
            
            # 3. users 디렉토리 전체에서 찾기 (user_id 모를 때)
            if not template_path:
                users_dir = Path(__file__).parents[3] / 'uploads' / 'templates' / 'users'
                if users_dir.exists():
                    for user_dir in users_dir.iterdir():
                        if user_dir.is_dir():
                            for pptx_file in user_dir.glob('*.pptx'):
                                if pptx_file.stem.lower().replace(' ', '_') == template_id:
                                    template_path = str(pptx_file)
                                    break
                        if template_path:
                            break
            
            if template_path and Path(template_path).exists():
                logger.info(f"템플릿에서 썸네일 생성: {template_id} from {template_path}")
                thumbnails = self.generate_template_thumbnails(template_path, template_id)
                
                # 생성 후 다시 캐시에서 찾기
                cache_files = list(self.thumbnail_cache_dir.glob(f"{template_id}_*_{slide_index}.png"))
                if cache_files:
                    with open(cache_files[0], 'rb') as f:
                        return f.read()
                
                # 메모리에서 직접 반환
                if slide_index < len(thumbnails) and 'thumbnail_data' in thumbnails[slide_index]:
                    import base64
                    return base64.b64decode(thumbnails[slide_index]['thumbnail_data'])
            
            return None
            
        except Exception as e:
            logger.error(f"썸네일 이미지 조회 실패: {e}")
            return None

# 전역 인스턴스
thumbnail_generator = ThumbnailGenerator()
