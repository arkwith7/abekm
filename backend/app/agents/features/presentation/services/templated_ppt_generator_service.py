"""Templated PPT Generator Service - 템플릿 기반 생성 전용"""
from __future__ import annotations

import json
import asyncio
import re
import os
import tempfile
from pathlib import Path
from typing import Any, Dict, List, Optional
from datetime import datetime

from loguru import logger
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

from app.core.config import settings
from app.services.core.ai_service import ai_service
from .ppt_models import ChartData, DiagramData, SlideSpec, DeckSpec
from .ppt_template_manager import PPTTemplateManager, template_manager
from .enhanced_object_processor import EnhancedPPTObjectProcessor
# Note: TemplateContentCleaner 사용 안 함 (스타일 손실 방지)


class TemplatedPPTGeneratorService:
    """템플릿 기반 PPT 생성 전용 서비스 - AI 생성, 템플릿 적용, 고급 기능"""
    
    def __init__(self):
        self.prompts_dir = Path(__file__).parents[3] / "prompts"
        self.upload_dir = settings.resolved_upload_dir
        self.upload_dir.mkdir(parents=True, exist_ok=True)
        self.template_manager = template_manager
        self.object_processor = EnhancedPPTObjectProcessor()
        # Note: content_cleaner 사용 안 함 (스타일 손실 방지)
        
        # 풍부한 색상 테마
        self.color_themes = {
            "corporate_blue": {"primary": RGBColor(0, 102, 204), "secondary": RGBColor(102, 153, 255), "accent": RGBColor(255, 153, 0), "text": RGBColor(51, 51, 51), "background": RGBColor(248, 249, 250)},
            "modern_green": {"primary": RGBColor(34, 139, 34), "secondary": RGBColor(144, 238, 144), "accent": RGBColor(255, 215, 0), "text": RGBColor(47, 79, 79), "background": RGBColor(248, 255, 248)},
            "professional_gray": {"primary": RGBColor(70, 70, 70), "secondary": RGBColor(169, 169, 169), "accent": RGBColor(220, 20, 60), "text": RGBColor(0, 0, 0), "background": RGBColor(245, 245, 245)},
            "playful_violet": {"primary": RGBColor(111, 45, 168), "secondary": RGBColor(181, 126, 220), "accent": RGBColor(255, 181, 71), "text": RGBColor(60, 60, 60), "background": RGBColor(250, 248, 255)},
        }

    # ---------------- Filename / Topic Normalization Helpers ----------------
    def _normalize_topic_for_filename(self, topic: str, max_chars: int = 50) -> str:
        """과도하게 긴 topic(여러 줄, 키메시지 포함 등)을 파일명용으로 정제.

        규칙:
        0. 🆕 요청 표현 제거 (PPT 작성해 주세요, 만들어줘 등)
        1. 줄 단위로 분리 후 첫 줄 우선. (첫 줄이 5자 미만이면 다음 줄 탐색)
        2. '키메시지', '키 메시지', '제품 개요' 이후 내용 잘라냄.
        3. 중복 연속 단어 제거.
        4. 허용 문자만 남기고 공백은 '_'로 치환.
        5. 길이 제한 (기본 50자) - 멀티바이트 안전하게 자르기.
        """
        if not topic:
            return "presentation"

        original = topic
        
        # 🆕 Step 0: 요청 표현 제거 (명사형으로 축약)
        topic = self._remove_request_expressions(topic)
        
        # 줄 분리 + 첫 적절한 라인
        lines = [ln.strip() for ln in re.split(r"[\r\n]+", topic) if ln.strip()]
        if lines:
            # 첫 줄이 지나치게 짧고 두 번째가 더 의미 있으면 교체
            if len(lines[0]) < 5 and len(lines) > 1:
                topic = lines[1]
            else:
                topic = lines[0]

        # 키메시지 및 추가 설명 트리거어 제거
        topic = re.split(r"키 ?메시지|Key Message|상세 설명|제품 개요", topic)[0].strip()

        # 연속 중복 단어 제거 (ex: '제품 소개 제품 개요' -> 앞 1~2개만)
        words = topic.split()
        dedup_words = []
        for w in words:
            if not dedup_words or dedup_words[-1] != w:
                dedup_words.append(w)
        topic = " ".join(dedup_words)

        # 길이 제한 (문자 기준)
        if len(topic) > max_chars:
            topic = topic[:max_chars].rstrip()
        # 최소 길이 보장
        if len(topic) < 2:
            topic = original[:max_chars] if original else "presentation"

        # 파일명용 정규화
        safe = re.sub(r"[^\w\s-]", "", topic)
        safe = re.sub(r"[\s-]+", "_", safe).strip("_")
        if not safe:
            safe = "presentation"

        logger.info(f"🧪 토픽 파일명 정규화: original='{original[:60]}', normalized='{safe}'")
        return safe
    
    def _remove_request_expressions(self, text: str) -> str:
        """요청 표현을 제거하고 명사형 제목으로 정제.
        
        예시:
        - '자동차 산업의 특허분석 방법론에 대해 PPT 작성해 주세요' → '자동차 산업의 특허분석 방법론'
        - 'AI 기술 트렌드 발표 자료 만들어줘' → 'AI 기술 트렌드'
        - '2024 마케팅 전략' → '2024 마케팅 전략' (이미 명사형)
        """
        if not text:
            return text
        
        original = text
        
        # 1. 후위 요청 표현 패턴 (끝에서부터 제거)
        suffix_patterns = [
            r'\s*(에 대해|에 대한|에 관한|에 관해|을 위한|를 위한)\s*(PPT|ppt|프레젠테이션|발표\s*자료|슬라이드).*$',
            r'\s*(PPT|ppt|프레젠테이션|발표\s*자료|슬라이드)\s*(작성|생성|만들|제작).*$',
            r'\s*(작성|생성|만들어|제작)\s*(해|좀)?\s*(주세요|줘|줘요|주십시오|부탁).*$',
            r'\s*(해|좀)?\s*(주세요|줘|줘요|주십시오|부탁).*$',
            r'\s+PPT\s*$',
            r'\s+ppt\s*$',
        ]
        
        for pattern in suffix_patterns:
            text = re.sub(pattern, '', text, flags=re.IGNORECASE).strip()
        
        # 2. 전위 요청 표현 패턴 (앞에서부터 제거)
        prefix_patterns = [
            r'^(다음|아래|위)\s*(내용|주제)(에 대해|으로|로)?\s*',
        ]
        
        for pattern in prefix_patterns:
            text = re.sub(pattern, '', text, flags=re.IGNORECASE).strip()
        
        # 3. 중간에 있는 불필요한 표현 제거
        mid_patterns = [
            r'\s+에 대해\s+',
            r'\s+에 대한\s+',
            r'\s+관련\s+',
        ]
        
        for pattern in mid_patterns:
            # 문맥상 필요할 수 있으므로 공백으로 대체
            text = re.sub(pattern, ' ', text).strip()
        
        # 4. 조사 정리 (끝에 '의', '에', '를' 등이 남으면 제거)
        text = re.sub(r'[의에를을가이]$', '', text).strip()
        
        # 결과가 너무 짧으면 원본 반환
        if len(text) < 3:
            text = original
        
        if text != original:
            logger.info(f"📝 요청 표현 제거: '{original[:50]}' → '{text[:50]}'")
        
        return text

    async def generate_pptx_from_data(
        self,
        template_id: str,
        slides_data: List[Dict[str, Any]],
        output_filename: str = "generated_presentation",
        user_id: Optional[str] = None
    ) -> str:
        """
        사용자가 편집한 데이터(slides_data)를 기반으로 PPT를 생성합니다.
        (Template-First Approach)
        """
        # 1. 템플릿 로드
        template_path = self.template_manager.get_template_path(template_id)
        
        # 시스템 템플릿에서 못 찾은 경우, 사용자 템플릿 검색
        if not template_path or not os.path.exists(template_path):
            try:
                from app.agents.features.presentation.services.user_template_manager import user_template_manager
                
                # 1. user_id가 있으면 해당 사용자의 템플릿 확인
                if user_id:
                    template_path = user_template_manager.get_template_path(user_id, template_id)
                
                # 2. 없으면 전체 사용자 템플릿에서 검색 (소유자 찾기)
                if not template_path:
                    owner_id = user_template_manager.find_template_owner(template_id)
                    if owner_id:
                        template_path = user_template_manager.get_template_path(owner_id, template_id)
            except Exception as e:
                logger.warning(f"User template lookup failed: {e}")

        if not template_path or not os.path.exists(template_path):
            raise ValueError(f"Template file not found: {template_id}")
            
        prs = Presentation(template_path)
        
        # 2. 데이터 적용
        for slide_data in slides_data:
            slide_index = slide_data.get("index", 0)
            # 1-based index to 0-based
            if slide_index < 1 or slide_index > len(prs.slides):
                continue
                
            slide = prs.slides[slide_index - 1]
            elements = slide_data.get("elements", [])
            
            for element in elements:
                el_id = element.get("id")
                text = element.get("text")
                
                if text is None: # Skip if text is None (keep original or empty)
                    continue

                # Find shape by name (id)
                for shape in slide.shapes:
                    if shape.name == el_id:
                        if hasattr(shape, "text_frame"):
                            # 텍스트 교체 (서식 유지를 위해 run 단위 교체 시도 가능하나, 일단 전체 교체)
                            shape.text_frame.text = text
                        break
                        
        # 3. 저장
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_filename = self._normalize_topic_for_filename(output_filename)
        final_filename = f"{safe_filename}_{timestamp}.pptx"
        output_path = self.upload_dir / final_filename
        prs.save(output_path)
        
        logger.info(f"✅ PPT Generated from data: {output_path}")
        return str(output_path)

    def _load_prompt(self) -> str:
        """프롬프트 파일 로드"""
        try:
            prompt_file = self.prompts_dir / "ppt_generation.txt"
            if prompt_file.exists():
                return prompt_file.read_text(encoding='utf-8')
        except Exception as e:
            logger.warning(f"프롬프트 파일 로드 실패: {e}")
        
        return (
            "당신은 전문 프레젠테이션 디자이너입니다. JSON만 출력. "
            "필드: topic,max_slides,slides[].title,key_message,bullets,layout,diagram,visual_suggestion,speaker_notes"
        )

    async def generate_enhanced_outline(self, topic: str, context_text: str, provider: Optional[str] = None,
                                        template_style: str = "business", include_charts: bool = True,
                                        retries: int = 2, document_filename: Optional[str] = None,
                                        custom_template_path: Optional[str] = None,
                                        presentation_type: str = "general",
                                        user_template_id: Optional[str] = None) -> DeckSpec:
        """AI 기반 향상된 아웃라인 생성"""
        
        logger.info(f"🚀 템플릿 기반 아웃라인 생성 시작: topic='{topic[:50]}', template_style='{template_style}'")
        logger.info(f"📝 파라미터: include_charts={include_charts}, presentation_type={presentation_type}")
        
        try:
            # 주제 개선
            improved_topic = self._improve_topic(topic, context_text)
            
            # AI 프롬프트 구성
            system = self._load_prompt()
            enhanced_requirements = [
                "- AI 응답 내용의 제목과 구조를 정확히 반영하여 슬라이드 생성",
                "- 번호가 있는 섹션(1. 제품 개요, 2. 기술 사양 등)은 각각 별도 슬라이드로 구성",
                "- 각 섹션의 세부 항목들은 bullets로 정확히 나열",
                "- 두 번째 슬라이드는 번호가 있는 섹션들을 목차로 구성",
                "- bullets 항목당 50자 이내로 간결하게 표현",
                f"- include_charts={include_charts} 이면 수치 데이터를 차트로 변환",
                f"- template_style={template_style} (business|minimal|modern|playful)",
                "- visual_suggestion: 관련 아이콘/이미지 아이디어 1줄",
                "- speaker_notes: 발표자 스크립트 2~4문장 한국어",
                "- 각 슬라이드 title은 섹션 번호와 제목을 포함 (예: '1. 제품 개요')",
                "- key_message는 해당 섹션의 핵심 설명문으로 구성"
            ]
            
            user_content = [
                f"주제: {improved_topic}",
                f"컨텍스트:\n{context_text[:8000]}",
                "요구사항:",
                *enhanced_requirements
            ]
            
            # AI 호출
            for attempt in range(retries + 1):
                try:
                    logger.info(f"🤖 AI 호출 시도 {attempt + 1}/{retries + 1}")
                    
                    # AI 메시지 구성
                    ai_message = f"{system}\n\n{chr(10).join(user_content)}"
                    logger.info(f"🔍 AI 메시지 길이: {len(ai_message)}문자")
                    logger.debug(f"🔍 AI 메시지 내용: {ai_message[:500]}...")
                    
                    # Provider 기본값 설정 (.env 기반 Settings 우선)
                    effective_provider = provider or settings.get_current_llm_provider()
                    logger.info(f"🔍 사용할 AI 제공자: {effective_provider}")
                    
                    response_generator = ai_service.chat_stream(
                        messages=[{"role": "user", "content": ai_message}],
                        provider=effective_provider
                    )
                    
                    # 스트림 응답 수집 (안전장치 추가)
                    full_response = ""
                    chunk_count = 0
                    max_chunks = 1000  # 최대 청크 수 제한
                    max_response_length = 50000  # 최대 응답 길이 제한
                    
                    # 타임아웃 처리를 위한 래퍼 함수
                    async def collect_response():
                        nonlocal full_response, chunk_count
                        async for chunk in response_generator:
                            chunk_count += 1
                            
                            # 안전장치: 최대 청크 수 초과 시 종료
                            if chunk_count > max_chunks:
                                logger.warning(f"⚠️ 최대 청크 수({max_chunks}) 초과로 스트림 종료")
                                break
                            
                            # 디버그 로깅 빈도 조절 (100개마다만 로깅)
                            if chunk_count % 100 == 0:
                                logger.debug(f"🔄 청크 {chunk_count}: {type(chunk)} - {str(chunk)[:50]}...")
                            
                            content = ""
                            if chunk:
                                if hasattr(chunk, 'text') and callable(getattr(chunk, 'text', None)):
                                    try:
                                        text_method = getattr(chunk, 'text')
                                        content = text_method()
                                    except Exception:
                                        content = str(chunk)
                                elif isinstance(chunk, str):
                                    content = chunk
                                else:
                                    content = str(chunk)
                            
                            # 빈 내용이나 메서드 바인딩 문자열 제외
                            if content and isinstance(content, str) and not content.startswith('<bound method'):
                                full_response += content
                                
                                # 안전장치: 최대 응답 길이 초과 시 종료
                                if len(full_response) > max_response_length:
                                    logger.warning(f"⚠️ 최대 응답 길이({max_response_length}) 초과로 스트림 종료")
                                    break
                                
                                # JSON 완료 패턴 감지
                                if content.strip().endswith('}') and '{"topic"' in full_response:
                                    # 기본적인 JSON 구조가 완성되었는지 확인
                                    brace_count = full_response.count('{') - full_response.count('}')
                                    if brace_count <= 0:
                                        logger.info(f"🔚 JSON 완료 패턴 감지, 스트림 종료 (청크: {chunk_count})")
                                        break
                                
                                # 진행상황 로깅 (100개마다)
                                if chunk_count % 100 == 0:
                                    logger.info(f"📊 진행상황: {chunk_count}청크, {len(full_response)}자")
                    
                    # 타임아웃과 함께 응답 수집 실행 (최대 60초)
                    try:
                        await asyncio.wait_for(collect_response(), timeout=60.0)
                    except asyncio.TimeoutError:
                        logger.warning(f"⚠️ AI 응답 수집 타임아웃 (60초), 현재까지 수집: {len(full_response)}문자")
                    
                    logger.info(f"📝 AI 응답 수집 완료: {len(full_response)}문자, {chunk_count}개 청크")
                    
                    # JSON 파싱
                    logger.info(f"🔍 AI 응답 파싱 시작: {full_response[:200]}...")
                    deck_spec = self._parse_ai_response(full_response, improved_topic, template_style)
                    
                    if deck_spec and len(deck_spec.slides) >= 2:
                        logger.info(f"✅ 아웃라인 생성 성공: {len(deck_spec.slides)}개 슬라이드")
                        for i, slide in enumerate(deck_spec.slides):
                            logger.info(f"  슬라이드 {i+1}: '{slide.title}' (bullets: {len(slide.bullets)}개)")
                        return deck_spec
                    else:
                        logger.warning(f"⚠️ AI 응답이 부적절함 (시도 {attempt + 1}): deck_spec={deck_spec}, slides={len(deck_spec.slides) if deck_spec else 0}")
                        
                except Exception as e:
                    logger.error(f"❌ AI 호출 실패 (시도 {attempt + 1}): {e}")
                    if attempt == retries:
                        break
                    await asyncio.sleep(1)
            
            # 폴백: 간단한 구조
            logger.warning("⚠️ AI 생성 실패, 폴백 구조 사용")
            return self._create_fallback_outline(improved_topic, context_text)
            
        except Exception as e:
            logger.error(f"generate_enhanced_outline 실패: {e}")
            return self._create_fallback_outline(topic, context_text)

    def _improve_topic(self, topic: str, context_text: str) -> str:
        """주제 개선 - AI 답변에서 실제 제목 추출"""
        if not context_text:
            return topic.strip()
        
        # AI 답변에서 실제 제목 추출 (quick PPT와 동일한 로직)
        lines = [ln.strip() for ln in context_text.split('\n') if ln.strip()]
        actual_title = topic  # 기본값
        
        for line in lines[:5]:  # 처음 5줄에서 찾기
            line = line.strip()
            if line.startswith('###') and not line.startswith('####'):
                # ### 헤딩에서 제목 추출
                actual_title = line.lstrip('#').strip()
                logger.info(f"🎯 템플릿 PPT 제목 추출 (###): '{actual_title}'")
                break
            elif line.startswith('##') and not line.startswith('###'):
                # ## 헤딩에서 제목 추출
                actual_title = line.lstrip('#').strip()
                logger.info(f"🎯 템플릿 PPT 제목 추출 (##): '{actual_title}'")
                break
            elif (not line.startswith('#') and len(line) > 5 and len(line) <= 50 and 
                  ('제품' in line or '소개' in line or '시스템' in line or '서비스' in line)):
                # 일반 텍스트에서 제목으로 보이는 라인 추출
                actual_title = line
                logger.info(f"🎯 템플릿 PPT 제목 추출 (텍스트): '{actual_title}'")
                break
        
        return actual_title.strip()

    def _extract_json(self, text: str) -> str:
        """텍스트에서 JSON 부분을 추출"""
        if text.strip().startswith('{'):
            return text
        block = re.search(r"```(?:json)?\n(.*)```", text, re.DOTALL)
        if block:
            return block.group(1)
        brace = re.search(r"{.*}", text, re.DOTALL)
        return brace.group(0) if brace else text

    def _parse_outline(self, text: str, fallback_topic: Optional[str] = None) -> DeckSpec:
        """아웃라인 텍스트를 파싱하여 DeckSpec으로 변환 (enhanced 서비스와 호환)"""
        try:
            data = json.loads(self._extract_json(text))
            max_slides = int(data.get("max_slides", 10))
            raw_slides = data.get("slides", [])[:max_slides]
            slides: List[SlideSpec] = []
            
            for s in raw_slides:
                diagram_info = s.get("diagram") or {}
                chart = None
                if diagram_info.get("chart"):
                    chart_raw = diagram_info.get("chart")
                    if isinstance(chart_raw, dict):
                        # 허용된 필드만 전달
                        allowed = {k: v for k, v in chart_raw.items() if k in ChartData.__fields__}
                        chart = ChartData(**allowed)
                raw_data = diagram_info.get("data")
                # Normalize list -> {'items': list} to avoid validation restrictions later
                if isinstance(raw_data, list):
                    raw_data = {"items": raw_data}
                diagram = DiagramData(type=diagram_info.get("type", "none"), data=raw_data, chart=chart)
                
                # role 정보를 style에 포함
                style_info = s.get("style") or {}
                role = s.get("role")
                if role:
                    style_info["role"] = role
                    # role에 따른 추가 플래그 설정
                    if role == "title":
                        style_info["title"] = True
                    elif role == "agenda":
                        style_info["agenda"] = True
                
                slides.append(SlideSpec(
                    title=s.get("title", ""),
                    key_message=s.get("key_message", ""),
                    bullets=s.get("bullets", []),
                    diagram=diagram,
                    layout=s.get("layout", "title-and-content"),
                    style=style_info if style_info else None,
                    visual_suggestion=s.get("visual_suggestion"),
                    speaker_notes=s.get("speaker_notes"),
                ))
            
            # 토픽 결정 (번호형 섹션 제목 금지)
            parsed_topic = data.get("topic") or fallback_topic or "발표자료"
            if re.match(r"^\d+\.\s+", parsed_topic):  # 번호형이면 폐기
                parsed_topic = fallback_topic or "발표자료"
            # fallback_topic이 더 구체적이고 질의형이 아니면 교체
            if fallback_topic and len(fallback_topic) > len(parsed_topic):
                parsed_topic = fallback_topic

            return DeckSpec(
                topic=parsed_topic,
                slides=slides,
                template_style=data.get("template_style", "business")
            )
        except Exception as e:
            logger.error(f"아웃라인 파싱 실패: {e}")
            # 기본 DeckSpec 반환
            return DeckSpec(
                topic=fallback_topic or "발표자료",
                slides=[],
                template_style="business"
            )

    def _parse_ai_response(self, response: str, topic: str, template_style: str) -> Optional[DeckSpec]:
        """AI 응답 JSON 파싱"""
        try:
            logger.info(f"🔍 JSON 추출 시작: 응답 길이={len(response)}")
            # JSON 추출
            json_match = re.search(r'\{.*\}', response, re.DOTALL)
            if not json_match:
                logger.error(f"❌ JSON 패턴을 찾을 수 없음. 응답 내용: {response[:500]}...")
                return None
            
            json_str = json_match.group()
            logger.info(f"🔍 JSON 문자열 추출 성공: {len(json_str)}자")
            
            data = json.loads(json_str)
            logger.info(f"🔍 JSON 파싱 성공: {list(data.keys())}")
            slides = []
            
            slides_data = data.get('slides', [])
            
            # slides가 비어있으면 slide_management에서 추출 시도
            if not slides_data and data.get('slide_management'):
                logger.warning("⚠️ 'slides' 필드가 비어있음, 'slide_management'에서 슬라이드 데이터 추출 시도")
                slide_mgmt = data.get('slide_management', [])
                slides_data = []
                for sm in slide_mgmt:
                    if sm.get('action') in ['keep', 'update'] and sm.get('content'):
                        # slide_management의 content를 slides 형식으로 변환
                        slides_data.append({
                            'title': sm.get('content', {}).get('title', '제목 없음'),
                            'key_message': sm.get('content', {}).get('key_message', ''),
                            'bullets': sm.get('content', {}).get('bullets', []),
                            'layout': sm.get('layout', 'title_and_content'),
                            'speaker_notes': sm.get('content', {}).get('speaker_notes', ''),
                            'visual_suggestion': ''
                        })
                logger.info(f"🔄 slide_management에서 {len(slides_data)}개 슬라이드 복원")
            
            logger.info(f"🔍 슬라이드 데이터 수: {len(slides_data)}")
            
            for i, slide_data in enumerate(slides_data):
                slide = SlideSpec(
                    title=slide_data.get('title', '제목 없음'),
                    key_message=slide_data.get('key_message', ''),
                    bullets=slide_data.get('bullets', []),
                    layout=slide_data.get('layout', 'title_and_content'),
                    speaker_notes=slide_data.get('speaker_notes', ''),
                    visual_suggestion=slide_data.get('visual_suggestion', '')
                )
                slides.append(slide)
                logger.info(f"  슬라이드 {i+1}: '{slide.title}', bullets={len(slide.bullets)}")
            
            deck = DeckSpec(
                topic=data.get('topic', topic),
                slides=slides,
                max_slides=len(slides),
                template_style=template_style
            )
            logger.info(f"✅ DeckSpec 생성 완료: topic='{deck.topic}', slides={len(deck.slides)}")
            return deck
            
        except Exception as e:
            logger.error(f"AI 응답 파싱 실패: {e}")
            return None

    def _create_fallback_outline(self, topic: str, context_text: str) -> DeckSpec:
        """폴백 아웃라인 생성"""
        slides = [
            SlideSpec(title=topic, key_message="", bullets=[], layout="title-only"),
            SlideSpec(title="목차", key_message="", bullets=["주요 내용 1", "주요 내용 2", "결론"], layout="title-and-content"),
            SlideSpec(title="주요 내용 1", key_message="첫 번째 주요 내용입니다.", bullets=["세부사항 1", "세부사항 2"], layout="title-and-content"),
            SlideSpec(title="주요 내용 2", key_message="두 번째 주요 내용입니다.", bullets=["세부사항 1", "세부사항 2"], layout="title-and-content"),
            SlideSpec(title="결론", key_message="결론입니다.", bullets=["요약", "향후 계획"], layout="title-and-content")
        ]
        return DeckSpec(topic=topic, slides=slides, max_slides=len(slides))

    def build_templated_pptx(self, spec: DeckSpec, template_path: Path, file_basename: Optional[str] = None,
                            text_box_mappings: Optional[List[Dict[str, Any]]] = None,
                            content_segments: Optional[List[Dict[str, Any]]] = None) -> str:
        """템플릿 기반 PPT 빌드"""
        
        logger.info(f"🏗️ 템플릿 기반 PPT 빌드 시작: {len(spec.slides)}개 슬라이드")
        logger.info(f"📄 템플릿 파일: {template_path}")
        
        if not template_path.exists():
            raise FileNotFoundError(f"템플릿을 찾을 수 없습니다: {template_path}")
        
        try:
            # 파일명 생성
            if not file_basename:
                safe_topic = self._normalize_topic_for_filename(spec.topic)
                file_basename = f"templated_presentation_{safe_topic}"
            
            filename = f"{file_basename}.pptx"
            output_path = self.upload_dir / filename
            
            # 템플릿 로드
            prs = Presentation(str(template_path))
            logger.info(f"📋 템플릿 로드 완료: {len(prs.slide_layouts)}개 레이아웃")
            
            # 기존 슬라이드 제거 (템플릿 슬라이드만 유지)
            slide_count = len(prs.slides)
            for i in range(slide_count - 1, -1, -1):
                if i > 0:  # 첫 번째 슬라이드는 유지
                    rId = prs.slides._sldIdLst[i].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[i]
            
            # 새 슬라이드 생성
            for i, slide_spec in enumerate(spec.slides):
                logger.info(f"📄 슬라이드 {i+1} 생성: '{slide_spec.title}'")
                
                if i == 0:
                    # 첫 번째 슬라이드 수정
                    slide = prs.slides[0]
                    self._update_title_slide(slide, slide_spec)
                else:
                    # 새 슬라이드 추가
                    layout_idx = 1 if i == 1 else 1  # 목차와 내용 모두 같은 레이아웃
                    slide_layout = prs.slide_layouts[layout_idx]
                    slide = prs.slides.add_slide(slide_layout)
                    self._populate_template_slide(slide, slide_spec, text_box_mappings)
                
                logger.info(f"✅ 슬라이드 {i+1} 완료")
            
            # 파일 저장
            prs.save(str(output_path))
            logger.info(f"✅ 템플릿 기반 PPT 빌드 완료: {output_path}")
            return str(output_path)
            
        except Exception as e:
            logger.error(f"build_templated_pptx 실패: {e}")
            raise

    def _update_title_slide(self, slide, spec: SlideSpec):
        """제목 슬라이드 업데이트"""
        try:
            if slide.shapes.title:
                slide.shapes.title.text = spec.title
                logger.info(f"✅ 제목 설정: '{spec.title}'")
        except Exception as e:
            logger.error(f"제목 슬라이드 업데이트 실패: {e}")

    def _populate_template_slide(self, slide, spec: SlideSpec, text_box_mappings: Optional[List[Dict[str, Any]]] = None):
        """템플릿 슬라이드에 콘텐츠 채우기"""
        try:
            # 제목 설정
            if slide.shapes.title:
                slide.shapes.title.text = spec.title
            
            # 콘텐츠 영역 찾기 및 채우기
            content_shape = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    content_shape = shape
                    break
            
            if content_shape and hasattr(content_shape, 'text_frame'):
                tf = content_shape.text_frame
                tf.clear()
                tf.word_wrap = True
                
                # 키 메시지 추가
                if spec.key_message:
                    p = tf.paragraphs[0]
                    p.text = spec.key_message
                    p.font.size = Pt(22)
                    p.font.bold = True
                
                # 불릿 포인트 추가
                for bullet in spec.bullets:
                    if bullet.strip():
                        p = tf.add_paragraph()
                        p.text = f"• {bullet.strip()}"
                        p.font.size = Pt(18)
                        p.level = 1
            
            logger.info(f"✅ 템플릿 슬라이드 콘텐츠 완료: '{spec.title}'")
            
        except Exception as e:
            logger.error(f"템플릿 슬라이드 채우기 실패: {e}")

    def build_enhanced_pptx_with_slide_management(self, spec: DeckSpec, file_basename: Optional[str] = None,
                                                 template_style: str = "business", include_charts: bool = True,
                                                 custom_template_path: Optional[str] = None,
                                                 user_template_id: Optional[str] = None,
                                                 text_box_mappings: Optional[List[Dict[str, Any]]] = None,
                                                 content_segments: Optional[List[Dict[str, Any]]] = None,
                                                 slide_management: Optional[List[Dict[str, Any]]] = None,
                                                 used_template_indices: Optional[List[int]] = None,
                                                 template_metadata: Optional[Dict[str, Any]] = None) -> str:
        """슬라이드 관리가 포함된 Enhanced PPT 빌드 (enhanced 서비스와 호환)
        
        Args:
            spec: DeckSpec - AI 생성 콘텐츠 슬라이드들
            file_basename: 출력 파일명
            template_style: 템플릿 스타일
            include_charts: 차트 포함 여부
            custom_template_path: 커스텀 템플릿 경로
            user_template_id: 사용자 템플릿 ID
            text_box_mappings: 텍스트박스 매핑 정보
            content_segments: 콘텐츠 세그먼트
            slide_management: 슬라이드 관리 정보
            used_template_indices: 🆕 사용할 템플릿 슬라이드 인덱스 (slide_type_matcher 결과)
            template_metadata: 🆕 템플릿 메타데이터 (매핑되지 않은 요소 클리어용)
        """
        
        logger.info(f"🏗️ Enhanced PPT 빌드 시작: {len(spec.slides)}개 슬라이드")
        logger.info(f"📋 매핑 정보: text_box_mappings={len(text_box_mappings or [])}, content_segments={len(content_segments or [])}, slide_management={len(slide_management or [])}")
        if used_template_indices:
            logger.info(f"📋 사용할 템플릿 슬라이드: {used_template_indices}")
        if template_metadata:
            logger.info(f"📋 템플릿 메타데이터 제공됨: {len(template_metadata.get('slides', []))}개 슬라이드")
        
        try:
            # 커스텀 템플릿 경로가 있으면 템플릿 기반 빌드 사용
            if custom_template_path and os.path.exists(custom_template_path):
                logger.info(f"📄 커스텀 템플릿 사용: {custom_template_path}")
                
                # 🆕 템플릿 메타데이터: 파라미터로 전달된 것이 없으면 로드 시도
                if not template_metadata:
                    try:
                        from app.agents.features.presentation.services.user_template_manager import user_template_manager
                        # user_template_id가 있으면 해당 템플릿의 메타데이터 로드
                        if user_template_id:
                            # 경로에서 user_id 추출 시도
                            import re
                            match = re.search(r'/users/(\d+)/', custom_template_path)
                            if match:
                                owner_id = match.group(1)
                                template_metadata = user_template_manager.get_template_metadata(owner_id, user_template_id)
                                if template_metadata:
                                    logger.info(f"📋 템플릿 메타데이터 로드됨: {len(template_metadata.get('slides', []))}개 슬라이드")
                    except Exception as meta_e:
                        logger.warning(f"⚠️ 템플릿 메타데이터 로드 실패 (무시): {meta_e}")
                else:
                    logger.info(f"📋 전달된 메타데이터 사용: {len(template_metadata.get('slides', []))}개 슬라이드")
                
                # 🆕 매핑 또는 used_template_indices가 있으면 _build_with_mappings 사용
                # (used_template_indices가 있으면 Strategy C: 슬라이드 복제/삭제 필요)
                has_mappings = bool(text_box_mappings) or bool(content_segments) or bool(slide_management)
                has_template_indices = bool(used_template_indices)
                
                if has_mappings or has_template_indices:
                    logger.info(f"🎯 매핑 기반 템플릿 빌드 실행 (mappings={has_mappings}, indices={has_template_indices})")
                    return self._build_with_mappings(
                        spec=spec,
                        template_path=Path(custom_template_path),
                        file_basename=file_basename,
                        text_box_mappings=text_box_mappings,
                        content_segments=content_segments,
                        slide_management=slide_management,
                        template_metadata=template_metadata,  # 🆕 메타데이터 전달
                        used_template_indices=used_template_indices,
                    )
                else:
                    # 매핑 없으면 기본 템플릿 빌드
                    logger.info(f"📄 기본 템플릿 빌드 (매핑 없음)")
                    return self.build_templated_pptx(
                        spec=spec,
                        template_path=Path(custom_template_path),
                        file_basename=file_basename,
                        text_box_mappings=text_box_mappings,
                        content_segments=content_segments
                    )
            else:
                # 기본 빌더 사용 (템플릿 없음) - 간단한 PPT 생성
                logger.info(f"📄 기본 빌더 사용 (템플릿 없음)")
                
                # 파일명 생성
                if not file_basename:
                    safe_topic = self._normalize_topic_for_filename(spec.topic)
                    file_basename = f"enhanced_presentation_{safe_topic}"
                
                filename = f"{file_basename}.pptx"
                output_path = self.upload_dir / filename
                
                # 간단한 PPT 생성
                prs = Presentation()
                
                # 각 슬라이드 생성
                for i, slide_spec in enumerate(spec.slides):
                    if i == 0:
                        # 첫 번째 슬라이드 (제목 슬라이드)
                        title_slide_layout = prs.slide_layouts[0]
                        slide = prs.slides.add_slide(title_slide_layout)
                        if slide.shapes.title:
                            slide.shapes.title.text = slide_spec.title
                        # 부제목 placeholder 처리
                        if len(slide.shapes.placeholders) > 1:
                            subtitle_placeholder = slide.shapes.placeholders[1]
                            if getattr(subtitle_placeholder, 'has_text_frame', False):
                                text_frame = getattr(subtitle_placeholder, 'text_frame', None)
                                if text_frame:
                                    text_frame.text = slide_spec.key_message
                    else:
                        # 내용 슬라이드
                        content_slide_layout = prs.slide_layouts[1]
                        slide = prs.slides.add_slide(content_slide_layout)
                        if slide.shapes.title:
                            slide.shapes.title.text = slide_spec.title
                        
                        # 콘텐츠 추가
                        if len(slide.shapes.placeholders) > 1:
                            content_placeholder = slide.shapes.placeholders[1]
                            if getattr(content_placeholder, 'has_text_frame', False):
                                tf = getattr(content_placeholder, 'text_frame', None)
                                if tf:
                                    tf.clear()
                                    if slide_spec.key_message:
                                        p = tf.paragraphs[0]
                                        p.text = slide_spec.key_message
                                    
                                    # 불릿 포인트 추가
                                    for bullet in slide_spec.bullets:
                                        p = tf.add_paragraph()
                                        p.text = bullet
                                        p.level = 1
                
                prs.save(str(output_path))
                logger.info(f"✅ Enhanced PPT 빌드 완료: {output_path}")
                
                return str(output_path)
                
        except Exception as e:
            logger.error(f"Enhanced PPT 빌드 실패: {e}")
            raise

    def _build_with_mappings(self, spec: DeckSpec, template_path: Path, file_basename: Optional[str] = None,
                            text_box_mappings: Optional[List[Dict[str, Any]]] = None,
                            content_segments: Optional[List[Dict[str, Any]]] = None,
                            slide_management: Optional[List[Dict[str, Any]]] = None,
                            template_metadata: Optional[Dict[str, Any]] = None,
                            used_template_indices: Optional[List[int]] = None) -> str:
        """매핑을 적용한 템플릿 기반 PPT 빌드
        
        🆕 전략 C: AI 슬라이드가 템플릿보다 많으면 content 슬라이드 복제
        - AI 슬라이드의 유형(title, toc, content, thanks)과 
        - 템플릿 슬라이드의 role(title, toc, content, section, thanks)을 매칭
        - 초과 AI 슬라이드는 content 레이아웃을 복제하여 추가
        
        Args:
            spec: AI 생성 DeckSpec
            template_path: 템플릿 파일 경로
            file_basename: 출력 파일명
            text_box_mappings: 텍스트박스 매핑 (slideIndex가 template index임)
            content_segments: 콘텐츠 세그먼트
            slide_management: 슬라이드 관리 정보
            template_metadata: 템플릿 메타데이터
            used_template_indices: 사용할 템플릿 슬라이드 인덱스 목록
        """
        
        logger.info(f"🎯 매핑 기반 PPT 빌드 시작")
        
        try:
            # 파일명 생성
            if not file_basename:
                safe_topic = self._normalize_topic_for_filename(spec.topic)
                file_basename = f"mapped_presentation_{safe_topic}"
            
            filename = f"{file_basename}.pptx"
            output_path = self.upload_dir / filename
            
            # 템플릿 로드
            prs = Presentation(str(template_path))
            template_slide_count = len(prs.slides)
            ai_slide_count = len(spec.slides)
            logger.info(f"📋 템플릿 로드 완료: {template_slide_count}개 슬라이드, AI 슬라이드: {ai_slide_count}개")
            
            # 🆕 스타일 보존을 위해 템플릿 텍스트 정리 생략
            # content_cleaner._clean_slide_content()를 호출하면 스타일이 손실됨
            # 대신 enhanced_object_processor가 매핑된 텍스트박스의 텍스트만 교체 (run.text 직접 변경)
            logger.info(f"📋 템플릿 스타일 보존 모드: 텍스트 정리 생략, 매핑만 적용")
            
            # 🆕 전략 C: AI 슬라이드가 템플릿보다 많은 경우 처리
            if ai_slide_count > template_slide_count and not used_template_indices:
                logger.info(f"📋 AI 슬라이드({ai_slide_count}) > 템플릿({template_slide_count}) - 슬라이드 복제 필요")
                
                # content 타입 슬라이드 찾기 (복제 대상)
                content_slide_idx = self._find_content_slide_index(prs, template_metadata)
                
                if content_slide_idx is not None:
                    # 추가로 필요한 슬라이드 수
                    slides_to_add = ai_slide_count - template_slide_count
                    logger.info(f"📋 content 슬라이드(idx={content_slide_idx}) {slides_to_add}개 복제")
                    
                    # 슬라이드 복제 (마지막 슬라이드 앞에 삽입)
                    for i in range(slides_to_add):
                        self._duplicate_slide(prs, content_slide_idx)
                        logger.info(f"✅ 슬라이드 복제 완료: {i+1}/{slides_to_add}")
                    
                    template_slide_count = len(prs.slides)
                    logger.info(f"📋 복제 후 슬라이드 수: {template_slide_count}개")
                else:
                    logger.warning(f"⚠️ content 슬라이드를 찾을 수 없음 - 기본 레이아웃으로 추가")
                    # 기본 레이아웃으로 슬라이드 추가
                    slides_to_add = ai_slide_count - template_slide_count
                    for i in range(slides_to_add):
                        self._add_blank_content_slide(prs)
                    template_slide_count = len(prs.slides)
            
            # used_template_indices가 제공되면 해당 슬라이드만 사용
            if used_template_indices:
                logger.info(f"📋 사용할 템플릿 슬라이드: {used_template_indices}")
                
                # 🆕 AI 슬라이드가 템플릿보다 많은 경우, 슬라이드 복제 필요
                # used_template_indices는 AI 슬라이드 수만큼 있지만,
                # 실제 템플릿 슬라이드 수보다 많을 수 있음 (같은 인덱스 재사용)
                unique_indices = set(used_template_indices)
                if ai_slide_count > template_slide_count:
                    # 복제할 슬라이드: content 타입 중 가장 많이 사용된 인덱스
                    from collections import Counter
                    idx_counts = Counter(used_template_indices)
                    # content 슬라이드 찾기 (title=0, toc=1, thanks=마지막 제외)
                    content_indices = [idx for idx in idx_counts.keys() if idx not in [0, 1, template_slide_count - 1]]
                    if content_indices:
                        content_slide_template_idx = max(content_indices, key=lambda x: idx_counts[x])
                    else:
                        content_slide_template_idx = self._find_content_slide_in_indices(list(unique_indices), template_metadata)
                    
                    if content_slide_template_idx is not None:
                        slides_to_add = ai_slide_count - template_slide_count
                        logger.info(f"📋 AI({ai_slide_count}) > 템플릿({template_slide_count}): {slides_to_add}개 복제 필요 (template idx: {content_slide_template_idx})")
                        
                        # 실제 슬라이드 복제
                        for i in range(slides_to_add):
                            self._duplicate_slide(prs, content_slide_template_idx)
                            logger.info(f"✅ 슬라이드 복제: {i+1}/{slides_to_add}")
                        
                        template_slide_count = len(prs.slides)
                        logger.info(f"📋 복제 후 슬라이드 수: {template_slide_count}개")
                
                # 🆕 사용하지 않는 템플릿 슬라이드 삭제 로직 비활성화
                # (AI > 템플릿인 경우 모든 템플릿 슬라이드 사용)
                # AI < 템플릿인 경우에만 삭제 필요
                if ai_slide_count < template_slide_count:
                    slides_to_delete = [
                        i for i in range(template_slide_count) 
                        if i not in unique_indices
                    ]
                    slides_to_delete.sort(reverse=True)
                    
                    if slides_to_delete:
                        logger.info(f"🗑️ 삭제할 슬라이드: {slides_to_delete}")
                        
                        for slide_idx in slides_to_delete:
                            if slide_idx < len(prs.slides):
                                try:
                                    rId = prs.slides._sldIdLst[slide_idx].rId
                                    prs.part.drop_rel(rId)
                                    del prs.slides._sldIdLst[slide_idx]
                                    logger.info(f"🗑️ 슬라이드 {slide_idx} 삭제 완료")
                                except Exception as del_e:
                                    logger.warning(f"⚠️ 슬라이드 {slide_idx} 삭제 실패: {del_e}")
                        
                        logger.info(f"📋 삭제 후 슬라이드 수: {len(prs.slides)}개")
                    
                    # 매핑의 slideIndex를 새 인덱스로 재조정
                    if text_box_mappings:
                        # 원본 template_index -> 삭제 후 새 index 매핑
                        old_to_new_idx = {}
                        new_idx = 0
                        for old_idx in range(template_slide_count):
                            if old_idx not in slides_to_delete:
                                old_to_new_idx[old_idx] = new_idx
                                new_idx += 1
                        
                        logger.info(f"📋 인덱스 매핑: {old_to_new_idx}")
                        
                        # 매핑 업데이트
                        updated_mappings = []
                        for mapping in text_box_mappings:
                            old_slide_idx = mapping.get('slideIndex', 0)
                            if old_slide_idx in old_to_new_idx:
                                new_mapping = {**mapping, 'slideIndex': old_to_new_idx[old_slide_idx]}
                                updated_mappings.append(new_mapping)
                            else:
                                # 삭제된 슬라이드에 대한 매핑은 제외
                                logger.warning(f"⚠️ 삭제된 슬라이드({old_slide_idx})에 대한 매핑 제외")
                        
                        text_box_mappings = updated_mappings
                        logger.info(f"📋 업데이트된 매핑 수: {len(text_box_mappings)}개")
            
            # 🆕 매핑되지 않은 요소들을 클리어하기 위한 추가 매핑 생성
            if template_metadata and text_box_mappings:
                clear_mappings = self._generate_clear_mappings(template_metadata, text_box_mappings, len(prs.slides))
                if clear_mappings:
                    logger.info(f"🧹 매핑되지 않은 요소 클리어 매핑 추가: {len(clear_mappings)}개")
                    text_box_mappings = text_box_mappings + clear_mappings
            
            # Enhanced Object Processor로 매핑 적용
            if hasattr(self, 'object_processor') and text_box_mappings:
                logger.info(f"🔧 Enhanced Object Processor로 {len(text_box_mappings)}개 매핑 적용")
                self.object_processor.apply_object_mappings(prs, text_box_mappings, content_segments)
            else:
                logger.info(f"📄 매핑 없음 또는 Object Processor 없음 - AI 콘텐츠만 적용")
                
                # 매핑이 없으면 AI 콘텐츠를 순차적으로 적용
                for ai_idx, slide_spec in enumerate(spec.slides):
                    if ai_idx < len(prs.slides):
                        slide = prs.slides[ai_idx]
                        self._apply_ai_content_to_slide(slide, slide_spec, ai_idx)
            
            # 파일 저장
            prs.save(str(output_path))
            logger.info(f"✅ 매핑 기반 PPT 빌드 완료: {output_path}")
            return str(output_path)
            
        except Exception as e:
            logger.error(f"매핑 기반 PPT 빌드 실패: {e}")
            raise

    def _generate_clear_mappings(
        self, 
        template_metadata: Dict[str, Any], 
        existing_mappings: List[Dict[str, Any]],
        slide_count: int
    ) -> List[Dict[str, Any]]:
        """매핑되지 않은 is_fixed=False 요소들을 클리어하는 매핑 생성
        
        Args:
            template_metadata: 템플릿 메타데이터
            existing_mappings: 기존 매핑 리스트
            slide_count: 현재 PPT 슬라이드 수
            
        Returns:
            클리어 매핑 리스트
        """
        clear_mappings = []
        
        if not template_metadata:
            return clear_mappings
        
        # 기존 매핑된 element_id와 originalName 수집
        mapped_element_ids = set()
        mapped_original_names = set()
        for m in existing_mappings:
            if m.get('elementId'):
                mapped_element_ids.add(m.get('elementId'))
            if m.get('originalName'):
                mapped_original_names.add(m.get('originalName'))
        
        logger.info(f"🔍 클리어 매핑 생성: 매핑된 요소 {len(mapped_element_ids)}개, 원본이름 {len(mapped_original_names)}개")
        
        # 메타데이터의 각 슬라이드 요소 확인
        for slide_meta in template_metadata.get('slides', []):
            slide_idx = slide_meta.get('index', 1) - 1  # 1-based to 0-based
            
            # 슬라이드 범위 확인
            if slide_idx >= slide_count:
                continue
            
            for element in slide_meta.get('elements', []):
                element_id = element.get('id', '')
                original_name = element.get('original_name', '')
                is_fixed = element.get('is_fixed', False)
                element_role = element.get('element_role', '')
                
                # is_fixed=True인 요소는 클리어하지 않음 (Company Name, Logo 등)
                if is_fixed:
                    continue
                
                # 이미 매핑된 요소는 클리어하지 않음
                if element_id in mapped_element_ids or original_name in mapped_original_names:
                    continue
                
                # 클리어 대상이 아닌 역할 제외 (이미지 플레이스홀더 등)
                skip_roles = ['image_placeholder', 'chart_placeholder', 'diagram', 'decorative']
                if element_role in skip_roles:
                    continue
                
                # 클리어 매핑 생성
                clear_mappings.append({
                    'slideIndex': slide_idx,
                    'elementId': element_id,
                    'originalName': original_name,
                    'objectType': 'textbox',
                    'action': 'replace_content',
                    'newContent': '',  # 빈 문자열로 클리어
                    'isEnabled': True,
                    'target_role': f'clear_{element_role}'  # target_role 설정하여 빈 문자열 클리어 허용
                })
                logger.debug(f"🧹 클리어 대상: slide[{slide_idx}] {element_id} ({original_name}) role={element_role}")
        
        return clear_mappings

    def _copy_font_style(self, src_font, dst_font):
        """폰트 스타일 복사 (완전한 스타일 보존)"""
        try:
            # 폰트 이름
            if src_font.name: 
                dst_font.name = src_font.name
            # 폰트 크기
            if src_font.size: 
                dst_font.size = src_font.size
            # 굵기
            if src_font.bold is not None: 
                dst_font.bold = src_font.bold
            # 기울임
            if src_font.italic is not None: 
                dst_font.italic = src_font.italic
            # 밑줄
            if src_font.underline is not None: 
                dst_font.underline = src_font.underline
            
            # 색상 복사 (상세)
            try:
                if hasattr(src_font, 'color') and src_font.color:
                    src_color = src_font.color
                    # RGB 색상
                    if hasattr(src_color, 'type'):
                        if src_color.type == 1:  # RGB
                            if src_color.rgb:
                                dst_font.color.rgb = src_color.rgb
                        elif src_color.type == 2:  # THEME
                            if hasattr(src_color, 'theme_color') and src_color.theme_color:
                                dst_font.color.theme_color = src_color.theme_color
                            if hasattr(src_color, 'brightness') and src_color.brightness is not None:
                                dst_font.color.brightness = src_color.brightness
                    elif hasattr(src_color, 'rgb') and src_color.rgb:
                        # type 속성이 없는 경우 직접 RGB 복사
                        dst_font.color.rgb = src_color.rgb
            except Exception as color_err:
                logger.debug(f"색상 복사 중 오류 (무시됨): {color_err}")
                
        except Exception as e:
            logger.warning(f"폰트 스타일 복사 중 오류: {e}")

    def _copy_paragraph_style(self, src_para, dst_para):
        """문단 스타일 복사 (정렬, 레벨, 간격 등)"""
        try:
            # 정렬
            if src_para.alignment is not None:
                dst_para.alignment = src_para.alignment
            # 레벨 (들여쓰기)
            if hasattr(src_para, 'level') and src_para.level is not None:
                dst_para.level = src_para.level
            # 줄 간격
            if hasattr(src_para, 'line_spacing') and src_para.line_spacing:
                dst_para.line_spacing = src_para.line_spacing
            # 공백
            if hasattr(src_para, 'space_before') and src_para.space_before:
                dst_para.space_before = src_para.space_before
            if hasattr(src_para, 'space_after') and src_para.space_after:
                dst_para.space_after = src_para.space_after
        except Exception as e:
            logger.debug(f"문단 스타일 복사 중 오류 (무시됨): {e}")

    def _replace_text_preserving_style(self, shape, new_text):
        """스타일을 유지하면서 텍스트 교체"""
        try:
            tf = shape.text_frame
            if not tf.paragraphs:
                tf.text = new_text
                return

            # 첫 번째 문단 사용
            p = tf.paragraphs[0]
            
            # 첫 번째 run의 스타일 유지
            if p.runs:
                # 첫 번째 run에 텍스트 설정
                p.runs[0].text = new_text
                # 나머지 run의 텍스트 제거 (스타일은 유지되지만 내용은 비움)
                for i in range(1, len(p.runs)):
                    p.runs[i].text = ""
            else:
                p.text = new_text
                
            # 나머지 문단 제거 (내용 비우기)
            # python-pptx에서 문단 삭제가 까다로우므로 텍스트만 비움
            for i in range(1, len(tf.paragraphs)):
                tf.paragraphs[i].clear()
                
        except Exception as e:
            logger.warning(f"텍스트 교체 중 오류: {e}")
            # 폴백
            shape.text_frame.text = new_text

    def _apply_content_preserving_style(self, shape, bullets):
        """스타일을 유지하면서 콘텐츠(불릿) 적용 (개선된 버전)
        
        템플릿의 기존 폰트, 크기, 색상, 문단 스타일을 완전히 보존합니다.
        """
        try:
            tf = shape.text_frame
            if not tf.paragraphs:
                tf.text = "\n".join(bullets)
                return

            # 첫 번째 문단의 스타일 참조
            ref_p = tf.paragraphs[0]
            ref_run = ref_p.runs[0] if ref_p.runs else None
            
            # 원본 문단 정렬 저장
            original_alignment = ref_p.alignment
            original_level = ref_p.level if hasattr(ref_p, 'level') else 0
            
            # 기존 문단들 내용 비우기 (첫 번째 제외)
            for i in range(1, len(tf.paragraphs)):
                tf.paragraphs[i].clear()
                
            # 불릿 내용 적용
            if bullets:
                first_bullet = bullets[0]
                # 첫 번째 문단 업데이트
                if ref_run:
                    ref_run.text = f"• {first_bullet}" if not first_bullet.startswith("•") else first_bullet
                    for i in range(1, len(ref_p.runs)):
                        ref_p.runs[i].text = ""
                else:
                    ref_p.text = f"• {first_bullet}" if not first_bullet.startswith("•") else first_bullet
                    
                # 나머지 불릿 추가
                for b in bullets[1:]:
                    text = f"• {b}" if not b.startswith("•") else b
                    new_p = tf.add_paragraph()
                    new_p.text = text
                    
                    # 문단 스타일 복사 (정렬, 레벨)
                    self._copy_paragraph_style(ref_p, new_p)
                    
                    # 레벨 설정 (bullet level)
                    if original_level is not None:
                        try:
                            new_p.level = original_level
                        except:
                            new_p.level = 1
                    else:
                        new_p.level = 1
                    
                    # 폰트 스타일 복사 시도
                    if ref_run and new_p.runs:
                        self._copy_font_style(ref_run.font, new_p.runs[0].font)
            else:
                # 내용이 없으면 첫 번째 문단도 비움
                if ref_run:
                    ref_run.text = ""
                else:
                    ref_p.text = ""
                    
        except Exception as e:
            logger.warning(f"콘텐츠 적용 중 오류: {e}")
            # 폴백
            shape.text_frame.text = "\n".join(bullets)

    def _replace_table_cell_text_preserving_style(self, cell, new_text: str):
        """테이블 셀의 텍스트를 스타일 보존하면서 교체
        
        Args:
            cell: 테이블 셀 (_Cell 객체)
            new_text: 새로운 텍스트
        """
        try:
            tf = cell.text_frame
            if not tf.paragraphs:
                tf.text = new_text
                return
            
            # 첫 번째 문단 사용
            p = tf.paragraphs[0]
            
            # 첫 번째 run의 스타일 유지
            if p.runs:
                # 첫 번째 run에 텍스트 설정
                p.runs[0].text = new_text
                # 나머지 run의 텍스트 제거
                for i in range(1, len(p.runs)):
                    p.runs[i].text = ""
            else:
                p.text = new_text
            
            # 나머지 문단 비우기
            for i in range(1, len(tf.paragraphs)):
                tf.paragraphs[i].clear()
                
        except Exception as e:
            logger.debug(f"테이블 셀 텍스트 교체 중 오류: {e}")
            try:
                cell.text = new_text
            except:
                pass

    def _apply_content_to_table_preserving_style(self, table, table_data: List[List[str]]):
        """테이블에 데이터를 스타일 보존하면서 적용
        
        Args:
            table: Table 객체
            table_data: 2D 리스트 형태의 테이블 데이터 [[row1_col1, row1_col2], [row2_col1, row2_col2], ...]
        """
        try:
            for row_idx, row_data in enumerate(table_data):
                if row_idx >= len(table.rows):
                    break
                for col_idx, cell_text in enumerate(row_data):
                    if col_idx >= len(table.columns):
                        break
                    cell = table.cell(row_idx, col_idx)
                    self._replace_table_cell_text_preserving_style(cell, cell_text)
                    
            logger.debug(f"테이블 데이터 적용 완료: {len(table_data)}행")
        except Exception as e:
            logger.warning(f"테이블 데이터 적용 중 오류: {e}")

    def _get_all_text_shapes(self, slide):
        """슬라이드 내의 모든 텍스트 가능 객체를 재귀적으로 수집 (그룹 포함)"""
        text_shapes = []
        
        def _collect_text_shapes(shapes):
            for shape in shapes:
                # 그룹인 경우 재귀 호출
                if shape.shape_type == MSO_SHAPE.GROUP:
                    _collect_text_shapes(shape.shapes)
                    continue
                
                # 텍스트 프레임이 있는 경우
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text_shapes.append(shape)
                            
        _collect_text_shapes(slide.shapes)
        return text_shapes

    def _apply_ai_content_to_slide(self, slide, slide_spec, slide_index: int):
        """AI가 생성한 콘텐츠를 슬라이드에 적용 (스타일 보존 개선)
        
        템플릿 슬라이드의 기존 텍스트 스타일(폰트, 색상, 크기)을 유지하면서 내용을 교체합니다.
        """
        try:
            logger.info(f"🔄 슬라이드 {slide_index + 1}에 AI 콘텐츠 적용: '{slide_spec.title}'")
            
            # 텍스트박스 수집 (재귀적 탐색으로 변경)
            all_shapes = self._get_all_text_shapes(slide)
            
            text_shapes = []
            for shape in all_shapes:
                try:
                    # 그룹 내부 shape는 top/left가 그룹 기준일 수 있으나, 
                    # python-pptx에서는 절대 좌표를 제공하는 경우가 많음.
                    top = shape.top if hasattr(shape, 'top') else 0
                    area = shape.width * shape.height if hasattr(shape, 'width') else 0
                except:
                    top = 0
                    area = 0
                
                # 로고나 저작권 문구 등 보존해야 할 패턴 확인
                text = shape.text_frame.text
                # PRESERVE_PATTERNS를 매우 보수적으로 설정 (사용자 요청: 템플릿 텍스트는 모두 클리어)
                # "company", "date", "page" 등 일반적인 단어는 제거하여 오탐 방지
                PRESERVE_PATTERNS = ["<logo>", "<copyright>", "confidential", "all rights reserved"]
                should_preserve = any(p in text.lower() for p in PRESERVE_PATTERNS)
                
                if not should_preserve:
                    text_shapes.append({
                        'shape': shape,
                        'top': top,
                        'area': area,
                        'original_text': text.strip()[:30]
                    })
            
            # 위치(top)로 정렬 - 위쪽이 제목, 아래쪽이 콘텐츠
            text_shapes.sort(key=lambda x: x['top'])
            
            logger.info(f"  📋 편집 대상 텍스트박스 {len(text_shapes)}개 발견")
            
            title_applied = False
            content_applied = False
            
            # Step 1: 제목 적용 (첫 번째 텍스트박스)
            if len(text_shapes) > 0:
                title_shape = text_shapes[0]['shape']
                self._replace_text_preserving_style(title_shape, slide_spec.title)
                title_applied = True
                logger.info(f"  ✅ 제목 적용: '{slide_spec.title}'")
            
            # Step 2: 콘텐츠 적용 (두 번째 텍스트박스)
            content_items = []
            if slide_spec.bullets:
                content_items = slide_spec.bullets
            elif hasattr(slide_spec, 'key_message') and slide_spec.key_message:
                content_items = [slide_spec.key_message]
            
            if len(text_shapes) > 1:
                content_shape = text_shapes[1]['shape']
                if content_items:
                    self._apply_content_preserving_style(content_shape, content_items)
                    content_applied = True
                    logger.info(f"  ✅ 콘텐츠 적용: {len(content_items)}개 항목")
                else:
                    # 콘텐츠가 없으면 해당 텍스트박스 비우기 (스타일 보존하며 내용 삭제)
                    logger.debug(f"  🗑️ 콘텐츠 없음, 텍스트박스 비우기: '{text_shapes[1]['original_text']}...'")
                    self._replace_text_preserving_style(content_shape, "")
            
            # Step 3: 나머지 텍스트박스 비우기 (사용하지 않는 영역)
            for i in range(2, len(text_shapes)):
                unused_shape = text_shapes[i]['shape']
                logger.debug(f"  🗑️ 미사용 텍스트박스 비우기: '{text_shapes[i]['original_text']}...'")
                # 스타일 보존하며 내용 삭제 (빈 문자열 적용)
                self._replace_text_preserving_style(unused_shape, "")
            
            # 로깅
            if not title_applied:
                logger.warning(f"  ⚠️ 슬라이드 {slide_index + 1}: 제목 적용할 텍스트박스를 찾지 못함")
            if content_items and not content_applied:
                # 콘텐츠가 있는데 적용할 박스가 없는 경우만 경고
                if len(text_shapes) <= 1:
                    logger.warning(f"  ⚠️ 슬라이드 {slide_index + 1}: 콘텐츠 적용할 텍스트박스를 찾지 못함 (박스 부족)")
                
        except Exception as e:
            logger.error(f"슬라이드 {slide_index + 1} 콘텐츠 적용 실패: {e}", exc_info=True)

    def _find_content_slide_index(self, prs: Presentation, template_metadata: Optional[Dict[str, Any]] = None) -> Optional[int]:
        """템플릿에서 content 타입 슬라이드 인덱스 찾기
        
        우선순위:
        1. template_metadata에서 role='content'인 슬라이드
        2. 레이아웃 이름에 'content', 'body' 포함된 슬라이드
        3. 가장 많은 텍스트박스를 가진 슬라이드 (title, thanks 제외)
        """
        try:
            # 1. template_metadata에서 찾기
            if template_metadata and 'slides' in template_metadata:
                for slide_info in template_metadata['slides']:
                    if slide_info.get('role') == 'content':
                        idx = slide_info.get('index', 0)
                        logger.info(f"📋 메타데이터에서 content 슬라이드 발견: index={idx}")
                        return idx
            
            # 2. 레이아웃 이름으로 찾기
            for idx, slide in enumerate(prs.slides):
                layout_name = slide.slide_layout.name.lower() if slide.slide_layout else ''
                if 'content' in layout_name or 'body' in layout_name:
                    logger.info(f"📋 레이아웃 이름으로 content 슬라이드 발견: index={idx}, layout='{layout_name}'")
                    return idx
            
            # 3. 텍스트박스 수로 찾기 (첫/마지막 슬라이드 제외)
            max_textbox_count = 0
            best_idx = None
            
            for idx, slide in enumerate(prs.slides):
                # 첫 번째(title)와 마지막(thanks) 슬라이드 제외
                if idx == 0 or idx == len(prs.slides) - 1:
                    continue
                
                textbox_count = sum(1 for shape in slide.shapes if hasattr(shape, 'text_frame'))
                if textbox_count > max_textbox_count:
                    max_textbox_count = textbox_count
                    best_idx = idx
            
            if best_idx is not None:
                logger.info(f"📋 텍스트박스 수로 content 슬라이드 발견: index={best_idx}, textbox_count={max_textbox_count}")
                return best_idx
            
            # 폴백: 두 번째 슬라이드 (index 1)
            if len(prs.slides) > 2:
                logger.info(f"📋 폴백: 두 번째 슬라이드(index=1)를 content로 사용")
                return 1
            
            return None
            
        except Exception as e:
            logger.error(f"content 슬라이드 찾기 실패: {e}")
            return None

    def _find_content_slide_in_indices(self, used_indices: List[int], template_metadata: Optional[Dict[str, Any]] = None) -> Optional[int]:
        """used_template_indices 중에서 content 타입 슬라이드 찾기"""
        try:
            if template_metadata and 'slides' in template_metadata:
                for slide_info in template_metadata['slides']:
                    idx = slide_info.get('index', 0)
                    if idx in used_indices and slide_info.get('role') == 'content':
                        logger.info(f"📋 used_indices에서 content 슬라이드 발견: index={idx}")
                        return idx
            
            # 메타데이터 없으면 첫 번째/마지막 제외한 인덱스 중 하나 반환
            for idx in used_indices:
                if idx != 0 and idx != max(used_indices):
                    return idx
            
            # 폴백: 첫 번째 인덱스 제외한 아무거나
            return used_indices[1] if len(used_indices) > 1 else used_indices[0]
            
        except Exception as e:
            logger.error(f"used_indices에서 content 슬라이드 찾기 실패: {e}")
            return used_indices[0] if used_indices else None

    def _duplicate_slide(self, prs: Presentation, source_idx: int) -> bool:
        """슬라이드 복제 (마지막 위치에 추가)
        
        python-pptx는 직접적인 슬라이드 복제를 지원하지 않으므로,
        소스 슬라이드의 레이아웃으로 새 슬라이드를 생성하고 내용을 복사합니다.
        """
        try:
            if source_idx >= len(prs.slides):
                logger.warning(f"⚠️ 잘못된 소스 인덱스: {source_idx}")
                return False
            
            source_slide = prs.slides[source_idx]
            
            # 같은 레이아웃으로 새 슬라이드 추가
            new_slide = prs.slides.add_slide(source_slide.slide_layout)
            
            # 소스 슬라이드의 shape들을 복사 (단순화된 복사)
            # 참고: 완벽한 복제는 복잡하므로, 레이아웃만 복제하고 내용은 AI가 채움
            logger.info(f"✅ 슬라이드 복제 완료: source={source_idx}, new_idx={len(prs.slides)-1}")
            
            return True
            
        except Exception as e:
            logger.error(f"슬라이드 복제 실패: {e}")
            return False

    def _add_blank_content_slide(self, prs: Presentation) -> bool:
        """빈 content 슬라이드 추가 (레이아웃 1번 사용)"""
        try:
            # 일반적으로 레이아웃 1은 'Title and Content'
            if len(prs.slide_layouts) > 1:
                layout = prs.slide_layouts[1]
            else:
                layout = prs.slide_layouts[0]
            
            prs.slides.add_slide(layout)
            logger.info(f"✅ 빈 content 슬라이드 추가 완료: new_idx={len(prs.slides)-1}")
            
            return True
            
        except Exception as e:
            logger.error(f"빈 슬라이드 추가 실패: {e}")
            return False


# 전역 인스턴스
templated_ppt_service = TemplatedPPTGeneratorService()
