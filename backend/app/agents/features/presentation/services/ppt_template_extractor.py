# PPT 템플릿 메타데이터 추출기
# v3.0 - 완전한 오브젝트 인식, 요소 역할 세분화, 시각화 스타일 정의
"""
주요 개선사항 (v3.0):
1. 모든 텍스트 포함 오브젝트 인식 (TEXT_BOX, AUTO_SHAPE, PLACEHOLDER, FREEFORM, GROUP 등)
2. 요소 역할 세분화 (title, key_message, body_content, bullet_item, caption, label 등)
3. 슬라이드 시각화 스타일 정의 (card_grid, timeline, process_flow, comparison, icon_boxes 등)
4. 완전한 스타일 정보 보존 (폰트, 크기, 색상, 정렬, 굵기, 기울임 등)
5. 그룹 내 텍스트 요소 재귀 추출
"""
import json
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pathlib import Path
from typing import Dict, List, Optional, Any, Tuple
from collections import defaultdict

EMU_PER_INCH = 914400
DPI = 96.0

# ============================================================================
# 슬라이드 역할 분류 키워드
# ============================================================================
SLIDE_ROLE_KEYWORDS = {
    "title": ["제목", "title", "표지", "cover", "시작"],
    "toc": ["목차", "contents", "index", "목록", "agenda", "순서"],
    "section": ["섹션", "section", "장", "chapter", "파트", "part"],
    "thanks": ["감사", "thank", "끝", "end", "q&a", "질문", "마무리", "contact", "watching"],
}

# ============================================================================
# 요소 역할 분류 (세분화)
# ============================================================================
ELEMENT_ROLE_PATTERNS = {
    "slide_title": {"position": "top", "font_size_min": 18, "width_ratio_min": 0.4},
    "key_message": {"position": "upper", "font_size_min": 14, "text_length_max": 100, "shape_types": ["AUTO_SHAPE"]},
    "body_content": {"position": "middle", "text_length_min": 50},
    "bullet_item": {"patterns": [r'^[\•\-\*\◦\▪\●\○]', r'^\d+[\.\)]\s']},
    "caption": {"font_size_max": 12, "text_length_max": 50},
    "label": {"text_length_max": 20, "shape_types": ["AUTO_SHAPE"]},
}

# 고정 요소 키워드 (편집하면 안 되는 디자인 요소)
FIXED_ELEMENT_KEYWORDS = ["logo", "회사", "company", "team name", "copyright", "©", "ⓒ"]

# ============================================================================
# 시각화 스타일 정의
# ============================================================================
VISUALIZATION_STYLES = {
    "card_grid": {"description": "카드 형태의 그리드 배열", "patterns": ["사각형: 둥근 모서리", "rounded rectangle"], "min_similar_shapes": 3},
    "numbered_cards": {"description": "번호가 매겨진 카드 형태", "patterns": [r"^\d{2}\n", r"^0[1-9]\n"], "min_count": 3},
    "icon_boxes": {"description": "아이콘과 텍스트 조합 박스", "patterns": ["🔹", "🔸", "💎", "⭐", "📱", "🔄", "🏥", "🛡️", "🔍", "📊"], "min_count": 2},
    "timeline": {"description": "타임라인/프로세스 흐름", "patterns": ["→", ">>", "➜", "step", "단계"], "min_count": 2},
    "comparison": {"description": "비교 레이아웃 (좌우 대비)", "layout_keywords": ["비교", "vs", "대비"]},
    "process_flow": {"description": "프로세스 플로우 다이어그램", "patterns": ["화살표", "arrow", "flow"]},
    "table_style": {"description": "표 형식 데이터", "has_table": True},
    "image_with_text": {"description": "이미지와 텍스트 조합", "has_image": True},
    "bullet_list": {"description": "불릿 목록 스타일", "patterns": [r'^[\•\-\*]', r'^\d+[\.\)]'], "min_count": 3},
    "simple_text": {"description": "단순 텍스트 레이아웃", "is_default": True},
}

EDITABLE_MIN_WIDTH = 80
EDITABLE_MIN_HEIGHT = 25


def emu_to_px(v):
    """EMU를 픽셀로 변환"""
    try:
        return round((v / EMU_PER_INCH) * DPI, 2)
    except Exception:
        return None


def color_to_hex(color_obj):
    """색상 객체를 16진수 색상 코드로 변환 (테마 색상 지원)"""
    if color_obj is None:
        return None
    try:
        if hasattr(color_obj, "rgb") and color_obj.rgb is not None:
            rgb = color_obj.rgb
            return "{:02X}{:02X}{:02X}".format(rgb[0], rgb[1], rgb[2])
        if hasattr(color_obj, "type") and hasattr(color_obj, "theme_color"):
            from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
            if color_obj.type == MSO_COLOR_TYPE.SCHEME:
                theme_color = color_obj.theme_color
                theme_colors = {
                    MSO_THEME_COLOR.BACKGROUND_1: "FFFFFF", MSO_THEME_COLOR.TEXT_1: "000000",
                    MSO_THEME_COLOR.BACKGROUND_2: "F2F2F2", MSO_THEME_COLOR.TEXT_2: "333333",
                    MSO_THEME_COLOR.ACCENT_1: "5B9BD5", MSO_THEME_COLOR.ACCENT_2: "70AD47",
                    MSO_THEME_COLOR.ACCENT_3: "A5A5A5", MSO_THEME_COLOR.ACCENT_4: "FFC000",
                    MSO_THEME_COLOR.ACCENT_5: "4472C4", MSO_THEME_COLOR.ACCENT_6: "C55911",
                }
                if theme_color in theme_colors:
                    return theme_colors[theme_color]
        if isinstance(color_obj, RGBColor):
            return "{:02X}{:02X}{:02X}".format(color_obj[0], color_obj[1], color_obj[2])
    except Exception:
        pass
    return None


def get_font_info(run) -> Dict:
    """폰트 정보 추출"""
    f = run.font
    return {
        "name": f.name,
        "size_pt": float(f.size.pt) if f.size is not None else None,
        "bold": f.bold,
        "italic": f.italic,
        "underline": f.underline,
        "color": color_to_hex(f.color) if hasattr(f, "color") and f.color else None,
    }


def get_paragraph_info(paragraph) -> Dict:
    """단락 정보 추출"""
    return {
        "level": paragraph.level,
        "alignment": getattr(paragraph.alignment, "name", None) if paragraph.alignment else None,
        "runs": [{"text": run.text, "font": get_font_info(run)} for run in paragraph.runs]
    }


def get_shape_fill_info(shape) -> Optional[Dict]:
    """Shape의 fill 정보 추출"""
    if not hasattr(shape, "fill"):
        return None
    try:
        fill = shape.fill
        fill_type = getattr(fill.type, "name", None) if fill.type is not None else None
        fore_color = None
        if fill.fore_color:
            fore_color = color_to_hex(getattr(fill.fore_color, "rgb", None))
        return {"type": fill_type, "fore_color": fore_color}
    except Exception:
        return None


def get_shape_line_info(shape) -> Optional[Dict]:
    """Shape의 line 정보 추출"""
    if not hasattr(shape, "line"):
        return None
    try:
        line = shape.line
        line_color = None
        if line and line.color:
            line_color = color_to_hex(getattr(line.color, "rgb", None))
        return {"width_pt": float(line.width.pt) if getattr(line, "width", None) else None, "color": line_color}
    except Exception:
        return None


def extract_text_from_shape(shape) -> Tuple[str, List[Dict]]:
    """Shape에서 텍스트와 단락 정보 추출"""
    if not hasattr(shape, "has_text_frame") or not shape.has_text_frame:
        return "", []
    tf = shape.text_frame
    paragraphs = [get_paragraph_info(p) for p in tf.paragraphs]
    return shape.text, paragraphs


def extract_table_data(table_shape) -> Dict:
    """테이블 Shape에서 셀 데이터 추출"""
    if not hasattr(table_shape, "table"):
        return None
    
    table = table_shape.table
    rows_data = []
    all_cell_texts = []
    
    for row_idx, row in enumerate(table.rows):
        row_data = []
        for col_idx, cell in enumerate(row.cells):
            cell_text = cell.text.strip() if cell.text else ""
            all_cell_texts.append(cell_text)
            
            # 셀의 폰트 정보 추출
            cell_font_info = {}
            if cell.text_frame and cell.text_frame.paragraphs:
                first_para = cell.text_frame.paragraphs[0]
                if first_para.runs:
                    cell_font_info = get_font_info(first_para.runs[0])
            
            row_data.append({
                "row": row_idx,
                "col": col_idx,
                "text": cell_text,
                "font": cell_font_info,
                "is_merged": getattr(cell, "is_merge_origin", False) or getattr(cell, "is_spanned", False)
            })
        rows_data.append(row_data)
    
    return {
        "rows": len(table.rows),
        "cols": len(table.columns),
        "cells": rows_data,
        "all_text": "\n".join([t for t in all_cell_texts if t]),
        "header_row": rows_data[0] if rows_data else []
    }


def extract_shapes_recursive(shape, slide_idx: int, counters: Dict, 
                             slide_width_px: float, slide_height_px: float,
                             parent_group: str = None) -> List[Dict]:
    """Shape에서 텍스트 요소를 재귀적으로 추출 (그룹, 테이블 내부 포함)"""
    elements = []
    shape_type_name = getattr(shape.shape_type, "name", str(shape.shape_type))
    
    # 그룹 Shape인 경우 내부 요소 재귀 처리
    if shape_type_name == "GROUP" and hasattr(shape, "shapes"):
        group_id = f"group-{slide_idx}-{counters['group']}"
        counters['group'] += 1
        for child_shape in shape.shapes:
            child_elements = extract_shapes_recursive(child_shape, slide_idx, counters, slide_width_px, slide_height_px, parent_group=group_id)
            elements.extend(child_elements)
        return elements
    
    # 테이블 Shape 처리
    if shape_type_name == "TABLE" and hasattr(shape, "table"):
        if 'table' not in counters:
            counters['table'] = 0
        element_id = f"table-{slide_idx}-{counters['table']}"
        counters['table'] += 1
        
        original_name = shape.name
        # 원본 파일 보존: shape.name을 변경하지 않음
        # shape.name = element_id  # REMOVED - 원본 보존
        
        position = {
            "left": emu_to_px(shape.left), "top": emu_to_px(shape.top),
            "width": emu_to_px(shape.width), "height": emu_to_px(shape.height),
            "rotation": getattr(shape, "rotation", 0),
        }
        
        # 테이블 데이터 추출
        table_data = extract_table_data(shape)
        
        # 헤더 행의 폰트 정보
        font_info = {}
        if table_data and table_data.get("header_row"):
            first_cell = table_data["header_row"][0] if table_data["header_row"] else {}
            cell_font = first_cell.get("font", {})
            font_info = {
                "fontFamily": cell_font.get("name"),
                "fontSize": cell_font.get("size_pt"),
                "fontColor": cell_font.get("color"),
                "fontBold": cell_font.get("bold"),
                "fontItalic": cell_font.get("italic"),
            }
        
        element = {
            "id": element_id, 
            "original_name": original_name, 
            "type": "table", 
            "shape_type": "TABLE",
            "content": table_data.get("all_text", "") if table_data else "",
            "position": position,
            "style": {**font_info, "alignment": None, "fill": None, "line": None},
            "table_data": table_data,
            "parent_group": parent_group,
        }
        elements.append(element)
        return elements
    
    # 텍스트가 있는 Shape 처리
    has_text = hasattr(shape, "text") and getattr(shape, "text", "").strip()
    processable_types = ["TEXT_BOX", "AUTO_SHAPE", "PLACEHOLDER", "FREEFORM", "CALLOUT", "CHEVRON", "OVAL", "RECTANGLE"]
    
    if has_text and (shape_type_name in processable_types or has_text):
        if shape_type_name == "TEXT_BOX":
            element_id = f"textbox-{slide_idx}-{counters['textbox']}"
            counters['textbox'] += 1
            element_type = "textbox"
        elif shape_type_name == "PLACEHOLDER":
            element_id = f"placeholder-{slide_idx}-{counters['placeholder']}"
            counters['placeholder'] += 1
            element_type = "placeholder"
        else:
            element_id = f"shape-{slide_idx}-{counters['shape']}"
            counters['shape'] += 1
            element_type = "auto_shape"
        
        original_name = shape.name
        # 원본 파일 보존: shape.name을 변경하지 않음
        # shape.name = element_id  # REMOVED - 원본 보존
        
        position = {
            "left": emu_to_px(shape.left), "top": emu_to_px(shape.top),
            "width": emu_to_px(shape.width), "height": emu_to_px(shape.height),
            "rotation": getattr(shape, "rotation", 0),
        }
        
        text_content, paragraphs = extract_text_from_shape(shape)
        
        font_info = {}
        if paragraphs and paragraphs[0].get("runs"):
            first_run_font = paragraphs[0]["runs"][0].get("font", {})
            font_info = {
                "fontFamily": first_run_font.get("name"),
                "fontSize": first_run_font.get("size_pt"),
                "fontColor": first_run_font.get("color"),
                "fontBold": first_run_font.get("bold"),
                "fontItalic": first_run_font.get("italic"),
            }
        
        alignment = paragraphs[0].get("alignment") if paragraphs else None
        
        element = {
            "id": element_id, "original_name": original_name, "type": element_type, "shape_type": shape_type_name,
            "content": text_content, "position": position,
            "style": {**font_info, "alignment": alignment, "fill": get_shape_fill_info(shape), "line": get_shape_line_info(shape)},
            "paragraphs": paragraphs, "parent_group": parent_group,
        }
        elements.append(element)
    
    return elements


def classify_slide_role(slide_index: int, total_slides: int, layout_name: str, slide_texts: List[str]) -> Tuple[str, float]:
    """슬라이드의 역할(role)을 분류"""
    combined_text = " ".join(slide_texts).lower()
    layout_lower = (layout_name or "").lower()

    if slide_index == 1:
        if any(kw in layout_lower or kw in combined_text for kw in SLIDE_ROLE_KEYWORDS["title"]):
            return ("title", 0.95)
        return ("title", 0.85)

    if slide_index == total_slides:
        if any(kw in layout_lower or kw in combined_text for kw in SLIDE_ROLE_KEYWORDS["thanks"]):
            return ("thanks", 0.95)
        return ("thanks", 0.70)

    for role, keywords in SLIDE_ROLE_KEYWORDS.items():
        if any(kw in layout_lower for kw in keywords):
            return (role, 0.90)
        if any(kw in combined_text for kw in keywords):
            return (role, 0.80)

    numbered_items = sum(1 for text in slide_texts if re.match(r'^\d+[\.\s]', text.strip()))
    if numbered_items >= 3:
        return ("toc", 0.85)

    if len(slide_texts) <= 2 and any(len(t) < 30 for t in slide_texts):
        total_text_len = sum(len(t) for t in slide_texts)
        if total_text_len < 50:
            return ("section", 0.70)

    return ("content", 0.90)


def classify_element_role(element: Dict, slide_role: str, slide_width_px: float, slide_height_px: float) -> str:
    """요소의 역할을 세분화하여 분류"""
    text = element.get("content", "").strip()
    text_lower = text.lower()
    position = element.get("position", {})
    shape_type = element.get("shape_type", "")
    element_type = element.get("type", "")
    font_size = element.get("style", {}).get("fontSize") or 0
    
    # 테이블 요소 처리
    if element_type == "table" or shape_type == "TABLE":
        table_data = element.get("table_data", {})
        rows = table_data.get("rows", 0) if table_data else 0
        cols = table_data.get("cols", 0) if table_data else 0
        
        # 테이블 역할 세분화
        if rows <= 2 and cols >= 3:
            return "spec_table"  # 사양/스펙 테이블
        elif rows >= 3 and cols == 2:
            return "comparison_table"  # 비교 테이블
        elif rows >= 3 and cols >= 3:
            return "data_table"  # 데이터 테이블
        else:
            return "info_table"  # 정보 테이블
    
    top = position.get("top", 0) or 0
    width = position.get("width", 0) or 0
    
    top_ratio = top / slide_height_px if slide_height_px else 0
    width_ratio = width / slide_width_px if slide_width_px else 0
    
    if slide_role == "title":
        if font_size and font_size >= 30:
            return "main_title"
        if top_ratio < 0.4 and width_ratio > 0.5:
            return "main_title"
        if top_ratio > 0.5:
            return "subtitle"
        return "metadata"
    
    if slide_role == "toc":
        if re.match(r'^\d+[\.\s]', text) or re.match(r'^0[1-9]$', text):
            return "toc_number"
        if "목차" in text_lower or "contents" in text_lower:
            return "slide_title"
        return "toc_item"
    
    if slide_role == "thanks":
        if any(kw in text_lower for kw in ["감사", "thank", "watching"]):
            return "thanks_message"
        return "contact_info"
    
    if top_ratio < 0.15:
        if (font_size and font_size >= 18) or width_ratio > 0.5:
            return "slide_title"
    
    if shape_type in ["AUTO_SHAPE", "FREEFORM", "CALLOUT"]:
        if top_ratio < 0.35 and len(text) < 100:
            return "key_message"
    
    if re.match(r'^0[1-9]\n', text) or re.match(r'^\d{2}\n', text):
        return "numbered_card"
    
    if re.match(r'^[\•\-\*\◦\▪\●\○🔹🔸💎⭐]', text):
        return "bullet_item"
    
    # 🆕 v3.2: 이모지+텍스트 카드 vs 순수 아이콘 구분
    # 이모지로 시작하는 경우: 실질적 텍스트 길이에 따라 분류
    emoji_start_pattern = r'^[📱🔄🏥🛡️🔍📊💡✅❌➜→🎯📈📋🚀💻🔒🌐⚡🔧🎨📞📧🏢👥💰📦🔬🎓🏆🌟💼📊📈📉📌📍🔗]'
    if re.match(emoji_start_pattern, text):
        # 이모지 제거 후 실질적 텍스트 길이 계산
        text_without_emoji = re.sub(r'[📱🔄🏥🛡️🔍📊💡✅❌➜→🎯📈📋🚀💻🔒🌐⚡🔧🎨📞📧🏢👥💰📦🔬🎓🏆🌟💼📊📈📉📌📍🔗\s]', '', text)
        if len(text_without_emoji) >= 15:
            # 실질적 콘텐츠가 있는 카드 → 편집 가능
            return "icon_card"
        else:
            # 아이콘만 있거나 짧은 라벨 → 고정
            return "icon_text"
    
    if font_size and font_size <= 12 and len(text) < 50:
        return "caption"
    
    if len(text) < 20 and shape_type in ["AUTO_SHAPE", "OVAL", "RECTANGLE"]:
        return "label"
    
    if len(text) > 50 or "\n" in text:
        return "body_content"
    
    return "content_item"


def is_fixed_element(element: Dict, slide_width_px: float, slide_height_px: float) -> Tuple[bool, str]:
    """요소가 고정 요소인지 판단 (v3.1: 빈 요소, 아이콘, placeholder 고정 처리)"""
    text = element.get("content", "").strip()
    text_lower = text.lower()
    position = element.get("position", {})
    element_role = element.get("element_role", "")
    original_name = element.get("original_name", "").lower()
    
    # 🆕 v3.1: 텍스트가 없는 요소는 고정 (장식용 도형)
    if not text:
        return (True, "fixed:empty_content")
    
    # 🆕 v3.1: 아이콘/이모지만 있는 요소 고정
    emoji_only = all(ord(c) > 127 or c in '→←↑↓↔•●○▶▷►◀◁◄' or c.isspace() for c in text)
    if emoji_only and len(text.strip()) <= 3:
        return (True, "fixed:icon_only")
    
    # 🆕 v3.1: 화살표, 특수문자만 있는 요소 고정
    special_chars_only = {'→', '←', '↑', '↓', '|', '/', '-', '•', '▶', '▷', '►', '◀', '◁', '◄', '»', '«', '>>', '<<'}
    if text.strip() in special_chars_only:
        return (True, "fixed:special_char")
    
    # 🆕 v3.1: placeholder 텍스트 고정
    placeholder_patterns = {'제품 이미지', '이미지', 'image', 'placeholder', '사진', '그림', 'photo', 'picture'}
    if text_lower.strip() in placeholder_patterns:
        return (True, "fixed:placeholder_text")
    
    # 🆕 v3.1: 장식용 도형 이름 패턴 (대괄호, 화살표 등)
    decoration_name_patterns = ['대괄호', '괄호', 'bracket', '화살표', 'arrow', '타원', 'ellipse', '선', 'line', 'connector']
    for pattern in decoration_name_patterns:
        if pattern in original_name:
            # 단, 내용이 있는 경우는 편집 가능으로 (예: 타원 안에 번호가 있는 경우)
            if len(text) > 3 and not emoji_only:
                break  # 편집 가능 검사 계속
            return (True, f"fixed:decoration:{pattern}")
    
    editable_roles = [
        "main_title", "subtitle", "slide_title", "key_message", "body_content", "bullet_item",
        "numbered_card", "toc_item", "toc_number", "thanks_message", "content_item",
        "spec_table", "comparison_table", "data_table", "info_table",  # 테이블 역할
        "icon_card",  # 🆕 v3.2: 아이콘+실질적 텍스트 카드 (편집 가능)
    ]
    # 🆕 v3.1: icon_text는 편집 가능에서 제외 (아이콘만 있거나 짧은 텍스트)
    if element_role in editable_roles:
        return (False, f"editable:role:{element_role}")
    
    # 🆕 v3.1: label 역할은 짧은 경우 고정 (도식 내 라벨)
    if element_role == 'label' and len(text) <= 15:
        return (True, "fixed:short_label")
    
    # 🆕 v3.1: icon_text 역할은 고정
    if element_role == 'icon_text':
        return (True, "fixed:icon_text")
    
    for keyword in FIXED_ELEMENT_KEYWORDS:
        if keyword in text_lower and len(text) <= 25:
            return (True, f"fixed:keyword:{keyword}")
    
    top = position.get("top", 0) or 0
    width = position.get("width", 0) or 0
    height = position.get("height", 0) or 0
    
    if top < slide_height_px * 0.02:
        if width < slide_width_px * 0.25 and height < 35:
            return (True, "fixed:position:header_small")
    
    if top > slide_height_px * 0.95:
        if height < 35:
            return (True, "fixed:position:footer")
    
    if re.match(r'^[\d]+$', text) and len(text) <= 2:
        return (True, "fixed:pattern:page_number")
    
    if width < EDITABLE_MIN_WIDTH * 0.3 and height < EDITABLE_MIN_HEIGHT * 0.3:
        if len(text) < 3:
            return (True, "fixed:size:tiny")
    
    return (False, "editable")


def detect_visualization_style(elements: List[Dict], has_table: bool, has_image: bool) -> Dict:
    """슬라이드의 시각화 스타일을 감지"""
    if not elements:
        return {"name": "empty", "description": "빈 슬라이드", "confidence": 1.0}
    
    texts = [e.get("content", "") for e in elements]
    shape_types = [e.get("shape_type", "") for e in elements]
    original_names = [e.get("original_name", "") for e in elements]
    
    style_scores = defaultdict(float)
    
    if has_table:
        style_scores["table_style"] = 0.9
    if has_image:
        style_scores["image_with_text"] = 0.7
    
    numbered_cards = sum(1 for t in texts if re.match(r'^0[1-9]\n', t) or re.match(r'^\d{2}\n', t))
    if numbered_cards >= 3:
        style_scores["numbered_cards"] = 0.9
    elif numbered_cards >= 2:
        style_scores["numbered_cards"] = 0.7
    
    icon_patterns = ["🔹", "🔸", "💎", "⭐", "📱", "🔄", "🏥", "🛡️", "🔍", "📊", "💡"]
    icon_count = sum(1 for t in texts if any(icon in t for icon in icon_patterns))
    if icon_count >= 3:
        style_scores["icon_boxes"] = 0.9
    elif icon_count >= 2:
        style_scores["icon_boxes"] = 0.7
    
    rounded_rect_count = sum(1 for name in original_names if "둥근 모서리" in name or "rounded" in name.lower())
    auto_shape_count = sum(1 for st in shape_types if st == "AUTO_SHAPE")
    if rounded_rect_count >= 3:
        style_scores["card_grid"] = 0.9
    elif auto_shape_count >= 4:
        style_scores["card_grid"] = 0.7
    
    arrow_patterns = ["→", ">>", "➜", "▶"]
    arrow_count = sum(1 for t in texts if any(p in t for p in arrow_patterns))
    if arrow_count >= 2:
        style_scores["timeline"] = 0.8
    
    bullet_patterns = [r'^[\•\-\*\◦\▪\●\○]', r'^\d+[\.\)]']
    bullet_count = sum(1 for t in texts if any(re.match(p, t) for p in bullet_patterns))
    if bullet_count >= 3:
        style_scores["bullet_list"] = 0.85
    
    if style_scores:
        best_style = max(style_scores, key=style_scores.get)
        best_score = style_scores[best_style]
        if best_score >= 0.5:
            return {
                "name": best_style,
                "description": VISUALIZATION_STYLES.get(best_style, {}).get("description", best_style),
                "confidence": best_score,
                "detected_patterns": dict(style_scores)
            }
    
    return {"name": "simple_text", "description": "단순 텍스트 레이아웃", "confidence": 0.5, "detected_patterns": {}}


def extract_shape_metadata(shape) -> Dict:
    """Shape의 기본 메타데이터 추출 (shapes 배열용)"""
    shape_type_name = getattr(shape.shape_type, "name", str(shape.shape_type))
    base = {
        "name": shape.name, "type": shape_type_name,
        "left_px": emu_to_px(shape.left), "top_px": emu_to_px(shape.top),
        "width_px": emu_to_px(shape.width), "height_px": emu_to_px(shape.height),
        "rotation_deg": getattr(shape, "rotation", None),
        "is_placeholder": getattr(shape, "is_placeholder", False),
    }

    if base["is_placeholder"] and hasattr(shape, "placeholder_format"):
        phf = shape.placeholder_format
        base["placeholder"] = {"type": getattr(getattr(phf, "type", None), "name", None), "idx": getattr(phf, "idx", None)}

    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        text_content, paragraphs = extract_text_from_shape(shape)
        base["text"] = {"raw": text_content, "paragraphs": paragraphs}

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            img = shape.image
            base["image"] = {"filename": Path(img.filename).name if img.filename else None, "content_type": img.content_type}
        except Exception:
            base["image"] = None

    if shape.shape_type == MSO_SHAPE_TYPE.TABLE and hasattr(shape, "table"):
        table = shape.table
        base["table"] = {"rows": len(table.rows), "cols": len(table.columns)}

    base["fill"] = get_shape_fill_info(shape)
    base["line"] = get_shape_line_info(shape)
    return base


def extract_template_metadata(path: str, output_path: str) -> str:
    """
    PPT 템플릿에서 메타데이터를 추출합니다. (v3.0)
    
    - 모든 텍스트 포함 오브젝트 인식
    - 요소 역할 세분화 (title, key_message, body_content 등)
    - 시각화 스타일 감지 (card_grid, numbered_cards, icon_boxes 등)
    - 완전한 스타일 정보 보존
    - 그룹 내 텍스트 요소 재귀 추출
    
    중요: 원본 템플릿 파일은 절대 수정하지 않습니다.
          메타데이터만 추출하여 JSON으로 저장합니다.
    """
    # 원본 파일을 직접 읽기 (수정하지 않음)
    prs = Presentation(path)
    slide_width_px = emu_to_px(prs.slide_width)
    slide_height_px = emu_to_px(prs.slide_height)
    total_slides = len(prs.slides)

    result = {
        "file": Path(path).name, "slide_width_px": slide_width_px, "slide_height_px": slide_height_px,
        "version": "3.0", "total_slides": total_slides, "slides": [],
        "summary": {
            "title_slides": 0, "toc_slides": 0, "section_slides": 0, "content_slides": 0, "thanks_slides": 0,
            "total_editable_elements": 0, "total_fixed_elements": 0, "visualization_styles": {}
        }
    }

    for idx, slide in enumerate(prs.slides, start=1):
        layout_name = getattr(slide.slide_layout, "name", None)
        slide_texts = []
        has_table = False
        has_image = False
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_texts.append(shape.text.strip())
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                has_table = True
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_image = True

        slide_role, role_confidence = classify_slide_role(idx, total_slides, layout_name, slide_texts)
        counters = {"textbox": 0, "placeholder": 0, "shape": 0, "group": 0}

        all_elements = []
        shapes_data = []
        
        # 먼저 shapes_recursive를 실행하여 shape.name을 표준화된 ID로 업데이트
        for shape in slide.shapes:
            elements = extract_shapes_recursive(shape, idx - 1, counters, slide_width_px, slide_height_px)
            all_elements.extend(elements)
        
        # 그 후 메타데이터 추출 (업데이트된 shape.name 사용)
        for shape in slide.shapes:
            shape_meta = extract_shape_metadata(shape)
            shapes_data.append(shape_meta)

        editable_elements = []
        fixed_elements = []
        
        for element in all_elements:
            element_role = classify_element_role(element, slide_role, slide_width_px, slide_height_px)
            element["element_role"] = element_role
            is_fixed, fixed_reason = is_fixed_element(element, slide_width_px, slide_height_px)
            element["is_fixed"] = is_fixed
            element["fixed_reason"] = fixed_reason
            
            if is_fixed:
                fixed_elements.append(element["id"])
                result["summary"]["total_fixed_elements"] += 1
            else:
                editable_elements.append(element["id"])
                result["summary"]["total_editable_elements"] += 1

        viz_style = detect_visualization_style(all_elements, has_table, has_image)

        slide_info = {
            "index": idx, "layout_name": layout_name, "shapes_count": len(slide.shapes),
            "role": slide_role, "role_confidence": role_confidence, "visualization_style": viz_style,
            "shapes": shapes_data, "elements": all_elements,
            "editable_elements": editable_elements, "fixed_elements": fixed_elements,
        }

        if slide_role == "title":
            result["summary"]["title_slides"] += 1
        elif slide_role == "toc":
            result["summary"]["toc_slides"] += 1
        elif slide_role == "section":
            result["summary"]["section_slides"] += 1
        elif slide_role == "thanks":
            result["summary"]["thanks_slides"] += 1
        else:
            result["summary"]["content_slides"] += 1

        style_name = viz_style.get("name", "unknown")
        result["summary"]["visualization_styles"][style_name] = result["summary"]["visualization_styles"].get(style_name, 0) + 1
        result["slides"].append(slide_info)

    # 원본 파일은 수정하지 않음 - prs.save() 호출 제거
    # 메타데이터에 element_id (textbox-0-0 등)가 저장되어 있고,
    # PPT 생성 시 메타데이터의 position 정보로 shape를 찾아 매핑함

    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2)

    print(f"✅ 메타데이터 추출 완료 (v3.0): {output_path}")
    print(f"   - 총 슬라이드: {total_slides}")
    print(f"   - 편집 가능 요소: {result['summary']['total_editable_elements']}")
    print(f"   - 고정 요소: {result['summary']['total_fixed_elements']}")
    print(f"   - 시각화 스타일: {result['summary']['visualization_styles']}")
    return output_path


# 하위 호환성
extract_presentation = extract_template_metadata


if __name__ == "__main__":
    import sys
    if len(sys.argv) >= 3:
        extract_template_metadata(sys.argv[1], sys.argv[2])
    else:
        print("Usage: python ppt_template_extractor.py <input.pptx> <output.json>")
