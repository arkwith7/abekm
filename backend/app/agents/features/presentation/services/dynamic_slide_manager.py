"""
Dynamic Slide Manager - 하이브리드 PPT 생성을 위한 동적 슬라이드 관리

핵심 기능:
1. TOC 동적 조정: 목차 항목 수에 따라 폰트/간격 자동 조정
2. 본문 슬라이드 복제/삭제: 콘텐츠 양에 따른 슬라이드 가감
3. 슬라이드 타입 분류: 고정(표지, TOC, 마무리) vs 가변(본문)

v1.0: 2025-12-11
"""

import logging
import copy
from typing import Any, Dict, List, Optional, Tuple
from enum import Enum

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.slide import Slide
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)


class SlideType(Enum):
    """슬라이드 유형 분류"""
    TITLE = "title"           # 표지
    TOC = "toc"              # 목차
    SECTION = "section"       # 섹션 구분
    CONTENT = "content"       # 본문 (복제 가능)
    THANKS = "thanks"         # 마무리
    UNKNOWN = "unknown"


class DynamicSlideManager:
    """
    동적 슬라이드 관리자
    
    - TOC 슬라이드의 항목 수 동적 조정
    - 본문 슬라이드 복제/삭제
    - 슬라이드 타입별 처리 분기
    """
    
    # TOC 폰트 크기 조정 규칙
    TOC_FONT_RULES = {
        # (최소 항목 수, 최대 항목 수): (폰트 크기 Pt, 줄간격 배수)
        (1, 6): (18, 1.5),      # 기본 크기
        (7, 9): (16, 1.3),      # 약간 축소
        (10, 12): (14, 1.2),    # 중간 축소
        (13, 15): (12, 1.1),    # 많이 축소
        (16, 20): (10, 1.0),    # 최소 크기
    }
    
    def __init__(self, template_metadata: Dict[str, Any]):
        """
        Args:
            template_metadata: 템플릿 분석 메타데이터
        """
        self.metadata = template_metadata
        self.slide_types = self._classify_slides()
        
    def _classify_slides(self) -> Dict[int, SlideType]:
        """메타데이터 기반 슬라이드 타입 분류"""
        slide_types = {}
        
        for slide in self.metadata.get('slides', []):
            idx = slide.get('index', 0)
            role = slide.get('role', 'unknown').lower()
            
            if role == 'title':
                slide_types[idx] = SlideType.TITLE
            elif role == 'toc':
                slide_types[idx] = SlideType.TOC
            elif role == 'section':
                slide_types[idx] = SlideType.SECTION
            elif role in ['content', 'body']:
                slide_types[idx] = SlideType.CONTENT
            elif role in ['thanks', 'closing', 'end']:
                slide_types[idx] = SlideType.THANKS
            else:
                slide_types[idx] = SlideType.CONTENT  # 기본값은 콘텐츠
                
        logger.info(f"📊 슬라이드 타입 분류: {len(slide_types)}개")
        for idx, stype in slide_types.items():
            logger.debug(f"  슬라이드 {idx}: {stype.value}")
            
        return slide_types
    
    def get_content_slide_indices(self) -> List[int]:
        """복제/삭제 가능한 본문 슬라이드 인덱스 목록"""
        return [
            idx for idx, stype in self.slide_types.items() 
            if stype == SlideType.CONTENT
        ]
    
    def get_toc_slide_index(self) -> Optional[int]:
        """TOC 슬라이드 인덱스 반환"""
        for idx, stype in self.slide_types.items():
            if stype == SlideType.TOC:
                return idx
        return None
    
    def calculate_toc_adjustments(
        self, 
        toc_item_count: int,
        original_toc_count: int
    ) -> Dict[str, Any]:
        """
        TOC 항목 수에 따른 조정 파라미터 계산
        
        Args:
            toc_item_count: 새로운 TOC 항목 수
            original_toc_count: 템플릿 원본 TOC 항목 수
            
        Returns:
            {
                "font_size": Pt(N),
                "line_spacing": float,
                "items_to_remove": int,  # 제거할 항목 수 (음수면 추가)
                "needs_adjustment": bool
            }
        """
        # 폰트 크기 결정
        font_size = Pt(18)
        line_spacing = 1.5
        
        for (min_items, max_items), (size, spacing) in self.TOC_FONT_RULES.items():
            if min_items <= toc_item_count <= max_items:
                font_size = Pt(size)
                line_spacing = spacing
                break
        else:
            # 20개 초과시 최소값 사용
            font_size = Pt(10)
            line_spacing = 1.0
        
        items_diff = original_toc_count - toc_item_count
        needs_adjustment = items_diff != 0 or toc_item_count > 6
        
        logger.info(
            f"📋 TOC 조정 계산: {original_toc_count}→{toc_item_count}개, "
            f"폰트={font_size.pt}pt, 줄간격={line_spacing}"
        )
        
        return {
            "font_size": font_size,
            "line_spacing": line_spacing,
            "items_to_remove": items_diff,
            "needs_adjustment": needs_adjustment,
            "toc_item_count": toc_item_count
        }
    
    def calculate_slide_adjustments(
        self,
        content_sections: List[Dict[str, Any]],
        template_content_count: int
    ) -> Dict[str, Any]:
        """
        콘텐츠 양에 따른 슬라이드 조정 계산
        
        Args:
            content_sections: AI가 생성한 콘텐츠 섹션 리스트
            template_content_count: 템플릿의 본문 슬라이드 수
            
        Returns:
            {
                "required_slides": int,
                "slides_to_add": int,     # 추가할 슬라이드 수
                "slides_to_remove": int,  # 제거할 슬라이드 수
                "slide_mapping": Dict[int, int]  # 콘텐츠→슬라이드 매핑
            }
        """
        required_slides = len(content_sections)
        
        slides_to_add = max(0, required_slides - template_content_count)
        slides_to_remove = max(0, template_content_count - required_slides)
        
        # 콘텐츠 인덱스 → 슬라이드 인덱스 매핑
        content_indices = self.get_content_slide_indices()
        slide_mapping = {}
        
        for i, section in enumerate(content_sections):
            if i < len(content_indices):
                slide_mapping[i] = content_indices[i]
            else:
                # 복제된 슬라이드용 (나중에 인덱스 부여)
                slide_mapping[i] = -1  # placeholder
        
        logger.info(
            f"📊 슬라이드 조정 계산: 필요={required_slides}, "
            f"추가={slides_to_add}, 삭제={slides_to_remove}"
        )
        
        return {
            "required_slides": required_slides,
            "slides_to_add": slides_to_add,
            "slides_to_remove": slides_to_remove,
            "slide_mapping": slide_mapping,
            "content_indices": content_indices
        }


class TOCAdjuster:
    """
    TOC 슬라이드 동적 조정기
    
    - 항목 수에 따른 폰트 크기 조정
    - 불필요한 항목 제거
    - 레이아웃 최적화
    """
    
    def __init__(self, slide: Slide, metadata: Dict[str, Any]):
        """
        Args:
            slide: python-pptx Slide 객체
            metadata: 해당 슬라이드의 메타데이터
        """
        self.slide = slide
        self.metadata = metadata
        self.toc_elements = self._find_toc_elements()
        
    def _find_toc_elements(self) -> List[Dict[str, Any]]:
        """TOC 관련 요소(toc_item, toc_number) 찾기"""
        toc_elements = []
        
        for elem in self.metadata.get('elements', []):
            role = elem.get('element_role', '')
            if role in ['toc_item', 'toc_number']:
                toc_elements.append(elem)
                
        return sorted(toc_elements, key=lambda x: x.get('id', ''))
    
    def get_toc_item_count(self) -> int:
        """현재 TOC 항목 수"""
        return len([e for e in self.toc_elements if e.get('element_role') == 'toc_item'])
    
    def adjust_for_items(
        self,
        new_items: List[str],
        font_size: Pt,
        line_spacing: float
    ) -> Dict[str, Any]:
        """
        새 TOC 항목에 맞게 슬라이드 조정
        
        Args:
            new_items: 새 TOC 항목 텍스트 리스트
            font_size: 적용할 폰트 크기
            line_spacing: 줄간격
            
        Returns:
            {"success": bool, "adjusted_count": int}
        """
        toc_item_elements = [
            e for e in self.toc_elements 
            if e.get('element_role') == 'toc_item'
        ]
        
        adjusted_count = 0
        
        # 각 TOC 항목 shape 처리
        for i, elem in enumerate(toc_item_elements):
            shape = self._find_shape_by_name(elem.get('original_name', ''))
            if not shape or not shape.has_text_frame:
                continue
            
            if i < len(new_items):
                # 항목이 있으면 텍스트 설정 + 폰트 조정
                self._set_text_with_style(
                    shape.text_frame,
                    new_items[i],
                    font_size,
                    line_spacing
                )
                adjusted_count += 1
            else:
                # 항목이 없으면 빈 텍스트로 설정 (제거 효과)
                self._clear_shape(shape)
        
        # TOC 번호도 조정
        toc_number_elements = [
            e for e in self.toc_elements 
            if e.get('element_role') == 'toc_number'
        ]
        
        for i, elem in enumerate(toc_number_elements):
            shape = self._find_shape_by_name(elem.get('original_name', ''))
            if not shape or not shape.has_text_frame:
                continue
            
            if i < len(new_items):
                # 번호 설정 (01, 02, ...)
                self._set_text_with_style(
                    shape.text_frame,
                    f"{i+1:02d}",
                    font_size,
                    line_spacing
                )
            else:
                self._clear_shape(shape)
        
        logger.info(f"✅ TOC 조정 완료: {adjusted_count}개 항목")
        return {"success": True, "adjusted_count": adjusted_count}
    
    def _find_shape_by_name(self, name: str):
        """이름으로 shape 찾기"""
        for shape in self.slide.shapes:
            if shape.name == name:
                return shape
        return None
    
    def _set_text_with_style(
        self, 
        text_frame, 
        text: str,
        font_size: Pt,
        line_spacing: float
    ):
        """텍스트 설정 + 스타일 적용"""
        if not text_frame.paragraphs:
            return
        
        para = text_frame.paragraphs[0]
        
        # 기존 스타일 보존하면서 텍스트 교체
        if para.runs:
            run = para.runs[0]
            run.text = text
            run.font.size = font_size
        else:
            para.clear()
            run = para.add_run()
            run.text = text
            run.font.size = font_size
        
        # 줄간격 설정
        para.line_spacing = line_spacing
    
    def _clear_shape(self, shape):
        """shape 내용 비우기 (삭제 효과)"""
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                para.clear()


class SlideReplicator:
    """
    슬라이드 복제/삭제 관리자
    
    - 본문 슬라이드 복제
    - 불필요한 슬라이드 삭제
    - 슬라이드 순서 관리
    """
    
    def __init__(self, presentation: Presentation):
        """
        Args:
            presentation: python-pptx Presentation 객체
        """
        self.prs = presentation
    
    def duplicate_slide(self, slide_index: int) -> int:
        """
        슬라이드 복제 (원본 다음 위치에 삽입)
        
        Args:
            slide_index: 복제할 슬라이드의 인덱스 (0-based)
            
        Returns:
            새 슬라이드의 인덱스
        """
        if slide_index >= len(self.prs.slides):
            raise ValueError(f"슬라이드 인덱스 초과: {slide_index}")
        
        # python-pptx에서 슬라이드 복제
        source_slide = self.prs.slides[slide_index]
        
        # 슬라이드 레이아웃 가져오기
        slide_layout = source_slide.slide_layout
        
        # 새 슬라이드 추가
        new_slide = self.prs.slides.add_slide(slide_layout)
        
        # 원본 슬라이드의 모든 shape 복사
        self._copy_shapes(source_slide, new_slide)
        
        # 슬라이드 위치 조정 (원본 다음으로)
        new_index = self._move_slide(len(self.prs.slides) - 1, slide_index + 1)
        
        logger.info(f"📋 슬라이드 복제: {slide_index} → {new_index}")
        return new_index
    
    def _copy_shapes(self, source_slide: Slide, target_slide: Slide):
        """
        슬라이드의 shape들을 복사
        
        주의: python-pptx는 완벽한 shape 복사가 어려움
        텍스트 프레임 위주로 복사
        """
        # 기존 shape 제거 (레이아웃에서 상속된 것 제외)
        # 실제로는 복잡해서, 텍스트만 복사하는 것이 안전
        
        for shape in source_slide.shapes:
            if shape.has_text_frame:
                # 대응하는 shape 찾기 (이름 기반)
                target_shape = None
                for ts in target_slide.shapes:
                    if ts.name == shape.name and ts.has_text_frame:
                        target_shape = ts
                        break
                
                if target_shape:
                    # 텍스트 복사
                    for i, para in enumerate(shape.text_frame.paragraphs):
                        if i < len(target_shape.text_frame.paragraphs):
                            target_para = target_shape.text_frame.paragraphs[i]
                            target_para.clear()
                            for run in para.runs:
                                new_run = target_para.add_run()
                                new_run.text = run.text
                                # 스타일 복사
                                if run.font.size:
                                    new_run.font.size = run.font.size
                                if run.font.bold is not None:
                                    new_run.font.bold = run.font.bold
    
    def _move_slide(self, from_index: int, to_index: int) -> int:
        """
        슬라이드 위치 이동
        
        python-pptx는 직접 이동 API가 없어서 XML 조작 필요
        현재는 맨 뒤에 추가되므로, 정확한 위치 이동은 추후 구현
        """
        # TODO: XML 조작으로 슬라이드 순서 변경
        # 현재는 맨 뒤에 추가된 상태로 유지
        return len(self.prs.slides) - 1
    
    def remove_slide(self, slide_index: int) -> bool:
        """
        슬라이드 삭제
        
        Args:
            slide_index: 삭제할 슬라이드 인덱스 (0-based)
            
        Returns:
            성공 여부
        """
        if slide_index >= len(self.prs.slides):
            logger.warning(f"⚠️ 삭제할 슬라이드 없음: {slide_index}")
            return False
        
        try:
            # python-pptx 슬라이드 삭제 (XML 조작)
            slide_id = self.prs.slides._sldIdLst[slide_index].rId
            self.prs.part.drop_rel(slide_id)
            del self.prs.slides._sldIdLst[slide_index]
            
            logger.info(f"🗑️ 슬라이드 삭제: {slide_index}")
            return True
            
        except Exception as e:
            logger.error(f"❌ 슬라이드 삭제 실패: {e}")
            return False
    
    def remove_slides_from_end(self, count: int) -> int:
        """
        뒤에서부터 N개 슬라이드 삭제 (Thanks 슬라이드 보존)
        
        Args:
            count: 삭제할 슬라이드 수
            
        Returns:
            실제 삭제된 수
        """
        removed = 0
        total_slides = len(self.prs.slides)
        
        # 마지막 슬라이드(Thanks)는 보존
        for i in range(count):
            # Thanks 슬라이드 직전까지만 삭제
            target_idx = total_slides - 2 - i  # -2: Thanks 보존
            if target_idx > 1:  # 최소 Title, TOC는 보존
                if self.remove_slide(target_idx):
                    removed += 1
                    total_slides -= 1
        
        return removed


def create_dynamic_manager(
    presentation: Presentation,
    metadata: Dict[str, Any]
) -> Tuple[DynamicSlideManager, SlideReplicator]:
    """
    동적 슬라이드 관리자 팩토리 함수
    
    Returns:
        (DynamicSlideManager, SlideReplicator) 튜플
    """
    manager = DynamicSlideManager(metadata)
    replicator = SlideReplicator(presentation)
    return manager, replicator
