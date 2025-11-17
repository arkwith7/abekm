"""
문서 처리 서비스 - 완전한 파이프라인 구현
"""

import os
import hashlib
import time
from pathlib import Path
from typing import Dict, Any, Optional
import logging

# Korean NLP Service 임포트
from app.services.core.korean_nlp_service import korean_nlp_service

logger = logging.getLogger(__name__)

class DocumentProcessorService:
    """
    문서 처리 서비스 - 설계 문서 기반 완전한 파이프라인 구현
    """
    
    def __init__(self):
        """문서 처리 서비스 초기화"""
        
        # 지원되는 파일 형식과 처리 함수 매핑
        self.supported_formats = {
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.pptx': self._process_pptx,
            '.xlsx': self._process_xlsx,
            '.txt': self._process_txt,
            '.hwp': self._process_hwp
        }
        
        logger.info("문서 처리 서비스 초기화 완료")

    async def process_document(
        self, 
        file_path: str, 
        file_name: str, 
        user_id: Optional[int] = None,
        container_id: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        문서 처리 메인 함수 - 완전한 파이프라인 구현
        설계 문서의 전체 워크플로우 실행
        """
        import time
        start_time = time.time()
        
        try:
            logger.info(f"문서 처리 시작: {file_name}")
            
            # 1단계: 파일 검증 및 기본 정보 수집
            validation_result = await self._validate_and_extract_file_info(file_path, file_name)
            if not validation_result["success"]:
                return validation_result
            
            file_info = validation_result["file_info"]
            
            # 2단계: 문서 내용 추출
            extraction_result = await self._extract_document_content(file_path, file_info["file_extension"])
            if not extraction_result["success"]:
                return {
                    "success": False,
                    "error": extraction_result["error"],
                    "file_info": file_info,
                    "processing_time": time.time() - start_time
                }
            
            raw_text = extraction_result["text"]
            document_metadata = extraction_result.get("metadata", {})
            
            # 3단계: 한국어 텍스트 전처리 및 분석
            logger.info("한국어 텍스트 처리 시작")
            korean_analysis = await korean_nlp_service.hybrid_process_korean_text(raw_text)
            
            # 4단계: 텍스트 청킹 및 임베딩 생성
            logger.info("텍스트 청킹 및 임베딩 생성 시작")
            chunk_embeddings = await korean_nlp_service.create_document_chunks_with_embeddings(raw_text)
            
            # 5단계: 문서 품질 검증
            quality_metrics = await self._assess_document_quality(raw_text, korean_analysis)
            
            # 6단계: 메타데이터 통합
            integrated_metadata = await self._integrate_document_metadata(
                file_info, document_metadata, korean_analysis, quality_metrics, user_id
            )
            
            # 7단계: 권한 및 컨테이너 분류 (컨테이너 ID가 지정된 경우 우선 사용)
            container_assignment = await self._assign_knowledge_container(
                integrated_metadata, korean_analysis, user_id, container_id
            )
            
            processing_time = time.time() - start_time
            
            # 최종 결과 구성
            result = {
                "success": True,
                "file_info": file_info,
                "content": {
                    "raw_text": raw_text,
                    "text_length": len(raw_text),
                    "chunk_count": len(chunk_embeddings),
                    "chunks_with_embeddings": chunk_embeddings
                },
                "metadata": integrated_metadata,
                "korean_analysis": korean_analysis,
                "quality_metrics": quality_metrics,
                "container_assignment": container_assignment,
                "processing_time": processing_time,
                "pipeline_status": {
                    "content_extraction": "success",
                    "korean_processing": korean_analysis.get("success", False),
                    "chunking_embedding": len(chunk_embeddings) > 0,
                    "quality_assessment": quality_metrics.get("overall_score", 0) > 0,
                    "container_assignment": container_assignment.get("assigned", False)
                }
            }
            
            logger.info(f"문서 처리 완료: {file_name}, 소요시간: {processing_time:.2f}초")
            return result
            
        except Exception as e:
            processing_time = time.time() - start_time
            logger.error(f"문서 처리 실패: {file_name}, 에러: {e}")
            return {
                "success": False,
                "error": str(e),
                "file_info": {"file_name": file_name, "file_path": file_path},
                "processing_time": processing_time
            }

    async def _validate_and_extract_file_info(self, file_path: str, file_name: str) -> Dict[str, Any]:
        """파일 검증 및 기본 정보 추출"""
        try:
            # 파일 확장자 확인
            file_ext = Path(file_name).suffix.lower()
            if file_ext not in self.supported_formats:
                return {
                    "success": False,
                    "error": f"지원하지 않는 파일 형식: {file_ext}",
                    "supported_formats": list(self.supported_formats.keys())
                }
            
            # 파일 존재 확인
            if not os.path.exists(file_path):
                return {
                    "success": False,
                    "error": f"파일을 찾을 수 없습니다: {file_path}"
                }
            
            # 파일 기본 정보
            file_stats = os.stat(file_path)
            file_hash = await self._calculate_file_hash(file_path)
            
            file_info = {
                "file_name": file_name,
                "file_path": file_path,
                "file_size": file_stats.st_size,
                "file_extension": file_ext,
                "modified_time": file_stats.st_mtime,
                "file_hash": file_hash,
                "processing_strategy": self._determine_processing_strategy(file_stats.st_size)
            }
            
            return {
                "success": True,
                "file_info": file_info
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": f"파일 검증 실패: {str(e)}"
            }

    def _determine_processing_strategy(self, file_size: int) -> str:
        """파일 크기에 따른 처리 전략 결정"""
        if file_size < 5 * 1024 * 1024:  # 5MB 미만
            return "immediate_small"
        elif file_size < 20 * 1024 * 1024:  # 20MB 미만
            return "immediate_medium"
        elif file_size < 50 * 1024 * 1024:  # 50MB 미만
            return "background_large"
        else:
            return "background_xlarge"

    async def _extract_document_content(self, file_path: str, file_ext: str) -> Dict[str, Any]:
        """문서 내용 추출"""
        try:
            processor = self.supported_formats[file_ext]
            return await processor(file_path)
        except Exception as e:
            return {
                "success": False,
                "error": f"문서 내용 추출 실패: {str(e)}"
            }

    async def _calculate_file_hash(self, file_path: str) -> str:
        """파일 해시 계산 (SHA-256)"""
        hash_sha256 = hashlib.sha256()
        try:
            with open(file_path, "rb") as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    hash_sha256.update(chunk)
            return hash_sha256.hexdigest()
        except Exception as e:
            logger.error(f"파일 해시 계산 실패: {e}")
            return ""

    async def _process_pdf(self, file_path: str) -> Dict[str, Any]:
        """PDF 파일 처리 - OCR 포함"""
        try:
            import fitz  # PyMuPDF
            from PIL import Image
            import io
            
            doc = fitz.open(file_path)
            text_content = []
            metadata = {}
            
            # 메타데이터 추출
            metadata.update({
                "title": doc.metadata.get("title", ""),
                "author": doc.metadata.get("author", ""),
                "creation_date": doc.metadata.get("creationDate", ""),
                "modification_date": doc.metadata.get("modDate", ""),
                "page_count": doc.page_count
            })
            
            # 페이지별 텍스트 추출
            for page_num in range(doc.page_count):
                page = doc[page_num]
                
                # 텍스트 추출 시도
                page_text = page.get_text()
                
                # 텍스트가 부족한 경우 OCR 수행
                if len(page_text.strip()) < 50:
                    try:
                        # 페이지를 이미지로 변환
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2배 확대
                        img_data = pix.tobytes("png")
                        img = Image.open(io.BytesIO(img_data))
                        
                        # OCR 수행
                        ocr_text = await self._perform_ocr(img)
                        page_text = ocr_text if len(ocr_text) > len(page_text) else page_text
                        
                    except Exception as ocr_error:
                        logger.warning(f"OCR 실패 (페이지 {page_num}): {ocr_error}")
                
                text_content.append(page_text)
            
            doc.close()
            
            final_text = "\n".join(text_content)
            
            return {
                "success": True,
                "text": final_text,
                "metadata": metadata
            }
            
        except Exception as e:
            logger.error(f"PDF 처리 실패: {e}")
            return {
                "success": False,
                "error": f"PDF 처리 중 오류 발생: {str(e)}"
            }

    async def _process_docx(self, file_path: str) -> Dict[str, Any]:
        """DOCX 파일 처리"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            text_content = []
            
            # 메타데이터 추출
            metadata = {
                "title": doc.core_properties.title or "",
                "author": doc.core_properties.author or "",
                "creation_date": str(doc.core_properties.created) if doc.core_properties.created else "",
                "modification_date": str(doc.core_properties.modified) if doc.core_properties.modified else "",
                "paragraph_count": len(doc.paragraphs)
            }
            
            # 텍스트 추출
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # 테이블 텍스트 추출
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    if row_text.strip():
                        text_content.append(row_text)
            
            final_text = "\n".join(text_content)
            
            return {
                "success": True,
                "text": final_text,
                "metadata": metadata
            }
            
        except Exception as e:
            logger.error(f"DOCX 처리 실패: {e}")
            return {
                "success": False,
                "error": f"DOCX 처리 중 오류 발생: {str(e)}"
            }

    async def _process_pptx(self, file_path: str) -> Dict[str, Any]:
        """PPTX 파일 처리"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text_content = []
            
            # 메타데이터 추출
            metadata = {
                "title": prs.core_properties.title or "",
                "author": prs.core_properties.author or "",
                "creation_date": str(prs.core_properties.created) if prs.core_properties.created else "",
                "modification_date": str(prs.core_properties.modified) if prs.core_properties.modified else "",
                "slide_count": len(prs.slides)
            }
            
            # 슬라이드별 텍스트 추출
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"=== 슬라이드 {slide_num} ==="]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                text_content.extend(slide_text)
            
            final_text = "\n".join(text_content)
            
            return {
                "success": True,
                "text": final_text,
                "metadata": metadata
            }
            
        except Exception as e:
            logger.error(f"PPTX 처리 실패: {e}")
            return {
                "success": False,
                "error": f"PPTX 처리 중 오류 발생: {str(e)}"
            }

    async def _process_xlsx(self, file_path: str) -> Dict[str, Any]:
        """XLSX 파일 처리"""
        try:
            import pandas as pd
            
            # 모든 시트 읽기
            excel_file = pd.ExcelFile(file_path)
            text_content = []
            
            metadata = {
                "sheet_count": len(excel_file.sheet_names),
                "sheet_names": excel_file.sheet_names
            }
            
            for sheet_name in excel_file.sheet_names:
                try:
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    
                    # 시트별 텍스트 구성
                    sheet_text = [f"=== 시트: {sheet_name} ==="]
                    
                    # 컬럼명 추가
                    if not df.empty:
                        sheet_text.append("컬럼: " + " | ".join(str(col) for col in df.columns))
                        
                        # 데이터 행 추가 (최대 100행)
                        for _, row in df.head(100).iterrows():
                            row_text = " | ".join(str(val) for val in row.values if pd.notna(val))
                            if row_text.strip():
                                sheet_text.append(row_text)
                    
                    text_content.extend(sheet_text)
                    
                except Exception as sheet_error:
                    logger.warning(f"시트 '{sheet_name}' 처리 중 오류: {sheet_error}")
            
            final_text = "\n".join(text_content)
            
            return {
                "success": True,
                "text": final_text,
                "metadata": metadata
            }
            
        except Exception as e:
            logger.error(f"XLSX 처리 실패: {e}")
            return {
                "success": False,
                "error": f"XLSX 처리 중 오류 발생: {str(e)}"
            }

    async def _process_txt(self, file_path: str) -> Dict[str, Any]:
        """TXT 파일 처리"""
        try:
            # 인코딩 자동 감지
            encodings = ['utf-8', 'cp949', 'euc-kr', 'utf-16']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    break
                except UnicodeDecodeError:
                    continue
            else:
                # 모든 인코딩 실패시 바이너리로 읽고 처리
                with open(file_path, 'rb') as f:
                    raw_data = f.read()
                content = raw_data.decode('utf-8', errors='ignore')
            
            metadata = {
                "encoding": encoding if 'encoding' in locals() else 'binary',
                "line_count": len(content.split('\n')),
                "char_count": len(content)
            }
            
            return {
                "success": True,
                "text": content,
                "metadata": metadata
            }
            
        except Exception as e:
            logger.error(f"TXT 처리 실패: {e}")
            return {
                "success": False,
                "error": f"TXT 처리 중 오류 발생: {str(e)}"
            }

    async def _process_hwp(self, file_path: str) -> Dict[str, Any]:
        """HWP 파일 처리 (한글 문서)"""
        try:
            # HWP 파일 처리를 위한 라이브러리 시도
            try:
                import olefile
                
                if olefile.isOleFile(file_path):
                    ole = olefile.OleFileIO(file_path)
                    
                    # HWP 문서 구조에서 텍스트 추출 시도
                    # 실제 구현은 HWP 문서 구조에 따라 복잡할 수 있음
                    text_streams = []
                    
                    # HWP 스트림 탐색
                    for stream_name in ole.listdir():
                        if 'BodyText' in stream_name or 'text' in str(stream_name).lower():
                            try:
                                stream = ole.opendir(stream_name)
                                # 기본적인 텍스트 추출 로직
                                # 실제로는 HWP 포맷 파싱이 필요
                            except:
                                continue
                    
                    ole.close()
                
            except ImportError:
                # olefile이 없는 경우 대체 방법
                pass
            
            # HWP 처리가 실패한 경우 경고 메시지와 함께 빈 결과 반환
            return {
                "success": True,
                "text": "",
                "metadata": {
                    "file_format": "hwp",
                    "extraction_method": "limited",
                    "note": "HWP 파일 처리가 제한적입니다. 전용 라이브러리가 필요합니다."
                }
            }
            
        except Exception as e:
            logger.error(f"HWP 처리 실패: {e}")
            return {
                "success": False,
                "error": f"HWP 처리 중 오류 발생: {str(e)}"
            }

    async def _perform_ocr(self, image) -> str:
        """OCR 수행 (이미지에서 텍스트 추출)"""
        try:
            import pytesseract
            
            # 한국어 + 영어 OCR
            custom_config = r'--oem 3 --psm 6 -l kor+eng'
            text = pytesseract.image_to_string(image, config=custom_config)
            
            return text.strip()
            
        except ImportError:
            logger.warning("pytesseract가 설치되지 않았습니다. OCR를 건너뜁니다.")
            return ""
        except Exception as e:
            logger.error(f"OCR 처리 실패: {e}")
            return ""

    async def _assess_document_quality(self, text: str, korean_analysis: Dict[str, Any]) -> Dict[str, Any]:
        """문서 품질 평가"""
        try:
            metrics = {
                "text_length": len(text),
                "word_count": len(text.split()),
                "line_count": len(text.split('\n')),
                "korean_ratio": 0.0,
                "content_density": 0.0,
                "readability_score": 0.0,
                "overall_score": 0.0
            }
            
            # 한국어 비율 계산
            korean_chars = len([c for c in text if '\uac00' <= c <= '\ud7af'])
            total_chars = len([c for c in text if c.isalnum()])
            metrics["korean_ratio"] = korean_chars / total_chars if total_chars > 0 else 0.0
            
            # 내용 밀도 (문장당 평균 단어 수)
            sentences = text.split('.')
            metrics["content_density"] = metrics["word_count"] / len(sentences) if sentences else 0.0
            
            # 가독성 점수 (간단한 휴리스틱)
            if korean_analysis.get("success", False):
                keywords_count = len(korean_analysis.get("keywords", []))
                proper_nouns_count = len(korean_analysis.get("proper_nouns", []))
                metrics["readability_score"] = min(10.0, (keywords_count + proper_nouns_count) / 10.0)
            
            # 전체 품질 점수 계산
            metrics["overall_score"] = (
                min(10.0, metrics["text_length"] / 1000) * 0.3 +  # 텍스트 길이
                metrics["korean_ratio"] * 10 * 0.3 +               # 한국어 비율
                min(10.0, metrics["content_density"] / 10) * 0.2 + # 내용 밀도
                metrics["readability_score"] * 0.2                 # 가독성
            )
            
            return metrics
            
        except Exception as e:
            logger.error(f"품질 평가 실패: {e}")
            return {"overall_score": 0.0, "error": str(e)}

    async def _integrate_document_metadata(self, file_info: Dict, doc_metadata: Dict, 
                                         korean_analysis: Dict, quality_metrics: Dict,
                                         user_id: Optional[int] = None) -> Dict[str, Any]:
        """문서 메타데이터 통합"""
        try:
            integrated = {
                # 파일 기본 정보
                "file_name": file_info["file_name"],
                "file_size": file_info["file_size"],
                "file_extension": file_info["file_extension"],
                "file_hash": file_info["file_hash"],
                
                # 문서 메타데이터
                "title": doc_metadata.get("title", file_info["file_name"]),
                "author": doc_metadata.get("author", ""),
                "creation_date": doc_metadata.get("creation_date", ""),
                "modification_date": doc_metadata.get("modification_date", ""),
                
                # 한국어 분석 결과
                "document_type": korean_analysis.get("document_type", "unknown"),
                "keywords": korean_analysis.get("keywords", []),
                "proper_nouns": korean_analysis.get("proper_nouns", []),
                "company_names": korean_analysis.get("company_names", []),
                
                # 품질 메트릭
                "quality_score": quality_metrics.get("overall_score", 0.0),
                "text_length": quality_metrics.get("text_length", 0),
                "korean_ratio": quality_metrics.get("korean_ratio", 0.0),
                
                # 처리 정보
                "processing_timestamp": time.time(),
                "uploaded_by": user_id,
                "processing_strategy": file_info.get("processing_strategy", "immediate_small")
            }
            
            return integrated
            
        except Exception as e:
            logger.error(f"메타데이터 통합 실패: {e}")
            return {"error": str(e)}

    async def _assign_knowledge_container(
        self, 
        metadata: Dict, 
        korean_analysis: Dict, 
        user_id: Optional[int] = None,
        container_id: Optional[str] = None
    ) -> Dict[str, Any]:
        """지식 컨테이너 할당 - 지정된 컨테이너 우선 사용"""
        try:
            # 컨테이너 ID가 명시적으로 지정된 경우 우선 사용
            if container_id:
                return {
                    "assigned": True,
                    "container_id": container_id,
                    "assignment_method": "user_specified",
                    "confidence": 1.0,
                    "reason": f"사용자가 지정한 컨테이너: {container_id}"
                }
            
            # 기존 자동 할당 로직
            doc_type = korean_analysis.get("document_type", "general")
            keywords = korean_analysis.get("keywords", [])
            company_names = korean_analysis.get("company_names", [])
            
            # 문서 유형별 컨테이너 분류
            container_mappings = {
                "정책": "policy_documents",
                "절차": "procedure_documents", 
                "매뉴얼": "manual_documents",
                "보고서": "report_documents",
                "회계": "finance_documents",
                "인사": "hr_documents",
                "기술": "technical_documents",
                "마케팅": "marketing_documents"
            }
            
            assigned_container = "general_documents"  # 기본값
            
            # 키워드 기반 컨테이너 분류
            for category, container in container_mappings.items():
                if any(category in keyword for keyword in keywords):
                    assigned_container = container
                    break
            
            # 회사명 기반 접근 권한 설정
            access_level = "internal"  # 기본값
            if any("웅진" in name for name in company_names):
                access_level = "company_wide"
            elif user_id:
                access_level = "department"
            
            return {
                "assigned": True,
                "container_id": assigned_container,
                "access_level": access_level,
                "auto_assigned": True,
                "assignment_reason": f"문서유형: {doc_type}, 키워드 기반 분류",
                "requires_review": metadata.get("quality_score", 0) < 5.0
            }
            
        except Exception as e:
            logger.error(f"컨테이너 할당 실패: {e}")
            return {
                "assigned": False,
                "container_id": "general_documents",
                "access_level": "internal",
                "error": str(e)
            }


# 싱글톤 인스턴스 생성
document_processor_service = DocumentProcessorService()
