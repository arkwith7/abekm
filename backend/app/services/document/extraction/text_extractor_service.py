"""
📄 텍스트 추출 서비스
===================

다양한 파일 형식에서 텍스트를 추출하는 서비스
- PDF, DOCX, TXT, HWP 등 지원
- 한국어 인코딩 처리
- 메타데이터 추출
"""

import os
import logging
from typing import Dict, Any, Optional
from pathlib import Path

from app.core.config import settings

logger = logging.getLogger(__name__)

class TextExtractorService:
    """다양한 파일 형식에서 텍스트를 추출하는 서비스"""
    
    def __init__(self):
        self.supported_extensions = {'.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.xml'}
        
    async def extract_text_from_file(self, file_path: str, file_extension: str = None) -> Dict[str, Any]:
        """
        파일에서 텍스트를 추출합니다.
        
        Args:
            file_path: 파일 경로
            file_extension: 파일 확장자 (None인 경우 파일 경로에서 추출)
            
        Returns:
            Dict containing extracted text and metadata
        """
        if file_extension is None:
            file_extension = Path(file_path).suffix
            
        return await self.extract_text(file_path, file_extension)
    
    async def extract_text(self, file_path: str, file_extension: str) -> Dict[str, Any]:
        """
        파일에서 텍스트를 추출합니다. (내부 메서드)
        
        Args:
            file_path: 파일 경로
            file_extension: 파일 확장자
            
        Returns:
            Dict containing extracted text and metadata
        """
        try:
            result = {
                "text": "",
                "metadata": {},
                "success": True,
                "error": None,
                "text_length": 0,
                "encoding": "utf-8"
            }
            
            # Azure Blob 경로인지 확인 (raw/ 또는 processed/로 시작)
            is_blob_path = file_path.startswith('raw/') or file_path.startswith('processed/')
            actual_file_path = file_path
            temp_file_path = None
            
            try:
                if is_blob_path:
                    # Azure Blob에서 임시 파일로 다운로드
                    from app.services.core.azure_blob_service import azure_blob_service
                    import tempfile
                    
                    logger.info(f"📥 Azure Blob에서 파일 다운로드 시작: {file_path}")
                    
                    # 임시 파일 생성
                    file_ext = os.path.splitext(file_path)[1]
                    temp_fd, temp_file_path = tempfile.mkstemp(suffix=file_ext)
                    os.close(temp_fd)
                    
                    # Blob에서 다운로드 (동기 메서드)
                    blob_data = azure_blob_service.download_blob_to_bytes(file_path, purpose='raw')
                    if not blob_data:
                        raise Exception(f"Azure Blob 다운로드 실패: {file_path}")
                    
                    # 임시 파일에 저장
                    with open(temp_file_path, 'wb') as f:
                        f.write(blob_data)
                    
                    actual_file_path = temp_file_path
                    logger.info(f"✅ Azure Blob 다운로드 완료: {file_path} → {temp_file_path}")
                
                if not os.path.exists(actual_file_path):
                    result["success"] = False
                    result["error"] = "파일을 찾을 수 없습니다."
                    return result
                
                # 파일 확장자별 처리
                if file_extension.lower() in self.supported_extensions:
                    result = await self._extract_text_file(actual_file_path, result)
                elif file_extension.lower() == '.pdf':
                    result = await self._extract_pdf_file(actual_file_path, result)
                elif file_extension.lower() in ['.docx', '.doc']:
                    result = await self._extract_docx_file(actual_file_path, result)
                elif file_extension.lower() in ['.pptx', '.ppt']:
                    result = await self._extract_pptx_file(actual_file_path, result)
                elif file_extension.lower() in ['.xlsx', '.xls']:
                    result = await self._extract_excel_file(actual_file_path, result)
                elif file_extension.lower() in ['.hwp', '.hwpx']:
                    result = await self._extract_hwp_file(actual_file_path, result)
                else:
                    result["text"] = f"지원하지 않는 파일 형식입니다: {file_extension}"
                    result["success"] = False
                    result["error"] = "Unsupported file format"
                
                # 텍스트 길이 계산
                result["text_length"] = len(result["text"])
                
                # 기본 메타데이터 추가
                file_stats = os.stat(actual_file_path)
                result["metadata"].update({
                    "file_size": file_stats.st_size,
                    "last_modified": file_stats.st_mtime,
                    "extraction_method": self._get_extraction_method(file_extension)
                })
                
                # 실제 파일 경로 추가 (멀티모달 파이프라인에서 이미지 추출용)
                result["actual_file_path"] = actual_file_path
                result["is_temp_file"] = temp_file_path is not None
                
                logger.info(f"텍스트 추출 완료 - 파일: {file_path}, 길이: {result['text_length']}자")
                return result
            
            finally:
                # 임시 파일 정리 (extraction_result 반환 후 multimodal pipeline에서 처리)
                # 멀티모달 파이프라인이 이미지를 추출할 수 있도록 여기서는 삭제하지 않음
                # multimodal_document_service에서 처리 완료 후 정리
                if temp_file_path and os.path.exists(temp_file_path):
                    logger.info(f"🔄 임시 파일 유지 (멀티모달 처리용): {temp_file_path}")
                    # os.remove는 호출자에게 위임
            
        except Exception as e:
            logger.error(f"텍스트 추출 실패 - 파일: {file_path}, 오류: {e}")
            return {
                "text": "",
                "metadata": {},
                "success": False,
                "error": str(e),
                "text_length": 0,
                "encoding": "utf-8"
            }
    
    async def _extract_text_file(self, file_path: str, result: Dict[str, Any]) -> Dict[str, Any]:
        """텍스트 파일 추출"""
        try:
            # 여러 인코딩 시도
            encodings = ['utf-8', 'cp949', 'euc-kr', 'latin-1']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        result["text"] = f.read()
                        result["encoding"] = encoding
                        break
                except UnicodeDecodeError:
                    continue
            else:
                # 모든 인코딩 실패시 바이너리로 읽어서 에러 무시
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    result["text"] = f.read()
                    result["encoding"] = "utf-8 (with errors ignored)"
            
        except Exception as e:
            result["success"] = False
            result["error"] = f"텍스트 파일 읽기 실패: {str(e)}"
        
        return result
    
    async def _extract_pdf_file(self, file_path: str, result: Dict[str, Any]) -> Dict[str, Any]:
        """PDF 파일 텍스트 추출 - Provider 기반 라우팅 + Fallback 로직"""
        
        provider = settings.document_processing_provider.lower().strip()
        fallback_provider = settings.document_processing_fallback.lower().strip() if settings.document_processing_fallback else None
        
        logger.info(f"📄 [PDF-EXTRACT] 문서 처리 Provider: {provider} (Fallback: {fallback_provider or 'None'})")
        logger.info(f"📄 [PDF-EXTRACT] 파일: {file_path}")
        
        # Primary Provider 시도
        primary_success = False
        
        # Azure Document Intelligence
        if provider == "azure_di":
            try:
                from .azure_document_intelligence_service import azure_document_intelligence_service
                
                logger.info(f"Azure Document Intelligence로 PDF 분석 시도: {file_path}")
                di_result = await azure_document_intelligence_service.analyze_pdf(file_path)
                
                if di_result.success:
                    logger.info(f"✅ Azure DI 성공: {file_path}")
                    converted_result = azure_document_intelligence_service.create_internal_extraction_result(di_result)
                    result.update(converted_result)
                    return result
                else:
                    logger.warning(f"⚠️ Azure DI 실패: {di_result.error}")
                    
            except Exception as e:
                logger.warning(f"⚠️ Azure DI 예외: {e}")
        
        # Upstage Document Parse
        elif provider == "upstage":
            try:
                logger.info(f"🔷 [UPSTAGE] Upstage Document Parse 사용 - 파일: {file_path}")
                from .upstage_document_service import upstage_document_service
                
                logger.info(f"🔷 [UPSTAGE] upstage_document_service 모듈 로드 완료")
                logger.info(f"🔷 [UPSTAGE] API 키 설정 여부: {bool(upstage_document_service.api_key)}")
                
                logger.info(f"🔷 [UPSTAGE] Document Parse 호출 시작: {file_path}")
                upstage_result = await upstage_document_service.parse_document(file_path)
                
                logger.info(f"🔷 [UPSTAGE] Document Parse 호출 완료 - success: {upstage_result.success}")
                
                if upstage_result.success:
                    logger.info(f"✅ [UPSTAGE] Upstage 성공: {file_path}")
                    logger.info(f"✅ [UPSTAGE] 추출된 텍스트 길이: {len(upstage_result.text)}")
                    logger.info(f"✅ [UPSTAGE] 페이지 수: {len(upstage_result.pages)}")
                    logger.info(f"✅ [UPSTAGE] 테이블 수: {len(upstage_result.tables)}")
                    logger.info(f"✅ [UPSTAGE] 이미지 수: {len(upstage_result.figures)}")
                    
                    converted_result = upstage_document_service.create_internal_extraction_result(upstage_result)
                    result.update(converted_result)
                    primary_success = True
                else:
                    logger.warning(f"⚠️ [UPSTAGE] Upstage 실패: {upstage_result.error}")
            
            except Exception as e:
                logger.error(f"❌ [UPSTAGE] Upstage 예외 발생: {e}", exc_info=True)
        
        # AWS Textract (향후 구현)
        elif provider == "aws_textract":
            logger.warning(f"⚠️ AWS Textract provider는 아직 구현되지 않았습니다.")
            # TODO: AWS Textract 구현
        
        # 기타 Provider (pdfplumber, tesseract 등)
        elif provider == "etc_other":
            logger.info(f"📚 기타 오픈소스 라이브러리(pdfplumber) 사용")
            primary_success = True  # pdfplumber는 아래에서 항상 실행
        
        # 알 수 없는 Provider
        else:
            logger.warning(f"⚠️ 알 수 없는 Provider '{provider}'")
        
        # Fallback Provider 시도
        if not primary_success and fallback_provider and fallback_provider != provider:
            logger.info(f"🔄 Fallback Provider로 재시도: {fallback_provider}")
            
            if fallback_provider == "upstage":
                try:
                    logger.info(f"🔷 [FALLBACK-UPSTAGE] Upstage Document Parse 사용 - 파일: {file_path}")
                    from .upstage_document_service import upstage_document_service
                    
                    logger.info(f"🔷 [FALLBACK-UPSTAGE] upstage_document_service 모듈 로드 완료")
                    logger.info(f"🔷 [FALLBACK-UPSTAGE] API 키 설정 여부: {bool(upstage_document_service.api_key)}")
                    
                    logger.info(f"🔷 [FALLBACK-UPSTAGE] Document Parse 호출 시작: {file_path}")
                    upstage_result = await upstage_document_service.parse_document(file_path)
                    
                    logger.info(f"🔷 [FALLBACK-UPSTAGE] Document Parse 호출 완료 - success: {upstage_result.success}")
                    
                    if upstage_result.success:
                        logger.info(f"✅ [FALLBACK-UPSTAGE] Upstage 성공: {file_path}")
                        logger.info(f"✅ [FALLBACK-UPSTAGE] 추출된 텍스트 길이: {len(upstage_result.text)}")
                        
                        converted_result = upstage_document_service.create_internal_extraction_result(upstage_result)
                        result.update(converted_result)
                        return result
                    else:
                        logger.warning(f"⚠️ [FALLBACK-UPSTAGE] Upstage 실패: {upstage_result.error}")
                        
                except Exception as e:
                    logger.error(f"❌ [FALLBACK-UPSTAGE] Upstage 예외 발생: {e}", exc_info=True)
            
            elif fallback_provider == "azure_di":
                try:
                    from .azure_document_intelligence_service import azure_document_intelligence_service
                    
                    logger.info(f"[Fallback] Azure DI로 PDF 분석 시도: {file_path}")
                    di_result = await azure_document_intelligence_service.analyze_pdf(file_path)
                    
                    if di_result.success:
                        logger.info(f"✅ [Fallback] Azure DI 성공: {file_path}")
                        converted_result = azure_document_intelligence_service.create_internal_extraction_result(di_result)
                        result.update(converted_result)
                        return result
                    else:
                        logger.warning(f"⚠️ [Fallback] Azure DI 실패: {di_result.error}")
                        
                except Exception as e:
                    logger.warning(f"⚠️ [Fallback] Azure DI 예외: {e}")
        
        # 최종 Fallback: pdfplumber (항상 사용 가능)
        
        # pdfplumber 폴백 또는 기본 방식
        return await self._extract_pdf_with_pdfplumber(file_path, result)
    
    async def _extract_pdf_with_pdfplumber(self, file_path: str, result: Dict[str, Any]) -> Dict[str, Any]:
        """PDF 파일 텍스트 추출 (pdfplumber 사용) - 페이지별 구조화"""
        try:
            import pdfplumber
            
            text_content = ""
            page_count = 0
            pages_data = []  # 페이지별 구조화 데이터
            
            with pdfplumber.open(file_path) as pdf:
                page_count = len(pdf.pages)
                for page_num, page in enumerate(pdf.pages):
                    page_text = page.extract_text() or ""
                    page_tables = page.extract_tables() or []
                    page_images = page.images or []
                    
                    # 이미지 메타데이터 추출 (멀티모달 검색용)
                    images_metadata = []
                    for img_idx, img in enumerate(page_images):
                        images_metadata.append({
                            'image_index': img_idx,
                            'x0': img.get('x0', 0),
                            'y0': img.get('y0', 0),
                            'x1': img.get('x1', 0),
                            'y1': img.get('y1', 0),
                            'width': img.get('width', 0),
                            'height': img.get('height', 0),
                            # 향후 확장: 'image_path', 'image_base64', 'ocr_text'
                        })
                    
                    # 페이지별 데이터 저장
                    pages_data.append({
                        'page_no': page_num + 1,
                        'text': page_text.strip(),
                        'tables_count': len(page_tables),
                        'images_count': len(page_images),
                        'images_metadata': images_metadata,  # ✅ 이미지 상세 정보
                        'char_count': len(page_text),
                        'has_content': bool(page_text.strip() or page_tables or page_images)
                    })
                    
                    if page_text:
                        text_content += f"\n[페이지 {page_num + 1}]\n{page_text}\n"
            
            if text_content.strip():
                result["text"] = text_content.strip()
                result["metadata"].update({
                    "page_count": page_count,
                    "pages": pages_data,  # ✅ 추가
                    "total_tables": sum(p['tables_count'] for p in pages_data),
                    "total_images": sum(p['images_count'] for p in pages_data),
                    "extraction_method": "pdfplumber_fallback" if settings.use_azure_document_intelligence_pdf else "pdfplumber",
                    "char_count": len(text_content),
                    "extraction_note": f"PDF 텍스트 추출 성공 (페이지별 구조화) - {'Azure DI 폴백' if settings.use_azure_document_intelligence_pdf else 'pdfplumber'}"
                })
                logger.info(f"PDF 텍스트 추출 성공: {len(text_content)}자, {page_count}페이지")
            else:
                result["text"] = f"PDF 파일: {Path(file_path).name}\n\n[텍스트를 추출할 수 없는 PDF 파일입니다]"
                result["metadata"]["extraction_note"] = "추출 가능한 텍스트가 없음"
                result["metadata"]["pages"] = pages_data  # 빈 페이지 정보라도 저장
                
        except Exception as e:
            result["success"] = False
            result["error"] = f"PDF 파일 처리 실패: {str(e)}"
            logger.error(f"PDF 텍스트 추출 실패: {e}")
        
        return result
    
    async def _extract_docx_file(self, file_path: str, result: Dict[str, Any]) -> Dict[str, Any]:
        """DOCX 파일 텍스트 및 이미지 추출 (python-docx 사용) - 구조화"""
        try:
            from docx import Document
            import zipfile
            import io
            from PIL import Image
            
            doc = Document(file_path)
            text_content = ""
            paragraph_count = 0
            paragraphs_data = []  # 문단 데이터
            tables_data = []  # 표 데이터
            images_metadata = []  # 이미지 메타데이터
            
            # 문단별로 텍스트 추출
            for para_idx, paragraph in enumerate(doc.paragraphs):
                if paragraph.text.strip():
                    text_content += paragraph.text + "\n"
                    paragraph_count += 1
                    paragraphs_data.append({
                        'paragraph_no': para_idx + 1,
                        'text': paragraph.text.strip(),
                        'char_count': len(paragraph.text)
                    })
            
            # 표 내용도 추출
            table_count = 0
            for table_idx, table in enumerate(doc.tables):
                table_count += 1
                table_text = f"\n[표 {table_count}]\n"
                table_rows = []
                table_cells = []  # 2D 셀 구조 (중복 저장 가능하나 검색/구조 확장 용도)
                for row in table.rows:
                    row_cell_texts = []
                    for cell in row.cells:
                        ctext = cell.text.strip() if cell.text else ""
                        row_cell_texts.append(ctext)
                    row_text = " | ".join(row_cell_texts)
                    if row_text.strip():
                        table_text += row_text + "\n"
                    table_rows.append(row_text)
                    table_cells.append(row_cell_texts)
                text_content += table_text
                
                tables_data.append({
                    'table_no': table_count,
                    'rows_count': len(table.rows),
                    'cols_count': len(table.columns) if hasattr(table, 'columns') else 0,
                    'content': table_rows,
                    'cells': table_cells,
                    'has_header': True if table_rows else False  # 단순 첫 행을 헤더로 간주 (추후 고도화 가능)
                })
            
            # 이미지 추출 (DOCX는 실제로는 ZIP 파일)
            image_count = 0
            try:
                with zipfile.ZipFile(file_path, 'r') as docx_zip:
                    # word/media/ 폴더에서 이미지 파일들 찾기
                    media_files = [name for name in docx_zip.namelist() if name.startswith('word/media/')]
                    
                    for media_file in media_files:
                        try:
                            # 이미지 파일인지 확인 (확장자 기준)
                            if any(media_file.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff']):
                                image_count += 1
                                
                                # 이미지 바이너리 데이터 읽기
                                image_data = docx_zip.read(media_file)
                                
                                # PIL로 이미지 정보 추출
                                try:
                                    with Image.open(io.BytesIO(image_data)) as img:
                                        width, height = img.size
                                        format_name = img.format or 'Unknown'
                                        
                                        images_metadata.append({
                                            'image_index': image_count,
                                            'filename': os.path.basename(media_file),
                                            'format': format_name,
                                            'width': width,
                                            'height': height,
                                            'size_bytes': len(image_data),
                                            'media_path': media_file,
                                            'binary_data': image_data  # 바이너리 데이터 포함 (일시적)
                                        })
                                        
                                except Exception as img_err:
                                    logger.warning(f"DOCX 이미지 정보 추출 실패 ({media_file}): {img_err}")
                                    # 기본 메타데이터라도 저장
                                    images_metadata.append({
                                        'image_index': image_count,
                                        'filename': os.path.basename(media_file),
                                        'format': 'Unknown',
                                        'size_bytes': len(image_data),
                                        'media_path': media_file,
                                        'binary_data': image_data  # 바이너리 데이터 포함 (일시적)
                                    })
                                    
                        except Exception as file_err:
                            logger.warning(f"DOCX 미디어 파일 처리 실패 ({media_file}): {file_err}")
                            
            except Exception as zip_err:
                logger.warning(f"DOCX ZIP 처리 실패: {zip_err}")
            
            # 결과 구성
            if text_content.strip():
                result["text"] = text_content.strip()
                
                # pages 구조로 변환 (멀티모달 서비스와의 호환성을 위해)
                pages_data = [{
                    'page_no': 1,
                    'text': text_content.strip(),
                    'images_metadata': images_metadata,
                    'tables_count': table_count,
                    'images_count': image_count,
                    'tables_metadata': tables_data  # 페이지 수준 표 상세 구조 제공
                }]
                
                result["metadata"].update({
                    "paragraph_count": paragraph_count,
                    "paragraphs": paragraphs_data[:50],  # 최대 50개까지만 (크기 제한)
                    "table_count": table_count,
                    "tables": tables_data,
                    "pages": pages_data,  # 멀티모달 파이프라인을 위한 pages 구조 추가
                    "extraction_method": "python-docx",
                    "char_count": len(text_content),
                    "extraction_note": f"DOCX 텍스트 및 이미지 추출 성공 (구조화) - 이미지 {image_count}개"
                })
                logger.info(f"DOCX 추출 성공: {len(text_content)}자, {paragraph_count}개 문단, {table_count}개 표, {image_count}개 이미지")
            else:
                result["text"] = f"DOCX 파일: {Path(file_path).name}\n\n[텍스트를 추출할 수 없는 DOCX 파일입니다]"
                result["metadata"]["extraction_note"] = "추출 가능한 텍스트가 없음"
            
        except Exception as e:
            result["success"] = False
            result["error"] = f"DOCX 파일 처리 실패: {str(e)}"
            logger.error(f"DOCX 텍스트 추출 실패: {e}")
        
        return result
    
    async def _extract_excel_file(self, file_path: str, result: Dict[str, Any]) -> Dict[str, Any]:
        """Excel 파일 텍스트 추출 (openpyxl 사용) - 시트별 구조화"""
        try:
            from openpyxl import load_workbook
            
            workbook = load_workbook(file_path, read_only=True)
            text_content = ""
            sheet_count = 0
            sheets_data = []  # 시트 데이터
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_count += 1
                text_content += f"\n[시트: {sheet_name}]\n"
                
                row_count = 0
                sheet_rows = []
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None and str(cell).strip() for cell in row):
                        row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                        text_content += row_text + "\n"
                        sheet_rows.append(row_text)
                        row_count += 1
                        
                        # 너무 많은 행은 제한
                        if row_count > 1000:
                            text_content += "[... 추가 행 생략 ...]\n"
                            break
                
                sheets_data.append({
                    'sheet_no': sheet_count,
                    'sheet_name': sheet_name,
                    'rows_count': row_count,
                    'content': sheet_rows[:100]  # 최대 100행까지만
                })
            
            workbook.close()
            
            if text_content.strip():
                result["text"] = text_content.strip()
                result["metadata"].update({
                    "sheet_count": sheet_count,
                    "sheets": sheets_data,
                    "extraction_method": "openpyxl",
                    "char_count": len(text_content),
                    "extraction_note": "Excel 텍스트 추출 성공 (시트별 구조화)"
                })
                logger.info(f"Excel 텍스트 추출 성공: {len(text_content)}자, {sheet_count}개 시트")
            else:
                result["text"] = f"Excel 파일: {Path(file_path).name}\n\n[텍스트를 추출할 수 없는 Excel 파일입니다]"
                result["metadata"]["extraction_note"] = "추출 가능한 텍스트가 없음"
            
        except Exception as e:
            result["success"] = False
            result["error"] = f"Excel 파일 처리 실패: {str(e)}"
            logger.error(f"Excel 텍스트 추출 실패: {e}")
        
        return result
    
    async def _extract_hwp_file(self, file_path: str, result: Dict[str, Any]) -> Dict[str, Any]:
        """HWP 파일 및 HWPX 파일 텍스트 추출 (개선된 방식)"""
        try:
            # 새로운 HWP 변환 서비스 사용
            from .hwp_converter_service import hwp_converter_service
            
            # LibreOffice 변환 방식 시도
            hwp_result = await hwp_converter_service.extract_text_from_hwp(file_path)
            
            if hwp_result['success']:
                result['success'] = True
                result['text'] = hwp_result['text']
                result['metadata'].update(hwp_result['metadata'])
                result['metadata']['extraction_method'] = 'libreoffice_conversion'
            else:
                # 기존 방식으로 폴백
                ext = Path(file_path).suffix.lower()
                if ext == '.hwpx':
                    # HWPX: ZIP + XML 구조
                    import zipfile
                    try:
                        import lxml.etree as ET
                    except ImportError:
                        import xml.etree.ElementTree as ET
                    
                    with zipfile.ZipFile(file_path, 'r') as z:
                        # 기본 섹션 파일 읽기
                        xml_names = [n for n in z.namelist() if n.endswith('.xml')]
                        text_content = ''
                        for name in xml_names:
                            try:
                                data = z.read(name)
                                tree = ET.fromstring(data)
                                # 모든 텍스트 노드 수집
                                texts = tree.xpath('//text()')
                                text_content += '\n'.join(texts) + '\n'
                            except Exception:
                                continue
                    
                    result['text'] = text_content.strip() or f'HWPX 파일입니다: {Path(file_path).name}\n[텍스트 추출 실패]'
                    result['metadata'].update({
                        'extraction_method': 'hwp5-xml-fallback',
                        'char_count': len(result['text']),
                        'extraction_note': 'HWPX 텍스트 추출 완료 (폴백 방식)'
                    })
                else:
                    # HWP: OLE 파일 PrvText 스트림 추출
                    import olefile
                    ole = olefile.OleFileIO(file_path)
                    if ole.exists('PrvText'):
                        raw = ole.openstream('PrvText').read()
                        try:
                            text = raw.decode('utf-16le')
                        except Exception:
                            text = raw.decode('cp949', errors='ignore')
                        result['text'] = text.strip() or f'HWP 파일입니다: {Path(file_path).name}\n[PrvText 빈 스트림]'
                        result['metadata'].update({
                            'extraction_method': 'olefile-PrvText-fallback',
                            'char_count': len(result['text']),
                            'extraction_note': 'HWP 텍스트 추출 완료 (폴백 방식)'
                        })
                    else:
                        result['text'] = f'HWP 파일입니다: {Path(file_path).name}\n[PrvText 스트림 없음]'
                        result['metadata']['extraction_note'] = 'HWP PrvText 스트림 없음 (폴백 방식)'
                        
        except Exception as e:
            result['success'] = False
            result['error'] = f'HWP/HWPX 처리 실패: {str(e)}'
        
        return result
    
    def _get_extraction_method(self, file_extension: str) -> str:
        """파일 확장자별 추출 방법 반환"""
        methods = {
            '.txt': 'direct_text_read',
            '.md': 'direct_text_read',
            '.py': 'direct_text_read',
            '.js': 'direct_text_read',
            '.html': 'direct_text_read',
            '.css': 'direct_text_read',
            '.json': 'direct_text_read',
            '.xml': 'direct_text_read',
            '.pdf': 'pdf_extraction_library',
            '.docx': 'python_docx_library',
            '.doc': 'python_docx_library',
            '.xlsx': 'openpyxl_library',
            '.xls': 'openpyxl_library',
            '.pptx': 'python_pptx_library',
            '.ppt': 'python_pptx_library',
            '.hwp': 'olefile_library'
        }
        return methods.get(file_extension.lower(), 'unknown')

    async def _extract_pptx_file(self, file_path: str, result: Dict[str, Any]) -> Dict[str, Any]:
        """PPTX 파일 텍스트 추출 (python-pptx 사용)"""
        try:
            from pptx import Presentation
            import io
            from PIL import Image
            
            presentation = Presentation(file_path)
            text_content = ""
            slide_count = 0
            shape_count = 0
            total_text_length = 0
            slides_data = []  # 슬라이드별 구조화 데이터
            
            # 슬라이드별로 텍스트 추출
            for slide_idx, slide in enumerate(presentation.slides):
                slide_count += 1
                slide_text = f"\n[슬라이드 {slide_idx + 1}]\n"
                slide_has_content = False
                slide_tables_count = 0
                slide_charts_count = 0
                slide_text_content = ""
                slide_tables_metadata = []  # 표 상세 구조 (cells 포함)
                
                # 슬라이드의 모든 shape에서 텍스트 추출
                for shape_idx, shape in enumerate(slide.shapes):
                    shape_content = ""
                    
                    try:
                        # Shape 유형 확인 (안전한 처리를 위해)
                        shape_type = getattr(shape, 'shape_type', None)
                        
                        # 텍스트 박스/도형의 텍스트 추출
                        if hasattr(shape, "text"):
                            text = shape.text.strip()
                            if text:
                                shape_content += text + "\n"
                                slide_has_content = True
                        
                        # 텍스트 프레임이 있는 경우 (더 상세한 텍스트 추출)
                        if hasattr(shape, "text_frame") and shape.text_frame is not None:
                            try:
                                for paragraph in shape.text_frame.paragraphs:
                                    paragraph_text = ""
                                    for run in paragraph.runs:
                                        if run.text and run.text.strip():
                                            paragraph_text += run.text
                                    if paragraph_text.strip():
                                        shape_content += paragraph_text.strip() + "\n"
                                        slide_has_content = True
                            except Exception as e:
                                logger.debug(f"텍스트 프레임 처리 중 오류 (무시됨): {e}")
                                pass
                        
                        # 표가 있는 경우 표 내용도 추출
                        if hasattr(shape, "table") and shape.table is not None:
                            try:
                                table = shape.table
                                slide_tables_count += 1  # ✅ 표 카운트
                                table_content = "\n[표]\n"
                                table_rows_raw = []
                                table_cells_2d = []
                                for row in table.rows:
                                    row_texts = []
                                    for cell in row.cells:
                                        cell_text = cell.text.strip() if cell.text else ""
                                        if cell_text:
                                            row_texts.append(cell_text)
                                        else:
                                            row_texts.append("")
                                    if row_texts:
                                        table_content += " | ".join(row_texts) + "\n"
                                    table_rows_raw.append(" | ".join(row_texts))
                                    table_cells_2d.append(row_texts)
                                if len(table_content) > 10:  # 실제 내용이 있는 경우만
                                    shape_content += table_content
                                    slide_has_content = True
                                # 표 메타데이터 저장 (내용 유무 상관없이 구조 보존)
                                slide_tables_metadata.append({
                                    'table_index': slide_tables_count - 1,
                                    'rows_count': len(table.rows),
                                    'cols_count': len(table.columns) if hasattr(table, 'columns') else 0,
                                    'content': table_rows_raw,
                                    'cells': table_cells_2d,
                                    'has_header': True if table_rows_raw else False
                                })
                            except Exception as e:
                                logger.debug(f"표 처리 중 오류 (무시됨): {e}")
                                pass
                        
                        # 차트의 제목이나 데이터 레이블 추출 시도
                        try:
                            if hasattr(shape, "chart"):
                                # 차트 객체가 실제로 존재하는지 더 안전하게 확인
                                chart = getattr(shape, "chart", None)
                                if chart is not None:
                                    slide_charts_count += 1  # ✅ 차트 카운트
                                    try:
                                        # 차트 제목 추출 시도
                                        if hasattr(chart, "chart_title") and chart.chart_title:
                                            if hasattr(chart.chart_title, "has_text_frame") and chart.chart_title.has_text_frame:
                                                chart_title = chart.chart_title.text_frame.text.strip()
                                                if chart_title:
                                                    shape_content += f"[차트 제목] {chart_title}\n"
                                                    slide_has_content = True
                                    except Exception as chart_title_error:
                                        logger.debug(f"차트 제목 처리 중 오류 (무시됨): {chart_title_error}")
                                        pass
                        except Exception as e:
                            # 차트 처리 중 오류가 발생해도 계속 진행
                            logger.debug(f"차트 처리 중 오류 (무시됨): {e}")
                            pass
                        
                        if shape_content:
                            shape_count += 1
                            slide_text += f"  {shape_content}"
                            total_text_length += len(shape_content)
                            
                    except Exception as shape_error:
                        # 개별 shape 처리 중 오류가 발생해도 다음 shape로 계속 진행
                        logger.debug(f"Shape 처리 중 오류 (무시됨): {shape_error}")
                        continue
                
                # 슬라이드 노트 추출 시도
                try:
                    if hasattr(slide, "notes_slide") and slide.notes_slide and slide.notes_slide.shapes:
                        notes_text = ""
                        for shape in slide.notes_slide.shapes:
                            if hasattr(shape, "text") and shape.text and shape.text.strip():
                                notes_text += shape.text.strip() + " "
                        if notes_text.strip():
                            slide_text += f"\n[노트] {notes_text.strip()}\n"
                            slide_has_content = True
                            total_text_length += len(notes_text)
                except Exception as e:
                    logger.debug(f"슬라이드 노트 처리 중 오류 (무시됨): {e}")
                    pass
                
                if slide_has_content:
                    text_content += slide_text
                    slide_text_content = slide_text
                else:
                    # 텍스트가 없는 슬라이드도 기록
                    text_content += f"\n[슬라이드 {slide_idx + 1}] (시각적 콘텐츠 - 텍스트 없음)\n"
                    slide_text_content = "(시각적 콘텐츠 - 텍스트 없음)"
                
                # 슬라이드별 구조화 데이터 저장
                
                # 이미지 메타데이터 추출 (멀티모달 검색용)
                images_metadata = []
                for shape in slide.shapes:
                    try:
                        # 이미지 또는 그림 shape 확인
                        if hasattr(shape, 'image'):
                            try:
                                img_blob = shape.image.blob  # 원본 이미지 바이너리
                                img_ext = shape.image.ext or 'png'
                                size_bytes = len(img_blob)
                                pixel_width = None
                                pixel_height = None
                                try:
                                    with Image.open(io.BytesIO(img_blob)) as im:
                                        pixel_width, pixel_height = im.size
                                except Exception as pil_err:
                                    logger.debug(f"PPTX 이미지 PIL 로드 실패 (무시): {pil_err}")
                                images_metadata.append({
                                    'image_index': len(images_metadata),
                                    'left': int(getattr(shape, 'left', 0)),
                                    'top': int(getattr(shape, 'top', 0)),
                                    'width': int(getattr(shape, 'width', 0)),
                                    'height': int(getattr(shape, 'height', 0)),
                                    'ext': img_ext,
                                    'size_bytes': size_bytes,
                                    'pixel_width': pixel_width,
                                    'pixel_height': pixel_height,
                                    'binary_data': img_blob  # 멀티모달 후처리 단계에서 제거 / 활용
                                })
                            except Exception as im_err:
                                logger.debug(f"PPTX 이미지 메타 수집 실패 (무시): {im_err}")
                    except:
                        pass
                
                slides_data.append({
                    'slide_no': slide_idx + 1,
                    'text': slide_text_content.strip(),
                    'tables_count': slide_tables_count,
                    'charts_count': slide_charts_count,
                    'images_count': len(images_metadata),  # ✅ 이미지 개수
                    'images_metadata': images_metadata,  # ✅ 이미지 상세 정보
                    'tables_metadata': slide_tables_metadata,  # 표 상세 구조
                    'char_count': len(slide_text_content),
                    'has_content': slide_has_content
                })
            
            # 결과 처리
            if total_text_length > 10:  # 최소 10자 이상의 의미있는 텍스트가 있을 때
                result["text"] = text_content.strip()
                result["metadata"].update({
                    "slide_count": slide_count,
                    "slides": slides_data,  # ✅ 추가
                    "shape_count": shape_count,
                    "total_tables": sum(s['tables_count'] for s in slides_data),
                    "total_charts": sum(s['charts_count'] for s in slides_data),
                    "total_images": sum(s.get('images_count', 0) for s in slides_data),  # ✅ 전체 이미지 개수
                    "extraction_method": "python-pptx-enhanced",
                    "char_count": len(text_content),
                    "meaningful_text_length": total_text_length,
                    "extraction_note": "PPTX 텍스트 추출 성공 (슬라이드별 구조화, 이미지 메타데이터 포함)"
                })
                logger.info(f"PPTX 텍스트 추출 성공: {len(text_content)}자 (의미있는 텍스트: {total_text_length}자), {slide_count}슬라이드")
            else:
                # 텍스트가 거의 없는 경우 기본 메시지와 함께 슬라이드 정보 제공
                fallback_text = f"PowerPoint 파일: {Path(file_path).name}\n\n"
                fallback_text += f"총 {slide_count}개 슬라이드가 포함되어 있습니다.\n"
                fallback_text += "주로 이미지, 도형, 차트 등의 시각적 콘텐츠로 구성되어 있어 텍스트 추출이 제한적입니다.\n\n"
                if text_content.strip():
                    fallback_text += "추출된 일부 내용:\n" + text_content.strip()
                
                result["text"] = fallback_text
                result["metadata"].update({
                    "slide_count": slide_count,
                    "slides": slides_data,  # ✅ 추가
                    "shape_count": shape_count,
                    "total_tables": sum(s['tables_count'] for s in slides_data),
                    "total_charts": sum(s['charts_count'] for s in slides_data),
                    "extraction_method": "python-pptx-enhanced",
                    "char_count": len(fallback_text),
                    "meaningful_text_length": total_text_length,
                    "extraction_note": "PPTX 파일이지만 추출 가능한 텍스트가 제한적임"
                })
                logger.info(f"PPTX 처리 완료: {slide_count}슬라이드, 추출된 텍스트 {total_text_length}자 (제한적)")
                
        except ImportError:
            result["success"] = False
            result["error"] = "python-pptx 라이브러리가 설치되지 않았습니다. pip install python-pptx"
            logger.error("python-pptx 라이브러리 없음")
        except Exception as e:
            result["success"] = False
            result["error"] = f"PPTX 파일 처리 실패: {str(e)}"
            logger.error(f"PPTX 텍스트 추출 실패: {e}")
        
        return result

# 전역 인스턴스
text_extractor_service = TextExtractorService()
