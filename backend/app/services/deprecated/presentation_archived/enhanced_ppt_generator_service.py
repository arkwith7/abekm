"""Enhanced PPT Generator Service (LLM-driven)"""
from __future__ import annotations

import json
import asyncio
import re
import os
import tempfile
from pathlib import Path
from typing import Any, Dict, List, Optional
from datetime import datetime

from loguru import logger
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pydantic import BaseModel, Field, validator

from app.core.config import settings
from app.services.core.ai_service import ai_service
from .ppt_models import ChartData, DiagramData, SlideSpec, DeckSpec
from .ppt_template_manager import PPTTemplateManager, template_manager
from .product_template_manager import ProductTemplateManager, product_template_manager
from .dynamic_template_manager import DynamicTemplateManager, dynamic_template_manager
from .enhanced_object_processor import EnhancedPPTObjectProcessor


class EnhancedPPTGeneratorService:
    def __init__(self):
        # backend/prompts 지속 사용 (root/prompts 제거 예정)
        self.prompts_dir = Path(__file__).parents[3] / "prompts"
        self.upload_dir = settings.resolved_upload_dir
        self.upload_dir.mkdir(parents=True, exist_ok=True)
        self.template_manager = template_manager  # 템플릿 관리자 추가
        self.object_processor = EnhancedPPTObjectProcessor()  # 🆕 확장된 오브젝트 처리기
        self.color_themes = {
            "corporate_blue": {"primary": RGBColor(0, 102, 204), "secondary": RGBColor(102, 153, 255), "accent": RGBColor(255, 153, 0), "text": RGBColor(51, 51, 51), "background": RGBColor(248, 249, 250)},
            "modern_green": {"primary": RGBColor(34, 139, 34), "secondary": RGBColor(144, 238, 144), "accent": RGBColor(255, 215, 0), "text": RGBColor(47, 79, 79), "background": RGBColor(248, 255, 248)},
            "professional_gray": {"primary": RGBColor(70, 70, 70), "secondary": RGBColor(169, 169, 169), "accent": RGBColor(220, 20, 60), "text": RGBColor(0, 0, 0), "background": RGBColor(245, 245, 245)},
            "playful_violet": {"primary": RGBColor(111, 45, 168), "secondary": RGBColor(181, 126, 220), "accent": RGBColor(255, 181, 71), "text": RGBColor(60, 60, 60), "background": RGBColor(250, 248, 255)},
        }

    # ---------------- Title Helpers (Query Detection & Normalization) ----------------
    def _is_query_like(self, text: str) -> bool:
        if not text:
            return False
        t = text.strip()
        # 물음표 / 명령형 / 요청형 표현 포함 여부
        patterns = [
            r"\?$", r"해주세요$", r"해 주세요$", r"알려줘$", r"알려주세요$", r"만들어줘$", r"작성해줘$", r"정리해줘$",
            r"설명해줘$", r"소개해줘$", r"소개해 주세요$", r"요약해줘$", r"정리해 주세요$"
        ]
        if any(re.search(p, t) for p in patterns):
            return True
        # 문장 길이 대비 동사/요청 비율 (간단 힌트)
        if len(t) <= 30 and any(x in t for x in ["해주세요", "해줘", "어떻게", "무엇", "?"]):
            return True
        return False

    def _normalize_topic_text(self, text: str) -> str:
        if not text:
            return text
        t = text.strip()
        # 요청형 접미사 제거
        t = re.sub(r"(을|를)?\s*(소개|설명|요약)?해\s*주세요$", "", t)
        t = re.sub(r"(을|를)?\s*(소개|설명|요약)?해줘$", "", t)
        t = re.sub(r"(알려줘|알려주세요)$", "", t)
        t = re.sub(r"(만들어줘|작성해줘|정리해줘)$", "", t)
        t = re.sub(r"\?$", "", t)
        t = t.strip(' ,.-')
        # 너무 짧아졌다면 원본 보존
        return t if len(t) >= 4 else text

    # ---------------- Topic & Content Analysis ----------------
    def _extract_clean_title(self, text: str, document_filename: Optional[str] = None) -> str:
        """제목 신뢰도 향상: AI 응답에서 실제 제목 추출"""
        if not text:
            return ""
        
        lines = text.strip().split('\n')
        title_candidates = []
        
        # 첫 번째 줄이 섹션 형태가 아니고 독립적인 제목인지 확인
        first_line = lines[0].strip() if lines else ""
        if first_line and not re.match(r'^\d+\.', first_line) and len(first_line) <= 50:
            # 다음 줄이 빈 줄이거나 섹션 시작이면 첫 줄이 제목일 가능성 높음
            if len(lines) > 1:
                second_line = lines[1].strip() if len(lines) > 1 else ""
                if not second_line or re.match(r'^\d+\.', second_line):
                    return first_line
        
        for i, line in enumerate(lines[:10]):  # 처음 10줄만 검사
            # 마크다운 데코레이터 제거 (# > * 공백)
            clean_line = re.sub(r'^[#>*\s]*', '', line).strip()
            
            if not clean_line or len(clean_line) <= 5:
                continue
                
            # 제목 후보 점수 계산
            score = self._title_score(clean_line)
            
            # 첫 번째 줄이고 구체적인 제목이면 높은 점수
            if i == 0 and not any(word in clean_line.lower() for word in ['질문', '문의', '해주세요', '알려주세요', '입니다', '합니다']):
                score += 50
            
            # 번호나 목차 형태가 아닌 제목이면 가점
            if not re.match(r'^\d+\.', clean_line.strip()) and not clean_line.lower().startswith('목차'):
                score += 10
                
            # 길이가 적절한 제목이면 가점 (10-50자)
            if 10 <= len(clean_line) <= 50:
                score += 20
            elif len(clean_line) > 100:
                score -= 30  # 너무 긴 설명문은 감점
                
            title_candidates.append((clean_line, score))
        
        # 가장 높은 점수의 제목 선택
        if title_candidates:
            best_title = max(title_candidates, key=lambda x: x[1])[0]
            
            # 문서명과 비교
            if document_filename:
                doc_title = re.sub(r'\.(docx?|pdf|txt)$', '', document_filename, flags=re.IGNORECASE)
                doc_score = self._title_score(doc_title)
                best_score = max(title_candidates, key=lambda x: x[1])[1]
                
                logger.debug(f"Title comparison: doc='{doc_title}'({doc_score}) vs best='{best_title}'({best_score})")
                return doc_title if doc_score > best_score + 10 else best_title
            
            return best_title
            
        # 문서명 fallback
        if document_filename:
            return re.sub(r'\.(docx?|pdf|txt)$', '', document_filename, flags=re.IGNORECASE)
        return ""

    def _title_score(self, title: str) -> int:
        """제목 품질 점수 계산 (길이 + 의미 키워드 + 질의형 패널티)"""
        if not title:
            return -100
        t = title.strip()
        score = 0
        ln = len(t)
        # 이상적인 길이(8~40자) 가산, 길이 벗어나면 완만 감점
        if 8 <= ln <= 40:
            score += 40
        else:
            score += max(0, 40 - abs(ln - 24))
        # 키워드 가중치
        high_value_keywords = ['제품', '시스템', '서비스', '개발', '분석', '보고서', '계획', '전략', '가이드', '로드맵']
        medium_value_keywords = ['개요', '소개', '설명', '정보', '플랫폼', '솔루션']
        low_value_keywords = ['발표자료', '문서', '자료', 'ppt', 'presentation']
        for kw in high_value_keywords:
            if kw in t:
                score += 16
        for kw in medium_value_keywords:
            if kw in t:
                score += 8
        for kw in low_value_keywords:
            if kw in t.lower():
                score -= 25
        # 명확한 제품/영문 토큰
        if re.search(r"[A-Z][A-Za-z0-9]{2,}", t):
            score += 6
        # 질의/요청형 패턴 패널티
        if re.search(r"(해주세요|해줘|알려줘|알려주세요|소개해줘|설명해줘|요약해줘|정리해줘)$", t):
            score -= 45
        if t.endswith('?'):
            score -= 30
        # 단어 수 (2~8 적정)
        words = re.split(r"\s+", t)
        if 2 <= len(words) <= 8:
            score += 10
        # 금지 기본 제목
        if t in ["발표자료", "프레젠테이션", "슬라이드"]:
            score -= 50
        return score

    def _extract_structured_sections(self, text: str) -> List[Dict[str, Any]]:
        """AI 답변에서 번호가 있는 구조화된 섹션들을 추출"""
        sections = []
        lines = text.strip().split('\n')
        current_section = None
        current_content = []
        
        for line in lines:
            # 번호가 있는 섹션 제목 찾기 (1. 제품 개요, 2. 기술 사양 등)
            section_match = re.match(r'^(\d+)\.\s*(.+)$', line.strip())
            if section_match:
                # 이전 섹션 저장
                if current_section:
                    sections.append({
                        'number': current_section['number'],
                        'title': current_section['title'],
                        'content': '\n'.join(current_content).strip()
                    })
                
                # 새 섹션 시작
                current_section = {
                    'number': int(section_match.group(1)),
                    'title': section_match.group(2).strip()
                }
                current_content = []
            elif current_section and line.strip():
                # 현재 섹션에 내용 추가 (빈 줄은 제외)
                current_content.append(line.strip())
        
        # 마지막 섹션 저장
        if current_section:
            sections.append({
                'number': current_section['number'],
                'title': current_section['title'],
                'content': '\n'.join(current_content).strip()
            })
        
        return sections

    def _parse_section_content(self, content: str) -> tuple[str, List[str]]:
        """섹션 내용을 key_message와 bullets로 분리"""
        lines = content.split('\n')
        key_message = ""
        bullets = []
        
        # 첫 문단은 key_message로 사용
        paragraph_lines = []
        for line in lines:
            if line.strip():
                paragraph_lines.append(line.strip())
            elif paragraph_lines:
                # 빈 줄을 만나면 문단 완료
                break
        
        if paragraph_lines:
            key_message = ' '.join(paragraph_lines)
        
        # 나머지는 bullets로 처리
        in_bullet_section = False
        for line in lines:
            stripped = line.strip()
            if not stripped:
                continue
                
            # key_message에 이미 포함된 내용은 스킵
            if not in_bullet_section and stripped in key_message:
                in_bullet_section = True
                continue
            elif not in_bullet_section and any(stripped in key_message for part in [stripped]):
                continue
            
            in_bullet_section = True
            
            # 불릿 포인트 정리 (-, *, 번호 등 제거)
            clean_bullet = re.sub(r'^[-*•]\s*', '', stripped)
            clean_bullet = re.sub(r'^\d+\.\s*', '', clean_bullet)
            
            if clean_bullet and len(clean_bullet) > 5:
                bullets.append(clean_bullet[:80])  # 길이 제한
        
        return key_message[:200], bullets[:6]  # 길이 제한

    def _extract_keyvalue_blocks(self, text: str) -> List[Dict[str, Any]]:
        """키-값 패턴 블록 추출"""
        blocks = []
        lines = text.split('\n')
        current_block = []
        
        for line in lines:
            # 키: 값 패턴 매칭
            kv_match = re.match(r'^([^:]{1,25}):\s*(.{1,100})$', line.strip())
            if kv_match:
                key, value = kv_match.groups()
                if len(key.strip()) <= 15 and len(value.strip()) <= 30:  # 규칙 준수
                    current_block.append({"key": key.strip(), "value": value.strip()})
            else:
                # 블록 완료
                if len(current_block) >= 3:
                    blocks.append({
                        "type": "keyvalue",
                        "items": current_block.copy(),
                        "is_chart_candidate": self._is_chart_candidate(current_block)
                    })
                current_block = []
        
        # 마지막 블록 처리
        if len(current_block) >= 3:
            blocks.append({
                "type": "keyvalue", 
                "items": current_block.copy(),
                "is_chart_candidate": self._is_chart_candidate(current_block)
            })
        
        return blocks

    def _is_chart_candidate(self, items: List[Dict[str, str]]) -> bool:
        """차트 후보 여부 판단"""
        if len(items) < 3:
            return False
            
        numeric_count = 0
        units = set()
        
        for item in items:
            value = item["value"]
            # 숫자 패턴 (정수, 소수, 숫자+단위)
            numeric_match = re.search(r'(\d+(?:\.\d+)?)\s*([a-zA-Z가-힣%]*)', value)
            if numeric_match:
                numeric_count += 1
                unit = numeric_match.group(2).strip()
                if unit:
                    units.add(unit)
        
        # 3개 이상 숫자형이고, 단위가 3개 미만 (혼합 단위 회피)
        return numeric_count >= 3 and len(units) <= 2

    # ---------------- Prompt Handling ----------------
    def _load_prompt(self) -> str:
        # 우선순위: 환경변수 > backend/prompts/presentation.prompt > 기본 내장
        candidates: List[Path] = []
        env_path = os.environ.get("PRESENTATION_PROMPT_PATH")
        if env_path:
            candidates.append(Path(env_path))
        pres_file = self.prompts_dir / "presentation.prompt"
        candidates.append(pres_file)
        for p in candidates:
            try:
                if p.exists() and p.is_file():
                    text = p.read_text(encoding="utf-8")
                    if len(text) > 500 and ('slides' in text or 'topic' in text):
                        logger.debug({"phase": "prompt_load", "file": str(p)})
                        return text
                    else:
                        logger.warning({"phase": "prompt_load_warn", "file": str(p), "reason": "missing expected tokens"})
            except Exception as e:  # pragma: no cover
                logger.warning({"phase": "prompt_load_error", "file": str(p), "error": str(e)})
        return (
            "당신은 전문 프레젠테이션 디자이너입니다. JSON만 출력. "
            "필드: topic,max_slides,slides[].title,key_message,bullets,layout,diagram,visual_suggestion,speaker_notes"
        )

    async def generate_enhanced_outline(self, topic: str, context_text: str, provider: Optional[str] = None,
                                        template_style: str = "business", include_charts: bool = True,
                                        retries: int = 2, document_filename: Optional[str] = None,
                                        custom_template_path: Optional[str] = None,
                                        presentation_type: str = "general",
                                        user_template_id: Optional[str] = None) -> DeckSpec:
        
        # 로깅 보강: 요청 단계
        first_lines = context_text[:300].replace('\n', '\\n') if context_text else ""
        logger.debug({
            "phase": "ppt_request",
            "raw_topic_candidate": topic,
            "first_lines": first_lines,
            "document_filename": document_filename,
            "template_style": template_style,
            "include_charts": include_charts,
            "presentation_type": presentation_type,
            "user_template_id": user_template_id
        })
        
        # 🎯 사용자 업로드 템플릿 우선 처리
        if user_template_id and user_template_id.startswith("user_"):
            logger.info(f"사용자 템플릿 기반 PPT 생성: {user_template_id}")
            return await self._generate_user_template_outline(
                topic, context_text, provider, template_style,
                include_charts, retries, document_filename, user_template_id
            )
        
        # 🎯 제품소개서 특화 처리
        if presentation_type == "product_introduction":
            logger.info(f"제품소개서 모드로 PPT 생성: {topic}")
            return await self._generate_product_introduction_outline(
                topic, context_text, provider, template_style, 
                include_charts, retries, document_filename
            )
        
        # 일반 모드는 기존 로직 사용
        return await self._generate_general_outline(
            topic, context_text, provider, template_style, 
            include_charts, retries, document_filename
        )
        
        # 구조화된 섹션 추출 및 힌트 생성
        structured_sections = self._extract_structured_sections(context_text)
        
        # 키-값 패턴 추출 (Pre-processing)
        kv_blocks = self._extract_keyvalue_blocks(context_text)
        table_hints = []
        chart_hints = []
        
        for block in kv_blocks:
            if block["is_chart_candidate"]:
                chart_hints.append({
                    "type": "chart_candidate",
                    "items": block["items"][:6]  # 최대 6개 항목
                })
            else:
                table_hints.append({
                    "type": "table_candidate", 
                    "items": block["items"][:8]  # 최대 8개 항목
                })
        
        # 프롬프트 강화 - 구조화된 섹션 생성 강조
        system = self._load_prompt()
        enhanced_requirements = [
            "- AI 응답 내용의 제목과 구조를 정확히 반영하여 슬라이드 생성",
            "- 번호가 있는 섹션(1. 제품 개요, 2. 기술 사양 등)은 각각 별도 슬라이드로 구성",
            "- 각 섹션의 세부 항목들은 bullets로 정확히 나열",
            "- 두 번째 슬라이드는 번호가 있는 섹션들을 목차로 구성",
            "- 키:값 패턴이 연속으로 3개 이상 나오면 표(table) 슬라이드로 변환",
            "- bullets 항목당 50자 이내로 간결하게 표현",
            f"- include_charts={include_charts} 이면 수치 데이터를 차트로 변환",
            f"- template_style={template_style} (business|minimal|modern|playful)",
            "- visual_suggestion: 관련 아이콘/이미지 아이디어 1줄",
            "- speaker_notes: 발표자 스크립트 2~4문장 한국어",
            "- 각 슬라이드 title은 섹션 번호와 제목을 포함 (예: '1. 제품 개요')",
            "- key_message는 해당 섹션의 핵심 설명문으로 구성"
        ]
        
        base_user_content = [
            f"주제: {improved_topic}",
            f"컨텍스트:\n{context_text[:8000]}",  # 길이 제한
            "요구사항:",
            *enhanced_requirements
        ]
        
        # 구조화된 섹션 힌트 추가
        if structured_sections:
            sections_info = []
            for section in structured_sections:
                sections_info.append({
                    "number": section["number"],
                    "title": section["title"],
                    "preview": section["content"][:100] + "..." if len(section["content"]) > 100 else section["content"]
                })
            base_user_content.append(f"\n감지된 섹션 구조: {json.dumps(sections_info, ensure_ascii=False)}")
        
        # 힌트 추가
        if table_hints:
            base_user_content.append(f"\n표 후보 힌트: {json.dumps(table_hints, ensure_ascii=False)}")
        if chart_hints:
            base_user_content.append(f"\n차트 후보 힌트: {json.dumps(chart_hints, ensure_ascii=False)}")
            
        base_user_content.append("\nJSON만 출력 (추가 설명 금지)")
        base_user = '\n'.join(base_user_content)
        
        last_err: Optional[str] = None
        for attempt in range(retries + 1):
            user = base_user if attempt == 0 else base_user + f"\n이전 오류: {last_err}. 유효 JSON만 다시 출력."
            try:
                logger.debug(f"PPT LLM 호출 시도 {attempt + 1}/{retries + 1}, topic='{improved_topic}'")
                resp = await ai_service.chat_completion([
                    {"role": "system", "content": system},
                    {"role": "user", "content": user},
                ], provider=provider)
                raw = resp.get("response", "").strip()
                
                # 응답 로깅
                logger.debug(f"LLM raw response (first 200 chars): {raw[:200]}")
                
                deck = self._parse_outline(raw, improved_topic)
                
                # 후처리 및 검증
                deck = self._post_process_deck(deck, improved_topic)
                
                # 로깅 보강: 파싱 후
                has_chart = any(slide.diagram and slide.diagram.type == "chart" for slide in deck.slides)
                logger.debug({
                    "phase": "ppt_parse",
                    "final_topic": deck.topic,
                    "slide_count": len(deck.slides),
                    "has_chart": has_chart,
                    "ppt_fallback": False
                })
                
                return deck
                
            except Exception as e:  # noqa: BLE001
                last_err = str(e)
                logger.warning(f"PPT 생성 시도 {attempt + 1} 실패: {last_err}")
                if attempt == retries:
                    logger.error(f"Outline 생성 실패: {last_err}")
                    # 구조화된 섹션이 있으면 이를 활용한 fallback 생성
                    if structured_sections:
                        fallback_slides = []
                        # 제목 슬라이드
                        fallback_slides.append(SlideSpec(
                            title="개요", 
                            key_message="핵심 내용 요약", 
                            bullets=["주요 포인트 정리"], 
                            layout="title-and-content"
                        ))
                        
                        # 목차 슬라이드
                        agenda_bullets = [f"{s['number']}. {s['title']}" for s in structured_sections[:6]]
                        fallback_slides.append(SlideSpec(
                            title="목차",
                            key_message="",
                            bullets=agenda_bullets,
                            layout="title-and-content"
                        ))
                        
                        # 각 섹션별 슬라이드
                        for section in structured_sections[:4]:  # 최대 4개 섹션
                            key_msg, bullets = self._parse_section_content(section['content'])
                            fallback_slides.append(SlideSpec(
                                title=f"{section['number']}. {section['title']}",
                                key_message=key_msg,
                                bullets=bullets,
                                layout="title-and-content"
                            ))
                        
                        logger.debug({
                            "phase": "ppt_parse", 
                            "final_topic": improved_topic or "발표자료",
                            "slide_count": len(fallback_slides),
                            "has_chart": False,
                            "ppt_fallback": True,
                            "structured_fallback": True
                        })
                        
                        return DeckSpec(
                            topic=improved_topic or "발표자료", 
                            max_slides=len(fallback_slides), 
                            slides=fallback_slides,
                            theme={"color_scheme": "corporate_blue", "font_style": "modern"}
                        )
                    
                    # 기본 fallback (구조화된 섹션이 없을 때)
                    logger.debug({
                        "phase": "ppt_parse", 
                        "final_topic": improved_topic or "발표자료",
                        "slide_count": 3,
                        "has_chart": False,
                        "ppt_fallback": True,
                        "structured_fallback": False
                    })
                    return DeckSpec(topic=improved_topic or "발표자료", max_slides=3, slides=[
                        SlideSpec(title="개요", key_message="핵심 개요", bullets=["배경", "목표"], layout="title-and-content"),
                        SlideSpec(title="주요 내용", key_message="포인트", bullets=["포인트1", "포인트2"], layout="title-and-content"),
                        SlideSpec(title="결론", key_message="요약", bullets=["요약", "다음 단계"], layout="title-and-content"),
                    ], theme={"color_scheme": "corporate_blue", "font_style": "modern"})
            await asyncio.sleep(0.4)
        # 안전장치 (논리적으로 도달하지 않음)
        return DeckSpec(topic=improved_topic or "발표자료", max_slides=1, slides=[SlideSpec(title="개요", key_message="요약", bullets=["포인트"], layout="title-and-content")])

    async def _generate_user_template_outline(self, topic: str, context_text: str, 
                                            provider: Optional[str] = None,
                                            template_style: str = "business", 
                                            include_charts: bool = True,
                                            retries: int = 2, 
                                            document_filename: Optional[str] = None,
                                            user_template_id: str = None) -> DeckSpec:
        """사용자 업로드 템플릿 기반 아웃라인 생성"""
        try:
            logger.info(f"사용자 템플릿 기반 생성 시작: {user_template_id}")
            
            # 1. 사용자 템플릿 메타데이터 가져오기
            template_metadata = dynamic_template_manager.get_template_for_ai(user_template_id)
            if not template_metadata:
                logger.warning(f"템플릿 메타데이터를 찾을 수 없음: {user_template_id}")
                # 폴백: 일반 모드로 처리
                return await self._generate_general_outline(
                    topic, context_text, provider, template_style, 
                    include_charts, retries, document_filename
                )
            
            # 2. 동적 프롬프트 생성
            base_prompt = self._load_prompt()
            dynamic_prompt = dynamic_template_manager.generate_dynamic_prompt(
                user_template_id, base_prompt
            )
            
            # 3. 템플릿 특화 요구사항 생성
            template_requirements = [
                f"- 이 템플릿은 {template_metadata['total_slides']}개 슬라이드 구조입니다.",
                f"- 템플릿명: {template_metadata['template_name']}",
                f"- 주 색상: {template_metadata['style_guide']['color_scheme'].get('primary', '#1f4e79')}",
                "- 업로드된 템플릿의 레이아웃 구조를 최대한 활용하세요:",
            ]
            
            # 4. 사용 가능한 레이아웃 정보 추가
            for layout in template_metadata['layout_options']:
                template_requirements.append(
                    f"  * {layout['name']}: {layout['description']}"
                )
            
            # 5. 슬라이드 구조 가이드 추가
            template_requirements.extend([
                "",
                "- 다음 슬라이드 구조를 참고하세요:",
            ])
            
            for slide_struct in template_metadata['slide_structure_template'][:8]:  # 최대 8개
                template_requirements.append(
                    f"  {slide_struct['slide_number']}. {slide_struct['title']} "
                    f"(레이아웃: {slide_struct['layout']})"
                )
            
            # 6. AI 생성 요청
            user_content = [
                f"주제: {topic}",
                f"컨텍스트:\n{context_text[:8000]}",
                "템플릿 특화 요구사항:",
                *template_requirements,
                "",
                "위 템플릿 구조에 맞춰 JSON 아웃라인을 생성하세요. JSON만 출력 (추가 설명 금지)"
            ]
            
            user_prompt = '\n'.join(user_content)
            
            # 7. LLM 호출
            for attempt in range(retries + 1):
                try:
                    logger.debug(f"사용자 템플릿 기반 LLM 호출 시도 {attempt + 1}")
                    resp = await ai_service.chat_completion([
                        {"role": "system", "content": dynamic_prompt},
                        {"role": "user", "content": user_prompt},
                    ], provider=provider)
                    
                    raw = resp.get("response", "").strip()
                    logger.debug(f"LLM raw response (first 200 chars): {raw[:200]}")
                    
                    # 8. 아웃라인 파싱
                    deck = self._parse_outline(raw, topic)
                    
                    # 9. 템플릿 메타데이터 적용
                    enhanced_deck = dynamic_template_manager.apply_template_to_outline(
                        user_template_id, deck.dict()
                    )
                    
                    # 10. DeckSpec으로 변환
                    final_deck = DeckSpec(**enhanced_deck)
                    
                    logger.info(f"사용자 템플릿 기반 아웃라인 생성 완료: {len(final_deck.slides)}개 슬라이드")
                    return final_deck
                    
                except Exception as e:
                    logger.warning(f"사용자 템플릿 생성 시도 {attempt + 1} 실패: {e}")
                    if attempt == retries:
                        raise
                    continue
            
        except Exception as e:
            logger.error(f"사용자 템플릿 기반 생성 실패: {e}")
            # 폴백: 일반 모드로 처리
            logger.info("일반 모드로 폴백")
            return await self._generate_general_outline(
                topic, context_text, provider, template_style, 
                include_charts, retries, document_filename
            )

    async def _generate_product_introduction_outline(self, topic: str, context_text: str, 
                                                   provider: Optional[str] = None,
                                                   template_style: str = "business", 
                                                   include_charts: bool = True,
                                                   retries: int = 2, 
                                                   document_filename: Optional[str] = None) -> DeckSpec:
        """제품소개서 전용 아웃라인 생성"""
        try:
            # 1. 제품소개서 템플릿 매니저를 사용하여 구조화된 아웃라인 생성
            logger.info("제품소개서 전용 파이프라인 시작")
            
            # 2. RAG 답변을 제품소개서 구조로 분석 및 변환
            product_outline = product_template_manager.generate_product_outline(
                context_text, product_type="medical_device"
            )
            
            logger.debug({
                "phase": "product_outline_generated",
                "slides_count": len(product_outline.get("slides", [])),
                "topic": product_outline.get("topic", "")
            })
            
            # 3. 기본 DeckSpec으로 변환
            slides = []
            for slide_data in product_outline.get("slides", []):
                diagram = None
                if slide_data.get("diagram"):
                    diagram_data = slide_data["diagram"]
                    diagram = DiagramData(
                        type=diagram_data.get("type", "none"),
                        data=diagram_data.get("data", {}),
                        chart=None
                    )
                
                slide_spec = SlideSpec(
                    title=slide_data.get("title", ""),
                    key_message=slide_data.get("key_message", ""),
                    bullets=slide_data.get("bullets", []),
                    layout=slide_data.get("layout", "title-and-content"),
                    style=slide_data.get("style", {}),
                    diagram=diagram,
                    visual_suggestion=slide_data.get("visual_suggestion"),
                    speaker_notes=slide_data.get("speaker_notes")
                )
                slides.append(slide_spec)
            
            # 4. 최종 DeckSpec 생성
            deck = DeckSpec(
                topic=product_outline.get("topic", topic),
                max_slides=len(slides),
                slides=slides,
                theme=product_outline.get("theme", {"color_scheme": "medical_blue"})
            )
            
            logger.info(f"제품소개서 아웃라인 생성 완료: {len(slides)}개 슬라이드")
            return deck
            
        except Exception as e:
            logger.error(f"제품소개서 생성 실패: {e}")
            # 폴백: 일반 모드로 처리
            logger.info("일반 모드로 폴백")
            return await self._generate_general_outline(
                topic, context_text, provider, template_style, 
                include_charts, retries, document_filename
            )
    
    async def _generate_general_outline(self, topic: str, context_text: str, 
                                      provider: Optional[str] = None,
                                      template_style: str = "business", 
                                      include_charts: bool = True,
                                      retries: int = 2, 
                                      document_filename: Optional[str] = None) -> DeckSpec:
        """일반 아웃라인 생성 (기존 로직)"""
        # 기존 로직을 여기로 이동
        # 제목 신뢰도 향상 / 사용자 질의 제거 로직
        provided_topic = (topic or "").strip()
        extracted_topic = self._extract_clean_title(context_text, document_filename)
        chosen: str
        if provided_topic:
            if self._is_query_like(provided_topic):
                # 질의형이면 추출한 제목 우선
                chosen = extracted_topic or provided_topic
            else:
                # 두 후보 점수 비교 (추출 제목이 현저히 좋으면 교체)
                p_score = self._title_score(provided_topic)
                e_score = self._title_score(extracted_topic)
                chosen = extracted_topic if e_score > p_score + 12 else provided_topic
        else:
            chosen = extracted_topic or "발표자료"
        improved_topic = self._normalize_topic_text(chosen) or "발표자료"
        # 방어: 아직도 질의형 패턴이면 한 번 더 추출 제목 적용
        if self._is_query_like(improved_topic) and extracted_topic:
            improved_topic = self._normalize_topic_text(extracted_topic)
        if not improved_topic.strip():  # 최종 안전장치
            improved_topic = "발표자료"
        logger.debug({
            "phase": "title_select",
            "provided_topic": provided_topic,
            "extracted_topic": extracted_topic,
            "final_topic": improved_topic
        })
        
        # 구조화된 섹션 추출 및 힌트 생성
        structured_sections = self._extract_structured_sections(context_text)
        
        # 키-값 패턴 추출 (Pre-processing)
        kv_blocks = self._extract_keyvalue_blocks(context_text)
        table_hints = []
        chart_hints = []
        
        for block in kv_blocks:
            if block["is_chart_candidate"]:
                chart_hints.append({
                    "type": "chart_candidate",
                    "items": block["items"][:6]  # 최대 6개 항목
                })
            else:
                table_hints.append({
                    "type": "table_candidate", 
                    "items": block["items"][:8]  # 최대 8개 항목
                })
        
        # 프롬프트 강화 - 구조화된 섹션 생성 강조
        system = self._load_prompt()
        enhanced_requirements = [
            "- AI 응답 내용의 제목과 구조를 정확히 반영하여 슬라이드 생성",
            "- 번호가 있는 섹션(1. 제품 개요, 2. 기술 사양 등)은 각각 별도 슬라이드로 구성",
            "- 각 섹션의 세부 항목들은 bullets로 정확히 나열",
            "- 두 번째 슬라이드는 번호가 있는 섹션들을 목차로 구성",
            "- 키:값 패턴이 연속으로 3개 이상 나오면 표(table) 슬라이드로 변환",
            "- bullets 항목당 50자 이내로 간결하게 표현",
            f"- include_charts={include_charts} 이면 수치 데이터를 차트로 변환",
            f"- template_style={template_style} (business|minimal|modern|playful)",
            "- visual_suggestion: 관련 아이콘/이미지 아이디어 1줄",
            "- speaker_notes: 발표자 스크립트 2~4문장 한국어",
            "- 각 슬라이드 title은 섹션 번호와 제목을 포함 (예: '1. 제품 개요')",
            "- key_message는 해당 섹션의 핵심 설명문으로 구성"
        ]
        
        base_user_content = [
            f"주제: {improved_topic}",
            f"컨텍스트:\n{context_text[:8000]}",  # 길이 제한
            "요구사항:",
            *enhanced_requirements
        ]
        
        # 구조화된 섹션 힌트 추가
        if structured_sections:
            sections_info = []
            for section in structured_sections:
                sections_info.append({
                    "number": section["number"],
                    "title": section["title"],
                    "preview": section["content"][:100] + "..." if len(section["content"]) > 100 else section["content"]
                })
            base_user_content.append(f"\n감지된 섹션 구조: {json.dumps(sections_info, ensure_ascii=False)}")
        
        # 힌트 추가
        if table_hints:
            base_user_content.append(f"\n표 후보 힌트: {json.dumps(table_hints, ensure_ascii=False)}")
        if chart_hints:
            base_user_content.append(f"\n차트 후보 힌트: {json.dumps(chart_hints, ensure_ascii=False)}")
            
        base_user_content.append("\nJSON만 출력 (추가 설명 금지)")
        base_user = '\n'.join(base_user_content)
        
        last_err: Optional[str] = None
        for attempt in range(retries + 1):
            user = base_user if attempt == 0 else base_user + f"\n이전 오류: {last_err}. 유효 JSON만 다시 출력."
            try:
                logger.debug(f"PPT LLM 호출 시도 {attempt + 1}/{retries + 1}, topic='{improved_topic}'")
                resp = await ai_service.chat_completion([
                    {"role": "system", "content": system},
                    {"role": "user", "content": user},
                ], provider=provider)
                raw = resp.get("response", "").strip()
                
                # 응답 로깅
                logger.debug(f"LLM raw response (first 200 chars): {raw[:200]}")
                
                deck = self._parse_outline(raw, improved_topic)
                
                # 후처리 및 검증
                deck = self._post_process_deck(deck, improved_topic)
                
                # 로깅 보강: 파싱 후
                has_chart = any(slide.diagram and slide.diagram.type == "chart" for slide in deck.slides)
                logger.debug({
                    "phase": "ppt_parse",
                    "final_topic": deck.topic,
                    "slide_count": len(deck.slides),
                    "has_chart": has_chart,
                    "ppt_fallback": False
                })
                
                return deck
                
            except Exception as e:  # noqa: BLE001
                last_err = str(e)
                logger.warning(f"PPT 생성 시도 {attempt + 1} 실패: {last_err}")
                if attempt == retries:
                    logger.error(f"Outline 생성 실패: {last_err}")
                    # 구조화된 섹션이 있으면 이를 활용한 fallback 생성
                    if structured_sections:
                        fallback_slides = []
                        # 제목 슬라이드
                        fallback_slides.append(SlideSpec(
                            title="개요", 
                            key_message="핵심 내용 요약", 
                            bullets=["주요 포인트 정리"], 
                            layout="title-and-content"
                        ))
                        
                        # 목차 슬라이드
                        agenda_bullets = [f"{s['number']}. {s['title']}" for s in structured_sections[:6]]
                        fallback_slides.append(SlideSpec(
                            title="목차",
                            key_message="",
                            bullets=agenda_bullets,
                            layout="title-and-content"
                        ))
                        
                        # 각 섹션별 슬라이드
                        for section in structured_sections[:4]:  # 최대 4개 섹션
                            key_msg, bullets = self._parse_section_content(section['content'])
                            fallback_slides.append(SlideSpec(
                                title=f"{section['number']}. {section['title']}",
                                key_message=key_msg,
                                bullets=bullets,
                                layout="title-and-content"
                            ))
                        
                        logger.debug({
                            "phase": "ppt_parse", 
                            "final_topic": improved_topic or "발표자료",
                            "slide_count": len(fallback_slides),
                            "has_chart": False,
                            "ppt_fallback": True,
                            "structured_fallback": True
                        })
                        
                        return DeckSpec(
                            topic=improved_topic or "발표자료", 
                            max_slides=len(fallback_slides), 
                            slides=fallback_slides,
                            theme={"color_scheme": "corporate_blue", "font_style": "modern"}
                        )
                    
                    # 기본 fallback (구조화된 섹션이 없을 때)
                    logger.debug({
                        "phase": "ppt_parse", 
                        "final_topic": improved_topic or "발표자료",
                        "slide_count": 3,
                        "has_chart": False,
                        "ppt_fallback": True,
                        "structured_fallback": False
                    })
                    return DeckSpec(topic=improved_topic or "발표자료", max_slides=3, slides=[
                        SlideSpec(title="개요", key_message="핵심 개요", bullets=["배경", "목표"], layout="title-and-content"),
                        SlideSpec(title="주요 내용", key_message="포인트", bullets=["포인트1", "포인트2"], layout="title-and-content"),
                        SlideSpec(title="결론", key_message="요약", bullets=["요약", "다음 단계"], layout="title-and-content"),
                    ], theme={"color_scheme": "corporate_blue", "font_style": "modern"})
            await asyncio.sleep(0.4)
        # 안전장치 (논리적으로 도달하지 않음)
        return DeckSpec(topic=improved_topic or "발표자료", max_slides=1, slides=[SlideSpec(title="개요", key_message="요약", bullets=["포인트"], layout="title-and-content")])

    def _extract_json(self, text: str) -> str:
        if text.strip().startswith('{'):
            return text
        block = re.search(r"```(?:json)?\n(.*)```", text, re.DOTALL)
        if block:
            return block.group(1)
        brace = re.search(r"{.*}", text, re.DOTALL)
        return brace.group(0) if brace else text

    def _parse_outline(self, text: str, fallback_topic: Optional[str] = None) -> DeckSpec:
        data = json.loads(self._extract_json(text))
        max_slides = int(data.get("max_slides", 10))
        raw_slides = data.get("slides", [])[:max_slides]
        slides: List[SlideSpec] = []
        
        for s in raw_slides:
            diagram_info = s.get("diagram") or {}
            chart = None
            if diagram_info.get("chart"):
                chart_raw = diagram_info.get("chart")
                if isinstance(chart_raw, dict):
                    # 허용된 필드만 전달
                    allowed = {k: v for k, v in chart_raw.items() if k in ChartData.__fields__}
                    chart = ChartData(**allowed)
            raw_data = diagram_info.get("data")
            # Normalize list -> {'items': list} to avoid validation restrictions later
            if isinstance(raw_data, list):
                raw_data = {"items": raw_data}
            diagram = DiagramData(type=diagram_info.get("type", "none"), data=raw_data, chart=chart)
            
            # role 정보를 style에 포함
            style_info = s.get("style") or {}
            role = s.get("role")
            if role:
                style_info["role"] = role
                # role에 따른 추가 플래그 설정
                if role == "title":
                    style_info["title"] = True
                elif role == "agenda":
                    style_info["agenda"] = True
            
            slides.append(SlideSpec(
                title=s.get("title", ""),
                key_message=s.get("key_message", ""),
                bullets=s.get("bullets", []),
                diagram=diagram,
                layout=s.get("layout", "title-and-content"),
                style=style_info if style_info else None,
                visual_suggestion=s.get("visual_suggestion"),
                speaker_notes=s.get("speaker_notes"),
            ))
        
        # 토픽 결정 (번호형 섹션 제목 금지)
        parsed_topic = data.get("topic") or fallback_topic or "발표자료"
        if re.match(r"^\d+\.\s+", parsed_topic):  # 번호형이면 폐기
            parsed_topic = fallback_topic or "발표자료"
        # fallback_topic이 더 구체적이고 질의형이 아니면 교체
        if fallback_topic and not self._is_query_like(fallback_topic) and len(fallback_topic) > len(parsed_topic):
            parsed_topic = fallback_topic
        logger.debug({"phase": "topic_after_parse", "parsed": parsed_topic})
            
        return DeckSpec(
            topic=parsed_topic,
            max_slides=max_slides,
            slides=slides,
            theme=data.get("theme", {"color_scheme": "corporate_blue", "font_style": "modern"})
        )

    def _ensure_structure(self, deck: DeckSpec) -> DeckSpec:
        """보장: [0]=Title, [1]=Agenda. 없으면 생성/재구성."""
        slides = deck.slides
        changed = False
        
        # 1) Title slide 확인 / 재구성
        need_title = True
        if slides:
            first = slides[0]
            # role 또는 layout 또는 style 기반으로 title 슬라이드 판단
            is_title = (
                (first.style and first.style.get('role') == 'title') or
                (first.style and first.style.get('title')) or
                first.layout == 'title-only'
            )
            if is_title:
                need_title = False
        
        if need_title:
            title_slide = SlideSpec(
                title=deck.topic, 
                key_message="", 
                bullets=[], 
                layout="title-only", 
                style={"title": True, "role": "title"}
            )
            slides.insert(0, title_slide)
            deck.max_slides += 1
            changed = True
        
        # 2) Agenda 존재 여부 확인
        has_agenda = False
        agenda_index = -1
        
        # Title 다음 슬라이드들에서 agenda 찾기
        for i in range(1, min(len(slides), 4)):  # 처음 몇 개 슬라이드만 체크
            slide = slides[i]
            is_agenda = (
                (slide.style and slide.style.get('role') == 'agenda') or
                (slide.style and slide.style.get('agenda')) or
                any(k in (slide.title or '').lower() for k in ['목차', 'agenda', 'contents']) and len(slide.bullets) >= 2
            )
            if is_agenda:
                has_agenda = True
                agenda_index = i
                break
        
        # agenda가 없거나 위치가 잘못된 경우 생성/이동
        if not has_agenda and len(slides) >= 2:
            # 번호형 섹션들로 목차 생성
            numbered_sections = []
            for slide in slides[1:]:  # Title 제외
                title = slide.title or ""
                # 번호형 섹션이면서 목차/agenda가 아닌 경우만 포함
                if re.match(r"^\d+\.\s+", title) and not any(k in title.lower() for k in ['목차', 'agenda', 'contents']):
                    numbered_sections.append(title)
            
            if numbered_sections:
                agenda_slide = SlideSpec(
                    title="목차",
                    key_message="",
                    bullets=numbered_sections[:12],
                    layout="title-and-content",
                    style={"agenda": True, "role": "agenda"}
                )
                slides.insert(1, agenda_slide)
                deck.max_slides += 1
                changed = True
        
        elif has_agenda and agenda_index != 1:
            # agenda가 있지만 위치가 잘못된 경우 이동
            agenda_slide = slides.pop(agenda_index)
            slides.insert(1, agenda_slide)
            changed = True
        
        if changed:
            logger.info({"phase": "structure_fix", "slide_count": len(slides), "has_title": True, "has_agenda": True})
        
        return deck

    def _post_process_deck(self, deck: DeckSpec, improved_topic: str) -> DeckSpec:
        """후처리: 슬라이드 병합, 압축, 자동 테이블 추가"""
        # 구조 보정 우선
        deck = self._ensure_structure(deck)

        # 슬라이드 개수가 과다한 경우 압축 (> 9개)
        if len(deck.slides) > 9:
            logger.debug(f"슬라이드 압축: {len(deck.slides)}개 → 압축 시도")
            deck = self._compress_slides(deck)

        # 빈약한 슬라이드 통합
        deck = self._merge_weak_slides(deck)

        # 자동 테이블 슬라이드 추가 (키-값 블록 감지)
        deck = self._add_auto_tables(deck)

        return deck
    
    def _compress_slides(self, deck: DeckSpec) -> DeckSpec:
        """슬라이드 개수 압축 (9개 → 6개 목표)"""
        if len(deck.slides) <= 6:
            return deck
            
        # 유사한 제목/키워드 기반 그룹핑 시도
        groups = []
        current_group = []
        
        for slide in deck.slides:
            if len(current_group) == 0:
                current_group.append(slide)
            elif len(current_group) >= 2:  # 그룹 완성
                groups.append(current_group)
                current_group = [slide]
            else:
                # 유사성 검사 (간단한 키워드 기반)
                if self._slides_similar(current_group[-1], slide):
                    current_group.append(slide)
                else:
                    groups.append(current_group)
                    current_group = [slide]
        
        if current_group:
            groups.append(current_group)
        
        # 그룹을 슬라이드로 병합
        merged_slides = []
        for group in groups:
            if len(group) == 1:
                merged_slides.append(group[0])
            else:
                merged = self._merge_slide_group(group)
                merged_slides.append(merged)
        
        deck.slides = merged_slides[:6]  # 최대 6개로 제한
        logger.debug(f"슬라이드 압축 완료: {len(merged_slides)}개")
        return deck
    
    def _slides_similar(self, slide1: SlideSpec, slide2: SlideSpec) -> bool:
        """슬라이드 유사성 판단"""
        # 간단한 키워드 기반 유사성
        keywords1 = set(re.findall(r'[가-힣A-Za-z]+', slide1.title + ' ' + slide1.key_message))
        keywords2 = set(re.findall(r'[가-힣A-Za-z]+', slide2.title + ' ' + slide2.key_message))
        
        if not keywords1 or not keywords2:
            return False
            
        intersection = keywords1.intersection(keywords2)
        union = keywords1.union(keywords2)
        
        return len(intersection) / len(union) > 0.3  # 30% 이상 겹치면 유사
    
    def _merge_slide_group(self, slides: List[SlideSpec]) -> SlideSpec:
        """슬라이드 그룹을 하나로 병합"""
        if len(slides) == 1:
            return slides[0]
        
        # 첫 번째 슬라이드 기반으로 병합
        base = slides[0]
        merged_bullets = base.bullets.copy()
        
        # 다른 슬라이드의 bullets 추가 (중복 제거)
        for slide in slides[1:]:
            for bullet in slide.bullets:
                if bullet not in merged_bullets and len(merged_bullets) < 8:
                    merged_bullets.append(bullet)
        
        # 제목 조합
        titles = [s.title for s in slides if s.title.strip()]
        merged_title = titles[0] if titles else "통합 섹션"
        
        return SlideSpec(
            title=merged_title,
            key_message=base.key_message,
            bullets=merged_bullets,
            diagram=base.diagram,
            layout=base.layout,
            visual_suggestion=base.visual_suggestion,
            speaker_notes=base.speaker_notes
        )
    
    def _merge_weak_slides(self, deck: DeckSpec) -> DeckSpec:
        """빈약한 슬라이드들 통합"""
        if len(deck.slides) <= 3:
            return deck
            
        strong_slides = []
        weak_slides = []
        
        for slide in deck.slides:
            # 빈약함 기준: bullets가 2개 이하이고 key_message가 짧음
            protected = (slide.layout == 'title-only') or (slide.style and (slide.style.get('agenda') or slide.style.get('title')))
            if not protected and len(slide.bullets) <= 2 and len(slide.key_message) < 20:
                weak_slides.append(slide)
            else:
                strong_slides.append(slide)
        
        # 빈약한 슬라이드들을 strong 슬라이드에 병합 또는 별도 그룹화
        if weak_slides and len(weak_slides) >= 2:
            merged_weak = self._merge_slide_group(weak_slides)
            merged_weak.title = "기타 주요 사항"
            strong_slides.append(merged_weak)
        elif weak_slides:
            # 1개뿐이면 병합 고려 (protected 아닌 경우만)
            ws = weak_slides[0]
            if strong_slides and ws.layout != 'title-only' and not (ws.style and ws.style.get('agenda')):
                weak_bullets = [b for b in ws.bullets if b not in strong_slides[0].bullets]
                strong_slides[0].bullets.extend(weak_bullets[:3])
            else:
                strong_slides.append(ws)
        
        deck.slides = strong_slides
        return deck
    
    def _add_auto_tables(self, deck: DeckSpec) -> DeckSpec:
        """자동 테이블 슬라이드 추가 (키-값 패턴 감지)"""
        # 현재 슬라이드들의 bullets에서 키-값 패턴 찾기
        for idx, slide in enumerate(deck.slides):
            if slide.diagram and slide.diagram.type != "none":
                continue  # 이미 다이어그램이 있으면 스킵
            # Title / Agenda 슬라이드는 제외
            if idx == 0 or (slide.style and slide.style.get('agenda')):
                continue
                
            # bullets에서 키:값 패턴 찾기
            kv_items = []
            for bullet in slide.bullets:
                kv_match = re.match(r'^([^:]{1,20}):\s*(.{1,50})$', bullet.strip())
                if kv_match:
                    key, value = kv_match.groups()
                    kv_items.append({"key": key.strip(), "value": value.strip()})
            
            # 변환 조건 강화: 3개 이상 & 전체 bullets의 60% 이상 & value 평균 길이 <= 25
            if kv_items:
                avg_val_len = sum(len(i['value']) for i in kv_items)/len(kv_items)
            else:
                avg_val_len = 0
            ratio = (len(kv_items) / max(1, len(slide.bullets)))
            # 짧은 라벨형 (value 평균 3 미만) 은 설명 bullet 가능성 → 제외
            if len(kv_items) >= 3 and ratio >= 0.6 and 3 <= avg_val_len <= 25:
                slide.diagram = DiagramData(
                    type="table",
                    data={"items": kv_items},
                    chart=None
                )
                slide.layout = "two-content"  # 테이블 + 제목 레이아웃
                # bullets는 테이블로 이동했으므로 요약으로 교체
                slide.bullets = [f"총 {len(kv_items)}개 항목", "세부사항은 표 참조"]
                logger.debug({
                    "phase": "auto_table",
                    "title": slide.title,
                    "count": len(kv_items),
                    "ratio": round(ratio,2),
                    "avg_val_len": round(avg_val_len,1)
                })
        
        return deck

    # ---------------- PPT Building ----------------
    def build_enhanced_pptx(self, spec: DeckSpec, file_basename: Optional[str] = None,
                            template_style: str = "business", include_charts: bool = True,
                            custom_template_path: Optional[str] = None,
                            user_template_id: Optional[str] = None,
                            text_box_mappings: Optional[List[Dict[str, Any]]] = None,
                            content_segments: Optional[List[Dict[str, Any]]] = None) -> str:
        """Enhanced PPT 빌드 (템플릿 지원)"""
        logger.debug({
            'phase': 'build_entry',
            'custom_template_path': custom_template_path,
            'custom_template_exists': bool(custom_template_path and Path(custom_template_path).exists()),
            'user_template_id': user_template_id,
            'topic': spec.topic,
            'slides': len(spec.slides),
            'text_box_mappings_count': len(text_box_mappings) if text_box_mappings else 0,
            'content_segments_count': len(content_segments) if content_segments else 0
        })
        
        # 🎯 동적 템플릿 ID 로깅
        if user_template_id:
            logger.info(f"✅ 사용자 동적 템플릿 적용: {user_template_id}")
        if custom_template_path:
            logger.info(f"📁 사용될 템플릿 경로: {custom_template_path}")
        # 커스텀 템플릿 경로가 있으면 템플릿 기반 빌드 사용
        if custom_template_path and Path(custom_template_path).exists():
            logger.info(f"템플릿 적용 시도: {custom_template_path}")
            return self._build_from_template(spec, Path(custom_template_path), file_basename, 
                                           text_box_mappings, content_segments)
        elif custom_template_path:
            logger.warning(f"지정된 템플릿 경로가 존재하지 않아 레거시 모드로 진행: {custom_template_path}")
        
        # 기존 로직 (레거시 호환)
        return self._build_legacy_pptx(spec, file_basename, template_style, include_charts)
    
    def _build_from_template(self, spec: DeckSpec, template_path: Path, file_basename: Optional[str] = None,
                           text_box_mappings: Optional[List[Dict[str, Any]]] = None,
                           content_segments: Optional[List[Dict[str, Any]]] = None) -> str:
        """템플릿 기반 PPT 빌드 (매핑 정보 활용)"""
        try:
            logger.info(f"템플릿 기반 빌드 시작: {template_path}")
            logger.info(f"매핑 정보: {len(text_box_mappings) if text_box_mappings else 0}개")
            
            # 원본 템플릿 로드
            prs = Presentation(str(template_path))
            
            # 매핑 정보가 있으면 확장된 오브젝트 처리기 사용
            if text_box_mappings:
                prs = self.object_processor.apply_object_mappings(prs, text_box_mappings, content_segments)
            else:
                # 기존 템플릿 매니저 로직 사용
                template_spec = self.template_manager.analyze_template(template_path)
                if template_spec:
                    adapted_spec = self.template_manager.map_deck_to_template(spec, template_spec)
                    prs = self.template_manager.build_from_template(adapted_spec, template_path)
            
            # 파일 저장
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            basename = file_basename or f"presentation_template_{timestamp}"
            if not basename.endswith('.pptx'):
                basename += '.pptx'
            
            output_path = self.upload_dir / basename
            prs.save(str(output_path))
            
            logger.info(f"템플릿 기반 PPT 생성 완료: {output_path}")
            return str(output_path)
            
        except Exception as e:
            logger.error(f"템플릿 기반 빌드 실패: {e}")
            # 폴백: 레거시 모드
            return self._build_legacy_pptx(spec, file_basename, "business", True)

    def _apply_text_box_mappings(self, prs, mappings: List[Dict[str, Any]], 
                                segments: Optional[List[Dict[str, Any]]] = None):
        """템플릿에 텍스트박스 매핑 적용"""
        try:
            logger.info(f"텍스트박스 매핑 적용 시작: {len(mappings)}개 매핑")
            
            # 매핑 정보 상세 로깅
            for i, mapping in enumerate(mappings):
                logger.info(f"매핑 {i}: {mapping}")
            
            # 매핑 정보를 슬라이드별로 그룹화
            mappings_by_slide = {}
            for mapping in mappings:
                slide_idx = mapping.get('slideIndex', 0)
                if slide_idx not in mappings_by_slide:
                    mappings_by_slide[slide_idx] = []
                mappings_by_slide[slide_idx].append(mapping)
            
            # 각 슬라이드에 매핑 적용
            for slide_idx, slide_mappings in mappings_by_slide.items():
                if slide_idx < len(prs.slides):
                    slide = prs.slides[slide_idx]
                    self._apply_mappings_to_slide(slide, slide_mappings)
                    logger.info(f"슬라이드 {slide_idx}에 {len(slide_mappings)}개 매핑 적용")
            
            return prs
            
        except Exception as e:
            logger.error(f"텍스트박스 매핑 적용 실패: {e}")
            return prs
    
    def _apply_mappings_to_slide(self, slide, mappings: List[Dict[str, Any]]):
        """개별 슬라이드에 매핑 적용"""
        try:
            for mapping in mappings:
                element_id = mapping.get('elementId')
                assigned_content = mapping.get('assignedContent', '')
                use_original = mapping.get('useOriginal', False)
                content_source = mapping.get('contentSource', '')
                
                # useOriginal이 True이거나 contentSource가 'keep_original'인 경우 원본 템플릿 내용을 유지
                # 단, assignedContent가 originalContent와 다른 경우에는 변경을 적용
                if use_original or content_source == 'keep_original':
                    original_content = mapping.get('originalContent', '')
                    if assigned_content == original_content:
                        logger.debug(f"원본 내용과 동일하여 건너뜀: {element_id} (useOriginal: {use_original}, contentSource: {content_source})")
                        continue
                    else:
                        logger.info(f"내용이 변경되어 적용: {element_id} - '{original_content}' -> '{assigned_content}'")
                        # 변경된 내용이 있으므로 계속 진행
                
                if not element_id or not assigned_content:
                    continue
                
                # 슬라이드에서 해당 ID를 가진 shape 찾기
                target_shape = None
                logger.info(f"슬라이드에서 elementId '{element_id}' 찾는 중...")
                logger.info(f"슬라이드의 shapes 목록:")
                for j, shape in enumerate(slide.shapes):
                    shape_name = getattr(shape, 'name', f'unnamed_{j}')
                    logger.info(f"  Shape {j}: name='{shape_name}', has_text_frame={hasattr(shape, 'text_frame')}")
                
                for shape in slide.shapes:
                    if hasattr(shape, 'name') and shape.name == element_id:
                        target_shape = shape
                        logger.info(f"Shape 이름으로 찾음: {element_id}")
                        break
                    # ID가 textbox-0-0 형식인 경우 인덱스로도 찾아보기
                    if element_id.startswith('textbox-'):
                        parts = element_id.split('-')
                        if len(parts) == 3 and parts[2].isdigit():
                            shape_idx = int(parts[2])
                            text_shapes = [s for s in slide.shapes if hasattr(s, 'text_frame')]
                            logger.info(f"텍스트박스 인덱스 검색: {shape_idx}, 총 텍스트 shape: {len(text_shapes)}")
                            if shape_idx < len(text_shapes):
                                target_shape = text_shapes[shape_idx]
                                logger.info(f"Shape 인덱스로 찾음: {element_id} -> shape {shape_idx}")
                                break
                
                # shape를 찾았으면 텍스트 적용
                if target_shape and hasattr(target_shape, 'text_frame'):
                    try:
                        # 기존 스타일 정보 백업
                        original_font_style = None
                        original_paragraph_style = None
                        
                        if target_shape.text_frame.paragraphs:
                            first_para = target_shape.text_frame.paragraphs[0]
                            if first_para.runs:
                                first_run = first_para.runs[0]
                                # 폰트 스타일 백업
                                original_font_style = {
                                    'size': first_run.font.size,
                                    'bold': first_run.font.bold,
                                    'name': first_run.font.name,
                                    'color': first_run.font.color,
                                }
                            # 단락 스타일 백업
                            original_paragraph_style = {
                                'alignment': first_para.alignment,
                                'level': first_para.level
                            }
                        
                        # 텍스트 내용 업데이트
                        target_shape.text_frame.clear()
                        p = target_shape.text_frame.paragraphs[0]
                        p.text = assigned_content
                        
                        # 백업된 스타일 복원
                        if p.runs and original_font_style:
                            run = p.runs[0]
                            try:
                                if original_font_style.get('size'):
                                    run.font.size = original_font_style['size']
                                if original_font_style.get('bold') is not None:
                                    run.font.bold = original_font_style['bold']
                                if original_font_style.get('name'):
                                    run.font.name = original_font_style['name']
                                
                                # 색상 복원 (개선됨)
                                if original_font_style.get('color'):
                                    try:
                                        from pptx.dml.color import RGBColor
                                        from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
                                        
                                        original_color = original_font_style['color']
                                        if hasattr(original_color, 'type'):
                                            if original_color.type == MSO_COLOR_TYPE.RGB:
                                                run.font.color.rgb = original_color.rgb
                                                logger.info(f"RGB 색상 복원: {original_color.rgb}")
                                            elif original_color.type == MSO_COLOR_TYPE.THEME:
                                                # 테마 색상 처리 개선
                                                try:
                                                    if hasattr(original_color, 'theme_color'):
                                                        run.font.color.theme_color = original_color.theme_color
                                                        logger.info(f"테마 색상 복원: {original_color.theme_color}")
                                                    else:
                                                        # 테마 색상이 BACKGROUND_1(흰색)인 경우 직접 RGB로 설정
                                                        run.font.color.rgb = RGBColor(255, 255, 255)
                                                        logger.info("테마 색상을 RGB 흰색으로 변환")
                                                except Exception as theme_error:
                                                    # 테마 색상 설정 실패시 RGB 흰색으로 설정
                                                    run.font.color.rgb = RGBColor(255, 255, 255)
                                                    logger.info(f"테마 색상 실패, RGB 흰색으로 설정: {theme_error}")
                                            else:
                                                # 기본적으로 흰색 설정
                                                run.font.color.rgb = RGBColor(255, 255, 255)
                                                logger.info("알 수 없는 색상 타입, 흰색으로 설정")
                                        else:
                                            # 색상 객체를 직접 복사
                                            run.font.color.rgb = RGBColor(255, 255, 255)
                                            logger.info("색상 타입 없음, 흰색으로 설정")
                                    except Exception as color_error:
                                        # 색상 복원 실패 시 흰색으로 설정
                                        from pptx.dml.color import RGBColor
                                        run.font.color.rgb = RGBColor(255, 255, 255)
                                        logger.info(f"색상 복원 실패, 흰색으로 설정: {color_error}")
                            except Exception as style_error:
                                logger.debug(f"스타일 복원 부분 실패: {style_error}")
                        
                        if original_paragraph_style:
                            try:
                                if original_paragraph_style.get('alignment') is not None:
                                    p.alignment = original_paragraph_style['alignment']
                                if original_paragraph_style.get('level') is not None:
                                    p.level = original_paragraph_style['level']
                            except Exception as para_error:
                                logger.debug(f"단락 스타일 복원 실패: {para_error}")
                        
                        logger.debug(f"텍스트 및 스타일 적용 성공: {element_id} -> '{assigned_content[:50]}...'")
                    except Exception as text_error:
                        logger.warning(f"텍스트 적용 실패 {element_id}: {text_error}")
                else:
                    logger.warning(f"Shape을 찾을 수 없음: {element_id}")
                    
        except Exception as e:
            logger.error(f"슬라이드 매핑 적용 실패: {e}")
    
    def _build_legacy_pptx(self, spec: DeckSpec, file_basename: Optional[str] = None,
                          template_style: str = "business", include_charts: bool = True) -> str:
        """기존 PPT 빌드 로직 (레거시)"""
        
        # --- DeckSpec 후처리 ---
        improved_topic = file_basename if file_basename and len(file_basename) > 3 else spec.topic
        processed_spec = self._post_process_deck(spec, improved_topic)
        logger.info(f"후처리된 슬라이드 수: {len(processed_spec.slides)}, 원본: {len(spec.slides)}")

        # --- 템플릿 로드 ---
        prs = None
        tpl_setting = getattr(settings, 'ppt_template_path', None)
        candidates: List[Path] = [Path(tpl_setting)] if tpl_setting else []
        # 올바른 루트 uploads/templates 위치 탐색 (parents[4] 우선)
        try:
            root_candidate = Path(__file__).parents[4]
        except Exception:
            root_candidate = Path(__file__).parents[3]
        root_tpl = root_candidate / 'uploads' / 'templates' / 'ppt_template.pptx'
        backend_tpl = Path(__file__).parents[3] / 'uploads' / 'templates' / 'ppt_template.pptx'
        if root_tpl.exists():
            candidates.append(root_tpl)
        candidates.append(backend_tpl)
        logger.debug({
            'phase': 'legacy_template_candidates',
            'candidates': [str(c) for c in candidates]
        })
        for c in candidates:
            if c and c.exists():
                try:
                    prs = Presentation(str(c))
                    logger.info(f"템플릿 사용: {c}")
                    break
                except Exception as e:  # noqa: BLE001
                    logger.warning(f"템플릿 로드 실패 {c}: {e}")
        if prs is None:
            prs = Presentation()
            logger.debug("No PPT base template found; using blank Presentation()")

        style_to_theme = {"business": "corporate_blue", "minimal": "professional_gray", "modern": "modern_green", "playful": "playful_violet"}
        chosen_theme_any = style_to_theme.get(template_style) or (processed_spec.theme.get("color_scheme") if processed_spec.theme else "corporate_blue")
        chosen_theme = str(chosen_theme_any or "corporate_blue")
        colors = self.color_themes.get(chosen_theme, self.color_themes["corporate_blue"])

        if not include_charts:
            for s in processed_spec.slides:
                if s.diagram and s.diagram.type == 'chart':
                    s.diagram.type = 'none'
                    s.diagram.chart = None

        chart_palettes = {
            "corporate_blue": [RGBColor(0,102,204), RGBColor(102,153,255), RGBColor(255,153,0), RGBColor(0,176,80), RGBColor(112,48,160)],
            "professional_gray": [RGBColor(70,70,70), RGBColor(120,120,120), RGBColor(169,169,169), RGBColor(220,20,60), RGBColor(100,149,237)],
            "modern_green": [RGBColor(34,139,34), RGBColor(60,179,113), RGBColor(143,188,143), RGBColor(255,215,0), RGBColor(46,139,87)],
            "playful_violet": [RGBColor(111,45,168), RGBColor(181,126,220), RGBColor(255,181,71), RGBColor(0,153,255), RGBColor(255,105,180)],
        }
        palette = chart_palettes.get(chosen_theme, chart_palettes["corporate_blue"])

        for idx, slide_spec in enumerate(processed_spec.slides[:processed_spec.max_slides]):
            # 레이아웃 선택
            if idx == 0:
                layout = prs.slide_layouts[0]
            else:
                if slide_spec.diagram and slide_spec.diagram.type == 'chart' and slide_spec.layout not in ['two-content', 'title-only']:
                    slide_spec.layout = 'two-content'
                layout_map = {"title-only":5, "title-and-content":1, "two-content":3, "section-header":2, "blank":6}
                layout = prs.slide_layouts[layout_map.get(slide_spec.layout, 1)]
            slide = prs.slides.add_slide(layout)

            # 모던/플레이풀 배경 강조
            if template_style in ['modern', 'playful']:
                try:
                    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
                    bg.fill.solid()
                    try:
                        bg.fill.fore_color.rgb = colors['background']
                    except Exception:
                        pass
                    try:
                        bg.line.fill.background()
                    except Exception:
                        pass
                except Exception:
                    pass

            # 제목 슬라이드
            if idx == 0:
                if slide.shapes.title:
                    clean_title = self._extract_clean_title(improved_topic or processed_spec.topic)
                    slide.shapes.title.text = clean_title
                    logger.info(f"제목 슬라이드 제목 설정: '{clean_title}' (원본: '{processed_spec.topic}')")
                    f = slide.shapes.title.text_frame.paragraphs[0].font
                    f.color.rgb = colors['primary']; f.size = Pt(42); f.bold = True
                # Subtitle: try to safely get a secondary placeholder or create one
                try:
                    subtitle_text = ""
                    first = processed_spec.slides[0] if processed_spec.slides else None
                    if first:
                        subtitle_text = first.key_message or (first.bullets[0] if first.bullets else "")

                    tf = None
                    try:
                        if len(slide.placeholders) > 1:
                            sub = slide.placeholders[1]
                            if getattr(sub, 'has_text_frame', False) and hasattr(sub, 'text_frame'):
                                tf = getattr(sub, 'text_frame')  # type: ignore[attr-defined]
                    except Exception:
                        tf = None

                    if tf is None:
                        # fallback: find any non-title text shape
                        for sh in slide.shapes:
                            try:
                                if getattr(sh, 'has_text_frame', False) and sh is not slide.shapes.title:
                                    tf = sh.text_frame
                                    break
                            except Exception:
                                continue

                    if tf is None:
                        # create a subtitle textbox
                        from pptx.util import Inches
                        try:
                            sub_box = slide.shapes.add_textbox(Inches(1), Inches(1.7), Inches(9.5), Inches(1))
                            tf = sub_box.text_frame
                        except Exception:
                            tf = None

                    if tf is not None:
                        try:
                            tf.clear()
                        except Exception:
                            pass
                        if subtitle_text:
                            p = tf.paragraphs[0]
                            p.text = subtitle_text[:120]
                            try:
                                p.font.size = Pt(22); p.font.color.rgb = colors['text']
                            except Exception:
                                pass
                except Exception:
                    pass
                continue

            # 일반 제목 슬라이드 설정
            if slide.shapes.title:
                slide.shapes.title.text = slide_spec.title or processed_spec.topic
                ft = slide.shapes.title.text_frame.paragraphs[0].font
                ft.color.rgb = colors['primary']; ft.size = Pt(30 if slide_spec.layout=='section-header' else 26); ft.bold = True

            # 섹션 헤더 정리
            if slide_spec.layout == 'section-header':
                try:
                    for ph in slide.placeholders:
                        if getattr(ph,'placeholder_format',None) and ph.placeholder_format.type not in (1,):
                            if getattr(ph,'has_text_frame',False) and hasattr(ph,'text_frame'):
                                try:
                                    getattr(ph,'text_frame').clear()  # type: ignore[attr-defined]
                                except Exception:
                                    pass
                except Exception:
                    pass

            # 컨텐츠 / 이미지
            if slide_spec.layout in ['title-and-content','two-content']:
                self._add_content(slide, slide_spec, colors, template_style)
                if slide_spec.visual_suggestion:
                    self._maybe_add_image(slide, slide_spec.visual_suggestion)

            # 차트/다이어그램
            if include_charts and slide_spec.diagram and slide_spec.diagram.type != 'none':
                self._add_diagram(slide, slide_spec.diagram, colors, palette)

            # footer & page numbers
            try:
                if idx >= 1:
                    from pptx.util import Inches  # 명시적 import로 스코프 문제 해결
                    margin_x = Inches(0.6)
                    footer_text_top = Inches(6.6)
                    width_main = Inches(9.5)
                    box = slide.shapes.add_textbox(margin_x, footer_text_top, width_main, Inches(0.25))
                    tf = box.text_frame; tf.clear(); p = tf.paragraphs[0]
                    p.text = f"{processed_spec.topic} | {datetime.now().strftime('%Y-%m-%d')}"
                    p.font.size = Pt(8); p.font.color.rgb = colors['text']
                    page_box_width = Inches(1.5)
                    page_left = Inches(11.2)
                    pbox = slide.shapes.add_textbox(page_left, footer_text_top, page_box_width, Inches(0.25))
                    ptf = pbox.text_frame; ptf.clear(); p2 = ptf.paragraphs[0]
                    p2.text = f"{idx+1}/{processed_spec.max_slides}"
                    p2.font.size = Pt(8); p2.font.color.rgb = colors['text']; p2.alignment = PP_ALIGN.RIGHT
            except Exception as e:
                logger.debug(f"Footer 위치 조정 실패: {e}")

        # --- 파일명 및 저장 ---
        def _sanitize_name(text: str) -> str:
            cleaned = re.sub(r"\[\[PPT_OPTS:.*?\]\]", "", text)
            cleaned = re.sub(r"[\n\r\t]", " ", cleaned).strip() or "presentation"
            cleaned = re.sub(r"[^0-9A-Za-z가-힣 _-]", "_", cleaned)
            cleaned = re.sub(r"\s+", "_", cleaned)[:40]
            return cleaned or "presentation"
        base_name = _sanitize_name(spec.topic)
        if file_basename:
            fb_clean = _sanitize_name(file_basename)
            if re.search(r"(해주세요|해줘|알려줘|알려주세요|\?$)", fb_clean):
                fname = f"enhanced_presentation_{base_name}"
            else:
                fname = fb_clean
        else:
            fname = f"enhanced_presentation_{base_name}"
        out_path = self.upload_dir / f"{fname}.pptx"
        prs.save(str(out_path))
        logger.info(f"PPTX 생성 완료: {out_path}")
        return str(out_path)

    # ================= Quick pipeline additions =================
    def generate_fixed_outline(self, topic: str, context_text: str, max_slides: int = 8) -> DeckSpec:
        """원클릭(디자인 무시)용 고정 구조 아웃라인 생성"""
        try:
            logger.info(f"🚀 원클릭 고정 구조 생성 시작: topic='{topic[:50]}', max_slides={max_slides}")
            logger.info(f"📝 입력 컨텍스트 길이: {len(context_text)} 문자")
            
            max_slides = max(3, min(max_slides, 20))
            
            # 더 강력한 섹션 추출 로직
            lines = [ln.strip() for ln in (context_text or "").split("\n") if ln.strip()]
            logger.info(f"📄 총 라인 수: {len(lines)}")
            
            # 1) 명시적 헤딩 패턴 찾기
            headings = []
            for ln in lines:
                if (ln.startswith(('#', '##', '###', '####')) or 
                    ln.endswith(':') or 
                    (len(ln) <= 50 and any(word in ln for word in ['배경', '목표', '현황', '과제', '방안', '결론', '요약']))):
                    headings.append(ln)
            
            # 2) 헤딩이 부족하면 문장 기반으로 섹션 생성
            if len(headings) < 2:
                sentences = [s.strip() for s in context_text.split('.') if s.strip() and len(s.strip()) > 10]
                headings = []
                for i, sent in enumerate(sentences[:max_slides-3]):
                    if len(sent) <= 60:
                        headings.append(sent)
                    else:
                        # 긴 문장은 요약해서 제목으로 사용
                        words = sent.split()[:6]
                        headings.append(' '.join(words) + '...')
            
            logger.info(f"🎯 추출된 헤딩 수: {len(headings)}")
            for i, h in enumerate(headings):
                logger.info(f"  헤딩 {i+1}: '{h[:30]}...'")
            
            sections = []
            content_lines = [ln for ln in lines if ln not in headings]
            
            for i, h in enumerate(headings[:max(0, max_slides-3)]):
                title = h.lstrip('#').strip(':').strip()
                
                # 각 섹션의 키 메시지와 불릿 생성
                key_msg = f"{title}에 대한 주요 내용입니다."
                bullets = []
                
                # 관련 컨텐츠 라인 찾기 (헤딩 다음 2-3개 라인)
                start_idx = i * 2
                for j in range(start_idx, min(start_idx + 3, len(content_lines))):
                    if j < len(content_lines) and content_lines[j]:
                        bullets.append(content_lines[j][:80])
                
                if not bullets:
                    bullets = [f"{title} 관련 세부 사항", "주요 포인트 및 고려사항", "실행 방안 및 기대효과"]
                
                sections.append({
                    "title": title or f"섹션 {i+1}", 
                    "key_message": key_msg, 
                    "bullets": bullets[:3]  # 최대 3개 불릿
                })
            
            logger.info(f"📋 생성된 섹션 수: {len(sections)}")
            
            slides: List[SlideSpec] = []
            
            # 1) 제목 슬라이드
            slides.append(SlideSpec(title=topic or "발표자료", key_message="", bullets=[]))
            logger.info("✅ 제목 슬라이드 생성")
            
            # 2) 목차 슬라이드
            slides.append(SlideSpec(title="목차", key_message="", bullets=[s["title"] for s in sections]))
            logger.info("✅ 목차 슬라이드 생성")
            
            # 3) 내용 슬라이드들
            for s in sections:
                slides.append(SlideSpec(title=s["title"], key_message=s.get("key_message", ""), bullets=s.get("bullets", [])))
                logger.info(f"✅ 내용 슬라이드 생성: '{s['title'][:20]}...'")
            
            # 4) 종료 슬라이드
            slides.append(SlideSpec(title="감사합니다", key_message="경청해 주셔서 감사합니다.", bullets=[]))
            logger.info("✅ 종료 슬라이드 생성")
            
            deck = DeckSpec(topic=topic or "발표자료", slides=slides, max_slides=len(slides))
            logger.info(f"🎉 고정 구조 DeckSpec 생성 완료: 총 {len(slides)}개 슬라이드")
            return deck
            
        except Exception as e:
            logger.error(f"generate_fixed_outline 실패: {e}")
            # 폴백: 최소한의 구조
            fallback_slides = [
                SlideSpec(title=topic or "발표자료", key_message="", bullets=[]),
                SlideSpec(title="내용", key_message="주요 내용을 다룹니다.", bullets=["세부사항 1", "세부사항 2", "세부사항 3"]),
                SlideSpec(title="감사합니다", key_message="경청해 주셔서 감사합니다.", bullets=[])
            ]
            logger.info(f"⚠️ 폴백 구조 사용: {len(fallback_slides)}개 슬라이드")
            return DeckSpec(topic=topic or "발표자료", slides=fallback_slides, max_slides=len(fallback_slides))

    def build_quick_pptx(self, spec: DeckSpec, file_basename: Optional[str] = None) -> str:
        """원클릭 전용 빌더: 템플릿/매핑 비적용, 레거시 빌드 강제"""
        logger.info(f"🏗️ 원클릭 PPT 빌드 시작: {len(spec.slides)}개 슬라이드, topic='{spec.topic}'")
        result = self._build_legacy_pptx(spec, file_basename=file_basename, template_style="business", include_charts=False)
        logger.info(f"✅ 원클릭 PPT 빌드 완료: {result}")
        return result

    # ---------------- Content Helpers ----------------
    def _add_content(self, slide, spec: SlideSpec, colors, template_style: str):
        # Be resilient: templates may have 0/1 placeholders; try to find a usable text_frame
        tf = None
        try:
            if len(slide.placeholders) > 1:
                ph = slide.placeholders[1]
                if getattr(ph, 'has_text_frame', False) and hasattr(ph, 'text_frame'):
                    tf = ph.text_frame
        except Exception:
            tf = None

        if tf is None:
            # fallback: find first non-title shape with text_frame
            try:
                for sh in slide.shapes:
                    if getattr(sh, 'has_text_frame', False):
                        # skip title shape if present
                        try:
                            if sh is slide.shapes.title:
                                continue
                        except Exception:
                            pass
                        tf = sh.text_frame
                        break
            except Exception:
                tf = None

        if tf is None:
            # If still no text frame, we'll create one later when needed
            created_box = None
        else:
            try:
                tf.clear()
                tf.word_wrap = True
            except Exception:
                pass
        agenda_mode = (spec.style and spec.style.get('agenda')) or (spec.title in ['목차','Agenda','Contents'] and not spec.key_message)
        if agenda_mode:
            # 목차: 동일 레벨 나열
            if tf:
                tf.clear()
                for i, b in enumerate(spec.bullets):
                    txt = b.strip()
                    if len(txt) > 60: txt = txt[:60] + '…'
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = txt
                    p.level = 0
                    p.font.size = Pt(20)
                    p.font.color.rgb = colors['text']
        else:
            # 일반 콘텐츠 슬라이드: 키 메시지 + 불릿 포인트
            if tf:
                tf.clear()
                logger.info(f"📝 콘텐츠 추가: '{spec.title}' - key_message='{spec.key_message}', bullets={len(spec.bullets)}개")
                
                # 1) 키 메시지 추가 (첫 번째 패러그래프)
                if spec.key_message and spec.key_message.strip():
                    p = tf.paragraphs[0]
                    p.text = spec.key_message.strip()
                    p.level = 0
                    try:
                        p.font.size = Pt(22)
                        p.font.bold = True
                        p.font.color.rgb = colors['text']
                    except Exception:
                        pass
                    logger.info(f"✅ 키 메시지 추가됨: '{spec.key_message[:30]}...'")
                
                # 2) 불릿 포인트 추가
                for i, bullet in enumerate(spec.bullets):
                    if bullet and bullet.strip():
                        txt = bullet.strip()
                        if len(txt) > 80: 
                            txt = txt[:80] + '…'
                        
                        # 키 메시지가 있으면 두 번째부터, 없으면 첫 번째부터
                        if spec.key_message and spec.key_message.strip():
                            p = tf.add_paragraph()
                        else:
                            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        
                        p.text = f"• {txt}"
                        p.level = 1
                        try:
                            p.font.size = Pt(18)
                            p.font.color.rgb = colors['text']
                        except Exception:
                            pass
                        logger.info(f"✅ 불릿 {i+1} 추가됨: '{txt[:30]}...'")
                
                logger.info(f"🎯 '{spec.title}' 슬라이드 콘텐츠 추가 완료")
            else:
                logger.warning(f"⚠️ '{spec.title}' 슬라이드에 텍스트 프레임을 찾을 수 없습니다")
        if spec.speaker_notes:
            try:
                notes = slide.part.notes_slide.notes_text_frame
                notes.text = spec.speaker_notes[:1500]
            except Exception:
                pass

    def _add_diagram(self, slide, diagram: DiagramData, colors, palette: List[RGBColor]):
        try:
            if diagram.type == 'chart' and diagram.chart:
                self._add_chart(slide, diagram.chart, colors, palette)
            elif diagram.type == 'table' and diagram.data:
                self._add_table(slide, diagram.data, colors)
            elif diagram.type == 'flow' and diagram.data:
                self._add_flow(slide, diagram.data, colors)
        except Exception as e:  # noqa: BLE001
            logger.warning(f"다이어그램 추가 실패: {e}")

    def _add_chart(self, slide, chart_data: ChartData, colors, palette: List[RGBColor]):
        if not chart_data.categories or not chart_data.series:
            return
        data = CategoryChartData(); data.categories = chart_data.categories
        for s in chart_data.series:
            data.add_series(s.get('name','Series'), s.get('values', []))
        chart_type = {
            'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'bar': XL_CHART_TYPE.BAR_CLUSTERED,
            'line': XL_CHART_TYPE.LINE,
            'pie': XL_CHART_TYPE.PIE
        }.get(chart_data.type, XL_CHART_TYPE.COLUMN_CLUSTERED)
        x,y,cx,cy = Inches(4), Inches(2), Inches(5), Inches(4)
        chart = slide.shapes.add_chart(chart_type, x,y,cx,cy, data).chart
        try:
            for i, s in enumerate(chart.series):
                f = s.format.fill; f.solid(); f.fore_color.rgb = palette[i % len(palette)]
        except Exception:
            pass
        if hasattr(chart,'chart_title') and chart_data.title:
            chart.chart_title.text_frame.text = chart_data.title

    def _add_table(self, slide, table_data: Dict[str, Any], colors):
        headers = table_data.get('headers', []); rows = table_data.get('rows', [])
        if not headers or not rows: return
        x,y,cx,cy = Inches(1), Inches(3), Inches(8), Inches(3)
        table = slide.shapes.add_table(len(rows)+1, len(headers), x,y,cx,cy).table
        for i,h in enumerate(headers):
            cell = table.cell(0,i); cell.text = h; cell.text_frame.paragraphs[0].font.bold = True
            cell.fill.solid(); cell.fill.fore_color.rgb = colors['secondary']
        for r_idx,row in enumerate(rows):
            for c_idx,val in enumerate(row):
                if c_idx < len(headers): table.cell(r_idx+1,c_idx).text = str(val)

    def _add_flow(self, slide, flow: Dict[str, Any], colors):
        steps = flow.get('steps', [])
        if not steps: return
        box_w, box_h = Inches(1.5), Inches(0.8)
        start_x, start_y, spacing = Inches(1), Inches(3), Inches(2)
        for i, step in enumerate(steps[:5]):
            x = start_x + i*spacing
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, start_y, box_w, box_h)
            shape.fill.solid(); shape.fill.fore_color.rgb = colors['primary']; shape.line.color.rgb = colors['text']
            tf = shape.text_frame; tf.text = step
            tf.paragraphs[0].font.color.rgb = RGBColor(255,255,255); tf.paragraphs[0].font.size = Pt(12); tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            if i < len(steps)-1 and i < 4:
                slide.shapes.add_connector(1, x+box_w+Inches(0.1), start_y+box_h/2, x+box_w+Inches(0.4), start_y+box_h/2)

    def _maybe_add_image(self, slide, suggestion: str):  # optional external image fetch (disabled by default)
        if not (os.environ.get('PPT_IMAGE_FETCH') == '1' and os.environ.get('UNSPLASH_ACCESS_KEY')):
            return
        try:
            import requests  # local import
            query = suggestion.split()[0][:40] if suggestion else 'technology'
            resp = requests.get('https://api.unsplash.com/photos/random', params={'query': query, 'content_filter':'high','orientation':'landscape'}, headers={'Authorization': f"Client-ID {os.environ['UNSPLASH_ACCESS_KEY']}"}, timeout=4)
            if resp.status_code != 200: return
            image_url = resp.json().get('urls', {}).get('small')
            if not image_url: return
            img = requests.get(image_url, timeout=4).content
            with tempfile.NamedTemporaryFile(delete=False, suffix='.jpg') as tmp:
                tmp.write(img); path = tmp.name
            slide.shapes.add_picture(path, Inches(8.2), Inches(3.0), width=Inches(4.5))
        except Exception as e:  # noqa: BLE001
            logger.debug(f"이미지 삽입 실패: {e}")

    # ================== 슬라이드 관리 기능 ==================
    def apply_slide_management(self, prs: Presentation, slide_management: List[Dict[str, Any]], 
                             original_spec: DeckSpec) -> Presentation:
        """슬라이드 관리 정보를 적용하여 슬라이드 순서 조정, 복사, 삭제"""
        if not slide_management:
            logger.info("슬라이드 관리 정보가 없어 원본 그대로 반환")
            return prs
            
        logger.info(f"슬라이드 관리 적용: {len(slide_management)}개 슬라이드 정의")
        
        # 현재는 슬라이드 관리 기능을 비활성화하고 원본을 그대로 반환
        # TODO: 향후 안전한 슬라이드 관리 구현
        logger.info("슬라이드 관리 기능은 현재 개발 중입니다. 원본 슬라이드를 사용합니다.")
        return prs

    def build_enhanced_pptx_with_slide_management(self, spec: DeckSpec, file_basename: Optional[str] = None,
                                                 template_style: str = "business", include_charts: bool = True,
                                                 custom_template_path: Optional[str] = None,
                                                 user_template_id: Optional[str] = None,
                                                 text_box_mappings: Optional[List[Dict[str, Any]]] = None,
                                                 content_segments: Optional[List[Dict[str, Any]]] = None,
                                                 slide_management: Optional[List[Dict[str, Any]]] = None) -> str:
        """슬라이드 관리가 포함된 Enhanced PPT 빌드"""
        
        # 기본 PPT 빌드
        temp_file_path = self.build_enhanced_pptx(
            spec=spec,
            file_basename=file_basename,
            template_style=template_style,
            include_charts=include_charts,
            custom_template_path=custom_template_path,
            user_template_id=user_template_id,
            text_box_mappings=text_box_mappings,
            content_segments=content_segments
        )
        
        # 슬라이드 관리 적용이 필요한 경우
        if slide_management:
            # ⚠️ 커스텀 템플릿 사용 시 슬라이드 관리 기능 비활성화 (템플릿 내용 보존)
            if custom_template_path:
                logger.info(f"🚨 커스텀 템플릿 사용 시 슬라이드 관리 기능 스킵: 템플릿 내용 보존 우선")
                logger.info(f"원본 템플릿 유지: {custom_template_path}")
                logger.info(f"생성된 PPT 반환: {temp_file_path}")
                return temp_file_path
            
            logger.info(f"슬라이드 관리 기능 적용: {len(slide_management)}개 슬라이드")
            logger.info("슬라이드 관리 모드: 매핑이 적용된 결과를 사용합니다.")
            
            # 매핑이 있든 없든 슬라이드 관리 적용
            try:
                managed_file_path = self._apply_slide_management_to_ppt(
                    temp_file_path, 
                    slide_management, 
                    custom_template_path
                )
                return managed_file_path
            except Exception as e:
                logger.error(f"슬라이드 관리 적용 실패: {e}")
                return temp_file_path
            else:
                logger.warning("슬라이드 관리 모드이지만 사용자 템플릿이 없습니다. 기본 PPT 사용.")
        
        return temp_file_path

    def _apply_slide_management_to_ppt(self, source_ppt_path: str, slide_management: List[dict], custom_template_path: Optional[str] = None) -> str:
        """슬라이드 관리 정보를 바탕으로 PPT의 슬라이드를 복사/삭제/순서 변경"""
        try:
            from pptx import Presentation
            import shutil

            # 새 파일명 생성
            managed_file_path = source_ppt_path.replace('.pptx', '_slide_managed.pptx')

            # 소스 PPT 열기 (매핑이 적용된 결과 또는 원본)
            src_prs = Presentation(source_ppt_path)
            logger.info(f"소스 PPT 로드: {len(src_prs.slides)}개 슬라이드")

            # 원본 슬라이드 스냅샷 (인덱스 변동 방지용)
            original_slides = list(src_prs.slides)
            original_slide_count = len(original_slides)

            # 최종 생성할 슬라이드 소스 목록 구성 (원본/복사/새 생성 모두 포함, 순서 보장)
            final_sources = []  # List[Tuple[source_slide, title_override]]
            for slide_info in slide_management:
                if not slide_info.get('is_enabled', True):
                    continue
                if slide_info.get('is_visible') is False:
                    continue

                index = slide_info.get('index', 0)
                original_index = slide_info.get('original_index')
                title = slide_info.get('title', f'슬라이드 {index + 1}')

                # 새 슬라이드 (original_index가 없거나 범위 밖)
                if original_index is None or original_index < 0 or original_index >= original_slide_count:
                    base_slide_index = slide_info.get('base_slide_index', 0) or 0
                    if base_slide_index < 0 or base_slide_index >= original_slide_count:
                        base_slide_index = 0
                    source_slide = original_slides[base_slide_index]
                    final_sources.append((source_slide, title))
                    logger.info(f"최종 순서에 새 슬라이드 추가: base={base_slide_index}, title={title}")
                    continue

                # 기존 슬라이드 또는 복사본 (항상 원본에서 복사해 새 프레젠테이션에 구성)
                source_slide = original_slides[original_index]
                final_sources.append((source_slide, title))
                logger.info(f"최종 순서에 원본/복사본 추가: original={original_index}, title={title}")

            # 새 프레젠테이션에 최종 순서대로 슬라이드 생성
            # 중요: 다른 프레젠테이션의 레이아웃을 재사용하지 말고, 새 프레젠테이션의 자체 레이아웃을 사용해야
            # 패키지 리소스(테마/마스터/레이아웃) 중복으로 인한 파일 손상을 방지할 수 있습니다.
            # 사용자 템플릿이 있으면 해당 템플릿을 기반으로 새 프레젠테이션 생성하여 테마/마스터 보존
            try:
                if custom_template_path and os.path.exists(custom_template_path):
                    new_prs = Presentation(custom_template_path)
                    logger.info(f"새 프레젠테이션 템플릿 적용: {custom_template_path}")
                else:
                    new_prs = Presentation()
            except Exception:
                new_prs = Presentation()
            # 템플릿에 포함된 기존 슬라이드를 모두 제거한 뒤 최종 순서를 구성
            try:
                removed = 0
                for idx in range(len(new_prs.slides) - 1, -1, -1):
                    rId = new_prs.slides._sldIdLst[idx].rId
                    new_prs.part.drop_rel(rId)
                    del new_prs.slides._sldIdLst[idx]
                    removed += 1
                if removed:
                    logger.info(f"기존 템플릿 슬라이드 제거: {removed}개")
            except Exception as e:
                logger.debug(f"기존 슬라이드 제거 스킵: {e}")

            # 레이아웃 이름 기반 매칭 유틸
            def _match_layout_index(src_layout_name: str) -> int:
                if not src_layout_name:
                    return 0
                for idx, layout in enumerate(new_prs.slide_layouts):
                    try:
                        if getattr(layout, 'name', None) == src_layout_name:
                            return idx
                    except Exception:
                        continue
                # 폴백 순서: 제목+내용(1) -> 제목(0) -> 빈(6) -> 0
                for idx in (1, 0, 6):
                    if idx < len(new_prs.slide_layouts):
                        return idx
                return 0

            for i, (source_slide, title_text) in enumerate(final_sources):
                try:
                    # 슬라이드 관리 정보 가져오기 (텍스트 클리어 필요 여부 확인)
                    slide_info = slide_management[i] if i < len(slide_management) else {}
                    
                    src_layout_name = None
                    try:
                        src_layout_name = getattr(source_slide.slide_layout, 'name', None)
                    except Exception:
                        src_layout_name = None
                    layout_idx = _match_layout_index(src_layout_name)
                    new_slide = new_prs.slides.add_slide(new_prs.slide_layouts[layout_idx])

                    # 원본 슬라이드의 모든 내용 복사
                    try:
                        self._copy_slide_content(source_slide, new_slide)
                    except Exception as e:
                        logger.warning(f"슬라이드 내용 복사 실패(i={i}): {e}")

                    # 🆕 "추가" 버튼으로 생성된 슬라이드는 텍스트 클리어
                    if slide_info.get('needsTextClear', False):
                        try:
                            cleared_count = 0
                            for shape in new_slide.shapes:
                                try:
                                    if hasattr(shape, 'text_frame') and shape.text_frame:
                                        if shape.text_frame.text.strip():
                                            shape.text_frame.text = ""
                                            cleared_count += 1
                                    elif hasattr(shape, 'text') and shape.text.strip():
                                        shape.text = ""
                                        cleared_count += 1
                                except Exception:
                                    continue
                            logger.info(f"슬라이드 {i+1} 텍스트 클리어: {cleared_count}개 요소")
                        except Exception as e:
                            logger.warning(f"슬라이드 {i+1} 텍스트 클리어 실패: {e}")

                    # 제목 텍스트 설정 시도 (필요할 경우에만)
                    # 프론트의 표시용 기본 제목(예: "슬라이드 1")은 실제 PPT에 주입하지 않습니다.
                    try:
                        safe_title = (title_text or "").strip()
                        if safe_title and not re.match(r"^슬라이드\s*\d+$", safe_title):
                            if hasattr(new_slide.shapes, 'title') and new_slide.shapes.title:
                                try:
                                    # 기존 제목이 이미 채워져 있으면 유지 (소스 placeholder에서 복사된 경우)
                                    current = (new_slide.shapes.title.text or "").strip()
                                except Exception:
                                    current = ""
                                if not current:
                                    new_slide.shapes.title.text = safe_title
                            else:
                                # 첫 번째 텍스트 박스에 제목을 설정 (있을 경우)
                                for shape in new_slide.shapes:
                                    if getattr(shape, 'has_text_frame', False):
                                        try:
                                            setattr(shape, 'text', safe_title)
                                        except Exception:
                                            tf = getattr(shape, 'text_frame', None)
                                            if tf and len(getattr(tf, 'paragraphs', [])):
                                                tf.paragraphs[0].text = safe_title
                                        break
                    except Exception as e:
                        logger.warning(f"제목 설정 실패(i={i}): {e}")

                    logger.info(f"슬라이드 생성 완료(i={i}): '{title_text}'")
                except Exception as e:
                    logger.warning(f"슬라이드 생성 실패(i={i}): {e}")

            logger.info(f"슬라이드 관리 완료: 총 {len(new_prs.slides)}개 슬라이드")

            # 파일 저장
            new_prs.save(managed_file_path)
            logger.info(f"슬라이드 관리 적용 완료: {managed_file_path}")

            # 원본 파일 정리
            if os.path.exists(source_ppt_path) and source_ppt_path != managed_file_path:
                os.unlink(source_ppt_path)

            return managed_file_path
            
        except Exception as e:
            logger.error(f"슬라이드 관리 적용 실패: {e}")
            return source_ppt_path
    
    def _copy_slide_content(self, source_slide, target_slide):
        """한 슬라이드의 내용을 다른 슬라이드로 복사
        - 먼저 커스텀 템플릿 여부를 확인
        - Placeholder가 있으면 기존 로직, 없으면 커스텀 템플릿용 로직 사용
        """
        try:
            # 템플릿 타입 감지
            template_type = self._detect_template_type(target_slide)
            
            if template_type == 'custom':
                logger.info("커스텀 템플릿 감지: 직접 shape 복사 모드 사용")
                self._copy_slide_content_custom_template(source_slide, target_slide)
            else:
                logger.info("표준 템플릿 감지: placeholder 매핑 모드 사용")
                self._copy_slide_content_standard_template(source_slide, target_slide)
                
        except Exception as e:
            logger.warning(f"슬라이드 내용 복사 중 오류: {e}")
            # 폴백: 기존 로직 사용
            self._copy_slide_content_standard_template(source_slide, target_slide)

    def _detect_template_type(self, slide) -> str:
        """템플릿 타입 감지 (🔵 최적화)"""
        try:
            placeholders = getattr(slide, 'placeholders', [])
            total_shapes = len(slide.shapes)
            
            # placeholder 비율로 템플릿 타입 결정
            if len(placeholders) == 0 and total_shapes > 0:
                return 'custom'  # 완전 커스텀 템플릿
            elif len(placeholders) / max(total_shapes, 1) > 0.3:
                return 'standard'  # 표준 템플릿 (30% 이상이 placeholder)
            else:
                return 'hybrid'  # 혼합형 (일부 placeholder 사용)
        except Exception:
            return 'standard'  # 기본값
    
    def _copy_slide_content_custom_template(self, source_slide, target_slide):
        """커스텀 템플릿용 복사 로직 (🟡 중요)
        - Shape name 기반 1:1 매핑
        - 텍스트 내용만 교체, 스타일은 완전 보존
        - 추가 shape 생성 방지
        """
        try:
            import copy as _copy
            
            # 1) 소스 슬라이드의 텍스트 내용을 수집
            source_text_map = {}
            for shape in source_slide.shapes:
                try:
                    shape_name = getattr(shape, 'name', None)
                    if shape_name and hasattr(shape, 'text_frame') and shape.text_frame:
                        text_content = shape.text_frame.text.strip()
                        if text_content:
                            source_text_map[shape_name] = text_content
                except Exception:
                    continue
            
            logger.info(f"소스 텍스트 수집: {len(source_text_map)}개 shape")
            
            # 2) 타겟 슬라이드의 해당 shape에 텍스트만 복사 (🟢 개선: 추가 shape 생성 방지)
            for target_shape in target_slide.shapes:
                try:
                    target_name = getattr(target_shape, 'name', None)
                    if target_name and target_name in source_text_map:
                        source_text = source_text_map[target_name]
                        
                        # 텍스트 프레임이 있으면 내용만 교체 (스타일 보존)
                        if hasattr(target_shape, 'text_frame') and target_shape.text_frame:
                            self._replace_text_preserving_style_simple(
                                target_shape.text_frame, source_text
                            )
                            logger.debug(f"텍스트 복사 완료: {target_name} -> '{source_text[:30]}...'")
                except Exception as e:
                    logger.debug(f"Shape 텍스트 복사 실패: {e}")
                    continue
            
            logger.info("커스텀 템플릿 슬라이드 복사 완료")
            
        except Exception as e:
            logger.error(f"커스텀 템플릿 복사 실패: {e}")
            # 폴백: 표준 로직 사용
            self._copy_slide_content_standard_template(source_slide, target_slide)
    
    def _replace_text_preserving_style_simple(self, text_frame, new_text: str):
        """스타일을 완전히 보존하면서 텍스트만 교체 (단순화 버전)"""
        try:
            # 기존 첫 번째 문단의 스타일 정보를 백업
            if text_frame.paragraphs:
                first_para = text_frame.paragraphs[0]
                
                # 스타일 백업
                para_style = {
                    'alignment': getattr(first_para, 'alignment', None),
                    'level': getattr(first_para, 'level', 0)
                }
                
                run_style = {}
                if first_para.runs:
                    first_run = first_para.runs[0]
                    run_style = {
                        'font_name': getattr(first_run.font, 'name', None),
                        'font_size': getattr(first_run.font, 'size', None),
                        'bold': getattr(first_run.font, 'bold', None),
                        'italic': getattr(first_run.font, 'italic', None),
                        'color': getattr(first_run.font, 'color', None)
                    }
                
                # 텍스트 교체
                text_frame.clear()
                new_para = text_frame.paragraphs[0]
                new_para.text = new_text
                
                # 스타일 복원
                if para_style.get('alignment') is not None:
                    new_para.alignment = para_style['alignment']
                new_para.level = para_style.get('level', 0)
                
                if new_para.runs and run_style:
                    new_run = new_para.runs[0]
                    try:
                        if run_style.get('font_name'):
                            new_run.font.name = run_style['font_name']
                        if run_style.get('font_size'):
                            new_run.font.size = run_style['font_size']
                        if run_style.get('bold') is not None:
                            new_run.font.bold = run_style['bold']
                        if run_style.get('italic') is not None:
                            new_run.font.italic = run_style['italic']
                        if run_style.get('color'):
                            new_run.font.color = run_style['color']
                    except Exception:
                        pass  # 스타일 적용 실패 시 텍스트만 유지
                        
        except Exception as e:
            # 최소한 텍스트만이라도 설정
            try:
                text_frame.clear()
                text_frame.paragraphs[0].text = new_text
            except Exception:
                pass
            logger.debug(f"스타일 보존 실패, 텍스트만 설정: {e}")

    def _copy_slide_content_standard_template(self, source_slide, target_slide):
        """표준 템플릿용 복사 로직 (기존 로직)
        - Placeholder(자리표시자)는 대상 슬라이드의 동일 placeholder에 텍스트만 주입
        - 그 외 일반 도형은 XML 복사로 추가
        """
        try:
            import copy as _copy

            # 1) Placeholder 텍스트 매핑 (스타일은 대상 레이아웃이 제공)
            try:
                # 소스의 placeholder 텍스트를 수집: {idx: text}
                src_ph_text = {}
                for sh in source_slide.shapes:
                    try:
                        if getattr(sh, 'is_placeholder', False) and getattr(sh, 'has_text_frame', False):
                            phf = getattr(sh, 'placeholder_format', None)
                            if phf is not None:
                                idx = getattr(phf, 'idx', None)
                                if idx is not None:
                                    # 전체 텍스트(문단 합치기)
                                    txt = getattr(sh, 'text', '')
                                    src_ph_text[idx] = txt
                    except Exception:
                        continue

                # 대상 슬라이드 placeholder에 텍스트 주입
                placeholders = getattr(target_slide, 'placeholders', [])
                for shp in placeholders:
                    try:
                        phf = getattr(shp, 'placeholder_format', None)
                        if phf is None:
                            continue
                        idx = getattr(phf, 'idx', None)
                        if idx in src_ph_text and getattr(shp, 'has_text_frame', False):
                            # 기존 문단 초기화 후 텍스트 설정
                            tf = getattr(shp, 'text_frame', None)
                            if tf:
                                try:
                                    # clear()가 없을 수 있어 안전하게 첫 문단만 사용
                                    if getattr(tf, 'paragraphs', None):
                                        tf.paragraphs[0].text = src_ph_text[idx]
                                    else:
                                        setattr(shp, 'text', src_ph_text[idx])
                                except Exception:
                                    setattr(shp, 'text', src_ph_text[idx])
                    except Exception:
                        continue
            except Exception as e:
                logger.debug(f"Placeholder 매핑 스킵: {e}")

            # 2) Placeholder가 아닌 일반 도형은 그대로 복사 추가 (🟢 개선: 중복 방지)
            copied_shapes = 0
            max_shapes_to_copy = 20  # 과도한 shape 복사 방지
            
            for shape in source_slide.shapes:
                try:
                    # placeholder는 이미 대상 placeholder에 주입했으므로 스킵
                    if getattr(shape, 'is_placeholder', False):
                        continue
                        
                    if copied_shapes >= max_shapes_to_copy:
                        logger.warning(f"Shape 복사 제한 도달: {max_shapes_to_copy}개")
                        break
                        
                    if hasattr(shape, 'element'):
                        shape_element = _copy.deepcopy(shape.element)
                        target_slide.shapes._spTree.insert_element_before(shape_element, 'p:extLst')
                        copied_shapes += 1
                except Exception:
                    continue
                    
            if copied_shapes > 0:
                logger.info(f"일반 도형 복사 완료: {copied_shapes}개")

        except Exception as e:
            logger.warning(f"표준 템플릿 슬라이드 내용 복사 중 오류: {e}")


enhanced_ppt_generator_service = EnhancedPPTGeneratorService()
