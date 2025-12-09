"""
PPTX 템플릿 파일 분석 서비스
사용자가 업로드한 PPTX 파일의 구조, 레이아웃, 스타일을 분석하여
AI 생성 시스템에서 활용할 수 있는 메타데이터를 추출
"""
import json
import os
from dataclasses import dataclass, asdict
from typing import Dict, List, Optional, Any, Tuple
from pathlib import Path
import logging

from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from pptx.shapes.placeholder import PlaceholderPicture, SlidePlaceholder
from pptx.shapes.autoshape import Shape
from pptx.shapes.picture import Picture
from pptx.shapes.graphfrm import GraphicFrame
from pptx.dml.color import RGBColor, ColorFormat
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

@dataclass
@dataclass
class LayoutBox:
    """레이아웃 박스 정보"""
    x: float  # 인치 단위
    y: float
    width: float
    height: float
    type: str  # title, content, image, chart, table, shape
    placeholder_type: Optional[str] = None
    text_align: Optional[str] = None
    font_size: Optional[float] = None
    font_color: Optional[str] = None
    text_content: Optional[str] = None  # 원본 텍스트 내용

@dataclass
class SlideTemplate:
    """개별 슬라이드 템플릿 정보"""
    slide_number: int
    title: str
    layout_type: str  # title-only, title-content, two-column, chart-focus, etc.
    layout_boxes: List[LayoutBox]
    background_color: Optional[str] = None
    suggested_content_type: Optional[str] = None  # text, bullet-list, chart, table, image
    
@dataclass
class PresentationTemplate:
    """전체 프레젠테이션 템플릿 정보"""
    template_name: str
    total_slides: int
    slide_templates: List[SlideTemplate]
    color_scheme: Dict[str, str]
    font_scheme: Dict[str, str]
    layout_types: List[str]
    
    def to_json(self) -> str:
        """JSON 문자열로 변환"""
        return json.dumps(asdict(self), ensure_ascii=False, indent=2)

class PPTXTemplateAnalyzer:
    """PPTX 템플릿 분석기"""
    
    def __init__(self):
        self.layout_type_map = {
            "title-only": ["title"],
            "title-content": ["title", "content"],
            "two-column": ["title", "content", "content"],
            "chart-focus": ["title", "content", "chart"],
            "table-focus": ["title", "table"],
            "image-focus": ["title", "content", "image"],
            "section-header": ["title"],
            "blank": []
        }
    
    def analyze_pptx_template(self, pptx_path: str) -> PresentationTemplate:
        """PPTX 파일을 분석하여 템플릿 메타데이터 생성"""
        try:
            logger.info(f"PPTX 템플릿 분석 시작: {pptx_path}")
            
            # PPTX 파일 로드
            presentation = Presentation(pptx_path)
            
            # 기본 정보 추출
            template_name = Path(pptx_path).stem
            total_slides = len(presentation.slides)
            
            # 슬라이드별 분석
            slide_templates = []
            for i, slide in enumerate(presentation.slides):
                slide_template = self._analyze_slide(slide, i + 1)
                slide_templates.append(slide_template)
            
            # 전체 색상/폰트 스키마 추출
            color_scheme = self._extract_color_scheme(presentation)
            font_scheme = self._extract_font_scheme(presentation)
            
            # 레이아웃 타입 목록 추출
            layout_types = list(set(st.layout_type for st in slide_templates))
            
            template = PresentationTemplate(
                template_name=template_name,
                total_slides=total_slides,
                slide_templates=slide_templates,
                color_scheme=color_scheme,
                font_scheme=font_scheme,
                layout_types=layout_types
            )
            
            logger.info(f"PPTX 템플릿 분석 완료: {total_slides}개 슬라이드, {len(layout_types)}개 레이아웃")
            return template
            
        except Exception as e:
            logger.error(f"PPTX 템플릿 분석 실패: {e}")
            raise
    
    def _analyze_slide(self, slide: Slide, slide_number: int) -> SlideTemplate:
        """개별 슬라이드 분석"""
        
        # 슬라이드 제목 추출
        title = self._extract_slide_title(slide)
        
        # 레이아웃 박스들 분석
        layout_boxes = []
        shape_types = []
        
        for shape in slide.shapes:
            layout_box = self._analyze_shape(shape)
            if layout_box:
                layout_boxes.append(layout_box)
                shape_types.append(layout_box.type)
        
        # 레이아웃 타입 추론
        layout_type = self._infer_layout_type(shape_types, layout_boxes)
        
        # 추천 콘텐츠 타입 추론
        suggested_content_type = self._infer_content_type(shape_types, layout_boxes)
        
        # 배경색 추출
        background_color = self._extract_background_color(slide)
        
        return SlideTemplate(
            slide_number=slide_number,
            title=title,
            layout_type=layout_type,
            layout_boxes=layout_boxes,
            background_color=background_color,
            suggested_content_type=suggested_content_type
        )
    
    def _extract_slide_title(self, slide: Slide) -> str:
        """슬라이드 제목 추출"""
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                # 첫 번째 텍스트를 제목으로 간주
                title = shape.text.strip()
                if len(title) > 0:
                    return title[:50]  # 최대 50자로 제한
        return f"슬라이드 {slide.slide_id}"
    
    def _analyze_shape(self, shape: BaseShape) -> Optional[LayoutBox]:
        """개별 도형/요소 분석"""
        try:
            # 위치와 크기 (인치 단위로 변환)
            x = shape.left.inches if hasattr(shape, 'left') else 0
            y = shape.top.inches if hasattr(shape, 'top') else 0
            width = shape.width.inches if hasattr(shape, 'width') else 0
            height = shape.height.inches if hasattr(shape, 'height') else 0
            
            # 도형 타입 분석
            shape_type = self._get_shape_type(shape)
            
            # 플레이스홀더 타입 확인 (더 관대하게)
            placeholder_type = None
            is_placeholder = False
            
            try:
                if hasattr(shape, 'placeholder_format'):
                    try:
                        placeholder_format = shape.placeholder_format
                        if placeholder_format:
                            placeholder_type = str(placeholder_format.type)
                            is_placeholder = True
                    except Exception:
                        # placeholder_format 접근 실패해도 속성 존재 자체로 플레이스홀더로 간주
                        placeholder_type = "unknown_placeholder"
                        is_placeholder = True
            except Exception as e:
                logger.debug(f"플레이스홀더 확인 중 오류: {e}")
            
            # 텍스트 내용 확인
            has_text = False
            text_content = ""
            
            try:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text = shape.text_frame.text
                    if text.strip():
                        has_text = True
                        text_content = text.strip()
                elif hasattr(shape, 'text'):
                    text = shape.text
                    if text and text.strip():
                        has_text = True
                        text_content = text.strip()
            except Exception as e:
                logger.debug(f"텍스트 확인 중 오류: {e}")
            
            # 텍스트 스타일 분석
            text_align, font_size, font_color = self._extract_text_style(shape)
            
            # 로그 출력
            shape_name = getattr(shape, 'name', 'Unnamed')
            logger.debug(f"도형 분석: {shape_name} (타입: {shape_type}, "
                        f"플레이스홀더: {is_placeholder}, 텍스트: {has_text})")
            
            if has_text:
                logger.debug(f"  텍스트 내용: '{text_content[:50]}...'")
            
            return LayoutBox(
                x=x, y=y, width=width, height=height,
                type=shape_type,
                placeholder_type=placeholder_type,
                text_align=text_align,
                font_size=font_size,
                font_color=font_color,
                text_content=text_content if has_text else None
            )
            
        except Exception as e:
            logger.warning(f"도형 분석 실패: {e}")
            return None
    
    def _get_shape_type(self, shape: BaseShape) -> str:
        """도형 타입 분류"""
        if hasattr(shape, 'text') and shape.text.strip():
            # 텍스트가 있는 도형
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                if 'title' in str(shape.placeholder_format.type).lower():
                    return "title"
                elif 'content' in str(shape.placeholder_format.type).lower():
                    return "content"
            return "text"
        
        elif hasattr(shape, 'table') or str(shape.shape_type) == "TABLE":
            return "table"
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            return "chart"
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return "image"
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            return "shape"
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            return "group"
        else:
            return "unknown"
    
    def _extract_text_style(self, shape: BaseShape) -> Tuple[Optional[str], Optional[float], Optional[str]]:
        """텍스트 스타일 추출"""
        text_align = None
        font_size = None
        font_color = None
        
        try:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                # 첫 번째 문단의 스타일 가져오기
                if shape.text_frame.paragraphs:
                    para = shape.text_frame.paragraphs[0]
                    
                    # 정렬
                    if para.alignment:
                        align_map = {
                            PP_ALIGN.LEFT: "left",
                            PP_ALIGN.CENTER: "center", 
                            PP_ALIGN.RIGHT: "right",
                            PP_ALIGN.JUSTIFY: "justify"
                        }
                        text_align = align_map.get(para.alignment, "left")
                    
                    # 폰트 크기와 색상
                    if para.runs:
                        run = para.runs[0]
                        if run.font.size:
                            font_size = run.font.size.pt
                        if run.font.color and run.font.color.rgb:
                            rgb = run.font.color.rgb
                            font_color = f"#{rgb.r:02x}{rgb.g:02x}{rgb.b:02x}"
        except Exception as e:
            logger.debug(f"텍스트 스타일 추출 실패: {e}")
        
        return text_align, font_size, font_color
    
    def _infer_layout_type(self, shape_types: List[str], layout_boxes: List[LayoutBox]) -> str:
        """도형 구성을 바탕으로 레이아웃 타입 추론"""
        
        # 제목만 있는 경우
        if len([t for t in shape_types if t == "title"]) >= 1 and len(shape_types) <= 2:
            return "title-only"
        
        # 차트가 있는 경우
        if "chart" in shape_types:
            return "chart-focus"
        
        # 테이블이 있는 경우
        if "table" in shape_types:
            return "table-focus"
        
        # 이미지가 있는 경우
        if "image" in shape_types:
            return "image-focus"
        
        # 콘텐츠 영역이 2개 이상인 경우
        content_count = len([t for t in shape_types if t in ["content", "text"]])
        if content_count >= 2:
            return "two-column"
        
        # 제목 + 콘텐츠
        if "title" in shape_types and content_count >= 1:
            return "title-content"
        
        # 기본값
        return "title-content"
    
    def _infer_content_type(self, shape_types: List[str], layout_boxes: List[LayoutBox]) -> str:
        """추천 콘텐츠 타입 추론"""
        if "chart" in shape_types:
            return "chart"
        elif "table" in shape_types:
            return "table"
        elif "image" in shape_types:
            return "image"
        elif len([t for t in shape_types if t in ["content", "text"]]) > 1:
            return "bullet-list"
        else:
            return "text"
    
    def _extract_background_color(self, slide: Slide) -> Optional[str]:
        """배경색 추출"""
        try:
            if hasattr(slide, 'background') and slide.background:
                # 배경 정보 추출 시도
                return None  # 일단 None으로 설정 (복잡한 배경 처리는 추후)
        except Exception as e:
            logger.debug(f"배경색 추출 실패: {e}")
        return None
    
    def _extract_color_scheme(self, presentation: Presentation) -> Dict[str, str]:
        """전체 색상 스키마 추출"""
        colors = {
            "primary": "#1f4e79",      # 기본 파란색
            "secondary": "#70ad47",    # 기본 초록색  
            "accent": "#ffc000",       # 기본 노란색
            "background": "#ffffff",   # 기본 흰색
            "text": "#000000"          # 기본 검은색
        }
        
        try:
            # 테마 색상 추출 시도 (복잡하므로 기본값 사용)
            # presentation.core_properties, theme 등에서 추출 가능
            pass
        except Exception as e:
            logger.debug(f"색상 스키마 추출 실패: {e}")
        
        return colors
    
    def _extract_font_scheme(self, presentation: Presentation) -> Dict[str, str]:
        """폰트 스키마 추출"""
        fonts = {
            "title": "맑은 고딕",
            "body": "맑은 고딕",
            "accent": "맑은 고딕"
        }
        
        try:
            # 테마 폰트 추출 시도
            pass
        except Exception as e:
            logger.debug(f"폰트 스키마 추출 실패: {e}")
        
        return fonts
    
    def save_template_metadata(self, template: PresentationTemplate, output_path: str) -> str:
        """템플릿 메타데이터를 JSON 파일로 저장"""
        try:
            # 🎯 통일된 네이밍: template_name을 template_id 형식으로 변환 (공백을 언더스코어로)
            template_id_for_metadata = template.template_name.replace(' ', '_')
            output_file = Path(output_path) / f"{template_id_for_metadata}_metadata.json"
            
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(template.to_json())
            
            logger.info(f"템플릿 메타데이터 저장 완료: {output_file}")
            return str(output_file)
            
        except Exception as e:
            logger.error(f"템플릿 메타데이터 저장 실패: {e}")
            raise
    
    def generate_outline_tabs(self, template: PresentationTemplate) -> List[Dict[str, Any]]:
        """템플릿 기반 아웃라인 탭 구조 생성"""
        tabs = []
        
        for slide_template in template.slide_templates:
            tab = {
                "id": f"slide_{slide_template.slide_number}",
                "title": slide_template.title,
                "layout_type": slide_template.layout_type,
                "content_type": slide_template.suggested_content_type,
                "layout_boxes": [asdict(box) for box in slide_template.layout_boxes],
                "editable_areas": self._get_editable_areas(slide_template),
                "preview_info": {
                    "background_color": slide_template.background_color,
                    "main_content_area": self._get_main_content_area(slide_template.layout_boxes)
                }
            }
            tabs.append(tab)
        
        return tabs
    
    def _get_editable_areas(self, slide_template: SlideTemplate) -> List[Dict[str, Any]]:
        """편집 가능한 영역 정의"""
        editable_areas = []
        
        for box in slide_template.layout_boxes:
            if box.type in ["title", "content", "text"]:
                area = {
                    "type": box.type,
                    "placeholder": f"여기에 {box.type} 내용을 입력하세요",
                    "max_length": 200 if box.type == "title" else 1000,
                    "style_guide": {
                        "font_size": box.font_size,
                        "text_align": box.text_align,
                        "font_color": box.font_color
                    }
                }
                editable_areas.append(area)
        
        return editable_areas
    
    def _get_main_content_area(self, layout_boxes: List[LayoutBox]) -> Optional[Dict[str, float]]:
        """메인 콘텐츠 영역 정보"""
        content_boxes = [box for box in layout_boxes if box.type == "content"]
        if content_boxes:
            main_box = content_boxes[0]  # 첫 번째 콘텐츠 박스
            return {
                "x": main_box.x,
                "y": main_box.y, 
                "width": main_box.width,
                "height": main_box.height
            }
        return None


# 전역 인스턴스
pptx_template_analyzer = PPTXTemplateAnalyzer()
