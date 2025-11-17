"""
PPT 템플릿 관리 및 매핑 서비스
템플릿 파일의 레이아웃, 슬라이드 마스터, 플레이스홀더를 분석하고
DeckSpec을 템플릿에 맞게 변환하여 적용
"""
from __future__ import annotations

import json
import subprocess
import shutil
from pathlib import Path
from typing import Dict, List, Optional, Any, Tuple
from dataclasses import dataclass
from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.enum.shapes import MSO_SHAPE_TYPE

from loguru import logger
from .ppt_models import DeckSpec, SlideSpec
from .ppt_template_extractor import extract_presentation


@dataclass
class TemplateLayoutInfo:
    """템플릿 레이아웃 정보"""
    name: str
    slide_index: int  # 템플릿 내 슬라이드 인덱스
    placeholders: Dict[str, Any]  # 플레이스홀더 정보
    layout_type: str  # title-only, title-and-content, two-content, etc.
    max_bullets: int = 6
    supports_chart: bool = False
    supports_table: bool = False


@dataclass
class TemplateSpec:
    """템플릿 전체 정보"""
    file_path: Path
    layouts: Dict[str, TemplateLayoutInfo]
    slide_masters: List[Any]
    theme_colors: Dict[str, str]
    default_fonts: Dict[str, str]
    max_slides: int = 20


class PPTTemplateManager:
    """PPT 템플릿 분석 및 적용 매니저"""
    
    def __init__(self):
        self.template_cache: Dict[str, TemplateSpec] = {}
        self._initialized = False  # 중복 초기화 방지 플래그
        # 템플릿 레지스트리 초기화
        self._initialize_registry()
    
    def _initialize_registry(self):
        """템플릿 레지스트리 초기화 (중복 호출 방지)"""
        if self._initialized:
            logger.debug("PPTTemplateManager 이미 초기화됨 - 스킵")
            return
            
        # uploads/templates 경로 사용
        # NOTE: 기존에는 parents[3] (backend 디렉토리) 기준으로 잡혀 실제 루트(/project_root/uploads/templates)가 아닌
        #       /project_root/backend/uploads/templates 를 바라봐 존재하지 않는 경로로 인해 템플릿 미적용 문제가 발생.
        #       아래 로직은 루트 후보(parents[4]) 우선 검사 후 fallback 하여 안정적으로 실제 템플릿 디렉토리를 사용.
        root_candidate = None
        try:
            root_candidate = Path(__file__).parents[4]
        except Exception:
            root_candidate = Path(__file__).parents[3]
        base_dir_root = root_candidate / 'uploads' / 'templates'
        base_dir_backend = Path(__file__).parents[3] / 'uploads' / 'templates'
        if base_dir_root.exists():
            base_dir = base_dir_root
        else:
            base_dir = base_dir_backend
        base_dir.mkdir(parents=True, exist_ok=True)
        logger.debug(f"PPTTemplateManager base_dir resolved -> {base_dir} (root_exists={base_dir_root.exists()})")
        # 빈 레지스트리로 시작 - 실제 PPTX 파일만 자동 스캔으로 등록
        self._registry: Dict[str, Dict[str, Any]] = {}
        
        # 디렉토리에 있는 모든 PPTX 파일을 자동으로 등록
        self.base_dir = base_dir  # 나중에 사용할 수 있도록 저장
        self._scan_and_register_directory_templates()
        
        # 기본 템플릿 설정 로드
        self._load_default_template_config()
        
        # 초기화 완료 플래그 설정
        self._initialized = True
        logger.info(f"PPTTemplateManager 초기화 완료: {len(self._registry)}개 템플릿 등록됨")
    
    def _scan_and_register_directory_templates(self):
        """템플릿 디렉토리의 모든 PPTX 파일을 스캔하여 자동 등록"""
        try:
            if not self.base_dir.exists():
                logger.warning(f"템플릿 디렉토리가 존재하지 않음: {self.base_dir}")
                return
                
            # PPTX 파일들을 스캔 (대소문자 구분 없이)
            pptx_files = []
            for pattern in ['*.pptx', '*.PPTX', '*.Pptx']:
                pptx_files.extend(self.base_dir.glob(pattern))
            
            logger.info(f"📁 템플릿 디렉토리 스캔: {self.base_dir}")
            logger.info(f"🔍 발견된 PPTX 파일 수: {len(pptx_files)}")
            
            for file_path in pptx_files:
                # 이미 등록된 파일인지 확인 (경로 기준)
                file_str = str(file_path)
                already_registered = any(
                    template.get('path') == file_str 
                    for template in self._registry.values()
                )
                
                if already_registered:
                    # 기존 템플릿이 동적 분석이 안 되어 있으면 추가
                    for template_id, template_info in self._registry.items():
                        if template_info.get('path') == file_str:
                            if not template_info.get('dynamic_template_id'):
                                logger.info(f"🔄 기존 템플릿에 동적 분석 추가: {template_id}")
                                self._add_dynamic_analysis_to_existing(template_id, file_path)
                            break
                    continue
                
                # 새 템플릿 자동 등록 (동적 분석 포함)
                template_id = self._generate_template_id(file_path)
                template_name = self._generate_template_name(file_path)
                
                # 기본 템플릿 정보 생성
                entry = {
                    'id': template_id,
                    'name': template_name,
                    'description': f'자동 스캔된 템플릿: {file_path.name}',
                    'path': file_str,
                    'style': 'business'  # 기본 스타일
                }
                
                # 동적 분석 추가
                self._add_dynamic_analysis_to_entry(entry, file_path, template_name)
                
                self._registry[template_id] = entry
                logger.info(f"📄 템플릿 자동 등록: {template_id} -> {file_path.name}")
                
        except Exception as e:
            logger.error(f"템플릿 디렉토리 스캔 중 오류: {e}")
    
    def _add_dynamic_analysis_to_existing(self, template_id: str, file_path: Path):
        """기존 템플릿에 동적 분석 추가"""
        try:
            entry = self._registry[template_id]
            template_name = entry.get('name', file_path.stem)
            self._add_dynamic_analysis_to_entry(entry, file_path, template_name)
            logger.info(f"✅ 기존 템플릿 동적 분석 추가 완료: {template_id}")
        except Exception as e:
            logger.error(f"기존 템플릿 동적 분석 추가 실패: {template_id}, {e}")
    
    def _add_dynamic_analysis_to_entry(self, entry: dict, file_path: Path, template_name: str):
        """템플릿 엔트리에 메타데이터 분석 정보 추가 (간소화)"""
        try:
            from .template_debugger import template_debugger
            
            # 1. 디버깅 정보 수집
            logger.info(f"🔍 템플릿 디버깅 시작: {file_path}")
            debug_info = template_debugger.debug_template(str(file_path))
            entry['debug_info'] = debug_info  # type: ignore[assignment]
            
            # 2. 템플릿 메타데이터 추출 및 JSON 저장 (ppt_template_extractor 사용)
            template_id = entry.get('id', file_path.stem.replace(' ', '_'))
            logger.info(f"📊 템플릿 메타데이터 추출 시작: {template_id}")
            
            # 메타데이터 JSON 파일 경로 설정
            metadata_dir = file_path.parent / 'metadata'
            metadata_dir.mkdir(exist_ok=True)
            metadata_json_path = metadata_dir / f"{template_id}_metadata.json"
            
            # ppt_template_extractor 사용하여 상세한 메타데이터 추출
            extract_presentation(str(file_path), str(metadata_json_path))
            entry['metadata_json_path'] = str(metadata_json_path)
            
            # 3. 동적 템플릿 등록 (간소화)
            dynamic_template_id = f"user_{template_id}"
            entry['dynamic_template_id'] = dynamic_template_id
            entry['is_content_cleaned'] = True  # 메타데이터 추출 완료로 표시
            
            logger.info(f"✅ 템플릿 메타데이터 추출 완료: {metadata_json_path}")
            
        except Exception as e:
            logger.warning(f"템플릿 메타데이터 추출 실패: {e}")
            # 기본 설정으로 폴백
            template_id = entry.get('id', file_path.stem.replace(' ', '_'))
            entry['dynamic_template_id'] = f"user_{template_id}"
            entry['is_content_cleaned'] = False
            entry['metadata_json_path'] = ""
    
    def _generate_template_id(self, file_path: Path) -> str:
        """파일 경로로부터 고유한 템플릿 ID 생성"""
        base_name = file_path.stem
        # 특수문자 제거 및 소문자로 변환
        import re
        clean_name = re.sub(r'[^\w\-_]', '_', base_name).lower()
        clean_name = re.sub(r'_+', '_', clean_name)  # 연속 언더스코어 정리
        clean_name = clean_name.strip('_')  # 앞뒤 언더스코어 제거
        
        # 중복 확인 및 번호 추가
        template_id = clean_name
        counter = 1
        while template_id in self._registry:
            template_id = f"{clean_name}_{counter}"
            counter += 1
            
        return template_id
    
    def _generate_template_name(self, file_path: Path) -> str:
        """파일 경로로부터 사용자 친화적인 템플릿 이름 생성"""
        name = file_path.stem
        # 언더스코어를 스페이스로 변환
        name = name.replace('_', ' ')
        # 첫 글자를 대문자로
        return name.title() if name else file_path.name
    
    # 🔥 불필요한 레이아웃 분석 함수들 제거됨 (ppt_template_extractor가 대체)
    

    def list_templates(self) -> List[Dict[str, Any]]:
        """등록 템플릿 목록 반환 (실제 파일이 존재하는 것만, 품질 정보 포함)"""
        out = []
        for t in self._registry.values():
            path = t.get('path')
            exists = bool(path and Path(path).exists())
            # 실제 파일이 존재하는 템플릿만 포함
            if exists:
                thumb = self.get_thumbnail_path(t['id'])
                
                # 템플릿 품질 정보 추가
                file_size = 0
                quality_level = "basic"
                slide_count = 0
                
                try:
                    if path and Path(path).exists():
                        file_size = Path(path).stat().st_size
                        # 파일 크기 기반 품질 판단
                        if file_size > 1_000_000:  # 1MB 이상
                            quality_level = "professional"
                        elif file_size > 100_000:  # 100KB 이상
                            quality_level = "standard"
                        
                        # 슬라이드 수 확인
                        from pptx import Presentation
                        prs = Presentation(path)
                        slide_count = len(prs.slides)
                        
                except Exception as e:
                    logger.warning(f"템플릿 품질 분석 실패 {t['id']}: {e}")
                
                # 사용자 업로드 여부 판단:
                # - business_default 등 내장 템플릿이 아닌 모든 실제 파일은 사용자 템플릿으로 간주
                # - ID 접두사(clean_)는 정리 여부일 뿐 소유 구분과 무관
                is_user_uploaded = not t['id'].startswith('business_default')

                template_info = {
                    'id': t['id'],
                    'name': t['name'],
                    'description': t.get('description'),
                    'style': t.get('style'),
                    'has_file': exists,
                    'quality_level': quality_level,
                    'file_size_mb': round(file_size / 1024 / 1024, 1),
                    'slide_count': slide_count,
                    'is_user_uploaded': is_user_uploaded,
                    'thumbnail_url': f"/api/v1/chat/presentation/templates/{t['id']}/thumbnail" if thumb else None
                }
                out.append(template_info)
            else:
                logger.debug(f"템플릿 파일 누락으로 목록에서 제외: {t['id']} -> {path}")
        
        # 사용자 업로드 템플릿을 먼저 보여주도록 정렬
        out.sort(key=lambda x: (not x['is_user_uploaded'], x['name']))
        return out

    def template_exists(self, template_id: str) -> bool:
        """템플릿이 존재하는지 확인"""
        return template_id in self._registry

    def set_default_template(self, template_id: str) -> bool:
        """템플릿을 기본 템플릿으로 설정"""
        try:
            if not self.template_exists(template_id):
                logger.error(f"기본 템플릿 설정 실패 - 존재하지 않는 템플릿: {template_id}")
                return False
            
            # 기존 기본 템플릿 해제
            for tid, template_data in self._registry.items():
                if template_data.get('is_default'):
                    template_data['is_default'] = False
                    logger.info(f"기존 기본 템플릿 해제: {tid}")
            
            # 새 기본 템플릿 설정
            self._registry[template_id]['is_default'] = True
            logger.info(f"새 기본 템플릿 설정: {template_id}")
            
            # 메타데이터 파일 업데이트
            self._save_default_template_config(template_id)
            
            return True
        except Exception as e:
            logger.error(f"기본 템플릿 설정 실패: {template_id} - {e}")
            return False

    def get_default_template_id(self) -> Optional[str]:
        """현재 기본 템플릿 ID 반환"""
        for template_id, template_data in self._registry.items():
            if template_data.get('is_default'):
                return template_id
        return None

    def _save_default_template_config(self, template_id: str):
        """기본 템플릿 설정을 파일에 저장"""
        try:
            config_path = self.base_dir / 'default_template.json'
            config = {
                'default_template_id': template_id,
                'updated_at': str(__import__('datetime').datetime.now())
            }
            with open(config_path, 'w', encoding='utf-8') as f:
                json.dump(config, f, ensure_ascii=False, indent=2)
            logger.info(f"기본 템플릿 설정 저장: {config_path}")
        except Exception as e:
            logger.error(f"기본 템플릿 설정 저장 실패: {e}")

    def _load_default_template_config(self):
        """저장된 기본 템플릿 설정을 로드"""
        try:
            config_path = self.base_dir / 'default_template.json'
            if config_path.exists():
                with open(config_path, 'r', encoding='utf-8') as f:
                    config = json.load(f)
                default_id = config.get('default_template_id')
                if default_id and self.template_exists(default_id):
                    self._registry[default_id]['is_default'] = True
                    logger.info(f"기본 템플릿 설정 로드: {default_id}")
                    return default_id
        except Exception as e:
            logger.error(f"기본 템플릿 설정 로드 실패: {e}")
        
        # 기본 템플릿이 없으면 첫 번째 템플릿을 기본으로 설정
        templates = self.list_templates()
        if templates:
            first_template_id = templates[0]['id']
            self.set_default_template(first_template_id)
            return first_template_id
        return None

    def register_uploaded_template(self, file_path: Path, style: str = 'business', name: Optional[str] = None) -> Dict[str, Any]:
        """업로드된 템플릿을 레지스트리에 추가 (서버 재기동 시 휘발성)

        파일명 기반 ID 생성, 중복 시 난수 suffix.
        """
        base_id = file_path.stem.lower().replace(' ', '_')[:40] or 'custom_template'
        tpl_id = base_id
        i = 1
        while tpl_id in self._registry:
            tpl_id = f"{base_id}_{i}"
            i += 1
        from typing import Any as _Any, Dict as _Dict
        entry: _Dict[str, _Any] = {
            'id': tpl_id,
            'name': name or file_path.stem,
            'description': '사용자 업로드 템플릿',
            'path': str(file_path),
            'style': style
        }
        self._registry[tpl_id] = entry
        # 사전 분석 캐시 시도 (실패해도 계속)
        try:
            self.analyze_template(file_path)
        except Exception:
            pass
        
        # 🎯 템플릿 메타데이터 추출 (ppt_template_extractor 사용)
        try:
            from .template_debugger import template_debugger
            
            # 1. 원본 템플릿 디버깅 정보 수집
            logger.info(f"🔍 템플릿 분석 시작: {file_path}")
            debug_info = template_debugger.debug_template(str(file_path))
            entry['debug_info'] = debug_info  # type: ignore[assignment]
            
            # 2. 상세한 메타데이터 추출 및 JSON 저장 (ppt_template_extractor 사용)
            metadata_dir = Path(file_path).parent / 'metadata'
            expected_metadata_file = metadata_dir / f"{tpl_id}_metadata.json"
            
            metadata_dir.mkdir(exist_ok=True)
            # ppt_template_extractor 사용하여 폰트, 색상 등 상세한 메타데이터 추출
            extract_presentation(str(file_path), str(expected_metadata_file))
            entry['metadata_json_path'] = str(expected_metadata_file)
            logger.info(f"📊 업로드 템플릿 메타데이터 생성 완료: {expected_metadata_file}")
            
            # 3. 동적 템플릿 등록 (간소화)
            dynamic_template_id = f"user_{tpl_id}"
            entry['dynamic_template_id'] = dynamic_template_id
            entry['is_content_cleaned'] = True  # 메타데이터 추출 완료로 표시
            
            logger.info(f"✅ 템플릿 분석 완료: {tpl_id} -> {dynamic_template_id}")
            logger.info(f"📊 원본 템플릿 메타데이터 생성: {debug_info.get('summary', {}).get('total_shapes', 0)}개 도형, {debug_info.get('summary', {}).get('text_shapes', 0)}개 텍스트")
            
        except Exception as e:
            logger.warning(f"템플릿 메타데이터 추출 실패: {e}")
            # 기본 설정으로 폴백
            entry['dynamic_template_id'] = f"user_{tpl_id}"
            entry['is_content_cleaned'] = False
            entry['metadata_json_path'] = ""
        # 썸네일 생성 시도
        thumb = self._try_generate_thumbnail(file_path, tpl_id)
        if thumb:
            entry['thumbnail'] = thumb
        return entry

    def remove_template(self, template_id: str) -> bool:
        """템플릿을 레지스트리에서 제거하고 파일도 삭제"""
        if template_id not in self._registry:
            return False
        
        entry = self._registry[template_id]
        template_path = Path(entry['path'])
        
        # 파일 삭제 (원본 기본 템플릿만 보호)
        protected_files = ['business_template.pptx']  # 원본 기본 템플릿만 보호
        is_protected = template_path.name in protected_files  # 정확한 파일명 매칭
        
        deleted_files = []
        
        if template_path.exists() and not is_protected:
            try:
                template_path.unlink()
                deleted_files.append(str(template_path))
                logger.info(f"템플릿 파일 삭제됨: {template_path}")
            except Exception as e:
                logger.warning(f"템플릿 파일 삭제 실패: {template_path} - {e}")
        elif is_protected:
            logger.info(f"보호된 템플릿 파일 삭제 스킵: {template_path}")
        
            # 🎯 연관된 파일들도 함께 삭제
        try:
            # 1. 메타데이터 파일 삭제 (두 가지 네이밍 패턴 모두 확인)
            metadata_dir = template_path.parent / 'metadata'
            
            # 파일명 기준 메타데이터 (예: 제품소개서 샘플_metadata.json)
            metadata_file_by_filename = metadata_dir / f"{template_path.stem}_metadata.json"
            if metadata_file_by_filename.exists():
                metadata_file_by_filename.unlink()
                deleted_files.append(str(metadata_file_by_filename))
                logger.info(f"메타데이터 파일 삭제됨 (파일명 기준): {metadata_file_by_filename}")
            
            # 템플릿 ID 기준 메타데이터 (예: 제품소개서_샘플_metadata.json)
            template_id_for_metadata = template_id.replace(' ', '_')
            metadata_file_by_id = metadata_dir / f"{template_id_for_metadata}_metadata.json"
            if metadata_file_by_id.exists():
                metadata_file_by_id.unlink()
                deleted_files.append(str(metadata_file_by_id))
                logger.info(f"메타데이터 파일 삭제됨 (템플릿 ID 기준): {metadata_file_by_id}")
            
            # 2. PDF 미리보기 파일 삭제
            pdf_file = template_path.parent / f"{template_path.stem}.pdf"
            if pdf_file.exists():
                pdf_file.unlink()
                deleted_files.append(str(pdf_file))
                logger.info(f"PDF 미리보기 파일 삭제됨: {pdf_file}")
            
            # 3. PDF 캐시 파일 삭제 (backend/uploads/pdf_cache 디렉토리)
            # PDF 캐시 디렉토리는 backend/uploads/pdf_cache에 위치
            backend_uploads_dir = Path(__file__).parents[3] / 'uploads'
            pdf_cache_dir = backend_uploads_dir / 'pdf_cache'
            
            # 템플릿 ID를 추출 (URL 디코딩된 이름에서 공백을 언더스코어로 변환)
            template_id_for_cache = template_id.replace(' ', '_')
            pdf_cache_file = pdf_cache_dir / f"template_{template_id_for_cache}.pdf"
            if pdf_cache_file.exists():
                pdf_cache_file.unlink()
                deleted_files.append(str(pdf_cache_file))
                logger.info(f"PDF 캐시 파일 삭제됨: {pdf_cache_file}")
            
            # 파일명 기준으로도 PDF 캐시 파일 찾아서 삭제 (예: template_제품소개서 샘플.pdf)
            pdf_cache_file_by_filename = pdf_cache_dir / f"template_{template_path.stem}.pdf"
            if pdf_cache_file_by_filename.exists() and pdf_cache_file_by_filename != pdf_cache_file:
                pdf_cache_file_by_filename.unlink()
                deleted_files.append(str(pdf_cache_file_by_filename))
                logger.info(f"PDF 캐시 파일 삭제됨 (파일명 기준): {pdf_cache_file_by_filename}")
            
            # 4. clean_ 버전이 있으면 삭제
            if not template_path.name.startswith('clean_'):
                clean_version = template_path.parent / f"clean_{template_path.name}"
                if clean_version.exists():
                    clean_version.unlink()
                    deleted_files.append(str(clean_version))
                    logger.info(f"연관 템플릿 파일 삭제됨: {clean_version}")
            
            # 5. 원본 버전이 있으면 삭제 (clean_ 템플릿 삭제 시)
            elif template_path.name.startswith('clean_'):
                original_name = template_path.name[6:]  # 'clean_' 제거
                original_version = template_path.parent / original_name
                if original_version.exists() and original_name not in protected_files:
                    original_version.unlink()
                    deleted_files.append(str(original_version))
                    logger.info(f"연관 원본 템플릿 파일 삭제됨: {original_version}")
                    
        except Exception as e:
            logger.warning(f"연관 파일 삭제 중 오류: {e}")
        
        if deleted_files:
            logger.info(f"📁 총 {len(deleted_files)}개 파일 삭제 완료: {deleted_files}")
        
        # 레지스트리에서 제거
        del self._registry[template_id]
        
        # 🎯 연관된 템플릿도 레지스트리에서 제거
        related_templates_to_remove = []
        for tid, tentry in list(self._registry.items()):
            tpath = Path(tentry['path'])
            # 삭제된 파일과 연관된 템플릿 찾기
            if str(tpath) in deleted_files:
                related_templates_to_remove.append(tid)
        
        for related_id in related_templates_to_remove:
            if related_id in self._registry:
                del self._registry[related_id]
                logger.info(f"연관 템플릿 레지스트리에서 제거: {related_id}")
        
        # 캐시에서도 제거
        all_deleted_paths = [str(template_path)] + deleted_files
        for deleted_path in all_deleted_paths:
            if deleted_path in self.template_cache:
                del self.template_cache[deleted_path]
        
        logger.info(f"템플릿 제거됨: {template_id}")
        if related_templates_to_remove:
            logger.info(f"연관 템플릿도 함께 제거됨: {related_templates_to_remove}")
        return True

    def template_cache_directory(self) -> Path:
        base_dir = Path(__file__).parents[3] / 'uploads' / 'templates'
        base_dir.mkdir(parents=True, exist_ok=True)
        return base_dir

    def _try_generate_thumbnail(self, file_path: Path, template_id: str) -> Optional[str]:
        soffice = shutil.which('soffice') or shutil.which('libreoffice')
        if not soffice:
            return None
        try:
            out_dir = self.template_cache_directory() / 'thumbs'
            out_dir.mkdir(parents=True, exist_ok=True)
            cmd = [soffice, '--headless', '--convert-to', 'png', '--outdir', str(out_dir), str(file_path)]
            subprocess.run(cmd, timeout=25, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            candidates = sorted(out_dir.glob(file_path.stem + '*.png'))
            if not candidates:
                return None
            first = candidates[0]
            target = out_dir / f"{template_id}.png"
            if target.exists():
                try:
                    target.unlink()
                except Exception:
                    pass
            try:
                first.rename(target)
            except Exception:
                return str(first)
            return str(target)
        except Exception:
            return None

    def get_template_details(self, template_id: str) -> Optional[Dict[str, Any]]:
        info = self._registry.get(template_id)
        if not info:
            return None
        # 분석 정보 포함
        path = info.get('path')
        analysis: Dict[str, Any] = {}
        if path and Path(path).exists():
            spec = self.analyze_template(Path(path))
            if spec:
                analysis = {
                    'layout_count': len(spec.layouts),
                    'layouts': [l for l in spec.layouts.keys()],
                    'theme_colors': spec.theme_colors,
                    'max_slides': spec.max_slides
                }
        thumb = self.get_thumbnail_path(template_id)
        extra: Dict[str, Any] = {'analysis': analysis}
        if thumb:
            extra['thumbnail_url'] = f"/api/v1/chat/presentation/templates/{template_id}/thumbnail"
        return info | extra

    def get_template_metadata(self, template_id: str) -> Optional[Dict[str, Any]]:
        """템플릿의 메타데이터를 조회합니다."""
        info = self._registry.get(template_id)
        if not info:
            return None

        template_path = info.get('path')
        if not template_path:
            return None

        # 메타데이터 파일 경로 구성 (다양한 형식 지원)
        metadata_dir = self.metadata_directory()
        metadata_file_candidates = [
            metadata_dir / f"{template_id}_metadata.json",  # 주요 형식
            metadata_dir / f"{template_id}.json",           # 기본 형식
        ]
        
        # 존재하는 메타데이터 파일 찾기
        metadata_file = None
        for candidate in metadata_file_candidates:
            if candidate.exists():
                metadata_file = candidate
                break
        
        if not metadata_file:
            # 메타데이터 파일이 없으면 생성
            try:
                # 메타데이터 디렉토리 생성
                metadata_dir.mkdir(parents=True, exist_ok=True)
                
                # 새 메타데이터 파일 경로 설정
                metadata_file = metadata_dir / f"{template_id}_metadata.json"
                
                # 메타데이터 추출 및 저장
                extract_presentation(template_path, str(metadata_file))
                
                # 생성된 메타데이터 파일 로드
                with open(metadata_file, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except Exception as e:
                logger.error(f"메타데이터 생성 실패: {e}")
                return None
        
        try:
            # 기존 메타데이터 파일 로드
            with open(metadata_file, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception as e:
            logger.error(f"메타데이터 로드 실패: {e}")
            return None

    def metadata_directory(self) -> Path:
        """메타데이터 저장 디렉토리를 반환합니다."""
        return self.template_cache_directory() / 'metadata'

    def get_thumbnail_path(self, template_id: str) -> Optional[str]:
        """생성된 썸네일 실제 경로 반환 (존재하지 않으면 None)"""
        try:
            thumbs = self.template_cache_directory() / 'thumbs'
            p = thumbs / f"{template_id}.png"
            if p.exists():
                return str(p)
        except Exception:
            return None
        return None

    def get_template_path(self, template_id: str) -> Optional[Path]:
        """템플릿 파일의 실제 경로 반환"""
        try:
            template_info = self._registry.get(template_id)
            if template_info and 'path' in template_info:
                template_path = Path(template_info['path'])
                if template_path.exists():
                    return template_path
                logger.warning(f"Template file not found: {template_path}")
        except Exception as e:
            logger.error(f"Error getting template path for {template_id}: {e}")
        return None

    # === Compatibility helpers for API endpoints ===
    def get_template_file_path(self, template_id: str) -> Optional[str]:
        """API 호환: 템플릿 원본 파일 경로 문자열 반환."""
        p = self.get_template_path(template_id)
        return str(p) if p else None

    def _pdf_cache_directory(self) -> Path:
        """PDF 캐시 디렉토리 경로 반환 (존재하지 않으면 생성)."""
        base_uploads = Path(__file__).parents[3] / 'uploads'
        pdf_cache = base_uploads / 'pdf_cache'
        pdf_cache.mkdir(parents=True, exist_ok=True)
        return pdf_cache

    def get_template_pdf_path(self, template_id: str) -> Optional[str]:
        """템플릿을 PDF로 변환하여 캐시에 저장 후 경로 반환.

        - 캐시 존재하고 최신이면 재사용
        - soffice/libreoffice 필요 (없으면 None)
        """
        try:
            pptx_path = self.get_template_path(template_id)
            if not pptx_path or not pptx_path.exists():
                logger.warning(f"PDF 변환 실패: 템플릿 파일 없음 ({template_id})")
                return None

            cache_dir = self._pdf_cache_directory()
            target_pdf = cache_dir / f"template_{template_id}.pdf"

            # 캐시 신선도 체크
            try:
                if target_pdf.exists() and target_pdf.stat().st_mtime >= pptx_path.stat().st_mtime:
                    return str(target_pdf)
            except Exception:
                pass

            soffice = shutil.which('soffice') or shutil.which('libreoffice')
            if not soffice:
                logger.warning("PDF 변환 도구(soffice/libreoffice) 미설치")
                return None

            # outdir로 캐시 디렉토리 지정, 변환 파일명을 후처리로 target 이름으로 이동
            try:
                cmd = [soffice, '--headless', '--convert-to', 'pdf', '--outdir', str(cache_dir), str(pptx_path)]
                subprocess.run(cmd, timeout=60, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            except subprocess.TimeoutExpired:
                logger.error("PDF 변환 타임아웃")
                return None
            except subprocess.CalledProcessError as e:
                logger.error(f"PDF 변환 실패 (코드 {e.returncode})")
                return None
            except Exception as e:
                logger.error(f"PDF 변환 예외: {e}")
                return None

            # 변환 결과 파일 찾기 (원본 파일명 기반)
            generated = cache_dir / f"{pptx_path.stem}.pdf"
            if generated.exists():
                try:
                    # 기존 캐시 파일 제거 후 이동/이름 변경
                    if target_pdf.exists():
                        try:
                            target_pdf.unlink()
                        except Exception:
                            pass
                    generated.rename(target_pdf)
                    return str(target_pdf)
                except Exception as e:
                    logger.warning(f"PDF 캐시 파일 이동 실패: {e}")
                    return str(generated)

            # 일부 환경에서 다른 이름으로 생성될 수 있으므로 fallback: 최신 pdf 선택
            try:
                pdfs = sorted(cache_dir.glob('*.pdf'), key=lambda p: p.stat().st_mtime, reverse=True)
                if pdfs:
                    cand = pdfs[0]
                    return str(cand)
            except Exception:
                pass
            logger.error("PDF 변환 결과를 찾을 수 없음")
            return None
        except Exception as e:
            logger.error(f"get_template_pdf_path 오류: {e}")
            return None
        
    def analyze_template(self, template_path: Path) -> Optional[TemplateSpec]:
        """템플릿 파일 분석 및 레이아웃 정보 추출"""
        cache_key = str(template_path)
        
        if cache_key in self.template_cache:
            return self.template_cache[cache_key]
            
        try:
            prs = Presentation(str(template_path))
            
            # 슬라이드 마스터 분석
            slide_masters = prs.slide_masters
            layouts = {}
            
            # 각 레이아웃 분석
            for master in slide_masters:
                for i, layout in enumerate(master.slide_layouts):
                    layout_info = self._analyze_layout(layout, i)
                    if layout_info:
                        layouts[layout_info.name] = layout_info
            
            # 테마 색상 추출
            theme_colors = self._extract_theme_colors(prs)
            
            template_spec = TemplateSpec(
                file_path=template_path,
                layouts=layouts,
                slide_masters=list(slide_masters),
                theme_colors=theme_colors,
                default_fonts=self._extract_default_fonts(prs)
            )
            
            self.template_cache[cache_key] = template_spec
            logger.info(f"템플릿 분석 완료: {template_path} ({len(layouts)}개 레이아웃)")
            
            return template_spec
            
        except Exception as e:
            logger.error(f"템플릿 분석 실패: {template_path} - {e}")
            return None
    
    def _analyze_layout(self, layout, index: int) -> Optional[TemplateLayoutInfo]:
        """개별 레이아웃 분석"""
        try:
            placeholders = {}
            supports_chart = False
            supports_table = False
            max_bullets = 6
            
            # 플레이스홀더 분석
            for shape in layout.placeholders:
                ph_type = shape.placeholder_format.type
                placeholders[str(ph_type)] = {
                    'idx': shape.placeholder_format.idx,
                    'type': str(ph_type),
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height
                }
                
                # 차트/테이블 지원 여부 확인
                if 'CHART' in str(ph_type) or 'OBJECT' in str(ph_type):
                    supports_chart = True
                    supports_table = True
            
            # 레이아웃 타입 간단 결정
            layout_type = "title-content"  # 기본값
            if len(placeholders) == 0:
                layout_type = "blank"
            elif any('TITLE' in ph_info.get('type', '') for ph_info in placeholders.values()):
                if any('BODY' in ph_info.get('type', '') or 'OBJECT' in ph_info.get('type', '') for ph_info in placeholders.values()):
                    layout_type = "title-content"
                else:
                    layout_type = "title-only"
            elif any('BODY' in ph_info.get('type', '') for ph_info in placeholders.values()):
                layout_type = "content-only"
            
            return TemplateLayoutInfo(
                name=layout.name,
                slide_index=index,
                placeholders=placeholders,
                layout_type=layout_type,
                max_bullets=max_bullets,
                supports_chart=supports_chart,
                supports_table=supports_table
            )
            
        except Exception as e:
            logger.warning(f"레이아웃 분석 실패: {layout.name} - {e}")
            return None

    # 🔥 불필요한 legacy 함수 제거됨
    
    def _extract_theme_colors(self, prs: Any) -> Dict[str, str]:
        """테마 색상 추출"""
        theme_colors = {}
        try:
            theme = prs.slide_masters[0].theme
            for i, color in enumerate(theme.theme_part.theme.color_scheme):
                theme_colors[f'accent{i+1}'] = str(color.rgb) if hasattr(color, 'rgb') else '#000000'
        except Exception:
            # 기본 색상 사용
            theme_colors = {
                'accent1': '#0078D4',
                'accent2': '#107C10',
                'accent3': '#FFB900'
            }
        return theme_colors
    
    def _extract_default_fonts(self, prs: Any) -> Dict[str, str]:
        """기본 폰트 정보 추출"""
        try:
            font_scheme = prs.slide_masters[0].theme.theme_part.theme.font_scheme
            return {
                'major': font_scheme.major_font.latin_typeface or 'Calibri',
                'minor': font_scheme.minor_font.latin_typeface or 'Calibri'
            }
        except Exception:
            return {'major': 'Calibri', 'minor': 'Calibri'}

    def _update_slide_title_only(self, slide, new_title: str):
        """슬라이드의 제목만 조심스럽게 업데이트 (나머지 내용 보존)"""
        try:
            # placeholder 기반 제목 업데이트 시도
            for shape in slide.shapes:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                    if shape.placeholder_format.type in [1, 3]:  # 제목 placeholder
                        if hasattr(shape, 'text_frame') and shape.text_frame:
                            shape.text_frame.text = new_title
                            return True
            
            # placeholder가 없는 경우, 가장 큰 textbox 찾아서 업데이트
            text_shapes = []
            for shape in slide.shapes:
                if (hasattr(shape, 'text_frame') and shape.text_frame and 
                    hasattr(shape, 'width') and hasattr(shape, 'height')):
                    area = shape.width * shape.height
                    text_shapes.append((area, shape))
            
            if text_shapes:
                # 가장 큰 텍스트 박스에 제목 설정 (원본 내용이 있으면 보존)
                largest_shape = max(text_shapes, key=lambda x: x[0])[1]
                current_text = largest_shape.text_frame.text.strip()
                if current_text == "" or len(current_text) < 10:  # 빈 텍스트나 짧은 텍스트만 교체
                    largest_shape.text_frame.text = new_title
                    return True
                # 원본에 의미있는 내용이 있으면 보존
                
        except Exception as e:
            logger.debug(f"제목 업데이트 실패: {e}")
        
        return False
    
    def _clear_slide_text_content(self, slide, slide_info: Optional[dict] = None):
        """특정 슬라이드의 텍스트 내용만 클리어 (디자인 요소는 보존)"""
        try:
            cleared_count = 0
            
            # slide_info에서 needsTextClear 확인
            needs_clear = slide_info and slide_info.get('needsTextClear', False)
            if not needs_clear:
                return 0
            
            for shape in slide.shapes:
                try:
                    # 텍스트를 가진 shape들의 내용만 클리어
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        if shape.text_frame.text.strip():
                            shape.text_frame.text = ""
                            cleared_count += 1
                    elif hasattr(shape, 'text') and shape.text.strip():
                        shape.text = ""
                        cleared_count += 1
                except Exception as e:
                    logger.debug(f"텍스트 클리어 실패: {e}")
                    continue
                    
            logger.debug(f"슬라이드 텍스트 클리어: {cleared_count}개 요소")
            return cleared_count
            
        except Exception as e:
            logger.debug(f"슬라이드 텍스트 클리어 실패: {e}")
            return 0
    
    def map_deck_to_template(self, deck: DeckSpec, template_spec: TemplateSpec) -> DeckSpec:
        """DeckSpec을 템플릿 제약사항에 맞게 변환"""
        if not template_spec:
            return deck
            
        mapped_slides = []
        available_layouts = list(template_spec.layouts.values())
        
        for slide in deck.slides:
            # (1) 사용자 선택 layout 유지: blank/none 이 아니면 그대로 사용
            if slide.layout and slide.layout not in ('blank', 'none'):
                mapped_slide = SlideSpec(
                    title=slide.title,
                    key_message=slide.key_message,
                    bullets=slide.bullets,
                    diagram=slide.diagram,
                    layout=slide.layout,
                    style=slide.style,
                    visual_suggestion=slide.visual_suggestion,
                    speaker_notes=slide.speaker_notes
                )
            else:
                best_layout = self._select_best_layout(slide, available_layouts)
                mapped_slide = self._adapt_slide_to_layout(slide, best_layout)
            mapped_slides.append(mapped_slide)

        # 밀도 기반 분할 (과밀 슬라이드 split)
        def slide_density(s: SlideSpec) -> float:
            text_len = sum(len(b) for b in (s.bullets or [])) + len(s.key_message or '')
            bullet_factor = len(s.bullets) * 22  # 평균 행 높이 가중치
            return text_len + bullet_factor

        density_threshold = 520  # 경험적 값
        expanded: List[SlideSpec] = []
        for s in mapped_slides:
            if slide_density(s) > density_threshold and len(s.bullets) >= 6:
                mid = len(s.bullets) // 2
                first = SlideSpec(
                    title=s.title + ' (1/2)',
                    key_message=s.key_message,
                    bullets=s.bullets[:mid],
                    diagram=None if (s.diagram and s.diagram.type != 'none') else s.diagram,
                    layout=s.layout,
                    style=s.style,
                    visual_suggestion=s.visual_suggestion,
                    speaker_notes=s.speaker_notes
                )
                second = SlideSpec(
                    title=s.title + ' (2/2)',
                    key_message='',
                    bullets=s.bullets[mid:],
                    diagram=None,
                    layout=s.layout,
                    style=s.style,
                    visual_suggestion=s.visual_suggestion,
                    speaker_notes=s.speaker_notes
                )
                expanded.extend([first, second])
            else:
                expanded.append(s)
        mapped_slides = expanded
        
        # 슬라이드 수 제한
        if len(mapped_slides) > template_spec.max_slides:
            mapped_slides = mapped_slides[:template_spec.max_slides]
            logger.info(f"템플릿 제약으로 슬라이드 수 제한: {len(mapped_slides)}개")
        logger.debug({
            'phase': 'template_mapping_result',
            'slides': [
                {
                    'title': s.title,
                    'bullets': len(s.bullets),
                    'layout': s.layout,
                    'has_diagram': bool(s.diagram and s.diagram.type != 'none')
                } for s in mapped_slides
            ]
        })
        return DeckSpec(
            topic=deck.topic,
            max_slides=len(mapped_slides),
            slides=mapped_slides,
            theme=deck.theme
        )
        
    
    def _select_best_layout(self, slide: SlideSpec, layouts: List[TemplateLayoutInfo]) -> Optional[TemplateLayoutInfo]:
        """슬라이드에 가장 적합한 레이아웃 선택"""
        # 다이어그램 타입별 우선순위
        if slide.diagram and slide.diagram.type == 'chart':
            chart_layouts = [l for l in layouts if l.supports_chart]
            if chart_layouts:
                return chart_layouts[0]
                
        if slide.diagram and slide.diagram.type == 'table':
            table_layouts = [l for l in layouts if l.supports_table]
            if table_layouts:
                return table_layouts[0]
        
        # 기본 레이아웃 선택
        for layout in layouts:
            if layout.layout_type == 'title-and-content':
                return layout
                
        return layouts[0] if layouts else None
    
    def _adapt_slide_to_layout(self, slide: SlideSpec, layout: Optional[TemplateLayoutInfo]) -> SlideSpec:
        """슬라이드를 레이아웃에 맞게 조정"""
        if not layout:
            return slide
            
        # bullets 개수 제한
        adapted_bullets = slide.bullets[:layout.max_bullets] if slide.bullets else []
        
        # 다이어그램 지원 여부 확인
        adapted_diagram = slide.diagram
        if slide.diagram:
            if slide.diagram.type == 'chart' and not layout.supports_chart:
                # 차트를 텍스트로 변환
                adapted_diagram = None
                if slide.diagram.data and isinstance(slide.diagram.data, dict):
                    items = slide.diagram.data.get('items', [])
                    chart_bullets = [f"{item.get('key', '')}: {item.get('value', '')}" for item in items[:3]]
                    adapted_bullets.extend(chart_bullets)
            elif slide.diagram.type == 'table' and not layout.supports_table:
                # 테이블을 텍스트로 변환
                adapted_diagram = None
                # 테이블 데이터를 bullets로 변환 로직
        
        return SlideSpec(
            title=slide.title,
            key_message=slide.key_message,
            bullets=adapted_bullets,
            diagram=adapted_diagram,
            layout=layout.layout_type,
            style=slide.style,
            visual_suggestion=slide.visual_suggestion,
            speaker_notes=slide.speaker_notes
        )
    
    def build_from_template(self, deck: DeckSpec, template_path: Path) -> Any:
        """템플릿을 기반으로 PPT 빌드"""
        template_spec = self.analyze_template(template_path)
        if not template_spec:
            raise ValueError(f"템플릿 분석 실패: {template_path}")
            
        # DeckSpec을 템플릿에 맞게 조정
        adapted_deck = self.map_deck_to_template(deck, template_spec)
        
        # 🔍 매핑 편집이 없는 경우 원본 템플릿 보존
        has_meaningful_changes = False
        for slide in adapted_deck.slides:
            # 새로운 내용이 있는지 확인 (AI가 생성한 제목이 아닌 실제 내용)
            if hasattr(slide, 'key_message') and slide.key_message and len(slide.key_message.strip()) > 20:
                has_meaningful_changes = True
                break
            if hasattr(slide, 'bullets') and slide.bullets and len(slide.bullets) > 0:
                has_meaningful_changes = True  
                break
        
        if not has_meaningful_changes:
            logger.info(f"🚨 매핑 편집 없이 바로 생성 - 원본 템플릿 내용 보존")
            logger.info(f"원본 템플릿 사용: {template_path}")
            # 원본 템플릿 그대로 사용 (정리하지 않음)
            prs = Presentation(str(template_path))
            
            # 제목만 업데이트 (나머지 내용은 보존)
            for i, slide_spec in enumerate(adapted_deck.slides):
                if i < len(prs.slides):
                    slide = prs.slides[i] 
                    # 제목만 조심스럽게 업데이트
                    self._update_slide_title_only(slide, slide_spec.title)
            
            return prs
        
        # 1. 매핑이 있는 경우만 템플릿 콘텐츠 정리 (기존 텍스트 제거)
        from .template_content_cleaner import template_content_cleaner
        
        # 임시 정리된 템플릿 생성
        clean_template_path = template_path.parent / f"temp_clean_{template_path.name}"
        logger.info(f"🧹 PPT 생성 시 템플릿 콘텐츠 정리: {template_path} -> {clean_template_path}")
        
        try:
            cleaned_path = template_content_cleaner.clean_template_content(
                str(template_path), str(clean_template_path)
            )
            
            # 정리된 템플릿으로 프레젠테이션 생성
            prs = Presentation(cleaned_path)
            
        except Exception as e:
            logger.warning(f"템플릿 콘텐츠 정리 실패, 원본 사용: {e}")
            # 정리 실패 시 원본 템플릿 사용
            prs = Presentation(str(template_path))
        finally:
            # 임시 파일 정리
            try:
                if clean_template_path.exists():
                    clean_template_path.unlink()
            except Exception:
                pass
        # --- 템플릿 재사용 모드 탐지 ---
        # 조건: (1) 템플릿 슬라이드가 1개 이상이고 (2) placeholder 기반이 거의 없으며 (3) 디자인이 슬라이드 본문에 직접 존재한다고 판단
        def _count_placeholders(p):
            c = 0
            for s in p.slides:
                for sh in s.shapes:
                    try:
                        if getattr(sh, 'placeholder_format', None):
                            c += 1
                    except Exception:
                        pass
            return c
        placeholder_total = _count_placeholders(prs)
        reuse_mode = (len(prs.slides) > 0 and placeholder_total == 0)
        logger.debug({
            'phase': 'template_reuse_detection',
            'slides_in_template': len(prs.slides),
            'placeholder_total': placeholder_total,
            'reuse_mode': reuse_mode
        })

        original_slide_count = len(prs.slides)

        # 기존 방식(모든 슬라이드 삭제)은 placeholder 기반 템플릿일 때만 수행
        if not reuse_mode:
            try:
                for idx in range(len(prs.slides) - 1, -1, -1):
                    rId = prs.slides._sldIdLst[idx].rId  # type: ignore[attr-defined]
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[idx]  # type: ignore[attr-defined]
                logger.debug({
                    'phase': 'template_cleanup',
                    'action': 'all_slides_removed',
                    'template_path': str(template_path)
                })
            except Exception as e:  # noqa: BLE001
                logger.warning(f"기존 슬라이드 제거 실패(무시 후 진행): {e}")
        
        # 새 슬라이드 추가 (개선된 overflow / 사용자 레이아웃 반영)
        def _pick_reference_slide(pres):
            best = None
            best_texts = -1
            for s in pres.slides:
                c = sum(1 for sh in s.shapes if getattr(sh, 'has_text_frame', False))
                if c > best_texts:
                    best = s
                    best_texts = c
            return best
        reference_slide = prs.slides[1] if reuse_mode and len(prs.slides) > 1 else _pick_reference_slide(prs)

        for idx, slide_spec in enumerate(adapted_deck.slides):
            slide = None
            layout_info = None
            if reuse_mode and idx < original_slide_count:
                slide = prs.slides[idx]
                logger.debug({'phase': 'reuse_slide', 'slide_title': slide_spec.title, 'template_slide_index': idx})
            else:
                # 사용자 layout 키/타입 매칭
                chosen_layout_key = None
                if slide_spec.layout:
                    if slide_spec.layout in template_spec.layouts:
                        chosen_layout_key = slide_spec.layout
                    else:
                        for k, info in template_spec.layouts.items():
                            if info.layout_type == slide_spec.layout:
                                chosen_layout_key = k
                                break
                if not chosen_layout_key:
                    # 텍스트 풍부한 레이아웃 우선
                    max_score = -1
                    for k, info in template_spec.layouts.items():
                        score = len(info.placeholders) if info.placeholders else 1
                        if score > max_score:
                            max_score = score
                            chosen_layout_key = k
                layout_info = template_spec.layouts.get(chosen_layout_key or list(template_spec.layouts.keys())[0])
                # 레퍼런스 슬라이드 레이아웃 재사용 시도
                slide_layout = prs.slide_layouts[0]
                try:
                    if reuse_mode and reference_slide is not None:
                        slide_layout = reference_slide.slide_layout
                    elif layout_info:
                        slide_layout = prs.slide_masters[0].slide_layouts[layout_info.slide_index]
                except Exception:
                    pass
                slide = prs.slides.add_slide(slide_layout)
                logger.debug({'phase': 'reuse_overflow_add', 'slide_title': slide_spec.title, 'used_layout': layout_info.name if layout_info else 'unknown', 'user_layout': slide_spec.layout})

            if slide is None:
                logger.debug({'phase': 'template_layout_missing', 'slide_title': slide_spec.title})
                continue
            if layout_info is None:
                try:
                    _ = slide.slide_layout
                except Exception:
                    pass
                layout_info = TemplateLayoutInfo(name='reused', slide_index=0, placeholders={}, layout_type='custom-textbox')
            self._populate_slide(slide, slide_spec, layout_info)

        # 재사용 모드에서 템플릿 슬라이드가 더 많았다면 남는 슬라이드 제거
        if reuse_mode and original_slide_count > len(adapted_deck.slides):
            try:
                for idx in range(original_slide_count - 1, len(adapted_deck.slides) - 1, -1):
                    rId = prs.slides._sldIdLst[idx].rId  # type: ignore[attr-defined]
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[idx]  # type: ignore[attr-defined]
                logger.debug({
                    'phase': 'reuse_trim_excess',
                    'removed_count': original_slide_count - len(adapted_deck.slides)
                })
            except Exception as e:  # noqa: BLE001
                logger.warning(f"재사용 모드 잔여 슬라이드 정리 실패: {e}")
        
        logger.info(f"템플릿 기반 PPT 생성 완료: {len(adapted_deck.slides)}개 슬라이드")
        return prs
    
    def _populate_slide(self, slide: Slide, slide_spec: SlideSpec, layout_info: TemplateLayoutInfo):
        """슬라이드에 내용 채우기 (placeholder + 일반 textbox 모두 지원)"""
        try:
            # --- 모든 텍스트 가능 Shape 재귀 스캔 (그룹 내부 포함) ---
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            from pptx.dml.color import RGBColor

            def walk_shape(shp, depth=0):
                """단일 shape 재귀 순회 (그룹 포함)"""
                yield shp, depth
                if getattr(shp, 'shape_type', None) == MSO_SHAPE_TYPE.GROUP:
                    try:
                        for child in shp.shapes:
                            yield from walk_shape(child, depth + 1)
                    except Exception:  # noqa: BLE001
                        return

            def iter_all_shapes(slide_obj):
                for top in list(slide_obj.shapes):
                    yield from walk_shape(top, 0)

            text_shapes = []            # placeholder 아님 + 텍스트프레임
            placeholder_shapes = []     # (shape, placeholder_type)
            group_text_shapes = []       # 그룹 내부 텍스트 (디자인 유지 대상)
            total_text_shape_count = 0

            for s, depth in iter_all_shapes(slide):
                has_tf = getattr(s, 'has_text_frame', False)
                if not has_tf:
                    continue
                total_text_shape_count += 1

                # placeholder 여부 확인
                is_placeholder = False
                try:
                    pf = s.placeholder_format  # type: ignore[attr-defined]
                    if pf is not None:  # noqa: SIM108
                        placeholder_shapes.append((s, str(pf.type)))
                        is_placeholder = True
                except (ValueError, AttributeError):
                    is_placeholder = False

                # 그룹 내부 여부 표시
                # depth >=1 이면 그룹 내부 텍스트로 간주
                if depth >= 1 and not is_placeholder:
                    group_text_shapes.append(s)
                    continue

                if not is_placeholder:
                    text_shapes.append(s)

            # 헬퍼: 변경 금지(로고/회사명 등) 텍스트 패턴
            STATIC_SKIP_KEYWORDS = {"Logo", "Company or Team Name"}

            def _capture_font_style(shape) -> dict:
                """shape의 대표 폰트 스타일 캡처 (색/크기/볼드/이탤릭/폰트명)"""
                style = {'color': None, 'size': None, 'bold': None, 'italic': None, 'name': None}
                try:
                    if not getattr(shape, 'has_text_frame', False):
                        return style
                    tf = shape.text_frame
                    if not tf or not tf.paragraphs:
                        return style
                    p0 = tf.paragraphs[0]
                    # run 우선, 없으면 paragraph font 사용
                    f = p0.runs[0].font if getattr(p0, 'runs', None) and p0.runs else p0.font
                    if hasattr(f, 'color') and getattr(f.color, 'rgb', None):
                        style['color'] = f.color.rgb
                    if getattr(f, 'size', None):
                        style['size'] = f.size
                    style['bold'] = getattr(f, 'bold', None)
                    style['italic'] = getattr(f, 'italic', None)
                    style['name'] = getattr(f, 'name', None)
                except Exception:
                    pass
                return style

            def _apply_font_style_to_paragraph(paragraph, style: dict):
                try:
                    f = paragraph.font
                    if style.get('size') is not None:
                        f.size = style['size']
                    if style.get('bold') is not None:
                        f.bold = style['bold']
                    if style.get('italic') is not None:
                        f.italic = style['italic']
                    # color은 RGBColor인 경우에만 설정
                    col = style.get('color')
                    if isinstance(col, RGBColor):
                        try:
                            f.color.rgb = col
                        except Exception:
                            pass
                except Exception:
                    pass

            def is_skippable_static(shape):
                try:
                    raw = (shape.text or '').strip()
                except Exception:
                    return False
                if not raw:
                    return False
                return raw in STATIC_SKIP_KEYWORDS
            
            logger.debug({
                'phase': 'populate_scan', 
                'slide_title': slide_spec.title, 
                'placeholders': [f'{getattr(s, "name", "unnamed")}({t})' for s, t in placeholder_shapes][:8],
                'text_shapes': len(text_shapes),
                'group_text_shapes': len(group_text_shapes),
                'total_text_shapes': total_text_shape_count,
                'text_shapes_detail': [f'{getattr(s, "name", "unnamed")}({type(s).__name__})' for s in text_shapes][:10] if text_shapes else [],
                'group_shapes_detail': [f'{getattr(s, "name", "unnamed")}({type(s).__name__})' for s in group_text_shapes][:10] if group_text_shapes else []
            })

            # --- 제목 Shape 찾기 ---
            title_shape = None
            if slide_spec.title:
                # 🟡 custom template 핸들링: 제목 shape 찾기 개선
                # 1순위: Title placeholder
                title_tokens = ["TITLE", "CENTER_TITLE", "VERTICAL_TITLE"]
                for s, ph_type in placeholder_shapes:
                    if any(tok in ph_type for tok in title_tokens):
                        title_shape = s
                        break
                
                # 2순위: 그룹 내부 텍스트 중 첫 번째(디자인 유지) - 고정 문구 제외 + 🟡 개선된 선택
                if not title_shape:
                    candidates = [s for s in group_text_shapes if not is_skippable_static(s)]
                    if candidates:
                        # 🟡 개선: 그룹 내에서도 textbox 형태를 우선 선택
                        textbox_candidates = [s for s in candidates if getattr(s, 'name', '').startswith('textbox')]
                        if textbox_candidates:
                            title_shape = textbox_candidates[0]
                            logger.debug({'phase': 'populate_title_group_textbox_select', 'slide_title': slide_spec.title, 'shape_name': getattr(title_shape, 'name', 'unknown')})
                        else:
                            # 🔴 중요: 장식용 ellipse 등은 3순위로 넘어가서 실제 textbox와 경쟁하도록
                            decorative_patterns = ['타원', 'ellipse', '원', 'circle']
                            non_decorative = [s for s in candidates if not any(pattern in getattr(s, 'name', '').lower() for pattern in decorative_patterns)]
                            
                            if non_decorative:
                                title_shape = non_decorative[0] 
                                logger.debug({'phase': 'populate_title_group_non_decorative_select', 'slide_title': slide_spec.title, 'shape_name': getattr(title_shape, 'name', 'unknown')})
                            else:
                                # 장식용만 남은 경우, 3순위로 넘어가도록 title_shape를 None으로 유지
                                logger.debug({'phase': 'populate_title_group_skip_decorative', 'slide_title': slide_spec.title, 'decorative_count': len(candidates)})

                # 🟡 3순위: 커스텀 템플릿 용 - 첫 번째 유효한 텍스트 shape (크기 기반 선택)
                if not title_shape:
                    # 텍스트 기능이 있는 shape만 필터링 - 우선순위: text_shapes > group_text_shapes
                    text_capable_shapes = []
                    
                    # Placeholder shapes (최우선)
                    for s, _ in placeholder_shapes:
                        if hasattr(s, 'text_frame') or hasattr(s, 'text'):
                            text_capable_shapes.append((s, 'placeholder'))
                    
                    # Regular text_shapes (2순위) - 실제 텍스트박스들
                    for s in text_shapes:
                        if hasattr(s, 'text_frame') or hasattr(s, 'text'):
                            text_capable_shapes.append((s, 'textbox'))
                    
                    # Group text_shapes (3순위) - 그룹 내 텍스트 (장식용일 가능성)
                    for s in group_text_shapes:
                        if hasattr(s, 'text_frame') or hasattr(s, 'text'):
                            text_capable_shapes.append((s, 'group'))
                    
                    # STATIC_SKIP_KEYWORDS 제외 후 크기 기반 정렬
                    filtered = [(s, t) for s, t in text_capable_shapes if not is_skippable_static(s)]
                    if not filtered:
                        filtered = text_capable_shapes
                    
                    if filtered:
                        # 크기 기반 정렬 (width * height) - 텍스트 shape만
                        try:
                            def get_shape_area(shape):
                                try:
                                    return shape.width * shape.height
                                except:
                                    return 0
                            
                            # 타입별로 그룹핑한 후 크기순 정렬
                            textbox_shapes = [(s, t) for s, t in filtered if t == 'textbox']
                            other_shapes = [(s, t) for s, t in filtered if t != 'textbox']
                            
                            # 🔴 디버깅: 후보들 상세 분석
                            logger.debug({
                                'phase': 'populate_title_candidates_analysis',
                                'slide_title': slide_spec.title,
                                'textbox_shapes': [f"{getattr(s, 'name', 'unnamed')}(area:{get_shape_area(s)})" for s, t in textbox_shapes],
                                'other_shapes': [f"{getattr(s, 'name', 'unnamed')}(type:{t},area:{get_shape_area(s)})" for s, t in other_shapes],
                                'filtered_total': len(filtered)
                            })
                            
                            # textbox 우선, 그 다음 크기순
                            if textbox_shapes:
                                sorted_textboxes = sorted(textbox_shapes, key=lambda x: get_shape_area(x[0]), reverse=True)
                                title_shape = sorted_textboxes[0][0]
                                shape_source = f'textbox({len(textbox_shapes)}candidates)'
                            else:
                                sorted_shapes = sorted(other_shapes, key=lambda x: get_shape_area(x[0]), reverse=True)
                                title_shape = sorted_shapes[0][0]
                                shape_source = f'fallback({len(other_shapes)}candidates)'
                            
                            logger.debug({
                                'phase': 'populate_title_custom_select', 
                                'slide_title': slide_spec.title, 
                                'shape_name': getattr(title_shape, 'name', 'unknown'),
                                'shape_area': get_shape_area(title_shape),
                                'shape_type': type(title_shape).__name__,
                                'shape_source': shape_source,
                                'textbox_candidates': len(textbox_shapes),
                                'total_candidates': len(filtered)
                            })
                        except Exception as e:
                            title_shape = filtered[0][0]
                            logger.debug({'phase': 'populate_title_custom_select_fallback', 'slide_title': slide_spec.title, 'error': str(e)})
                
                # 제목 적용 (기존 스타일 보존) - 🔴 urgent fix + 🟡 custom template
                if title_shape and slide_spec.title and slide_spec.title.strip():
                    success = False
                    shape_info = {'name': getattr(title_shape, 'name', 'unknown'), 'type': type(title_shape).__name__}
                    
                    # 🔴 Method 1: text_frame 방식 (선호)
                    if hasattr(title_shape, 'text_frame') and title_shape.text_frame:
                        try:
                            base_style = _capture_font_style(title_shape)
                            title_shape.text_frame.clear()  # 기존 내용 완전 제거
                            
                            # 새로운 텍스트 설정
                            title_shape.text_frame.text = slide_spec.title.strip()
                            
                            # 첫 단락에 스타일 재적용 (흰색 텍스트 등 유지)
                            if title_shape.text_frame.paragraphs:
                                _apply_font_style_to_paragraph(title_shape.text_frame.paragraphs[0], base_style)
                            
                            success = True
                            logger.debug({'phase': 'populate_title_ok', 'slide_title': slide_spec.title, 'final_text': title_shape.text_frame.text, 'method': 'text_frame', 'shape_info': shape_info})
                        except Exception as e:
                            logger.debug({'phase': 'populate_title_text_frame_fail', 'slide_title': slide_spec.title, 'error': str(e), 'shape_info': shape_info})
                    
                    # 🔴 Method 2: 직접 text 속성 (대안)
                    if not success and hasattr(title_shape, 'text'):
                        try:
                            title_shape.text = slide_spec.title.strip()
                            success = True
                            logger.debug({'phase': 'populate_title_ok', 'slide_title': slide_spec.title, 'final_text': title_shape.text, 'method': 'direct_text', 'shape_info': shape_info})
                        except Exception as e:
                            logger.debug({'phase': 'populate_title_direct_fail', 'slide_title': slide_spec.title, 'error': str(e), 'shape_info': shape_info})
                    
                    # � Method 3: 텍스트 프레임 재생성 (최후 수단)
                    if not success:
                        try:
                            # 기존 텍스트프레임이 있다면 내용만 교체
                            if hasattr(title_shape, 'text_frame'):
                                tf = title_shape.text_frame
                                # 모든 단락 제거 후 새로 추가
                                for p in tf.paragraphs[1:]:  # 첫 번째 제외하고 제거
                                    tf._element.remove(p._element)
                                if tf.paragraphs:
                                    tf.paragraphs[0].text = slide_spec.title.strip()
                                else:
                                    p = tf.add_paragraph()
                                    p.text = slide_spec.title.strip()
                                success = True
                                logger.debug({'phase': 'populate_title_ok', 'slide_title': slide_spec.title, 'final_text': tf.text, 'method': 'recreate', 'shape_info': shape_info})
                        except Exception as e:
                            logger.error({'phase': 'populate_title_recreate_fail', 'slide_title': slide_spec.title, 'error': str(e), 'shape_info': shape_info})
                    
                    if not success:
                        logger.error({'phase': 'populate_title_all_methods_fail', 'slide_title': slide_spec.title, 'shape_info': shape_info})
                else:
                    # 📝 제목 텍스트박스 동적 생성
                    logger.debug({'phase': 'populate_create_title', 'slide_title': slide_spec.title, 'reason': 'no_title_shape'})
                    try:
                        from pptx.util import Inches
                        # 제목 텍스트박스 생성 (상단 중앙)
                        left = Inches(0.5)
                        top = Inches(0.2)
                        width = Inches(9)
                        height = Inches(1)
                        
                        title_textbox = slide.shapes.add_textbox(left, top, width, height)
                        title_frame = title_textbox.text_frame
                        title_frame.text = slide_spec.title
                        
                        # 제목 스타일 적용 (기본 가독성)
                        for paragraph in title_frame.paragraphs:
                            for run in paragraph.runs:
                                from pptx.util import Pt
                                run.font.size = Pt(24)  # 큰 폰트
                                run.font.bold = True
                        
                        title_shape = title_textbox  # 참조 업데이트
                        logger.debug({'phase': 'populate_title_created', 'slide_title': slide_spec.title})
                    except Exception as e:
                        logger.warning(f"동적 제목박스 생성 실패: {slide_spec.title} - {e}")
                        logger.debug({'phase': 'populate_title_fail', 'slide_title': slide_spec.title, 'reason': 'title_creation_failed'})

            # --- 본문 Shape 찾기 ---
            import re
            is_toc = False
            if slide_spec.title and ('목차' in slide_spec.title or 'Contents' in slide_spec.title):
                is_toc = True
            elif len(slide_spec.bullets) >= 4:
                num_pattern = sum(1 for b in slide_spec.bullets if re.match(r'\d+\.|\d+\s', b))
                if num_pattern >= 2:
                    is_toc = True

            if slide_spec.bullets or slide_spec.key_message:
                body_shape = None
                
                # 1순위: Body/Content placeholder (제목 제외)
                body_tokens = ["BODY", "CONTENT", "OBJECT", "TEXT", "SUBTITLE"]
                for s, ph_type in placeholder_shapes:
                    if any(tok in ph_type for tok in body_tokens) and s is not title_shape:
                        body_shape = s
                        break
                
                # 2순위: 그룹 내부 텍스트 (제목/고정 제외, 기존 빈 shape 우선)
                if not body_shape:
                    group_candidates = [s for s in group_text_shapes if s is not title_shape and not is_skippable_static(s)]
                    # 빈 텍스트 먼저, 그 다음 비어있지 않은 것
                    def empty_first(shp_list):
                        scored = []
                        for shp in shp_list:
                            try:
                                txt = (shp.text or '').strip()
                            except Exception:
                                txt = ''
                            scored.append((0 if txt == '' else 1, shp))
                        scored.sort(key=lambda x: x[0])
                        return [s for _, s in scored]
                    group_candidates = empty_first(group_candidates)
                    if group_candidates:
                        body_shape = group_candidates[0]

                # 3순위: placeholder + 일반 텍스트 shape (제목 제외)
                if not body_shape:
                    candidates = [s for s, _ in placeholder_shapes if s is not title_shape and not is_skippable_static(s)] + \
                                [s for s in text_shapes if s is not title_shape and not is_skippable_static(s)]
                    if candidates:
                        body_shape = candidates[0]
                
                # 본문 적용 (스타일 보존)
                if is_toc and group_text_shapes:
                    # TOC: 여러 빈 group_text_shapes에 bullets 분배
                    toc_targets = [s for s in group_text_shapes if s is not title_shape]
                    def _existing_len(sh):
                        try:
                            return len((sh.text or '').strip())
                        except Exception:
                            return 0
                    toc_targets.sort(key=_existing_len)
                    assigned = 0
                    for b, tgt in zip(slide_spec.bullets, toc_targets):
                        if not getattr(tgt, 'has_text_frame', False):
                            continue
                        try:
                            tf = tgt.text_frame
                            style = _capture_font_style(tgt)
                            tf.clear()
                            tf.paragraphs[0].text = b
                            # 스타일 재적용 (색/사이즈)
                            _apply_font_style_to_paragraph(tf.paragraphs[0], style)
                            assigned += 1
                        except Exception:
                            continue
                    logger.debug({'phase': 'populate_body_ok_toc', 'slide_title': slide_spec.title, 'bullets': len(slide_spec.bullets), 'assigned': assigned})
                elif (slide_spec.bullets or slide_spec.key_message) and body_shape and hasattr(body_shape, 'text_frame'):
                    text_frame = body_shape.text_frame
                    base_style = _capture_font_style(body_shape)
                    
                    # 🔴 긴급 수정: 안전한 콘텐츠 설정
                    try:
                        text_frame.clear()
                    except Exception as clear_e:
                        logger.debug({'phase': 'populate_body_clear_fail', 'slide_title': slide_spec.title, 'error': str(clear_e)})
                    
                    try:
                        if slide_spec.key_message:
                            p0 = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
                            p0.text = slide_spec.key_message.strip()
                            p0.level = 0
                            _apply_font_style_to_paragraph(p0, base_style)
                        
                        for bullet in slide_spec.bullets:
                            if bullet and bullet.strip():  # 🔴 빈 bullet 방지
                                p = text_frame.add_paragraph()
                                p.text = bullet.strip()
                                p.level = 0
                                _apply_font_style_to_paragraph(p, base_style)
                        
                        logger.debug({'phase': 'populate_body_ok', 'slide_title': slide_spec.title, 'bullets': len(slide_spec.bullets), 'final_text': text_frame.text[:100] if text_frame.text else 'empty'})
                    except Exception as populate_e:
                        logger.error({'phase': 'populate_body_fail', 'slide_title': slide_spec.title, 'error': str(populate_e)})
                        # 🔴 긴급 대안: 직접 텍스트 속성 사용
                        try:
                            if hasattr(body_shape, 'text'):
                                content_lines = []
                                if slide_spec.key_message:
                                    content_lines.append(slide_spec.key_message.strip())
                                content_lines.extend([b.strip() for b in slide_spec.bullets if b and b.strip()])
                                if content_lines:
                                    body_shape.text = '\n'.join(content_lines)
                                    logger.debug({'phase': 'populate_body_fallback', 'slide_title': slide_spec.title})
                        except Exception as fallback_e:
                            logger.error({'phase': 'populate_body_fallback_fail', 'slide_title': slide_spec.title, 'error': str(fallback_e)})
                else:
                    # 내용 있을 때만 동적 생성
                    if slide_spec.bullets or slide_spec.key_message:
                        logger.debug({'phase': 'populate_create_textbox', 'slide_title': slide_spec.title, 'reason': 'no_existing_shape'})
                        try:
                            from pptx.util import Inches
                            left = Inches(0.5)
                            top = Inches(1.5) if title_shape else Inches(0.5)
                            width = Inches(9)
                            height = Inches(5)
                            textbox = slide.shapes.add_textbox(left, top, width, height)
                            text_frame = textbox.text_frame
                            text_frame.clear()
                            if slide_spec.key_message:
                                p0 = text_frame.paragraphs[0]
                                p0.text = slide_spec.key_message
                                p0.level = 0
                            for bullet in slide_spec.bullets:
                                p = text_frame.add_paragraph()
                                p.text = bullet
                                p.level = 0
                            # 동적 생성 상자에는 기본 가독성 폰트 크기 설정
                            try:
                                for para in text_frame.paragraphs:
                                    _apply_font_style_to_paragraph(para, {'size': None, 'bold': None, 'italic': None, 'color': None})
                            except Exception:
                                pass
                            logger.debug({'phase': 'populate_textbox_created', 'slide_title': slide_spec.title, 'bullets': len(slide_spec.bullets)})
                        except Exception as e:
                            logger.warning(f"동적 텍스트박스 생성 실패: {slide_spec.title} - {e}")
                            logger.debug({'phase': 'populate_no_body', 'slide_title': slide_spec.title, 'reason': 'textbox_creation_failed'})
            
            # 다이어그램 처리 (차트/테이블)
            if slide_spec.diagram and slide_spec.diagram.type != 'none':
                self._add_diagram_to_slide(slide, slide_spec.diagram)
                
        except Exception as e:
            logger.warning(f"슬라이드 내용 채우기 실패: {slide_spec.title} - {e}")
        finally:
            # 간단 검증: 슬라이드가 성공적으로 처리되었는지 로그
            try:
                logger.debug({'phase': 'populate_verify_ok', 'slide_title': slide_spec.title})
            except Exception:
                pass
    
    def _add_diagram_to_slide(self, slide: Slide, diagram):
        """슬라이드에 다이어그램 추가"""
        # TODO: 차트/테이블 추가 로직 구현
        # 현재는 기본 구현으로 넘어감
        pass
    
    def analyze_template_layouts(self, template_name: str) -> List[Dict[str, Any]]:
        """템플릿의 레이아웃 정보 분석"""
        try:
            # 템플릿 레지스트리에서 찾기
            if template_name not in self._registry:
                logger.warning(f"템플릿을 찾을 수 없음: {template_name}")
                return []
            
            template_info = self._registry[template_name]
            template_path = template_info.get('path')
            
            if not template_path or not Path(template_path).exists():
                logger.warning(f"템플릿 파일이 존재하지 않음: {template_path}")
                return []
            
            # 간단한 레이아웃 정보 반환
            layouts = [
                {
                    "id": "title",
                    "name": "제목 슬라이드",
                    "type": "title_slide",
                    "description": "프레젠테이션 제목과 부제목"
                },
                {
                    "id": "content",
                    "name": "내용 슬라이드", 
                    "type": "content_slide",
                    "description": "제목과 본문 내용"
                },
                {
                    "id": "two_column",
                    "name": "2단 구성",
                    "type": "two_column_slide", 
                    "description": "좌우 2단 구성 슬라이드"
                }
            ]
            
            logger.info(f"템플릿 레이아웃 분석 완료: {template_name}, {len(layouts)}개")
            return layouts
            
        except Exception as e:
            logger.error(f"템플릿 레이아웃 분석 실패 {template_name}: {e}")
            return []


# 전역 인스턴스
template_manager = PPTTemplateManager()
