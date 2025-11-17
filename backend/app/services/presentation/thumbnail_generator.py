"""
PPT 슬라이드 썸네일 생성 서비스
사용자 업로드 템플릿의 각 페이지를 썸네일 이미지로 변환
"""
import os
import logging
from typing import Optional, List, Dict, Any
from pathlib import Path
import hashlib
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
import io
import base64

logger = logging.getLogger(__name__)

class ThumbnailGenerator:
    """슬라이드 썸네일 생성기"""
    
    def __init__(self):
        # 썸네일 캐시 디렉토리 - uploads/templates 디렉토리 사용
        self.thumbnail_cache_dir = Path("uploads/templates/thumbnails")
        self.thumbnail_cache_dir.mkdir(parents=True, exist_ok=True)
        
        # 썸네일 설정
        self.thumbnail_width = 320
        self.thumbnail_height = 240
        self.cache_duration = 86400 * 7  # 7일
    
    def generate_template_thumbnails(self, pptx_file_path: str, template_id: str) -> List[Dict[str, Any]]:
        """템플릿의 모든 슬라이드 썸네일 생성"""
        try:
            logger.info(f"썸네일 생성 시작: {pptx_file_path}")
            
            # 캐시 키 생성 (파일 내용 기반)
            cache_key = self._generate_cache_key(pptx_file_path, template_id)
            
            # 캐시 확인
            cached_thumbnails = self._get_cached_thumbnails(cache_key)
            if cached_thumbnails:
                logger.info(f"캐시된 썸네일 사용: {template_id}")
                return cached_thumbnails
            
            # 새로 생성
            thumbnails = self._create_thumbnails(pptx_file_path, template_id, cache_key)
            
            # 캐시 저장
            self._save_thumbnail_cache(cache_key, thumbnails)
            
            logger.info(f"썸네일 생성 완료: {len(thumbnails)}개 슬라이드")
            return thumbnails
            
        except Exception as e:
            logger.error(f"썸네일 생성 실패: {e}")
            return self._create_fallback_thumbnails(pptx_file_path, template_id)
    
    def _create_thumbnails(self, pptx_file_path: str, template_id: str, cache_key: str) -> List[Dict[str, Any]]:
        """실제 썸네일 생성"""
        try:
            # PowerPoint 파일 로드
            prs = Presentation(pptx_file_path)
            thumbnails = []
            
            for slide_idx, slide in enumerate(prs.slides):
                try:
                    # 슬라이드 정보 추출
                    slide_info = self._extract_slide_info(slide, slide_idx)
                    
                    # 썸네일 이미지 생성
                    thumbnail_data = self._create_slide_thumbnail(slide, slide_info, cache_key, slide_idx)
                    
                    thumbnail = {
                        "slide_index": slide_idx,
                        "slide_title": slide_info["title"],
                        "layout_name": slide_info["layout_name"],
                        "description": slide_info["description"],
                        "thumbnail_data": thumbnail_data,
                        "thumbnail_url": f"/api/v1/chat/presentation/templates/{template_id}/thumbnails/{slide_idx}",
                        "shape_count": slide_info["shape_count"],
                        "text_shapes": slide_info["text_shapes"],
                        "has_images": slide_info["has_images"],
                        "has_charts": slide_info["has_charts"]
                    }
                    
                    thumbnails.append(thumbnail)
                    
                except Exception as e:
                    logger.warning(f"슬라이드 {slide_idx} 썸네일 생성 실패: {e}")
                    # 기본 썸네일 생성
                    thumbnails.append(self._create_fallback_slide_thumbnail(slide_idx, template_id))
            
            return thumbnails
            
        except Exception as e:
            logger.error(f"썸네일 생성 실패: {e}")
            return []
    
    def _extract_slide_info(self, slide, slide_idx: int) -> Dict[str, Any]:
        """슬라이드 정보 추출"""
        info = {
            "title": f"슬라이드 {slide_idx + 1}",
            "layout_name": "사용자 정의",
            "description": "",
            "shape_count": len(slide.shapes),
            "text_shapes": 0,
            "has_images": False,
            "has_charts": False
        }
        
        # 제목 추출 시도
        title_texts = []
        for shape in slide.shapes:
            try:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text = shape.text_frame.text.strip()
                    if text and len(text) < 100:  # 제목으로 보이는 짧은 텍스트
                        title_texts.append(text)
                        info["text_shapes"] += 1
                
                # 이미지 확인
                if shape.shape_type == 13:  # PICTURE
                    info["has_images"] = True
                
                # 차트 확인
                if shape.shape_type == 3:  # CHART
                    info["has_charts"] = True
                    
            except Exception:
                continue
        
        # 가장 적절한 제목 선택
        if title_texts:
            # 가장 짧고 의미있는 텍스트를 제목으로 선택
            info["title"] = min(title_texts, key=len)
            if len(info["title"]) > 30:
                info["title"] = info["title"][:27] + "..."
        
        # 설명 생성
        features = []
        if info["text_shapes"] > 0:
            features.append(f"텍스트 {info['text_shapes']}개")
        if info["has_images"]:
            features.append("이미지")
        if info["has_charts"]:
            features.append("차트")
        
        info["description"] = ", ".join(features) if features else "기본 레이아웃"
        
        return info
    
    def _create_slide_thumbnail(self, slide, slide_info: Dict, cache_key: str, slide_idx: int) -> str:
        """슬라이드 썸네일 이미지 생성"""
        try:
            # 기본 이미지 생성 (실제로는 더 복잡한 렌더링 필요)
            img = Image.new('RGB', (self.thumbnail_width, self.thumbnail_height), color='white')
            draw = ImageDraw.Draw(img)
            
            # 기본 폰트 (시스템에 따라 조정 필요)
            try:
                font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 14)
                title_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 18)
            except:
                font = ImageFont.load_default()
                title_font = ImageFont.load_default()
            
            # 배경 그라데이션 효과
            for y in range(self.thumbnail_height):
                color_val = int(240 + (y / self.thumbnail_height) * 15)
                draw.line([(0, y), (self.thumbnail_width, y)], fill=(color_val, color_val, color_val))
            
            # 테두리
            draw.rectangle([0, 0, self.thumbnail_width-1, self.thumbnail_height-1], outline='#CCCCCC', width=2)
            
            # 제목 표시
            title = slide_info["title"]
            title_bbox = draw.textbbox((0, 0), title, font=title_font)
            title_width = title_bbox[2] - title_bbox[0]
            title_x = (self.thumbnail_width - title_width) // 2
            draw.text((title_x, 20), title, fill='#333333', font=title_font)
            
            # 레이아웃 정보 표시
            layout_text = f"레이아웃: {slide_info['layout_name']}"
            draw.text((10, 60), layout_text, fill='#666666', font=font)
            
            # 구성 요소 정보
            desc_text = f"구성: {slide_info['description']}"
            draw.text((10, 80), desc_text, fill='#666666', font=font)
            
            # 도형 개수 표시
            shape_text = f"도형: {slide_info['shape_count']}개"
            draw.text((10, 100), shape_text, fill='#666666', font=font)
            
            # 기능 아이콘 표시 (간단한 점으로 표현)
            icon_y = 130
            if slide_info["has_images"]:
                draw.rectangle([10, icon_y, 30, icon_y+15], fill='#4CAF50', outline='#388E3C')
                draw.text((35, icon_y+2), "이미지", fill='#4CAF50', font=font)
                icon_y += 25
            
            if slide_info["has_charts"]:
                draw.rectangle([10, icon_y, 30, icon_y+15], fill='#2196F3', outline='#1976D2')
                draw.text((35, icon_y+2), "차트", fill='#2196F3', font=font)
            
            # 슬라이드 번호
            num_text = f"#{slide_idx + 1}"
            draw.text((self.thumbnail_width - 40, self.thumbnail_height - 25), num_text, fill='#999999', font=font)
            
            # Base64 인코딩
            buffer = io.BytesIO()
            img.save(buffer, format='PNG')
            buffer.seek(0)
            
            thumbnail_data = base64.b64encode(buffer.getvalue()).decode('utf-8')
            
            # 파일로도 저장 (캐싱용)
            thumbnail_path = self.thumbnail_cache_dir / f"{cache_key}_{slide_idx}.png"
            img.save(thumbnail_path)
            
            return thumbnail_data
            
        except Exception as e:
            logger.error(f"썸네일 이미지 생성 실패: {e}")
            return self._create_fallback_image_data(slide_idx)
    
    def _create_fallback_image_data(self, slide_idx: int) -> str:
        """기본 썸네일 이미지 데이터 생성"""
        img = Image.new('RGB', (self.thumbnail_width, self.thumbnail_height), color='#F5F5F5')
        draw = ImageDraw.Draw(img)
        
        # 테두리
        draw.rectangle([0, 0, self.thumbnail_width-1, self.thumbnail_height-1], outline='#DDDDDD', width=2)
        
        # 텍스트
        text = f"슬라이드 {slide_idx + 1}"
        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 16)
        except:
            font = ImageFont.load_default()
        
        bbox = draw.textbbox((0, 0), text, font=font)
        text_width = bbox[2] - bbox[0]
        text_height = bbox[3] - bbox[1]
        
        x = (self.thumbnail_width - text_width) // 2
        y = (self.thumbnail_height - text_height) // 2
        
        draw.text((x, y), text, fill='#666666', font=font)
        
        buffer = io.BytesIO()
        img.save(buffer, format='PNG')
        buffer.seek(0)
        
        return base64.b64encode(buffer.getvalue()).decode('utf-8')
    
    def _create_fallback_thumbnails(self, pptx_file_path: str, template_id: str) -> List[Dict[str, Any]]:
        """기본 썸네일 목록 생성"""
        try:
            prs = Presentation(pptx_file_path)
            thumbnails = []
            
            for slide_idx in range(len(prs.slides)):
                thumbnail = self._create_fallback_slide_thumbnail(slide_idx, template_id)
                thumbnails.append(thumbnail)
            
            return thumbnails
            
        except Exception as e:
            logger.error(f"기본 썸네일 생성 실패: {e}")
            return []
    
    def _create_fallback_slide_thumbnail(self, slide_idx: int, template_id: str) -> Dict[str, Any]:
        """기본 슬라이드 썸네일 생성"""
        return {
            "slide_index": slide_idx,
            "slide_title": f"슬라이드 {slide_idx + 1}",
            "layout_name": "기본 레이아웃",
            "description": "썸네일 생성 중...",
            "thumbnail_data": self._create_fallback_image_data(slide_idx),
            "thumbnail_url": f"/api/v1/chat/presentation/templates/{template_id}/thumbnails/{slide_idx}",
            "shape_count": 0,
            "text_shapes": 0,
            "has_images": False,
            "has_charts": False
        }
    
    def _generate_cache_key(self, pptx_file_path: str, template_id: str) -> str:
        """캐시 키 생성"""
        try:
            with open(pptx_file_path, 'rb') as f:
                file_hash = hashlib.md5(f.read()).hexdigest()[:8]
            return f"{template_id}_{file_hash}"
        except Exception:
            return f"{template_id}_{hash(pptx_file_path) % 100000000}"
    
    def _get_cached_thumbnails(self, cache_key: str) -> Optional[List[Dict[str, Any]]]:
        """캐시된 썸네일 조회"""
        # 구현 시 JSON 파일이나 데이터베이스에서 읽기
        return None
    
    def _save_thumbnail_cache(self, cache_key: str, thumbnails: List[Dict[str, Any]]):
        """썸네일 캐시 저장"""
        # 구현 시 JSON 파일이나 데이터베이스에 저장
        pass
    
    def get_slide_thumbnail(self, template_id: str, slide_index: int) -> Optional[bytes]:
        """특정 슬라이드 썸네일 이미지 반환"""
        try:
            # 캐시에서 찾기
            cache_files = list(self.thumbnail_cache_dir.glob(f"*_{slide_index}.png"))
            if cache_files:
                with open(cache_files[0], 'rb') as f:
                    return f.read()
            
            return None
            
        except Exception as e:
            logger.error(f"썸네일 이미지 조회 실패: {e}")
            return None

# 전역 인스턴스
thumbnail_generator = ThumbnailGenerator()
