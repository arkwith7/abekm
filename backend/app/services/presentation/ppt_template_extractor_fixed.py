# PPT 템플릿 메타데이터 추출기 (수정된 버전)
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pathlib import Path

EMU_PER_INCH = 914400
DPI = 96.0

def emu_to_px(v):
    try:
        return round((v/EMU_PER_INCH)*DPI, 2)
    except Exception:
        return None

def color_to_hex(color_obj):
    """색상 객체를 16진수 색상 코드로 변환 (테마 색상 지원)"""
    if color_obj is None:
        return None
    
    try:
        # RGB 색상인 경우
        if hasattr(color_obj, "rgb") and color_obj.rgb is not None:
            rgb = color_obj.rgb
            return "{:02X}{:02X}{:02X}".format(rgb[0], rgb[1], rgb[2])
        
        # 테마 색상인 경우
        if hasattr(color_obj, "type") and hasattr(color_obj, "theme_color"):
            from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
            
            if color_obj.type == MSO_COLOR_TYPE.SCHEME:
                theme_color = color_obj.theme_color
                
                # 일반적인 테마 색상들을 기본값으로 매핑
                theme_colors = {
                    MSO_THEME_COLOR.BACKGROUND_1: "FFFFFF",  # 흰색
                    MSO_THEME_COLOR.TEXT_1: "000000",       # 검정색
                    MSO_THEME_COLOR.BACKGROUND_2: "F2F2F2", # 연한 회색
                    MSO_THEME_COLOR.TEXT_2: "333333",       # 진한 회색
                    MSO_THEME_COLOR.ACCENT_1: "5B9BD5",     # 파란색
                    MSO_THEME_COLOR.ACCENT_2: "70AD47",     # 초록색
                    MSO_THEME_COLOR.ACCENT_3: "A5A5A5",     # 회색
                    MSO_THEME_COLOR.ACCENT_4: "FFC000",     # 노란색
                    MSO_THEME_COLOR.ACCENT_5: "4472C4",     # 진한 파란색
                    MSO_THEME_COLOR.ACCENT_6: "C55911",     # 주황색
                }
                
                if theme_color in theme_colors:
                    return theme_colors[theme_color]
        
        # RGBColor 객체인 경우 (직접)
        if isinstance(color_obj, RGBColor):
            return "{:02X}{:02X}{:02X}".format(color_obj[0], color_obj[1], color_obj[2])
            
    except Exception as e:
        print(f"색상 변환 오류: {e}")
        pass
    
    return None

def get_font_info(run):
    f = run.font
    return {
        "name": f.name,
        "size_pt": float(f.size.pt) if f.size is not None else None,
        "bold": f.bold,
        "italic": f.italic,
        "underline": f.underline,
        "color": color_to_hex(f.color) if hasattr(f, "color") and f.color else None,
    }

def extract_shape(shape):
    shape_type_name = getattr(shape.shape_type, "name", str(shape.shape_type))
    base = {
        "name": shape.name,
        "type": shape_type_name,
        "left_px": emu_to_px(shape.left),
        "top_px": emu_to_px(shape.top),
        "width_px": emu_to_px(shape.width),
        "height_px": emu_to_px(shape.height),
        "rotation_deg": getattr(shape, "rotation", None),
        "zorder": getattr(shape, "zorder", None),
        "is_placeholder": getattr(shape, "is_placeholder", False),
    }

    if base["is_placeholder"] and hasattr(shape, "placeholder_format"):
        phf = shape.placeholder_format
        base["placeholder"] = {
            "type": getattr(getattr(phf, "type", None), "name", None),
            "idx": getattr(phf, "idx", None)
        }

    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        text_data = []
        tf = shape.text_frame
        for p in tf.paragraphs:
            p_info = {
                "level": p.level,
                "alignment": getattr(p.alignment, "name", None) if p.alignment else None,
                "runs": []
            }
            for run in p.runs:
                p_info["runs"].append({
                    "text": run.text,
                    "font": get_font_info(run)
                })
            text_data.append(p_info)
        base["text"] = {
            "raw": shape.text,
            "paragraphs": text_data
        }

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            img = shape.image
            base["image"] = {
                "filename": Path(img.filename).name if img.filename else None,
                "content_type": img.content_type,
                "orig_width_px": emu_to_px(img.width),
                "orig_height_px": emu_to_px(img.height),
            }
        except Exception:
            base["image"] = None

    if shape.shape_type == MSO_SHAPE_TYPE.TABLE and hasattr(shape, "table"):
        table = shape.table
        base["table"] = {
            "rows": len(table.rows),
            "cols": len(table.columns),
        }

    if hasattr(shape, "fill"):
        try:
            fc_rgb = getattr(shape.fill.fore_color, "rgb", None) if shape.fill.fore_color else None
            base["fill"] = {
                "type": getattr(shape.fill.type, "name", None) if shape.fill.type is not None else None,
                "fore_color": color_to_hex(fc_rgb)
            }
        except Exception:
            base["fill"] = None
    if hasattr(shape, "line"):
        try:
            lc_rgb = getattr(shape.line.color, "rgb", None) if shape.line and shape.line.color else None
            base["line"] = {
                "width_pt": float(shape.line.width.pt) if getattr(shape.line, "width", None) else None,
                "color": color_to_hex(lc_rgb)
            }
        except Exception:
            base["line"] = None

    return base

def extract_presentation(path, output_path):
    prs = Presentation(path)
    result = {
        "file": Path(path).name,
        "slide_width_px": emu_to_px(prs.slide_width),
        "slide_height_px": emu_to_px(prs.slide_height),
        "slides": []
    }
    for idx, slide in enumerate(prs.slides, start=1):
        slide_info = {
            "index": idx,
            "layout_name": getattr(slide.slide_layout, "name", None),
            "shapes_count": len(slide.shapes),
            "shapes": [],
            "elements": []  # UI 호환 요소들 추가
        }
        
        textbox_index = 0  # 텍스트박스 인덱스 카운터
        
        for shape in slide.shapes:
            shape_data = extract_shape(shape)
            slide_info["shapes"].append(shape_data)
            
            # 텍스트박스인 경우 elements에도 추가 (UI 호환성)
            if (shape.shape_type.name == "TEXT_BOX" and 
                hasattr(shape, "text") and getattr(shape, "text", "").strip()):
                
                # UI 호환 element 생성
                element = {
                    "id": f"textbox-{idx-1}-{textbox_index}",
                    "content": getattr(shape, "text", ""),
                    "type": "textbox",
                    "position": {
                        "left": emu_to_px(shape.left),
                        "top": emu_to_px(shape.top),
                        "width": emu_to_px(shape.width), 
                        "height": emu_to_px(shape.height)
                    }
                }
                
                # 폰트 정보 추출 (첫 번째 run에서)
                if hasattr(shape, "text_frame"):
                    text_frame = getattr(shape, "text_frame", None)
                    if text_frame and text_frame.paragraphs:
                        first_para = text_frame.paragraphs[0]
                        if first_para.runs:
                            first_run = first_para.runs[0]
                            font_info = get_font_info(first_run)
                            element.update({
                                "fontFamily": font_info.get("name"),
                                "fontSize": font_info.get("size_pt"),
                                "color": font_info.get("color"),
                                "fontWeight": "bold" if font_info.get("bold") else "normal"
                            })
                
                slide_info["elements"].append(element)
                textbox_index += 1
        
        result["slides"].append(slide_info)
        
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2)
    return output_path
