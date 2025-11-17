"""
템플릿 메타데이터 추출 및 JSON 저장 시스템
템플릿의 모든 디자인 요소를 상세히 분석하여 재사용 가능한 메타데이터 생성
"""

import json
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple
from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from pptx.shapes.connector import Connector
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Picture
from pptx.shapes.placeholder import PlaceholderPicture, SlidePlaceholder
from pptx.shapes.table import Table
from pptx.text.text import TextFrame
from pptx.dml.color import RGBColor, ColorFormat
from loguru import logger
import datetime


class TemplateMetadataExtractor:
    """템플릿 메타데이터 추출기 - 상세한 디자인 정보 JSON 저장"""
    
    def __init__(self):
        self.metadata_dir = Path(__file__).parents[3] / 'uploads' / 'templates' / 'metadata'
        self.metadata_dir.mkdir(parents=True, exist_ok=True)
        
    def extract_template_metadata(self, template_path: str, template_id: str) -> Dict[str, Any]:
        """템플릿의 모든 메타데이터를 추출하여 JSON 형태로 반환"""
        try:
            prs = Presentation(template_path)
            
            metadata = {
                "template_id": template_id,
                "template_path": template_path,
                "extracted_at": datetime.datetime.now().isoformat(),
                "slide_size": {
                    "width": prs.slide_width,
                    "height": prs.slide_height
                },
                "total_slides": len(prs.slides),
                "slides": [],
                "master_slides": [],
                "color_scheme": {},
                "font_scheme": {},
                "theme_info": {}
            }
            
            # 슬라이드별 메타데이터 추출
            for i, slide in enumerate(prs.slides):
                slide_metadata = self._extract_slide_metadata(slide, i)
                metadata["slides"].append(slide_metadata)
            
            # 마스터 슬라이드 정보 추출
            for i, master in enumerate(prs.slide_masters):
                master_metadata = self._extract_master_slide_metadata(master, i)
                metadata["master_slides"].append(master_metadata)
            
            # 테마 정보 추출
            metadata["theme_info"] = self._extract_theme_info(prs)
            
            # 색상 스킴 추출
            metadata["color_scheme"] = self._extract_color_scheme(prs)
            
            # 폰트 스킴 추출
            metadata["font_scheme"] = self._extract_font_scheme(prs)
            
            logger.info(f"✅ 템플릿 메타데이터 추출 완료: {template_id} - {metadata['total_slides']}개 슬라이드")
            return metadata
            
        except Exception as e:
            logger.error(f"❌ 템플릿 메타데이터 추출 실패 {template_id}: {e}")
            return {}
    
    def _extract_slide_metadata(self, slide: Slide, slide_index: int) -> Dict[str, Any]:
        """개별 슬라이드 메타데이터 추출"""
        slide_data = {
            "slide_index": slide_index,
            "slide_id": slide.slide_id,
            "layout_name": slide.slide_layout.name if slide.slide_layout else "Unknown",
            "background": self._extract_background_info(slide),
            "shapes": [],
            "placeholders": [],
            "text_styles": {},
            "design_elements": {}
        }
        
        # 도형별 메타데이터 추출
        for shape_idx, shape in enumerate(slide.shapes):
            shape_metadata = self._extract_shape_metadata(shape, shape_idx)
            slide_data["shapes"].append(shape_metadata)
            
            # 플레이스홀더 정보 수집
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                placeholder_info = {
                    "placeholder_type": shape.placeholder_format.type,
                    "shape_index": shape_idx,
                    "position": {
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height
                    }
                }
                slide_data["placeholders"].append(placeholder_info)
        
        return slide_data
    
    def _extract_shape_metadata(self, shape: BaseShape, shape_index: int) -> Dict[str, Any]:
        """도형 메타데이터 추출"""
        shape_data = {
            "shape_index": shape_index,
            "shape_type": str(type(shape).__name__),
            "name": getattr(shape, 'name', f"Shape_{shape_index}"),
            "position": {
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height
            },
            "rotation": getattr(shape, 'rotation', 0),
            "fill": self._extract_fill_info(shape),
            "line": self._extract_line_info(shape),
            "text": None,
            "is_placeholder": hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None
        }
        
        # 텍스트 정보 추출
        if hasattr(shape, 'text_frame') and shape.text_frame:
            shape_data["text"] = self._extract_text_metadata(shape.text_frame)
        
        # 특별한 도형 타입별 추가 정보
        if isinstance(shape, Picture):
            shape_data["image_info"] = self._extract_image_info(shape)
        elif isinstance(shape, Table):
            shape_data["table_info"] = self._extract_table_info(shape)
        elif isinstance(shape, GroupShape):
            shape_data["group_info"] = self._extract_group_info(shape)
            
        return shape_data
    
    def _extract_text_metadata(self, text_frame: TextFrame) -> Dict[str, Any]:
        """텍스트 프레임 메타데이터 추출"""
        text_data = {
            "content": text_frame.text,
            "paragraphs": []
        }
        
        for para_idx, paragraph in enumerate(text_frame.paragraphs):
            para_data = {
                "paragraph_index": para_idx,
                "text": paragraph.text,
                "alignment": str(paragraph.alignment) if paragraph.alignment else None,
                "runs": []
            }
            
            for run_idx, run in enumerate(paragraph.runs):
                run_data = {
                    "run_index": run_idx,
                    "text": run.text,
                    "font": {
                        "name": run.font.name,
                        "size": run.font.size.pt if run.font.size else None,
                        "bold": run.font.bold,
                        "italic": run.font.italic,
                        "underline": run.font.underline,
                        "color": self._extract_color_info(run.font.color) if run.font.color else None
                    }
                }
                para_data["runs"].append(run_data)
            
            text_data["paragraphs"].append(para_data)
        
        return text_data
    
    def _extract_color_info(self, color_format: ColorFormat) -> Dict[str, Any]:
        """색상 정보 추출"""
        if not color_format:
            return None
            
        color_data = {
            "type": str(color_format.type) if color_format.type else None
        }
        
        try:
            if color_format.type == 1:  # RGB 색상
                rgb = color_format.rgb
                color_data.update({
                    "rgb": {
                        "r": rgb.r,
                        "g": rgb.g, 
                        "b": rgb.b,
                        "hex": f"#{rgb.r:02x}{rgb.g:02x}{rgb.b:02x}"
                    }
                })
            elif color_format.type == 2:  # 테마 색상
                color_data["theme_color"] = str(color_format.theme_color)
            elif color_format.type == 3:  # 색상 스킴
                color_data["scheme_color"] = str(color_format.scheme_color)
        except Exception as e:
            logger.warning(f"색상 정보 추출 오류: {e}")
        
        return color_data
    
    def _extract_fill_info(self, shape: BaseShape) -> Dict[str, Any]:
        """채우기 정보 추출"""
        fill_data = {"type": "none"}
        
        try:
            if hasattr(shape, 'fill'):
                fill = shape.fill
                if fill.type is not None:
                    fill_data["type"] = str(fill.type)
                    
                    if fill.type == 1:  # 단색 채우기
                        fill_data["solid_color"] = self._extract_color_info(fill.fore_color)
                    elif fill.type == 2:  # 그래디언트
                        fill_data["gradient"] = "gradient_fill"
                    elif fill.type == 6:  # 패턴
                        fill_data["pattern"] = "pattern_fill"
        except Exception as e:
            logger.warning(f"채우기 정보 추출 오류: {e}")
        
        return fill_data
    
    def _extract_line_info(self, shape: BaseShape) -> Dict[str, Any]:
        """선 정보 추출"""
        line_data = {"type": "none"}
        
        try:
            if hasattr(shape, 'line'):
                line = shape.line
                if line.color:
                    line_data.update({
                        "color": self._extract_color_info(line.color),
                        "width": line.width.pt if line.width else None,
                        "dash_style": str(line.dash_style) if line.dash_style else None
                    })
        except Exception as e:
            logger.warning(f"선 정보 추출 오류: {e}")
        
        return line_data
    
    def _extract_background_info(self, slide: Slide) -> Dict[str, Any]:
        """배경 정보 추출"""
        background_data = {"type": "default"}
        
        try:
            if hasattr(slide, 'background'):
                bg = slide.background
                background_data["fill"] = self._extract_fill_info(bg)
        except Exception as e:
            logger.warning(f"배경 정보 추출 오류: {e}")
        
        return background_data
    
    def _extract_master_slide_metadata(self, master: any, master_index: int) -> Dict[str, Any]:
        """마스터 슬라이드 메타데이터 추출"""
        return {
            "master_index": master_index,
            "name": getattr(master, 'name', f"Master_{master_index}"),
            "layouts": [layout.name for layout in master.slide_layouts]
        }
    
    def _extract_theme_info(self, prs: Presentation) -> Dict[str, Any]:
        """테마 정보 추출"""
        return {
            "has_theme": hasattr(prs, 'core_properties'),
            "slide_masters_count": len(prs.slide_masters)
        }
    
    def _extract_color_scheme(self, prs: Presentation) -> Dict[str, Any]:
        """색상 스킴 추출"""
        # PowerPoint 색상 스킴 추출 (고급 기능)
        return {"extracted": False, "reason": "Advanced color scheme extraction needed"}
    
    def _extract_font_scheme(self, prs: Presentation) -> Dict[str, Any]:
        """폰트 스킴 추출"""
        # PowerPoint 폰트 스킴 추출 (고급 기능)
        return {"extracted": False, "reason": "Advanced font scheme extraction needed"}
    
    def _extract_image_info(self, picture: Picture) -> Dict[str, Any]:
        """이미지 정보 추출"""
        return {
            "image_type": "picture",
            "has_image": True
        }
    
    def _extract_table_info(self, table: Table) -> Dict[str, Any]:
        """테이블 정보 추출"""
        return {
            "rows": len(table.table.rows),
            "columns": len(table.table.columns)
        }
    
    def _extract_group_info(self, group: GroupShape) -> Dict[str, Any]:
        """그룹 도형 정보 추출"""
        return {
            "shape_count": len(group.shapes)
        }
    
    def save_metadata_to_json(self, metadata: Dict[str, Any], template_id: str) -> str:
        """메타데이터를 JSON 파일로 저장"""
        json_path = self.metadata_dir / f"{template_id}_metadata.json"
        
        try:
            with open(json_path, 'w', encoding='utf-8') as f:
                json.dump(metadata, f, ensure_ascii=False, indent=2, default=str)
            
            logger.info(f"📄 메타데이터 JSON 저장 완료: {json_path}")
            return str(json_path)
            
        except Exception as e:
            logger.error(f"❌ JSON 저장 실패 {template_id}: {e}")
            return ""
    
    def load_metadata_from_json(self, template_id: str) -> Optional[Dict[str, Any]]:
        """JSON 파일에서 메타데이터 로드"""
        json_path = self.metadata_dir / f"{template_id}_metadata.json"
        
        if not json_path.exists():
            return None
            
        try:
            with open(json_path, 'r', encoding='utf-8') as f:
                metadata = json.load(f)
            
            logger.info(f"📄 메타데이터 JSON 로드 완료: {template_id}")
            return metadata
            
        except Exception as e:
            logger.error(f"❌ JSON 로드 실패 {template_id}: {e}")
            return None
    
    def extract_and_save_template_metadata(self, template_path: str, template_id: str) -> str:
        """템플릿 메타데이터 추출하고 JSON으로 저장하는 통합 메서드"""
        metadata = self.extract_template_metadata(template_path, template_id)
        
        if metadata:
            json_path = self.save_metadata_to_json(metadata, template_id)
            return json_path
        
        return ""


# 전역 인스턴스
template_metadata_extractor = TemplateMetadataExtractor()
