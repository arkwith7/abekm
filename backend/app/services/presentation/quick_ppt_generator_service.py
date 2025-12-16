"""Quick PPT Generator Service - 원클릭 생성 전용 (템플릿 미적용)

복원 히스토리:
- 원본: quick_ppt_generator_service.py.backup_20250902_151051
- 복원일: 2025-12-09
- 복원자: AI Assistant
- 사유: quick_pptx_builder_tool 도구에서 사용 (프롬프트 참조 도구 복원)
"""
from __future__ import annotations

import re
import os
from pathlib import Path
from typing import Any, Dict, List, Optional
from datetime import datetime

from loguru import logger
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from app.core.config import settings
from .ppt_models import SlideSpec, DeckSpec


class QuickPPTGeneratorService:
    """원클릭 PPT 생성 전용 서비스 - 디자인 무시, 고정 구조
    
    주요 기능:
    - 템플릿 없이 고정 구조로 PPT 생성
    - 3단계 레이아웃: 제목 + 키메시지 + 내용
    - 슬라이드 유형: 표지, 목차, 내용, 감사인사
    - 마크다운/구조화 텍스트 파싱 지원
    
    스타일 기준: Quick_PPT_Generator_Sample.pptx (2024-12-09)
    """
    
    def __init__(self):
        self.upload_dir = Path(settings.resolved_upload_dir)
        self.upload_dir.mkdir(parents=True, exist_ok=True)
        
        # 폰트 설정 (Noto Sans KR 기준)
        self.font_name = 'Noto Sans KR'
        
        # 색상 테마 (Quick_PPT_Generator_Sample.pptx 기준)
        self.colors = {
            "title_main": RGBColor(0x1E, 0x3A, 0x8A),      # #1E3A8A - 대제목 (진한 파란색)
            "title_slide": RGBColor(0x0F, 0x17, 0x2A),     # #0F172A - 슬라이드 제목 (진한 남색)
            "subtitle": RGBColor(0x64, 0x74, 0x8B),        # #64748B - 부제목/설명 (연한 회색)
            "body": RGBColor(0x33, 0x41, 0x55),            # #334155 - 본문 텍스트 (중간 회색)
            "accent": RGBColor(0x1D, 0x4E, 0xD8),          # #1D4ED8 - 강조/번호 (파란색)
            "divider": RGBColor(0xE2, 0xE8, 0xF0),         # #E2E8F0 - 구분선 (연한 회색)
        }
        
        # 폰트 크기 설정 (pt)
        self.font_sizes = {
            "main_title": 48,      # 표지 대제목
            "main_subtitle": 18,   # 표지 부제목
            "slide_title": 27,     # 슬라이드 제목
            "section_title": 15,   # 섹션 제목 (번호 포함)
            "key_message": 13.5,   # 키 메시지/강조
            "body": 12,            # 본문 텍스트
        }
        
        # 레이아웃 설정 (inches)
        self.layout = {
            "margin": 0.42,        # 좌우/상하 마진
            "title_top": 0.42,     # 제목 상단 위치
            "divider_top": 1.11,   # 구분선 위치
            "content_top": 1.46,   # 콘텐츠 시작 위치
        }
    
    def _remove_request_expressions(self, text: str) -> str:
        """요청 표현을 제거하고 명사형 제목으로 정제.
        
        예시:
        - '자동차 산업의 특허분석 방법론에 대해 PPT 작성해 주세요' → '자동차 산업의 특허분석 방법론'
        - 'AI 기술 트렌드 발표 자료 만들어줘' → 'AI 기술 트렌드'
        """
        if not text:
            return text
        
        original = text
        
        # 1. 후위 요청 표현 패턴 (끝에서부터 제거)
        suffix_patterns = [
            r'\s*(에 대해|에 대한|에 관한|에 관해|을 위한|를 위한)\s*(PPT|ppt|프레젠테이션|발표\s*자료|슬라이드).*$',
            r'\s*(PPT|ppt|프레젠테이션|발표\s*자료|슬라이드)\s*(작성|생성|만들|제작).*$',
            r'\s*(작성|생성|만들어|제작)\s*(해|좀)?\s*(주세요|줘|줘요|주십시오|부탁).*$',
            r'\s*(해|좀)?\s*(주세요|줘|줘요|주십시오|부탁).*$',
            r'\s+PPT\s*$',
            r'\s+ppt\s*$',
        ]
        
        for pattern in suffix_patterns:
            text = re.sub(pattern, '', text, flags=re.IGNORECASE).strip()
        
        # 2. 전위 요청 표현 패턴 (앞에서부터 제거)
        prefix_patterns = [
            r'^(다음|아래|위)\s*(내용|주제)(에 대해|으로|로)?\s*',
        ]
        
        for pattern in prefix_patterns:
            text = re.sub(pattern, '', text, flags=re.IGNORECASE).strip()
        
        # 3. 조사 정리 (끝에 '의', '에', '를' 등이 남으면 제거)
        text = re.sub(r'[의에를을가이]$', '', text).strip()
        
        # 결과가 너무 짧으면 원본 반환
        if len(text) < 3:
            text = original
        
        return text

    def generate_fixed_outline(self, topic: str, context_text: str, max_slides: int = 8) -> DeckSpec:
        """원클릭(디자인 무시)용 고정 구조 아웃라인 생성"""
        try:
            logger.info(f"🚀 원클릭 고정 구조 생성 시작: topic='{topic[:50]}', max_slides={max_slides}")
            logger.info(f"📝 입력 컨텍스트 길이: {len(context_text)} 문자")
            logger.info(f"📝 입력 컨텍스트 앞 200자: '{context_text[:200]}'")
            
            max_slides = max(3, min(max_slides, 20))
            
            # 더 강력한 섹션 추출 로직
            lines = [ln.strip() for ln in (context_text or "").split("\n") if ln.strip()]
            logger.info(f"📄 총 라인 수: {len(lines)}")
            
            # 실제 문서 제목 추출 (첫 번째 헤딩이나 제목 라인에서)
            actual_title = topic  # 기본값
            logger.info(f"🔍 제목 추출 시작 - 기본값: '{topic}'")
            logger.info(f"🔍 분석할 라인 수: {len(lines[:5])}")
            for i, line in enumerate(lines[:5]):  # 처음 5줄에서 찾기
                line = line.strip()
                logger.info(f"🔍 라인 {i+1}: '{line}'")
                if line.startswith('###') and not line.startswith('####'):
                    # ### 헤딩에서 제목 추출
                    actual_title = line.lstrip('#').strip()
                    logger.info(f"🎯 문서 제목 추출 (###): '{actual_title}'")
                    break
                elif line.startswith('##') and not line.startswith('###'):
                    # ## 헤딩에서 제목 추출
                    actual_title = line.lstrip('#').strip()
                    logger.info(f"🎯 문서 제목 추출 (##): '{actual_title}'")
                    break
                elif (not line.startswith('#') and len(line) > 5 and len(line) <= 50 and 
                      ('제품' in line or '소개' in line or '시스템' in line or '서비스' in line)):
                    # 일반 텍스트에서 제목으로 보이는 라인 추출
                    actual_title = line
                    logger.info(f"🎯 문서 제목 추출 (텍스트): '{actual_title}'")
                    break
            logger.info(f"🎯 최종 제목: '{actual_title}'")
            
            # 1) 더 정교한 섹션 추출 로직 - 문서 제목은 제외
            sections = self._parse_structured_content(lines, max_slides-3, exclude_title=actual_title)
            
            logger.info(f"🎯 추출된 섹션 수: {len(sections)}")
            for i, section in enumerate(sections):
                logger.info(f"  섹션 {i+1}: '{section['title'][:30]}...' (bullets: {len(section.get('bullets', []))}개)")
            
            slides: List[SlideSpec] = []
            
            # 1) 제목 슬라이드 - 추출된 실제 제목 사용
            slides.append(SlideSpec(title=actual_title or "발표자료", key_message="", bullets=[], layout="title-slide"))
            logger.info("✅ 제목 슬라이드 생성")
            
            # 2) 목차 슬라이드 - 내용 슬라이드만 포함 (제목 슬라이드 제외)
            toc_items = []
            # 섹션들은 1번부터 시작 (제목 슬라이드는 목차에서 제외)
            for i, s in enumerate(sections, start=1):
                section_title = s["title"]
                # 기존 번호 제거 (1., 2., ### 등)
                clean_title = re.sub(r'^\s*(\d+\.|\#+)\s*', '', section_title).strip()
                toc_items.append(f"{i}. {clean_title}")
            
            slides.append(SlideSpec(title="목차", key_message="", bullets=toc_items, layout="title-and-content"))
            logger.info("✅ 목차 슬라이드 생성")
            
            # 3) 내용 슬라이드들 - 페이지 제목은 번호 없이 깔끔하게
            for s in sections:
                # 페이지 제목에서는 번호 제거하여 깔끔하게 표시
                page_title = re.sub(r'^\s*(\d+\.|\#+)\s*', '', s["title"]).strip()
                slides.append(SlideSpec(
                    title=page_title, 
                    key_message=s.get("key_message", ""), 
                    bullets=s.get("bullets", []), 
                    layout="title-and-content"
                ))
                logger.info(f"✅ 내용 슬라이드 생성: '{page_title[:20]}...'")
            
            # 4) 종료 슬라이드
            slides.append(SlideSpec(title="감사합니다", key_message="경청해 주셔서 감사합니다.", bullets=[], layout="title-slide"))
            logger.info("✅ 종료 슬라이드 생성")
            
            deck = DeckSpec(topic=actual_title or "발표자료", slides=slides, max_slides=len(slides))
            logger.info(f"🎉 고정 구조 DeckSpec 생성 완료: 총 {len(slides)}개 슬라이드")
            return deck
            
        except Exception as e:
            logger.error(f"generate_fixed_outline 실패: {e}")
            # 폴백: 최소한의 구조
            fallback_slides = [
                SlideSpec(title=topic or "발표자료", key_message="", bullets=[], layout="title-slide"),
                SlideSpec(title="내용", key_message="주요 내용을 다룹니다.", bullets=["세부사항 1", "세부사항 2", "세부사항 3"], layout="title-and-content"),
                SlideSpec(title="감사합니다", key_message="경청해 주셔서 감사합니다.", bullets=[], layout="title-slide")
            ]
            logger.info(f"⚠️ 폴백 구조 사용: {len(fallback_slides)}개 슬라이드")
            return DeckSpec(topic=topic or "발표자료", slides=fallback_slides, max_slides=len(fallback_slides))

    def build_quick_pptx(self, spec: DeckSpec, file_basename: Optional[str] = None) -> str:
        """원클릭 전용 빌더: Quick_PPT_Generator_Sample.pptx 스타일 적용"""
        # 🆕 topic에서 요청 표현 제거 (명사형으로 축약)
        refined_topic = self._remove_request_expressions(spec.topic)
        if refined_topic != spec.topic:
            logger.info(f"📝 Quick PPT 제목 정제: '{spec.topic[:40]}' → '{refined_topic[:40]}'")
            spec = DeckSpec(
                topic=refined_topic,
                slides=spec.slides,
                max_slides=spec.max_slides
            )
        
        logger.info(f"🏗️ 원클릭 PPT 빌드 시작: {len(spec.slides)}개 슬라이드, topic='{spec.topic}'")
        
        try:
            # 파일명 생성
            if not file_basename:
                safe_topic = re.sub(r'[^\w\s-]', '', spec.topic).strip()
                safe_topic = re.sub(r'[-\s]+', '_', safe_topic)
                # 파일명 길이 제한 추가 (OS 제한 초과 방지)
                safe_topic = safe_topic[:100]
                file_basename = f"quick_presentation_{safe_topic}"
            
            filename = f"{file_basename}.pptx"
            output_path = self.upload_dir / filename
            
            # 새 프레젠테이션 생성 (16:9 비율)
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.50)
            
            for i, slide_spec in enumerate(spec.slides):
                logger.info(f"📄 슬라이드 {i+1} 생성 중: '{slide_spec.title}'")
                
                if i == 0:
                    # 제목 슬라이드 (커스텀 스타일)
                    slide_layout = prs.slide_layouts[6]  # Blank layout
                    slide = prs.slides.add_slide(slide_layout)
                    self._create_title_slide(slide, slide_spec)
                elif slide_spec.title == "감사합니다":
                    # 마지막 슬라이드 (커스텀 스타일)
                    slide_layout = prs.slide_layouts[6]  # Blank layout
                    slide = prs.slides.add_slide(slide_layout)
                    self._create_closing_slide(slide, slide_spec)
                else:
                    # 내용 슬라이드 - 커스텀 3단계 구조
                    slide_layout = prs.slide_layouts[6]  # Blank layout
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # 3단계 구조 생성
                    self._create_three_tier_layout(slide, slide_spec)
                
                logger.info(f"✅ 슬라이드 {i+1} 완료")
            
            # 파일 저장
            prs.save(str(output_path))
            logger.info(f"✅ 원클릭 PPT 빌드 완료: {output_path}")
            return str(output_path)
            
        except Exception as e:
            logger.error(f"build_quick_pptx 실패: {e}")
            raise

    def _create_title_slide(self, slide, slide_spec: SlideSpec):
        """표지 슬라이드 생성 (Quick_PPT_Generator_Sample.pptx 스타일)"""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        
        try:
            slide_width = Inches(13.33)
            slide_height = Inches(7.50)
            
            # 1. 대제목 (중앙 상단)
            title_left = Inches(0.83)
            title_top = Inches(2.24)
            title_width = Inches(6.46)
            title_height = Inches(1.92)
            
            title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            title_frame = title_box.text_frame
            title_frame.text = slide_spec.title
            title_frame.word_wrap = True
            
            # 대제목 스타일링 (Noto Sans KR, 48pt, 볼드, #1E3A8A)
            title_para = title_frame.paragraphs[0]
            title_para.alignment = PP_ALIGN.LEFT
            title_font = title_para.font
            title_font.name = self.font_name
            title_font.size = Pt(self.font_sizes["main_title"])
            title_font.bold = True
            title_font.color.rgb = self.colors["title_main"]
            
            # 2. 부제목/설명 (있는 경우)
            subtitle_text = slide_spec.key_message or "AI 기반 자동 생성 프레젠테이션"
            
            subtitle_left = Inches(0.83)
            subtitle_top = Inches(4.61)
            subtitle_width = Inches(7.16)
            subtitle_height = Inches(0.40)
            
            subtitle_box = slide.shapes.add_textbox(subtitle_left, subtitle_top, subtitle_width, subtitle_height)
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle_text
            subtitle_frame.word_wrap = True
            
            # 부제목 스타일링 (Noto Sans KR, 18pt, #64748B)
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.alignment = PP_ALIGN.LEFT
            subtitle_font = subtitle_para.font
            subtitle_font.name = self.font_name
            subtitle_font.size = Pt(self.font_sizes["main_subtitle"])
            subtitle_font.bold = False
            subtitle_font.color.rgb = self.colors["subtitle"]
            
            logger.info(f"✅ 표지 슬라이드 생성 완료: '{slide_spec.title}'")
            
        except Exception as e:
            logger.error(f"표지 슬라이드 생성 실패: {e}")

    def _create_closing_slide(self, slide, slide_spec: SlideSpec):
        """종료/감사 슬라이드 생성 (Quick_PPT_Generator_Sample.pptx 스타일)"""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        
        try:
            slide_width = Inches(13.33)
            slide_height = Inches(7.50)
            
            # 1. 감사합니다 대제목 (중앙)
            title_left = Inches(0.83)
            title_top = Inches(2.80)  # 더 중앙에 배치
            title_width = Inches(11.67)
            title_height = Inches(1.5)
            
            title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            title_frame = title_box.text_frame
            title_frame.text = slide_spec.title
            title_frame.word_wrap = True
            
            # 대제목 스타일링 (Noto Sans KR, 48pt, 볼드, #1E3A8A)
            title_para = title_frame.paragraphs[0]
            title_para.alignment = PP_ALIGN.CENTER
            title_font = title_para.font
            title_font.name = self.font_name
            title_font.size = Pt(self.font_sizes["main_title"])
            title_font.bold = True
            title_font.color.rgb = self.colors["title_main"]
            
            # 2. 부제목 (있는 경우)
            if slide_spec.key_message and slide_spec.key_message.strip():
                subtitle_left = Inches(0.83)
                subtitle_top = Inches(4.50)
                subtitle_width = Inches(11.67)
                subtitle_height = Inches(0.60)
                
                subtitle_box = slide.shapes.add_textbox(subtitle_left, subtitle_top, subtitle_width, subtitle_height)
                subtitle_frame = subtitle_box.text_frame
                subtitle_frame.text = slide_spec.key_message
                subtitle_frame.word_wrap = True
                
                # 부제목 스타일링 (Noto Sans KR, 18pt, #64748B)
                subtitle_para = subtitle_frame.paragraphs[0]
                subtitle_para.alignment = PP_ALIGN.CENTER
                subtitle_font = subtitle_para.font
                subtitle_font.name = self.font_name
                subtitle_font.size = Pt(self.font_sizes["main_subtitle"])
                subtitle_font.bold = False
                subtitle_font.color.rgb = self.colors["subtitle"]
            
            logger.info(f"✅ 종료 슬라이드 생성 완료: '{slide_spec.title}'")
            
        except Exception as e:
            logger.error(f"종료 슬라이드 생성 실패: {e}")

    def _add_simple_content(self, slide, spec: SlideSpec):
        """간단한 콘텐츠 추가 (목차 구분) - Quick_PPT_Generator_Sample.pptx 스타일"""
        try:
            # 콘텐츠 영역 찾기
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:  # 일반적으로 콘텐츠 플레이스홀더
                    content_placeholder = shape
                    break
            
            if not content_placeholder:
                logger.warning(f"⚠️ '{spec.title}' 슬라이드에 콘텐츠 플레이스홀더를 찾을 수 없습니다")
                return
            
            tf = content_placeholder.text_frame
            tf.clear()
            tf.word_wrap = True
            
            # 목차 슬라이드 구분
            is_agenda = spec.title in ['목차', 'Agenda', 'Contents']
            
            if is_agenda:
                # 목차: 불릿만 표시
                for i, bullet in enumerate(spec.bullets):
                    if bullet and bullet.strip():
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.text = bullet.strip()
                        p.level = 0
                        p.font.name = self.font_name
                        p.font.size = Pt(self.font_sizes["section_title"])
                        p.font.color.rgb = self.colors["accent"]
                        logger.info(f"✅ 목차 항목 추가: '{bullet[:30]}...'")
            else:
                # 일반 슬라이드: 키 메시지 + 불릿
                paragraph_added = False
                
                # 키 메시지 추가
                if spec.key_message and spec.key_message.strip():
                    p = tf.paragraphs[0]
                    p.text = spec.key_message.strip()
                    p.level = 0
                    p.font.name = self.font_name
                    p.font.size = Pt(self.font_sizes["key_message"])
                    p.font.bold = True
                    p.font.color.rgb = self.colors["body"]
                    paragraph_added = True
                    logger.info(f"✅ 키 메시지 추가: '{spec.key_message[:30]}...'")
                
                # 불릿 포인트 추가
                for i, bullet in enumerate(spec.bullets):
                    if bullet and bullet.strip():
                        if paragraph_added:
                            p = tf.add_paragraph()
                        else:
                            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                            paragraph_added = True
                        
                        p.text = f"• {bullet.strip()}"
                        p.level = 1
                        p.font.name = self.font_name
                        p.font.size = Pt(self.font_sizes["body"])
                        p.font.color.rgb = self.colors["body"]
                        logger.info(f"✅ 불릿 추가: '{bullet[:30]}...'")
            
            logger.info(f"🎯 '{spec.title}' 슬라이드 콘텐츠 완료")
            
        except Exception as e:
            logger.error(f"_add_simple_content 실패: {e}")

    def _parse_structured_content(self, lines: List[str], max_sections: int, exclude_title: Optional[str] = None) -> List[Dict[str, Any]]:
        """구조화된 컨텍스트에서 섹션별 상세 내용 추출"""
        sections = []
        current_section = None
        current_bullets = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # 헤딩 패턴 감지 (더 정교한 조건)
            is_heading = False
            if line.startswith(('#', '##', '###')):
                is_heading = True
            elif (re.match(r'^\d+\.\s+[가-힣A-Za-z]', line) and len(line) <= 30):  # 짧은 제목일 때만 헤딩으로 인식
                is_heading = True
            elif (line.endswith(':') and len(line) <= 50 and 
                  not re.match(r'.*[0-9]+.*[x×].*[0-9]+', line) and  # 크기/측정값이 아닌 경우
                  not re.search(r'[0-9]+\s*(units?|mg/dL|mm|g|일)', line)):  # 단위가 없는 경우
                is_heading = True
                
            if is_heading:
                # 새 섹션 시작
                title = re.sub(r'^#+\s*|\d+\.\s*|:$', '', line).strip()
                
                # 문서 제목과 동일한 섹션은 제외 (중복 방지)
                if exclude_title and title == exclude_title:
                    logger.info(f"🚫 문서 제목과 동일한 섹션 제외: '{title}'")
                    continue
                
                # 이전 섹션 저장
                if current_section:
                    current_section['bullets'] = current_bullets[:6]
                    sections.append(current_section)
                
                current_section = {
                    'title': title,
                    'key_message': f"{title}의 핵심 내용입니다.",
                    'bullets': []
                }
                current_bullets = []
                
            elif current_section:
                # 현재 섹션의 내용 라인 처리
                if line.startswith('-'):
                    # 불릿 포인트
                    bullet_text = line.lstrip('- ').strip()
                    if bullet_text and len(bullet_text) > 5:
                        current_bullets.append(bullet_text)
                elif re.match(r'^\d+\.\s+', line) and len(line) > 30:  # 긴 numbered list는 bullet으로 처리
                    # 번호 목록 (1. 2. 3. 등)
                    bullet_text = re.sub(r'^\d+\.\s*', '', line).strip()
                    if bullet_text and len(bullet_text) > 5:
                        current_bullets.append(bullet_text)
                elif line.endswith(':') and len(line) <= 60:
                    # 소제목 (콜론으로 끝나는 짧은 라인)
                    subtitle = line.rstrip(':')
                    current_bullets.append(f"**{subtitle}**")
                elif ':' in line and len(line.split(':')) == 2:
                    # 키-값 쌍 (예: "크기: 60mm x 45mm x 15mm")
                    key, value = line.split(':', 1)
                    if len(key.strip()) <= 20 and len(value.strip()) > 0:
                        current_bullets.append(f"{key.strip()}: {value.strip()}")
                elif len(line) > 20 and not line.startswith('**'):
                    # 일반 텍스트 (키 메시지로 사용하거나 불릿으로 변환)
                    if (len(current_bullets) == 0 and len(line) <= 200 and
                        not re.search(r'^[가-힣A-Za-z]+:', line)):  # 키-값 형태가 아닌 경우
                        # 긴 문단은 키메시지로 사용
                        current_section['key_message'] = line
                    elif len(line) <= 200:  # 더 긴 텍스트도 bullet으로 허용
                        current_bullets.append(line)
                elif len(line) > 10:  # 짧은 텍스트도 bullet으로 추가
                    current_bullets.append(line)
                elif line.startswith('**') and line.endswith('**'):
                    # 굵은 텍스트 (소제목)
                    clean_text = line.strip('*')
                    if len(clean_text) <= 80:
                        current_bullets.append(clean_text)
        
        # 마지막 섹션 저장
        if current_section:
            current_section['bullets'] = current_bullets[:6]
            sections.append(current_section)
        
        # 빈 섹션이나 너무 적은 내용의 섹션 필터링 및 보완
        valid_sections = []
        for section in sections[:max_sections]:
            if section.get('title'):
                # 불릿이 없으면 키 메시지라도 있는지 확인
                if not section.get('bullets') and section.get('key_message'):
                    # 키 메시지를 불릿으로 변환
                    key_msg = section['key_message']
                    if len(key_msg) > 100:
                        # 긴 메시지는 문장 단위로 분할
                        sentences = [s.strip() for s in key_msg.split('.') if s.strip()]
                        section['bullets'] = sentences[:3]
                        section['key_message'] = f"{section['title']}에 대한 핵심 내용입니다."
                    else:
                        section['bullets'] = [key_msg]
                        section['key_message'] = f"{section['title']}에 대한 핵심 내용입니다."
                
                valid_sections.append(section)
        
        logger.info(f"📊 구조화 파싱 완료: {len(valid_sections)}개 유효 섹션")
        return valid_sections

    def _create_three_tier_layout(self, slide, slide_spec: SlideSpec):
        """3단계 구조 레이아웃 생성: 제목 + 구분선 + 내용 (Quick_PPT_Generator_Sample.pptx 스타일)"""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        
        try:
            # 슬라이드 크기 (16:9 비율)
            slide_width = Inches(13.33)
            slide_height = Inches(7.50)
            margin = Inches(self.layout["margin"])
            
            # 1. 상단: 슬라이드 타이틀 (Step 제목 스타일)
            title_left = margin
            title_top = Inches(self.layout["title_top"])
            title_width = slide_width - (margin * 2)
            title_height = Inches(0.57)
            
            title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            title_frame = title_box.text_frame
            title_frame.text = slide_spec.title
            title_frame.word_wrap = True
            
            # 제목 스타일링 (Noto Sans KR, 27pt, 볼드, #0F172A)
            title_para = title_frame.paragraphs[0]
            title_para.alignment = PP_ALIGN.LEFT
            title_font = title_para.font
            title_font.name = self.font_name
            title_font.size = Pt(self.font_sizes["slide_title"])
            title_font.bold = True
            title_font.color.rgb = self.colors["title_slide"]
            
            # 2. 구분선 (제목 아래)
            divider_left = margin
            divider_top = Inches(self.layout["divider_top"])
            divider_width = Inches(12.50)
            divider_height = Inches(0.03)
            
            divider = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                divider_left, divider_top, divider_width, divider_height
            )
            divider.fill.solid()
            divider.fill.fore_color.rgb = self.colors["divider"]
            divider.line.fill.background()  # 선 없음
            
            # 3. 키 메시지 (있는 경우)
            content_top = Inches(self.layout["content_top"])
            
            if slide_spec.key_message and slide_spec.key_message.strip():
                key_msg_left = margin
                key_msg_top = content_top
                key_msg_width = slide_width - (margin * 2)
                key_msg_height = Inches(0.6)
                
                key_msg_box = slide.shapes.add_textbox(key_msg_left, key_msg_top, key_msg_width, key_msg_height)
                key_msg_frame = key_msg_box.text_frame
                key_msg_frame.text = slide_spec.key_message
                key_msg_frame.word_wrap = True
                
                # 키 메시지 스타일링 (13.5pt, 볼드, #334155)
                key_msg_para = key_msg_frame.paragraphs[0]
                key_msg_para.alignment = PP_ALIGN.LEFT
                key_msg_font = key_msg_para.font
                key_msg_font.name = self.font_name
                key_msg_font.size = Pt(self.font_sizes["key_message"])
                key_msg_font.bold = True
                key_msg_font.color.rgb = self.colors["body"]
                
                content_top = content_top + Inches(0.7)
            
            # 4. 불릿 포인트 내용
            if slide_spec.bullets and len(slide_spec.bullets) > 0:
                content_left = margin
                content_width = slide_width - (margin * 2)
                content_height = slide_height - content_top - margin
                
                content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
                content_frame = content_box.text_frame
                content_frame.word_wrap = True
                
                # 불릿 포인트 추가
                for i, bullet in enumerate(slide_spec.bullets[:10]):  # 최대 10개
                    if bullet and bullet.strip():
                        if i == 0:
                            para = content_frame.paragraphs[0]
                        else:
                            para = content_frame.add_paragraph()
                        
                        para.text = f"• {bullet.strip()}"
                        para.alignment = PP_ALIGN.LEFT
                        para.level = 0
                        para.space_after = Pt(8)  # 줄 간격
                        
                        # 불릿 스타일링 (12pt, #334155)
                        bullet_font = para.font
                        bullet_font.name = self.font_name
                        bullet_font.size = Pt(self.font_sizes["body"])
                        bullet_font.color.rgb = self.colors["body"]
            
            logger.info(f"✅ Quick PPT 레이아웃 생성 완료: '{slide_spec.title}'")
            
        except Exception as e:
            logger.error(f"Quick PPT 레이아웃 생성 실패: {e}")


# 전역 인스턴스
quick_ppt_service = QuickPPTGeneratorService()
