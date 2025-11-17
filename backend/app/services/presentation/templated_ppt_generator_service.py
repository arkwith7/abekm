"""Templated PPT Generator Service - 템플릿 기반 생성 전용"""
from __future__ import annotations

import json
import asyncio
import re
import os
import tempfile
from pathlib import Path
from typing import Any, Dict, List, Optional
from datetime import datetime

from loguru import logger
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

from app.core.config import settings
from app.services.core.ai_service import ai_service
from .ppt_models import ChartData, DiagramData, SlideSpec, DeckSpec
from .ppt_template_manager import PPTTemplateManager, template_manager
from .enhanced_object_processor import EnhancedPPTObjectProcessor


class TemplatedPPTGeneratorService:
    """템플릿 기반 PPT 생성 전용 서비스 - AI 생성, 템플릿 적용, 고급 기능"""
    
    def __init__(self):
        self.prompts_dir = Path(__file__).parents[3] / "prompts"
        self.upload_dir = Path(settings.file_upload_path or settings.upload_dir)
        self.upload_dir.mkdir(parents=True, exist_ok=True)
        self.template_manager = template_manager
        self.object_processor = EnhancedPPTObjectProcessor()
        
        # 풍부한 색상 테마
        self.color_themes = {
            "corporate_blue": {"primary": RGBColor(0, 102, 204), "secondary": RGBColor(102, 153, 255), "accent": RGBColor(255, 153, 0), "text": RGBColor(51, 51, 51), "background": RGBColor(248, 249, 250)},
            "modern_green": {"primary": RGBColor(34, 139, 34), "secondary": RGBColor(144, 238, 144), "accent": RGBColor(255, 215, 0), "text": RGBColor(47, 79, 79), "background": RGBColor(248, 255, 248)},
            "professional_gray": {"primary": RGBColor(70, 70, 70), "secondary": RGBColor(169, 169, 169), "accent": RGBColor(220, 20, 60), "text": RGBColor(0, 0, 0), "background": RGBColor(245, 245, 245)},
            "playful_violet": {"primary": RGBColor(111, 45, 168), "secondary": RGBColor(181, 126, 220), "accent": RGBColor(255, 181, 71), "text": RGBColor(60, 60, 60), "background": RGBColor(250, 248, 255)},
        }

    # ---------------- Filename / Topic Normalization Helpers ----------------
    def _normalize_topic_for_filename(self, topic: str, max_chars: int = 50) -> str:
        """과도하게 긴 topic(여러 줄, 키메시지 포함 등)을 파일명용으로 정제.

        규칙:
        1. 줄 단위로 분리 후 첫 줄 우선. (첫 줄이 5자 미만이면 다음 줄 탐색)
        2. '키메시지', '키 메시지', '제품 개요' 이후 내용 잘라냄.
        3. 중복 연속 단어 제거.
        4. 허용 문자만 남기고 공백은 '_'로 치환.
        5. 길이 제한 (기본 50자) - 멀티바이트 안전하게 자르기.
        """
        if not topic:
            return "presentation"

        original = topic
        # 줄 분리 + 첫 적절한 라인
        lines = [ln.strip() for ln in re.split(r"[\r\n]+", topic) if ln.strip()]
        if lines:
            # 첫 줄이 지나치게 짧고 두 번째가 더 의미 있으면 교체
            if len(lines[0]) < 5 and len(lines) > 1:
                topic = lines[1]
            else:
                topic = lines[0]

        # 키메시지 및 추가 설명 트리거어 제거
        topic = re.split(r"키 ?메시지|Key Message|상세 설명|제품 개요", topic)[0].strip()

        # 연속 중복 단어 제거 (ex: '제품 소개 제품 개요' -> 앞 1~2개만)
        words = topic.split()
        dedup_words = []
        for w in words:
            if not dedup_words or dedup_words[-1] != w:
                dedup_words.append(w)
        topic = " ".join(dedup_words)

        # 길이 제한 (문자 기준)
        if len(topic) > max_chars:
            topic = topic[:max_chars].rstrip()
        # 최소 길이 보장
        if len(topic) < 2:
            topic = original[:max_chars] if original else "presentation"

        # 파일명용 정규화
        safe = re.sub(r"[^\w\s-]", "", topic)
        safe = re.sub(r"[\s-]+", "_", safe).strip("_")
        if not safe:
            safe = "presentation"

        logger.info(f"🧪 토픽 파일명 정규화: original='{original[:60]}', normalized='{safe}'")
        return safe

    def _load_prompt(self) -> str:
        """프롬프트 파일 로드"""
        try:
            prompt_file = self.prompts_dir / "ppt_generation.txt"
            if prompt_file.exists():
                return prompt_file.read_text(encoding='utf-8')
        except Exception as e:
            logger.warning(f"프롬프트 파일 로드 실패: {e}")
        
        return (
            "당신은 전문 프레젠테이션 디자이너입니다. JSON만 출력. "
            "필드: topic,max_slides,slides[].title,key_message,bullets,layout,diagram,visual_suggestion,speaker_notes"
        )

    async def generate_enhanced_outline(self, topic: str, context_text: str, provider: Optional[str] = None,
                                        template_style: str = "business", include_charts: bool = True,
                                        retries: int = 2, document_filename: Optional[str] = None,
                                        custom_template_path: Optional[str] = None,
                                        presentation_type: str = "general",
                                        user_template_id: Optional[str] = None) -> DeckSpec:
        """AI 기반 향상된 아웃라인 생성"""
        
        logger.info(f"🚀 템플릿 기반 아웃라인 생성 시작: topic='{topic[:50]}', template_style='{template_style}'")
        logger.info(f"📝 파라미터: include_charts={include_charts}, presentation_type={presentation_type}")
        
        try:
            # 주제 개선
            improved_topic = self._improve_topic(topic, context_text)
            
            # AI 프롬프트 구성
            system = self._load_prompt()
            enhanced_requirements = [
                "- AI 응답 내용의 제목과 구조를 정확히 반영하여 슬라이드 생성",
                "- 번호가 있는 섹션(1. 제품 개요, 2. 기술 사양 등)은 각각 별도 슬라이드로 구성",
                "- 각 섹션의 세부 항목들은 bullets로 정확히 나열",
                "- 두 번째 슬라이드는 번호가 있는 섹션들을 목차로 구성",
                "- bullets 항목당 50자 이내로 간결하게 표현",
                f"- include_charts={include_charts} 이면 수치 데이터를 차트로 변환",
                f"- template_style={template_style} (business|minimal|modern|playful)",
                "- visual_suggestion: 관련 아이콘/이미지 아이디어 1줄",
                "- speaker_notes: 발표자 스크립트 2~4문장 한국어",
                "- 각 슬라이드 title은 섹션 번호와 제목을 포함 (예: '1. 제품 개요')",
                "- key_message는 해당 섹션의 핵심 설명문으로 구성"
            ]
            
            user_content = [
                f"주제: {improved_topic}",
                f"컨텍스트:\n{context_text[:8000]}",
                "요구사항:",
                *enhanced_requirements
            ]
            
            # AI 호출
            for attempt in range(retries + 1):
                try:
                    logger.info(f"🤖 AI 호출 시도 {attempt + 1}/{retries + 1}")
                    
                    # AI 메시지 구성
                    ai_message = f"{system}\n\n{chr(10).join(user_content)}"
                    logger.info(f"🔍 AI 메시지 길이: {len(ai_message)}문자")
                    logger.debug(f"🔍 AI 메시지 내용: {ai_message[:500]}...")
                    
                    # Provider 기본값 설정 (.env 기반 Settings 우선)
                    effective_provider = provider or settings.get_current_llm_provider()
                    logger.info(f"🔍 사용할 AI 제공자: {effective_provider}")
                    
                    response_generator = ai_service.chat_stream(
                        messages=[{"role": "user", "content": ai_message}],
                        provider=effective_provider
                    )
                    
                    # 스트림 응답 수집 (안전장치 추가)
                    full_response = ""
                    chunk_count = 0
                    max_chunks = 1000  # 최대 청크 수 제한
                    max_response_length = 50000  # 최대 응답 길이 제한
                    
                    # 타임아웃 처리를 위한 래퍼 함수
                    async def collect_response():
                        nonlocal full_response, chunk_count
                        async for chunk in response_generator:
                            chunk_count += 1
                            
                            # 안전장치: 최대 청크 수 초과 시 종료
                            if chunk_count > max_chunks:
                                logger.warning(f"⚠️ 최대 청크 수({max_chunks}) 초과로 스트림 종료")
                                break
                            
                            # 디버그 로깅 빈도 조절 (100개마다만 로깅)
                            if chunk_count % 100 == 0:
                                logger.debug(f"🔄 청크 {chunk_count}: {type(chunk)} - {str(chunk)[:50]}...")
                            
                            content = ""
                            if chunk:
                                if hasattr(chunk, 'text') and callable(getattr(chunk, 'text', None)):
                                    try:
                                        text_method = getattr(chunk, 'text')
                                        content = text_method()
                                    except Exception:
                                        content = str(chunk)
                                elif isinstance(chunk, str):
                                    content = chunk
                                else:
                                    content = str(chunk)
                            
                            # 빈 내용이나 메서드 바인딩 문자열 제외
                            if content and isinstance(content, str) and not content.startswith('<bound method'):
                                full_response += content
                                
                                # 안전장치: 최대 응답 길이 초과 시 종료
                                if len(full_response) > max_response_length:
                                    logger.warning(f"⚠️ 최대 응답 길이({max_response_length}) 초과로 스트림 종료")
                                    break
                                
                                # JSON 완료 패턴 감지
                                if content.strip().endswith('}') and '{"topic"' in full_response:
                                    # 기본적인 JSON 구조가 완성되었는지 확인
                                    brace_count = full_response.count('{') - full_response.count('}')
                                    if brace_count <= 0:
                                        logger.info(f"🔚 JSON 완료 패턴 감지, 스트림 종료 (청크: {chunk_count})")
                                        break
                                
                                # 진행상황 로깅 (100개마다)
                                if chunk_count % 100 == 0:
                                    logger.info(f"📊 진행상황: {chunk_count}청크, {len(full_response)}자")
                    
                    # 타임아웃과 함께 응답 수집 실행 (최대 60초)
                    try:
                        await asyncio.wait_for(collect_response(), timeout=60.0)
                    except asyncio.TimeoutError:
                        logger.warning(f"⚠️ AI 응답 수집 타임아웃 (60초), 현재까지 수집: {len(full_response)}문자")
                    
                    logger.info(f"📝 AI 응답 수집 완료: {len(full_response)}문자, {chunk_count}개 청크")
                    
                    # JSON 파싱
                    logger.info(f"🔍 AI 응답 파싱 시작: {full_response[:200]}...")
                    deck_spec = self._parse_ai_response(full_response, improved_topic, template_style)
                    
                    if deck_spec and len(deck_spec.slides) >= 2:
                        logger.info(f"✅ 아웃라인 생성 성공: {len(deck_spec.slides)}개 슬라이드")
                        for i, slide in enumerate(deck_spec.slides):
                            logger.info(f"  슬라이드 {i+1}: '{slide.title}' (bullets: {len(slide.bullets)}개)")
                        return deck_spec
                    else:
                        logger.warning(f"⚠️ AI 응답이 부적절함 (시도 {attempt + 1}): deck_spec={deck_spec}, slides={len(deck_spec.slides) if deck_spec else 0}")
                        
                except Exception as e:
                    logger.error(f"❌ AI 호출 실패 (시도 {attempt + 1}): {e}")
                    if attempt == retries:
                        break
                    await asyncio.sleep(1)
            
            # 폴백: 간단한 구조
            logger.warning("⚠️ AI 생성 실패, 폴백 구조 사용")
            return self._create_fallback_outline(improved_topic, context_text)
            
        except Exception as e:
            logger.error(f"generate_enhanced_outline 실패: {e}")
            return self._create_fallback_outline(topic, context_text)

    def _improve_topic(self, topic: str, context_text: str) -> str:
        """주제 개선 - AI 답변에서 실제 제목 추출"""
        if not context_text:
            return topic.strip()
        
        # AI 답변에서 실제 제목 추출 (quick PPT와 동일한 로직)
        lines = [ln.strip() for ln in context_text.split('\n') if ln.strip()]
        actual_title = topic  # 기본값
        
        for line in lines[:5]:  # 처음 5줄에서 찾기
            line = line.strip()
            if line.startswith('###') and not line.startswith('####'):
                # ### 헤딩에서 제목 추출
                actual_title = line.lstrip('#').strip()
                logger.info(f"🎯 템플릿 PPT 제목 추출 (###): '{actual_title}'")
                break
            elif line.startswith('##') and not line.startswith('###'):
                # ## 헤딩에서 제목 추출
                actual_title = line.lstrip('#').strip()
                logger.info(f"🎯 템플릿 PPT 제목 추출 (##): '{actual_title}'")
                break
            elif (not line.startswith('#') and len(line) > 5 and len(line) <= 50 and 
                  ('제품' in line or '소개' in line or '시스템' in line or '서비스' in line)):
                # 일반 텍스트에서 제목으로 보이는 라인 추출
                actual_title = line
                logger.info(f"🎯 템플릿 PPT 제목 추출 (텍스트): '{actual_title}'")
                break
        
        return actual_title.strip()

    def _extract_json(self, text: str) -> str:
        """텍스트에서 JSON 부분을 추출"""
        if text.strip().startswith('{'):
            return text
        block = re.search(r"```(?:json)?\n(.*)```", text, re.DOTALL)
        if block:
            return block.group(1)
        brace = re.search(r"{.*}", text, re.DOTALL)
        return brace.group(0) if brace else text

    def _parse_outline(self, text: str, fallback_topic: Optional[str] = None) -> DeckSpec:
        """아웃라인 텍스트를 파싱하여 DeckSpec으로 변환 (enhanced 서비스와 호환)"""
        try:
            data = json.loads(self._extract_json(text))
            max_slides = int(data.get("max_slides", 10))
            raw_slides = data.get("slides", [])[:max_slides]
            slides: List[SlideSpec] = []
            
            for s in raw_slides:
                diagram_info = s.get("diagram") or {}
                chart = None
                if diagram_info.get("chart"):
                    chart_raw = diagram_info.get("chart")
                    if isinstance(chart_raw, dict):
                        # 허용된 필드만 전달
                        allowed = {k: v for k, v in chart_raw.items() if k in ChartData.__fields__}
                        chart = ChartData(**allowed)
                raw_data = diagram_info.get("data")
                # Normalize list -> {'items': list} to avoid validation restrictions later
                if isinstance(raw_data, list):
                    raw_data = {"items": raw_data}
                diagram = DiagramData(type=diagram_info.get("type", "none"), data=raw_data, chart=chart)
                
                # role 정보를 style에 포함
                style_info = s.get("style") or {}
                role = s.get("role")
                if role:
                    style_info["role"] = role
                    # role에 따른 추가 플래그 설정
                    if role == "title":
                        style_info["title"] = True
                    elif role == "agenda":
                        style_info["agenda"] = True
                
                slides.append(SlideSpec(
                    title=s.get("title", ""),
                    key_message=s.get("key_message", ""),
                    bullets=s.get("bullets", []),
                    diagram=diagram,
                    layout=s.get("layout", "title-and-content"),
                    style=style_info if style_info else None,
                    visual_suggestion=s.get("visual_suggestion"),
                    speaker_notes=s.get("speaker_notes"),
                ))
            
            # 토픽 결정 (번호형 섹션 제목 금지)
            parsed_topic = data.get("topic") or fallback_topic or "발표자료"
            if re.match(r"^\d+\.\s+", parsed_topic):  # 번호형이면 폐기
                parsed_topic = fallback_topic or "발표자료"
            # fallback_topic이 더 구체적이고 질의형이 아니면 교체
            if fallback_topic and len(fallback_topic) > len(parsed_topic):
                parsed_topic = fallback_topic

            return DeckSpec(
                topic=parsed_topic,
                slides=slides,
                template_style=data.get("template_style", "business")
            )
        except Exception as e:
            logger.error(f"아웃라인 파싱 실패: {e}")
            # 기본 DeckSpec 반환
            return DeckSpec(
                topic=fallback_topic or "발표자료",
                slides=[],
                template_style="business"
            )

    def _parse_ai_response(self, response: str, topic: str, template_style: str) -> Optional[DeckSpec]:
        """AI 응답 JSON 파싱"""
        try:
            logger.info(f"🔍 JSON 추출 시작: 응답 길이={len(response)}")
            # JSON 추출
            json_match = re.search(r'\{.*\}', response, re.DOTALL)
            if not json_match:
                logger.error(f"❌ JSON 패턴을 찾을 수 없음. 응답 내용: {response[:500]}...")
                return None
            
            json_str = json_match.group()
            logger.info(f"🔍 JSON 문자열 추출 성공: {len(json_str)}자")
            
            data = json.loads(json_str)
            logger.info(f"🔍 JSON 파싱 성공: {list(data.keys())}")
            slides = []
            
            slides_data = data.get('slides', [])
            logger.info(f"🔍 슬라이드 데이터 수: {len(slides_data)}")
            
            for i, slide_data in enumerate(slides_data):
                slide = SlideSpec(
                    title=slide_data.get('title', '제목 없음'),
                    key_message=slide_data.get('key_message', ''),
                    bullets=slide_data.get('bullets', []),
                    layout=slide_data.get('layout', 'title_and_content'),
                    speaker_notes=slide_data.get('speaker_notes', ''),
                    visual_suggestion=slide_data.get('visual_suggestion', '')
                )
                slides.append(slide)
                logger.info(f"  슬라이드 {i+1}: '{slide.title}', bullets={len(slide.bullets)}")
            
            deck = DeckSpec(
                topic=data.get('topic', topic),
                slides=slides,
                max_slides=len(slides),
                template_style=template_style
            )
            logger.info(f"✅ DeckSpec 생성 완료: topic='{deck.topic}', slides={len(deck.slides)}")
            return deck
            
        except Exception as e:
            logger.error(f"AI 응답 파싱 실패: {e}")
            return None

    def _create_fallback_outline(self, topic: str, context_text: str) -> DeckSpec:
        """폴백 아웃라인 생성"""
        slides = [
            SlideSpec(title=topic, key_message="", bullets=[], layout="title-only"),
            SlideSpec(title="목차", key_message="", bullets=["주요 내용 1", "주요 내용 2", "결론"], layout="title-and-content"),
            SlideSpec(title="주요 내용 1", key_message="첫 번째 주요 내용입니다.", bullets=["세부사항 1", "세부사항 2"], layout="title-and-content"),
            SlideSpec(title="주요 내용 2", key_message="두 번째 주요 내용입니다.", bullets=["세부사항 1", "세부사항 2"], layout="title-and-content"),
            SlideSpec(title="결론", key_message="결론입니다.", bullets=["요약", "향후 계획"], layout="title-and-content")
        ]
        return DeckSpec(topic=topic, slides=slides, max_slides=len(slides))

    def build_templated_pptx(self, spec: DeckSpec, template_path: Path, file_basename: Optional[str] = None,
                            text_box_mappings: Optional[List[Dict[str, Any]]] = None,
                            content_segments: Optional[List[Dict[str, Any]]] = None) -> str:
        """템플릿 기반 PPT 빌드"""
        
        logger.info(f"🏗️ 템플릿 기반 PPT 빌드 시작: {len(spec.slides)}개 슬라이드")
        logger.info(f"📄 템플릿 파일: {template_path}")
        
        if not template_path.exists():
            raise FileNotFoundError(f"템플릿을 찾을 수 없습니다: {template_path}")
        
        try:
            # 파일명 생성
            if not file_basename:
                safe_topic = self._normalize_topic_for_filename(spec.topic)
                file_basename = f"templated_presentation_{safe_topic}"
            
            filename = f"{file_basename}.pptx"
            output_path = self.upload_dir / filename
            
            # 템플릿 로드
            prs = Presentation(str(template_path))
            logger.info(f"📋 템플릿 로드 완료: {len(prs.slide_layouts)}개 레이아웃")
            
            # 기존 슬라이드 제거 (템플릿 슬라이드만 유지)
            slide_count = len(prs.slides)
            for i in range(slide_count - 1, -1, -1):
                if i > 0:  # 첫 번째 슬라이드는 유지
                    rId = prs.slides._sldIdLst[i].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[i]
            
            # 새 슬라이드 생성
            for i, slide_spec in enumerate(spec.slides):
                logger.info(f"📄 슬라이드 {i+1} 생성: '{slide_spec.title}'")
                
                if i == 0:
                    # 첫 번째 슬라이드 수정
                    slide = prs.slides[0]
                    self._update_title_slide(slide, slide_spec)
                else:
                    # 새 슬라이드 추가
                    layout_idx = 1 if i == 1 else 1  # 목차와 내용 모두 같은 레이아웃
                    slide_layout = prs.slide_layouts[layout_idx]
                    slide = prs.slides.add_slide(slide_layout)
                    self._populate_template_slide(slide, slide_spec, text_box_mappings)
                
                logger.info(f"✅ 슬라이드 {i+1} 완료")
            
            # 파일 저장
            prs.save(str(output_path))
            logger.info(f"✅ 템플릿 기반 PPT 빌드 완료: {output_path}")
            return str(output_path)
            
        except Exception as e:
            logger.error(f"build_templated_pptx 실패: {e}")
            raise

    def _update_title_slide(self, slide, spec: SlideSpec):
        """제목 슬라이드 업데이트"""
        try:
            if slide.shapes.title:
                slide.shapes.title.text = spec.title
                logger.info(f"✅ 제목 설정: '{spec.title}'")
        except Exception as e:
            logger.error(f"제목 슬라이드 업데이트 실패: {e}")

    def _populate_template_slide(self, slide, spec: SlideSpec, text_box_mappings: Optional[List[Dict[str, Any]]] = None):
        """템플릿 슬라이드에 콘텐츠 채우기"""
        try:
            # 제목 설정
            if slide.shapes.title:
                slide.shapes.title.text = spec.title
            
            # 콘텐츠 영역 찾기 및 채우기
            content_shape = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    content_shape = shape
                    break
            
            if content_shape and hasattr(content_shape, 'text_frame'):
                tf = content_shape.text_frame
                tf.clear()
                tf.word_wrap = True
                
                # 키 메시지 추가
                if spec.key_message:
                    p = tf.paragraphs[0]
                    p.text = spec.key_message
                    p.font.size = Pt(22)
                    p.font.bold = True
                
                # 불릿 포인트 추가
                for bullet in spec.bullets:
                    if bullet.strip():
                        p = tf.add_paragraph()
                        p.text = f"• {bullet.strip()}"
                        p.font.size = Pt(18)
                        p.level = 1
            
            logger.info(f"✅ 템플릿 슬라이드 콘텐츠 완료: '{spec.title}'")
            
        except Exception as e:
            logger.error(f"템플릿 슬라이드 채우기 실패: {e}")

    def build_enhanced_pptx_with_slide_management(self, spec: DeckSpec, file_basename: Optional[str] = None,
                                                 template_style: str = "business", include_charts: bool = True,
                                                 custom_template_path: Optional[str] = None,
                                                 user_template_id: Optional[str] = None,
                                                 text_box_mappings: Optional[List[Dict[str, Any]]] = None,
                                                 content_segments: Optional[List[Dict[str, Any]]] = None,
                                                 slide_management: Optional[List[Dict[str, Any]]] = None) -> str:
        """슬라이드 관리가 포함된 Enhanced PPT 빌드 (enhanced 서비스와 호환)"""
        
        logger.info(f"🏗️ Enhanced PPT 빌드 시작: {len(spec.slides)}개 슬라이드")
        logger.info(f"📋 매핑 정보: text_box_mappings={len(text_box_mappings or [])}, content_segments={len(content_segments or [])}, slide_management={len(slide_management or [])}")
        
        try:
            # 커스텀 템플릿 경로가 있으면 템플릿 기반 빌드 사용
            if custom_template_path and os.path.exists(custom_template_path):
                logger.info(f"📄 커스텀 템플릿 사용: {custom_template_path}")
                
                # 매핑이 있으면 Enhanced Object Processor를 사용해서 매핑 적용
                if text_box_mappings or content_segments or slide_management:
                    logger.info(f"🎯 매핑 기반 템플릿 빌드 실행")
                    return self._build_with_mappings(
                        spec=spec,
                        template_path=Path(custom_template_path),
                        file_basename=file_basename,
                        text_box_mappings=text_box_mappings,
                        content_segments=content_segments,
                        slide_management=slide_management
                    )
                else:
                    # 매핑 없으면 기본 템플릿 빌드
                    logger.info(f"📄 기본 템플릿 빌드 (매핑 없음)")
                    return self.build_templated_pptx(
                        spec=spec,
                        template_path=Path(custom_template_path),
                        file_basename=file_basename,
                        text_box_mappings=text_box_mappings,
                        content_segments=content_segments
                    )
            else:
                # 기본 빌더 사용 (템플릿 없음) - 간단한 PPT 생성
                logger.info(f"📄 기본 빌더 사용 (템플릿 없음)")
                
                # 파일명 생성
                if not file_basename:
                    safe_topic = self._normalize_topic_for_filename(spec.topic)
                    file_basename = f"enhanced_presentation_{safe_topic}"
                
                filename = f"{file_basename}.pptx"
                output_path = self.upload_dir / filename
                
                # 간단한 PPT 생성
                prs = Presentation()
                
                # 각 슬라이드 생성
                for i, slide_spec in enumerate(spec.slides):
                    if i == 0:
                        # 첫 번째 슬라이드 (제목 슬라이드)
                        title_slide_layout = prs.slide_layouts[0]
                        slide = prs.slides.add_slide(title_slide_layout)
                        if slide.shapes.title:
                            slide.shapes.title.text = slide_spec.title
                        # 부제목 placeholder 처리
                        if len(slide.shapes.placeholders) > 1:
                            subtitle_placeholder = slide.shapes.placeholders[1]
                            if getattr(subtitle_placeholder, 'has_text_frame', False):
                                text_frame = getattr(subtitle_placeholder, 'text_frame', None)
                                if text_frame:
                                    text_frame.text = slide_spec.key_message
                    else:
                        # 내용 슬라이드
                        content_slide_layout = prs.slide_layouts[1]
                        slide = prs.slides.add_slide(content_slide_layout)
                        if slide.shapes.title:
                            slide.shapes.title.text = slide_spec.title
                        
                        # 콘텐츠 추가
                        if len(slide.shapes.placeholders) > 1:
                            content_placeholder = slide.shapes.placeholders[1]
                            if getattr(content_placeholder, 'has_text_frame', False):
                                tf = getattr(content_placeholder, 'text_frame', None)
                                if tf:
                                    tf.clear()
                                    if slide_spec.key_message:
                                        p = tf.paragraphs[0]
                                        p.text = slide_spec.key_message
                                    
                                    # 불릿 포인트 추가
                                    for bullet in slide_spec.bullets:
                                        p = tf.add_paragraph()
                                        p.text = bullet
                                        p.level = 1
                
                prs.save(str(output_path))
                logger.info(f"✅ Enhanced PPT 빌드 완료: {output_path}")
                
                return str(output_path)
                
        except Exception as e:
            logger.error(f"Enhanced PPT 빌드 실패: {e}")
            raise

    def _build_with_mappings(self, spec: DeckSpec, template_path: Path, file_basename: Optional[str] = None,
                            text_box_mappings: Optional[List[Dict[str, Any]]] = None,
                            content_segments: Optional[List[Dict[str, Any]]] = None,
                            slide_management: Optional[List[Dict[str, Any]]] = None) -> str:
        """매핑을 적용한 템플릿 기반 PPT 빌드"""
        
        logger.info(f"🎯 매핑 기반 PPT 빌드 시작")
        
        try:
            # 파일명 생성
            if not file_basename:
                safe_topic = self._normalize_topic_for_filename(spec.topic)
                file_basename = f"mapped_presentation_{safe_topic}"
            
            filename = f"{file_basename}.pptx"
            output_path = self.upload_dir / filename
            
            # 템플릿 로드
            prs = Presentation(str(template_path))
            logger.info(f"📋 템플릿 로드 완료: {len(prs.slide_layouts)}개 레이아웃")
            
            # Enhanced Object Processor를 사용해서 매핑 적용
            if hasattr(self, 'object_processor'):
                logger.info(f"🔧 Enhanced Object Processor로 매핑 적용")
                
                # 슬라이드 관리 정보가 있으면 적용
                slides_to_process = spec.slides
                if slide_management:
                    logger.info(f"📋 슬라이드 관리 적용: {len(slide_management)}개 슬라이드")
                    # 슬라이드 순서나 가시성 조정 등을 여기서 처리할 수 있음
                
                # 각 슬라이드에 대해 매핑 적용
                for i, slide_spec in enumerate(slides_to_process):
                    if i < len(prs.slides):
                        slide = prs.slides[i]
                    else:
                        # 새 슬라이드 추가
                        layout_idx = min(1, len(prs.slide_layouts) - 1)
                        slide_layout = prs.slide_layouts[layout_idx]
                        slide = prs.slides.add_slide(slide_layout)
                    
                    # 🔧 AI 생성 콘텐츠를 먼저 적용
                    self._apply_ai_content_to_slide(slide, slide_spec, i)
                    
                    # Enhanced Object Processor로 매핑 적용
                    # 슬라이드별 매핑 준비
                    slide_mappings = []
                    for mapping in (text_box_mappings or []):
                        if mapping.get('slideIndex', 0) == i:
                            slide_mappings.append(mapping)
                    
                    if slide_mappings:
                        self.object_processor.apply_object_mappings(
                            prs, slide_mappings, content_segments
                        )
                        logger.info(f"✅ 슬라이드 {i+1} 매핑 적용 완료: '{slide_spec.title}' ({len(slide_mappings)}개 매핑)")
                    else:
                        logger.info(f"✅ 슬라이드 {i+1} AI 콘텐츠 적용 완료: '{slide_spec.title}' ({len(slide_spec.bullets)}개 bullets)")
            else:
                logger.warning(f"⚠️ Enhanced Object Processor가 없어 기본 빌드로 폴백")
                return self.build_templated_pptx(spec, template_path, file_basename, text_box_mappings, content_segments)
            
            # 파일 저장
            prs.save(str(output_path))
            logger.info(f"✅ 매핑 기반 PPT 빌드 완료: {output_path}")
            return str(output_path)
            
        except Exception as e:
            logger.error(f"매핑 기반 PPT 빌드 실패: {e}")
            raise

    def _apply_ai_content_to_slide(self, slide, slide_spec, slide_index: int):
        """AI가 생성한 콘텐츠를 슬라이드에 적용"""
        try:
            # 제목 적용
            if hasattr(slide, 'shapes') and slide.shapes.title:
                slide.shapes.title.text = slide_spec.title
                logger.debug(f"제목 적용: '{slide_spec.title}'")
            
            # 콘텐츠 적용 - 텍스트 박스나 콘텐츠 placeholder 찾기
            content_applied = False
            
            # bullets가 있으면 bullet points로 적용
            if slide_spec.bullets:
                bullet_text = "\n".join([f"• {bullet}" for bullet in slide_spec.bullets])
                
                # placeholder나 텍스트 박스를 찾아서 콘텐츠 적용
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        # 빈 텍스트 박스이거나 placeholder인 경우
                        if (not shape.text_frame.text.strip() or 
                            hasattr(shape, 'placeholder_format')):
                            shape.text_frame.text = bullet_text
                            content_applied = True
                            logger.debug(f"bullets 적용: {len(slide_spec.bullets)}개")
                            break
            
            # bullets가 없고 key_message가 있으면 적용
            elif hasattr(slide_spec, 'key_message') and slide_spec.key_message:
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        if (not shape.text_frame.text.strip() or 
                            hasattr(shape, 'placeholder_format')):
                            shape.text_frame.text = slide_spec.key_message
                            content_applied = True
                            logger.debug(f"key_message 적용: '{slide_spec.key_message[:50]}...'")
                            break
            
            if not content_applied:
                logger.debug(f"슬라이드 {slide_index + 1}: 콘텐츠 적용할 텍스트 박스를 찾지 못함")
                
        except Exception as e:
            logger.error(f"슬라이드 {slide_index + 1} 콘텐츠 적용 실패: {e}")


# 전역 인스턴스
templated_ppt_service = TemplatedPPTGeneratorService()
