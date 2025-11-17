import logging
from typing import List, Dict, Any, Optional
from pptx.util import Inches, Pt
from enum import Enum

logger = logging.getLogger(__name__)

class PPTObjectType(Enum):
    TEXTBOX = "textbox"
    IMAGE = "image"
    SHAPE = "shape"
    CHART = "chart"
    TABLE = "table"
    DIAGRAM = "diagram"
    ICON = "icon"
    LOGO = "logo"
    BACKGROUND = "background"

class ObjectAction(Enum):
    KEEP_ORIGINAL = "keep_original"     # 원본 유지
    REPLACE_CONTENT = "replace_content" # 내용 교체
    HIDE_OBJECT = "hide_object"         # 오브젝트 제거

class EnhancedPPTObjectProcessor:
    """확장된 PPT 오브젝트 처리 클래스"""
    
    def __init__(self):
        self.logger = logger
    
    def _map_ppt_type_to_object_type(self, ppt_type: str) -> str:
        """PPT 내부 타입을 표준 오브젝트 타입으로 매핑"""
        type_map = {
            'TEXT_BOX': 'textbox',
            'textbox': 'textbox',  # 소문자 버전도 지원
            'AUTO_SHAPE': 'shape',
            'LINE': 'shape',
            'PICTURE': 'image',
            'image': 'image',  # 소문자 버전도 지원
            'TABLE': 'table',
            'table': 'table',  # 소문자 버전도 지원
            'CHART': 'chart',
            'chart': 'chart',  # 소문자 버전도 지원
            'GROUP': 'shape',  # 그룹도 도형으로 분류
        }
        
        mapped_type = type_map.get(ppt_type, 'shape')
        self.logger.debug(f"타입 매핑: {ppt_type} -> {mapped_type}")
        return mapped_type
    
    def apply_object_mappings(self, prs, mappings: List[Dict[str, Any]], 
                             segments: Optional[List[Dict[str, Any]]] = None):
        """모든 타입의 오브젝트 매핑 적용"""
        try:
            # pptx 모듈을 실행 시점에 import
            from pptx import Presentation
            
            self.logger.info(f"오브젝트 매핑 적용 시작: {len(mappings)}개 매핑")
            
            # 활성화된 매핑만 필터링 + 콘텐츠 정규화(assignedContent -> newContent)
            active_mappings = []
            for m in mappings:
                if not m.get('isEnabled', True):
                    continue
                # assignedContent를 사용하는 구버전 입력을 newContent로 정규화
                if 'assignedContent' in m and not m.get('newContent'):
                    m = {**m, 'newContent': m.get('assignedContent')}
                # 빈 문자열 교체는 건너뜀 (원본 보존)
                new_content = m.get('newContent')
                if m.get('action', 'replace_content') == 'replace_content' and isinstance(new_content, str) and new_content.strip() == "":
                    self.logger.info(f"⏭️ 빈 내용 교체 스킵: elementId={m.get('elementId')}")
                    continue
                active_mappings.append(m)
            self.logger.info(f"활성화된 매핑: {len(active_mappings)}개")
            
            # 매핑 정보를 슬라이드별로 그룹화
            mappings_by_slide = {}
            for mapping in active_mappings:
                slide_idx = mapping.get('slideIndex', 0)
                if slide_idx not in mappings_by_slide:
                    mappings_by_slide[slide_idx] = []
                mappings_by_slide[slide_idx].append(mapping)
            
            # 각 슬라이드에 매핑 적용
            for slide_idx, slide_mappings in mappings_by_slide.items():
                if slide_idx < len(prs.slides):
                    slide = prs.slides[slide_idx]
                    self._apply_mappings_to_slide(slide, slide_mappings, segments)
                    self.logger.info(f"슬라이드 {slide_idx}에 {len(slide_mappings)}개 매핑 적용")
            
            return prs
            
        except Exception as e:
            self.logger.error(f"오브젝트 매핑 적용 실패: {e}")
            return prs
    
    def _apply_mappings_to_slide(self, slide, mappings: List[Dict[str, Any]], 
                                segments: Optional[List[Dict[str, Any]]] = None):
        """개별 슬라이드에 매핑 적용"""
        try:
            # 중복 매핑 제거 (동일한 elementId에 대해서는 마지막 매핑만 사용)
            unique_mappings = {}
            used_content = set()  # 이미 사용된 content 추적
            
            for mapping in mappings:
                element_id = mapping.get('elementId')
                new_content = mapping.get('newContent', '')
                
                if element_id:
                    # 동일한 content가 이미 사용되었다면 스킵 (긴급 수정)
                    if new_content and new_content.strip() and new_content.strip() in used_content:
                        self.logger.warning(f"🔄 중복 content 사용 방지: elementId={element_id}, content='{new_content[:30]}...'")
                        continue
                    
                    unique_mappings[element_id] = mapping
                    if new_content and new_content.strip():
                        used_content.add(new_content.strip())
            
            filtered_mappings = list(unique_mappings.values())
            
            if len(mappings) != len(filtered_mappings):
                self.logger.info(f"� 중복 매핑 제거: {len(mappings)}개 → {len(filtered_mappings)}개")
            
            self.logger.info(f"�🔍 슬라이드 매핑 처리 시작: {len(filtered_mappings)}개 매핑")
            
            for i, mapping in enumerate(filtered_mappings):
                element_id = mapping.get('elementId')
                object_type = mapping.get('objectType', PPTObjectType.TEXTBOX.value)
                action = mapping.get('action', ObjectAction.KEEP_ORIGINAL.value)
                is_enabled = mapping.get('isEnabled', True)
                
                self.logger.info(f"🔍 매핑 {i}: elementId={element_id}, objectType={object_type}, action={action}, isEnabled={is_enabled}")
                
                # 비활성화된 매핑은 건너뜀
                if not is_enabled:
                    self.logger.info(f"⏸️ 비활성화된 매핑 건너뜀: {element_id}")
                    continue
                
                # 기존 TextBoxMapping 형식 지원 (하위 호환성)
                if 'assignedContent' in mapping and 'contentSource' in mapping:
                    # action 필드가 있으면 우선 사용
                    if mapping.get('action'):
                        action = mapping.get('action')
                        self.logger.info(f"🔍 action 필드에서 가져온 액션: {action}")
                    # 기존 형식 데이터를 새 형식으로 변환
                    elif mapping.get('contentSource') == 'keep_original':
                        action = ObjectAction.KEEP_ORIGINAL.value
                        self.logger.info(f"🔍 contentSource에서 변환된 액션: {action}")
                    else:
                        action = ObjectAction.REPLACE_CONTENT.value
                        self.logger.info(f"🔍 기본 변환된 액션: {action}")
                    # object_type이 명시되지 않은 경우만 textbox로 설정
                    if not mapping.get('objectType'):
                        object_type = PPTObjectType.TEXTBOX.value
                
                if not element_id:
                    self.logger.warning(f"⚠️ elementId가 없는 매핑 건너뜀")
                    continue
                
                # 타겟 shape 찾기
                target_shape = self._find_shape_by_id(slide, element_id)
                if not target_shape:
                    self.logger.warning(f"Shape not found: {element_id}")
                    continue
                
                self.logger.info(f"🎯 액션 실행: {element_id} -> {action}")
                
                # 액션별 처리
                if action == ObjectAction.HIDE_OBJECT.value:
                    self.logger.info(f"🫥 오브젝트 숨김 처리: {element_id}")
                    self._hide_object(target_shape)
                elif action == ObjectAction.REPLACE_CONTENT.value:
                    self.logger.info(f"🔄 내용 교체 처리: {element_id}")
                    self._replace_content(target_shape, mapping, object_type)
                elif action == ObjectAction.KEEP_ORIGINAL.value:
                    self.logger.info(f"📋 원본 유지 처리: {element_id}")
                    # 원본 유지의 경우 텍스트박스는 기존 로직 사용
                    if object_type == PPTObjectType.TEXTBOX.value:
                        self._apply_textbox_content(target_shape, mapping)
                else:
                    self.logger.warning(f"⚠️ 알 수 없는 액션: {action}")
                
        except Exception as e:
            self.logger.error(f"슬라이드 매핑 적용 실패: {e}")
            import traceback
            self.logger.error(traceback.format_exc())
    
    def _extract_original_id(self, element_id: str) -> str:
        """복사된 오브젝트 ID에서 원본 ID 추출
        예: 표 4_copy_1756367342492_2tojx4o4d_3 -> 표 4
        예: element_0_copy_1756367342492_2tojx4o4d_0 -> element_0
        """
        if '_copy_' not in element_id:
            return element_id
        
        # _copy_ 앞부분이 원본 ID
        original_id = element_id.split('_copy_')[0]
        
        # element_X 형태의 경우 인덱스를 기반으로 원본 매핑
        if original_id.startswith('element_'):
            try:
                index = int(original_id.split('_')[1])
                # element_0 -> 첫 번째 오브젝트, element_1 -> 두 번째 오브젝트 등
                # 실제 슬라이드에서 순서를 따라 매핑
                return f'element_{index}'
            except (ValueError, IndexError):
                pass
        
        return original_id

    def _find_shape_by_id(self, slide, element_id: str):
        """슬라이드에서 element_id로 shape 찾기"""
        try:
            self.logger.info(f"🔍 Shape 찾기 시작: element_id='{element_id}'")
            
            # 복사된 오브젝트 ID 처리 (예: 표 4_copy_1756367342492_2tojx4o4d_3 -> 표 4)
            original_element_id = self._extract_original_id(element_id)
            if original_element_id != element_id:
                self.logger.info(f"🔄 복사된 ID 감지: {element_id} -> {original_element_id}")
                # 원본 ID로 먼저 시도
                result = self._find_shape_by_id(slide, original_element_id)
                if result:
                    return result
                # 원본 ID로 찾지 못하면 복사된 ID로 계속 진행
            
            # 다양한 방식으로 shape 찾기
            for i, shape in enumerate(slide.shapes):
                shape_name = getattr(shape, 'name', '(no name)')
                shape_type = self._get_shape_type(shape)
                
                # 텍스트가 있는 경우 내용도 표시
                text_content = ""
                try:
                    if hasattr(shape, 'text_frame') and shape.text_frame and shape.text_frame.text:
                        text_content = shape.text_frame.text.strip()[:30] + "..."
                except:
                    pass
                
                self.logger.info(f"  - Shape {i}: name='{shape_name}', type='{shape_type}', text='{text_content}'")
            
            # 1차: name 속성으로 매칭 (정확한 매칭 우선)
            for i, shape in enumerate(slide.shapes):
                if hasattr(shape, 'name') and shape.name == element_id:
                    self.logger.info(f"✅ Shape 매칭 성공 (name): {element_id} -> Shape {i}")
                    return shape
            
            # 표 오브젝트인 경우 특별 처리
            if element_id == '표 4' or 'table' in element_id.lower() or '표' in element_id:
                table_count = 0
                for i, shape in enumerate(slide.shapes):
                    if hasattr(shape, 'table') and shape.table is not None:
                        table_count += 1
                        # 표 4 -> 4번째 표 (1부터 시작)
                        if '표 4' in element_id and table_count == 4:
                            self.logger.info(f"✅ 표 매칭 성공 (순서): {element_id} -> Shape {i} (4번째 표)")
                            return shape
                        elif table_count == 1:  # 첫 번째 표라면 매칭
                            self.logger.info(f"✅ 표 매칭 성공 (첫번째): {element_id} -> Shape {i}")
                            return shape
            
            # 2차: textbox 타입에 한해 순차적 인덱스 매칭 (textbox-0-2 -> 3번째 textbox)
            if element_id.startswith('textbox-'):
                try:
                    # textbox-0-2에서 마지막 숫자 추출
                    parts = element_id.split('-')
                    if len(parts) >= 3:
                        target_index = int(parts[-1])
                        textbox_count = 0
                        
                        for i, shape in enumerate(slide.shapes):
                            shape_type = self._get_shape_type(shape)
                            if shape_type == 'textbox':
                                if textbox_count == target_index:
                                    self.logger.info(f"✅ Shape 매칭 성공 (textbox-sequential): {element_id} -> Shape {i} ({textbox_count}번째 텍스트박스)")
                                    return shape
                                textbox_count += 1
                except (ValueError, IndexError):
                    pass
            
            # 3차: 기존 인덱스 기반 매칭 (fallback)
            for i, shape in enumerate(slide.shapes):
                if element_id.endswith(f'-{i}'):
                    shape_type = self._get_shape_type(shape)
                    # textbox를 찾는데 shape가 매칭되는 경우 건너뜀
                    if element_id.startswith('textbox-') and shape_type != 'textbox':
                        self.logger.info(f"⚠️ 타입 불일치로 인덱스 매칭 건너뜀: {element_id} (type={shape_type})")
                        continue
                    self.logger.info(f"✅ Shape 매칭 성공 (index): {element_id} -> Shape {i}")
                    return shape
            
            # 4차: element_N 형태 매칭 (순서 기반) - 개선된 로직
            if element_id.startswith('element_'):
                try:
                    target_index = int(element_id.split('_')[1])
                    
                    # 텍스트가 있는 shape들만 필터링하여 순서 매칭 (🔴 긴급 수정)
                    text_shapes = []
                    for i, shape in enumerate(slide.shapes):
                        if hasattr(shape, 'text_frame') and shape.text_frame:
                            try:
                                text = shape.text_frame.text.strip()
                                if text:  # 의미있는 텍스트가 있는 shape만
                                    text_shapes.append((i, shape))
                            except:
                                pass
                    
                    if target_index < len(text_shapes):
                        original_index, target_shape = text_shapes[target_index]
                        self.logger.info(f"✅ Shape 매칭 성공 (element-text-sequential): {element_id} -> Shape {original_index} (텍스트가 있는 {target_index}번째)")
                        return target_shape
                    
                    # 폴백: 전체 shape 순서
                    elif target_index < len(slide.shapes):
                        target_shape = slide.shapes[target_index]
                        self.logger.info(f"✅ Shape 매칭 성공 (element-sequential-fallback): {element_id} -> Shape {target_index}")
                        return target_shape
                except (ValueError, IndexError):
                    pass
            
            # 5차: 타입-인덱스 매칭
            for i, shape in enumerate(slide.shapes):
                shape_type = self._get_shape_type(shape)
                if element_id == f"{shape_type}-{i}":
                    self.logger.info(f"✅ Shape 매칭 성공 (type-index): {element_id} -> Shape {i}")
                    return shape
            
            # 매칭 실패 시 사용 가능한 모든 element ID 출력
            available_ids = []
            for i, shape in enumerate(slide.shapes):
                shape_type = self._get_shape_type(shape)
                available_ids.append(f"{shape_type}-{i}")
                if hasattr(shape, 'name') and shape.name:
                    available_ids.append(shape.name)
            
            self.logger.warning(f"❌ Shape 매칭 실패: element_id='{element_id}' 를 찾을 수 없음")
            self.logger.warning(f"🔍 사용 가능한 element ID들: {available_ids[:10]}{'...' if len(available_ids) > 10 else ''}")
            return None
            
        except Exception as e:
            self.logger.error(f"Shape 찾기 실패: {e}")
            return None
    
    def _get_shape_type(self, shape) -> str:
        """Shape의 타입 결정"""
        try:
            # python-pptx의 실제 shape_type 사용
            if hasattr(shape, 'shape_type'):
                shape_type_name = getattr(shape.shape_type, 'name', str(shape.shape_type))
                return self._map_ppt_type_to_object_type(shape_type_name)
            
            # Fallback: 기존 방식
            if hasattr(shape, 'text_frame') and shape.text_frame:
                return PPTObjectType.TEXTBOX.value
            elif hasattr(shape, 'image'):
                return PPTObjectType.IMAGE.value
            elif hasattr(shape, 'chart'):
                return PPTObjectType.CHART.value
            elif hasattr(shape, 'table'):
                return PPTObjectType.TABLE.value
            else:
                return PPTObjectType.SHAPE.value
                
        except Exception as e:
            self.logger.warning(f"Shape 타입 결정 오류: {e}, 기본값 shape 사용")
            return PPTObjectType.SHAPE.value
    
    def _capture_text_style(self, text_frame):
        """텍스트 프레임의 스타일 정보를 수집"""
        style_info = {
            'paragraphs': []
        }
        
        try:
            for para in text_frame.paragraphs:
                para_info = {
                    'alignment': para.alignment,
                    'level': para.level,
                    'runs': []
                }
                
                for run in para.runs:
                    run_info = {
                        'font_name': run.font.name,
                        'font_size': run.font.size,
                        'bold': run.font.bold,
                        'italic': run.font.italic,
                        'underline': run.font.underline,
                        'color_info': None
                    }
                    
                    # 색상 정보 상세 수집
                    try:
                        if run.font.color:
                            color_info = {
                                'type': None,
                                'rgb': None,
                                'theme_color': None,
                                'brightness': None
                            }
                            
                            # 색상 타입 확인
                            if hasattr(run.font.color, 'type'):
                                color_info['type'] = run.font.color.type
                                self.logger.debug(f"색상 타입: {run.font.color.type}")
                            
                            # RGB 색상 정보
                            if hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                                color_info['rgb'] = run.font.color.rgb
                                self.logger.debug(f"RGB 색상: {run.font.color.rgb}")
                            
                            # 테마 색상 정보
                            if hasattr(run.font.color, 'theme_color'):
                                color_info['theme_color'] = run.font.color.theme_color
                                self.logger.debug(f"테마 색상: {run.font.color.theme_color}")
                            
                            # 밝기 정보
                            if hasattr(run.font.color, 'brightness'):
                                color_info['brightness'] = run.font.color.brightness
                                self.logger.debug(f"밝기: {run.font.color.brightness}")
                            
                            run_info['color_info'] = color_info
                            
                    except Exception as e:
                        self.logger.debug(f"색상 정보 수집 실패: {e}")
                    
                    para_info['runs'].append(run_info)
                
                style_info['paragraphs'].append(para_info)
        
        except Exception as e:
            self.logger.warning(f"스타일 캡처 실패: {e}")
        
        return style_info
    
    def _replace_text_preserving_style(self, text_frame, new_content, original_style):
        """스타일을 보존하면서 텍스트 내용만 교체"""
        try:
            # 새 내용으로 텍스트 설정
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = new_content
            
            # 원본 스타일이 있으면 적용
            if original_style and original_style.get('paragraphs'):
                first_para_style = original_style['paragraphs'][0]
                
                # 첫 번째 paragraph 스타일 적용
                if hasattr(p, 'alignment'):
                    p.alignment = first_para_style.get('alignment')
                if hasattr(p, 'level'):
                    p.level = first_para_style.get('level', 0)
                
                # Run 스타일 적용
                if first_para_style.get('runs') and p.runs:
                    first_run = p.runs[0]
                    first_run_style = first_para_style['runs'][0]
                    
                    # 폰트 정보 적용
                    if first_run_style.get('font_name'):
                        first_run.font.name = first_run_style['font_name']
                    if first_run_style.get('font_size'):
                        first_run.font.size = first_run_style['font_size']
                    if first_run_style.get('bold') is not None:
                        first_run.font.bold = first_run_style['bold']
                    if first_run_style.get('italic') is not None:
                        first_run.font.italic = first_run_style['italic']
                    if first_run_style.get('underline') is not None:
                        first_run.font.underline = first_run_style['underline']
                    
                    # 색상 적용 (개선된 로직)
                    if first_run_style.get('color_info'):
                        try:
                            color_info = first_run_style['color_info']
                            self.logger.debug(f"색상 적용 시도: {color_info}")
                            
                            # 색상 타입에 따른 적용
                            if color_info.get('type') is not None:
                                from pptx.enum.dml import MSO_COLOR_TYPE
                                
                                if color_info['type'] == MSO_COLOR_TYPE.RGB and color_info.get('rgb'):
                                    # RGB 색상 적용
                                    first_run.font.color.rgb = color_info['rgb']
                                    self.logger.debug(f"RGB 색상 적용 완료: {color_info['rgb']}")
                                
                                elif color_info['type'] == MSO_COLOR_TYPE.SCHEME and color_info.get('theme_color'):
                                    # 테마 색상 적용
                                    first_run.font.color.theme_color = color_info['theme_color']
                                    if color_info.get('brightness') is not None:
                                        first_run.font.color.brightness = color_info['brightness']
                                    self.logger.debug(f"테마 색상 적용 완료: {color_info['theme_color']}")
                                
                                else:
                                    # RGB가 있으면 RGB로 폴백
                                    if color_info.get('rgb'):
                                        first_run.font.color.rgb = color_info['rgb']
                                        self.logger.debug(f"RGB 폴백 적용: {color_info['rgb']}")
                            
                            else:
                                # 타입 정보가 없으면 RGB로 시도
                                if color_info.get('rgb'):
                                    first_run.font.color.rgb = color_info['rgb']
                                    self.logger.debug(f"RGB 직접 적용: {color_info['rgb']}")
                                    
                        except Exception as e:
                            self.logger.warning(f"색상 적용 실패: {e}")
                            # 색상 적용에 실패해도 다른 스타일은 유지
        
        except Exception as e:
            self.logger.warning(f"스타일 적용 실패: {e}, 텍스트만 교체됨")
    
    def _hide_object(self, shape):
        """오브젝트 숨기기"""
        try:
            # 숨기기 전 정보 로깅
            shape_name = getattr(shape, 'name', '(no name)')
            text_content = ""
            try:
                if hasattr(shape, 'text_frame') and shape.text_frame and shape.text_frame.text:
                    text_content = shape.text_frame.text.strip()[:30] + "..."
            except:
                pass
                
            self.logger.info(f"🫥 숨김 대상: name='{shape_name}', text='{text_content}'")
            
            # 방법 1: 슬라이드 밖으로 이동하지 않고 현 위치에서 최소 크기로 축소
            # - 음수 좌표로 이동하면 뷰포트가 왼쪽으로 확장되어 슬라이드가 화면에서 오른쪽으로 치우쳐 보일 수 있음
            # - 따라서 위치는 유지하고 크기만 최소화하여 시각적으로 제거
            shape.width = Inches(0.01)
            shape.height = Inches(0.01)
            
            # 방법 3: 가능한 경우 투명도 설정
            try:
                if hasattr(shape, 'fill'):
                    shape.fill.solid()
                    # 흰색으로 설정하고 투명도 최대로
                    from pptx.dml.color import RGBColor
                    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    shape.fill.transparency = 1.0
            except:
                pass  # fill 설정에 실패해도 위치/크기 조정으로 숨김 효과
            
            # 방법 4: 텍스트가 있는 경우 빈 텍스트로 설정
            try:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    shape.text_frame.clear()
            except:
                pass
            
            self.logger.info(f"오브젝트 숨김 처리 완료: 최소 크기 및 투명화 (좌표 고정)")
            
        except Exception as e:
            self.logger.error(f"오브젝트 숨김 실패: {e}")
            import traceback
            self.logger.error(traceback.format_exc())
    
    def _replace_content(self, shape, mapping: Dict[str, Any], object_type: str):
        """오브젝트 내용 교체"""
        try:
            if object_type == PPTObjectType.TEXTBOX.value:
                self._replace_textbox_content(shape, mapping)
            elif object_type == PPTObjectType.SHAPE.value:
                self._replace_shape_content(shape, mapping)
            elif object_type == PPTObjectType.IMAGE.value:
                self._replace_image_content(shape, mapping)
            elif object_type == PPTObjectType.CHART.value:
                self._replace_chart_content(shape, mapping)
            elif object_type == PPTObjectType.TABLE.value:
                self._replace_table_content(shape, mapping)
                
        except Exception as e:
            self.logger.error(f"내용 교체 실패: {e}")
    
    def _replace_textbox_content(self, shape, mapping: Dict[str, Any]):
        """텍스트박스 내용 교체 (스타일 보존)"""
        try:
            if not hasattr(shape, 'text_frame') or not shape.text_frame:
                self.logger.warning("텍스트박스에 text_frame이 없음")
                return
            
            new_content = mapping.get('newContent', '')
            old_content = shape.text_frame.text if shape.text_frame.text else ""
            
            self.logger.info(f"📝 텍스트박스 내용 교체: '{old_content}' -> '{new_content}'")
            
            if new_content:
                # 기존 스타일 정보 수집
                original_style = self._capture_text_style(shape.text_frame)
                
                # 텍스트만 교체 (스타일 유지)
                self._replace_text_preserving_style(shape.text_frame, new_content, original_style)
                
                self.logger.info(f"✅ 텍스트박스 내용 교체 완료 (스타일 보존): '{new_content}'")
            else:
                self.logger.warning(f"⚠️ 새 내용이 비어있음: newContent='{new_content}'")
                
        except Exception as e:
            self.logger.error(f"텍스트박스 내용 교체 실패: {e}")
            import traceback
            self.logger.error(traceback.format_exc())
    
    def _replace_shape_content(self, shape, mapping: Dict[str, Any]):
        """도형 내용 교체 (텍스트가 있는 도형, 스타일 보존)"""
        try:
            if not hasattr(shape, 'text_frame') or not shape.text_frame:
                self.logger.warning(f"도형에 text_frame이 없음: name='{shape.name if hasattr(shape, 'name') else 'Unknown'}'")
                return
            
            new_content = mapping.get('newContent', '')
            old_content = shape.text_frame.text if shape.text_frame.text else ""
            
            self.logger.info(f"🎨 도형 내용 교체: '{old_content}' -> '{new_content}'")
            
            if new_content:
                # 기존 스타일 정보 수집
                original_style = self._capture_text_style(shape.text_frame)
                
                # 텍스트만 교체 (스타일 유지)
                self._replace_text_preserving_style(shape.text_frame, new_content, original_style)
                
                self.logger.info(f"✅ 도형 내용 교체 완료 (스타일 보존): '{new_content}'")
            else:
                self.logger.warning(f"⚠️ 새 내용이 비어있음: newContent='{new_content}'")
                
        except Exception as e:
            self.logger.error(f"도형 내용 교체 실패: {e}")
            import traceback
            self.logger.error(traceback.format_exc())
    
    def _replace_image_content(self, shape, mapping: Dict[str, Any]):
        """이미지 내용 교체"""
        try:
            new_image_url = mapping.get('newImageUrl', '')
            if not new_image_url:
                return
            
            # 이미지 다운로드 및 교체 로직
            # 실제 구현에서는 URL에서 이미지를 다운로드하여 교체
            self.logger.info(f"이미지 교체 요청: {new_image_url}")
            # TODO: 실제 이미지 교체 구현
            
        except Exception as e:
            self.logger.error(f"이미지 교체 실패: {e}")
    
    def _replace_chart_content(self, shape, mapping: Dict[str, Any]):
        """차트 내용 교체"""
        try:
            # 차트 데이터 교체 로직
            self.logger.info("차트 내용 교체 요청")
            # TODO: 차트 데이터 교체 구현
            
        except Exception as e:
            self.logger.error(f"차트 교체 실패: {e}")
    
    def _replace_table_content(self, shape, mapping: Dict[str, Any]):
        """테이블 내용 교체 (열/행 동적 추가 지원, 폰트 색상 유지)"""
        try:
            from pptx.util import Emu
            from pptx.dml.color import RGBColor

            # 입력 데이터 확보
            metadata = mapping.get('metadata') or {}
            table_data = metadata.get('tableData') or {}
            headers = table_data.get('headers') or []
            rows = table_data.get('rows') or []

            # 텍스트 형태(newContent)로 전달된 경우 간단 파싱은 생략하고 로그만 남김
            if not headers and not rows:
                self.logger.warning("테이블 교체 데이터가 비어있음 (headers/rows 없음)")
                return

            if not hasattr(shape, 'table') or not shape.table:
                self.logger.warning("타겟 shape에 table 속성이 없음")
                return

            tbl = shape.table
            original_rows = len(tbl.rows)
            original_cols = len(tbl.columns)

            # 헤더가 없으면 rows의 첫 행을 헤더로 간주
            if not headers and rows:
                headers = rows[0]
                rows = rows[1:]

            # 전체 데이터를 [헤더] + [행들] 형태로 구성
            data_matrix = []
            if headers:
                data_matrix.append([str(x) if x is not None else '' for x in headers])
            for r in rows:
                data_matrix.append([str(x) if x is not None else '' for x in r])

            # 필요한 크기 계산
            needed_rows = len(data_matrix)
            needed_cols = max(len(r) for r in data_matrix) if data_matrix else 0

            self.logger.info(f"📏 테이블 크기 분석: 기존({original_rows}x{original_cols}) → 필요({needed_rows}x{needed_cols})")

            # 원본 폭/높이/열/행 스타일 백업 (재생성 폴백 대비)
            orig_left, orig_top, orig_width, orig_height = shape.left, shape.top, shape.width, shape.height
            orig_name = getattr(shape, 'name', None)
            orig_col_widths = []
            try:
                for c in range(original_cols):
                    orig_col_widths.append(tbl.columns[c].width)
            except Exception:
                pass
            orig_row_heights = []
            try:
                for r in range(original_rows):
                    orig_row_heights.append(tbl.rows[r].height)
            except Exception:
                pass

            # 🆕 열 추가 (필요시)
            if needed_cols > original_cols:
                cols_to_add = needed_cols - original_cols
                self.logger.info(f"➕ {cols_to_add}개 열 추가 중...")
                
                try:
                    # python-pptx의 테이블은 열을 동적으로 추가하는 공식 API가 제한적
                    # 대신 테이블을 재생성하는 방식으로 처리
                    self.logger.info(f"테이블 크기 변경 필요 (열 추가): {original_cols} → {needed_cols}")
                    
                    # 열 추가가 필요한 경우 테이블 재생성으로 처리
                    self._recreate_table_with_data(
                        shape=shape,
                        data_matrix=data_matrix,
                        headers=headers,
                        name=orig_name,
                        left=orig_left,
                        top=orig_top,
                        width=orig_width,
                        height=orig_height,
                        orig_col_widths=orig_col_widths,
                        orig_row_heights=orig_row_heights
                    )
                    self.logger.info("✅ 열 추가를 위한 테이블 재생성 완료")
                    return
                    
                except Exception as e:
                    self.logger.error(f"❌ 열 추가 실패: {e}")
                    # 폴백: 테이블 재생성으로 시도
                    self._recreate_table_with_data(
                        shape=shape,
                        data_matrix=data_matrix,
                        headers=headers,
                        name=orig_name,
                        left=orig_left,
                        top=orig_top,
                        width=orig_width,
                        height=orig_height,
                        orig_col_widths=orig_col_widths,
                        orig_row_heights=orig_row_heights
                    )
                    self.logger.info("✅ 재생성 폴백으로 테이블 교체 완료")
                    return

            # 🆕 행 추가 (필요시)
            if needed_rows > original_rows:
                rows_to_add = needed_rows - original_rows
                self.logger.info(f"➕ {rows_to_add}개 행 추가 중...")
                
                try:
                    # 기존 첫 번째 행의 높이 정보 수집
                    first_row_height = tbl.rows[0].height if tbl.rows else None
                    
                    for _ in range(rows_to_add):
                        # 올바른 API 사용: add_row (높이는 별도로 설정)
                        new_row = tbl.rows.add_row()
                        
                        # 새로 추가된 행의 높이 설정
                        if first_row_height:
                            new_row.height = first_row_height
                        
                        # 새로 추가된 행의 각 셀에 첫 번째 행의 스타일 복사
                        if len(tbl.rows) > 1:  # 기존 행이 있는 경우
                            first_row = tbl.rows[0]
                            for col_idx in range(len(new_row.cells)):
                                if col_idx < len(first_row.cells):
                                    # 기존 첫 번째 행의 셀 스타일을 새 행에 복사
                                    try:
                                        source_cell = first_row.cells[col_idx]
                                        target_cell = new_row.cells[col_idx]
                                        
                                        # 텍스트 프레임이 있는 경우 폰트 스타일 복사
                                        if (source_cell.text_frame and source_cell.text_frame.paragraphs and
                                            target_cell.text_frame and target_cell.text_frame.paragraphs):
                                            
                                            source_para = source_cell.text_frame.paragraphs[0]
                                            target_para = target_cell.text_frame.paragraphs[0]
                                            
                                            if source_para.runs and target_para.runs:
                                                source_run = source_para.runs[0]
                                                target_run = target_para.runs[0]
                                                
                                                # 폰트 속성 복사
                                                if source_run.font.name:
                                                    target_run.font.name = source_run.font.name
                                                if source_run.font.size:
                                                    target_run.font.size = source_run.font.size
                                                target_run.font.bold = source_run.font.bold
                                                target_run.font.italic = source_run.font.italic
                                                
                                                # 색상 복사
                                                try:
                                                    if source_run.font.color.rgb:
                                                        target_run.font.color.rgb = source_run.font.color.rgb
                                                    elif source_run.font.color.theme_color:
                                                        target_run.font.color.theme_color = source_run.font.color.theme_color
                                                        if hasattr(source_run.font.color, 'brightness'):
                                                            target_run.font.color.brightness = source_run.font.color.brightness
                                                except:
                                                    pass  # 색상 복사 실패 시 무시
                                                    
                                    except Exception as cell_style_error:
                                        self.logger.warning(f"새 행 셀 스타일 복사 실패: {cell_style_error}")
                    
                    self.logger.info(f"✅ 행 추가 완료: {original_rows} → {len(tbl.rows)}")
                except Exception as e:
                    self.logger.error(f"❌ 행 추가 실패: {e}")
                    # 폴백: 테이블 재생성으로 시도
                    self._recreate_table_with_data(
                        shape=shape,
                        data_matrix=data_matrix,
                        headers=headers,
                        name=orig_name,
                        left=orig_left,
                        top=orig_top,
                        width=orig_width,
                        height=orig_height,
                        orig_col_widths=orig_col_widths,
                        orig_row_heights=orig_row_heights
                    )
                    self.logger.info("✅ 재생성 폴백으로 테이블 교체 완료")
                    return

            # 최종 크기 설정
            final_rows = len(tbl.rows)
            final_cols = len(tbl.columns)
            
            rows_to_fill = min(final_rows, len(data_matrix))
            cols_to_fill = min(final_cols, needed_cols)

            self.logger.info(f"📝 테이블 데이터 채우기: {rows_to_fill}x{cols_to_fill}")

            # 🎨 원본 셀 스타일 정보 수집 (색상 유지를 위해)
            original_styles = {}
            try:
                for r in range(original_rows):
                    for c in range(original_cols):
                        cell = tbl.cell(r, c)
                        if cell.text_frame and cell.text_frame.paragraphs:
                            para = cell.text_frame.paragraphs[0]
                            if para.runs:
                                run = para.runs[0]
                                if run.font:
                                    # 완전한 색상 정보 수집
                                    color_info = {
                                        'type': None,
                                        'rgb': None,
                                        'theme_color': None,
                                        'brightness': None
                                    }
                                    
                                    try:
                                        if run.font.color:
                                            # 색상 타입 확인
                                            if hasattr(run.font.color, 'type'):
                                                color_info['type'] = run.font.color.type
                                            
                                            # RGB 색상 정보
                                            if hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                                                color_info['rgb'] = run.font.color.rgb
                                            
                                            # 테마 색상 정보
                                            if hasattr(run.font.color, 'theme_color') and run.font.color.theme_color is not None:
                                                color_info['theme_color'] = run.font.color.theme_color
                                            
                                            # 밝기 정보
                                            if hasattr(run.font.color, 'brightness') and run.font.color.brightness is not None:
                                                color_info['brightness'] = run.font.color.brightness
                                    except Exception:
                                        pass
                                        
                                    original_styles[(r, c)] = {
                                        'color_info': color_info,
                                        'font_name': run.font.name,
                                        'font_size': run.font.size,
                                        'bold': run.font.bold,
                                        'italic': run.font.italic,
                                        'underline': run.font.underline
                                    }
                                    self.logger.debug(f"셀({r},{c}) 스타일 수집: font={run.font.name}, size={run.font.size}, color_type={color_info['type']}")
            except Exception as e:
                self.logger.warning(f"스타일 수집 실패: {e}")

            # 📝 데이터 채우기 및 스타일 적용
            for r in range(rows_to_fill):
                row_vals = data_matrix[r]
                for c in range(cols_to_fill):
                    val = row_vals[c] if c < len(row_vals) else ''
                    cell = tbl.cell(r, c)
                    
                    # 텍스트 설정
                    try:
                        # 기존 텍스트 삭제하고 새 텍스트 설정
                        if cell.text_frame and cell.text_frame.paragraphs:
                            para = cell.text_frame.paragraphs[0]
                            para.clear()
                            para.text = val
                            
                            # 🎨 원본 스타일 복원 (우선순위 1: 해당 위치 스타일)
                            if (r, c) in original_styles and para.runs:
                                run = para.runs[0]
                                style = original_styles[(r, c)]
                                
                                try:
                                    # 색상 정보 적용
                                    color_info = style.get('color_info', {})
                                    if color_info and color_info.get('type'):
                                        from pptx.enum.dml import MSO_COLOR_TYPE
                                        
                                        if color_info['type'] == MSO_COLOR_TYPE.SCHEME:
                                            # 테마 색상 적용
                                            if color_info.get('theme_color') is not None:
                                                run.font.color.theme_color = color_info['theme_color']
                                                if color_info.get('brightness') is not None:
                                                    run.font.color.brightness = color_info['brightness']
                                                self.logger.debug(f"셀({r},{c}) 테마 색상 적용: {color_info['theme_color']}")
                                        
                                        elif color_info['type'] == MSO_COLOR_TYPE.RGB:
                                            # RGB 색상 적용
                                            if color_info.get('rgb'):
                                                run.font.color.rgb = color_info['rgb']
                                                self.logger.debug(f"셀({r},{c}) RGB 색상 적용: {color_info['rgb']}")
                                    
                                    elif color_info.get('rgb'):
                                        # 타입 정보가 없으면 RGB로 폴백
                                        run.font.color.rgb = color_info['rgb']
                                        self.logger.debug(f"셀({r},{c}) RGB 폴백 적용: {color_info['rgb']}")
                                    
                                    # 폰트 속성 적용
                                    if style.get('font_name'):
                                        run.font.name = style['font_name']
                                    if style.get('font_size'):
                                        run.font.size = style['font_size']
                                    if style.get('bold') is not None:
                                        run.font.bold = style['bold']
                                    if style.get('italic') is not None:
                                        run.font.italic = style['italic']
                                    if style.get('underline') is not None:
                                        run.font.underline = style['underline']
                                        
                                    self.logger.debug(f"셀({r},{c}) 완전 스타일 적용 완료")
                                except Exception as se:
                                    self.logger.warning(f"셀({r},{c}) 스타일 적용 실패: {se}")
                                    
                            # 새로 추가된 셀의 경우 헤더 또는 첫 번째 셀 스타일 적용 (우선순위 2)
                            elif para.runs and original_styles:
                                run = para.runs[0]
                                
                                # 스타일 참조 우선순위: 같은 열 > 같은 행 > 첫 번째 셀
                                style = None
                                if r < original_rows and (r, 0) in original_styles:
                                    # 같은 행의 첫 번째 셀 스타일
                                    style = original_styles[(r, 0)]
                                elif c < original_cols and (0, c) in original_styles:
                                    # 같은 열의 첫 번째 셀 스타일
                                    style = original_styles[(0, c)]
                                elif (0, 0) in original_styles:
                                    # 폴백: 첫 번째 셀 스타일
                                    style = original_styles[(0, 0)]
                                
                                if style:
                                    try:
                                        # 색상 정보 적용
                                        color_info = style.get('color_info', {})
                                        if color_info and color_info.get('type'):
                                            from pptx.enum.dml import MSO_COLOR_TYPE
                                            
                                            if color_info['type'] == MSO_COLOR_TYPE.SCHEME:
                                                # 테마 색상 적용
                                                if color_info.get('theme_color') is not None:
                                                    run.font.color.theme_color = color_info['theme_color']
                                                    if color_info.get('brightness') is not None:
                                                        run.font.color.brightness = color_info['brightness']
                                            
                                            elif color_info['type'] == MSO_COLOR_TYPE.RGB:
                                                # RGB 색상 적용
                                                if color_info.get('rgb'):
                                                    run.font.color.rgb = color_info['rgb']
                                        
                                        elif color_info.get('rgb'):
                                            # 타입 정보가 없으면 RGB로 폴백
                                            run.font.color.rgb = color_info['rgb']
                                        
                                        # 폰트 속성 적용
                                        if style.get('font_name'):
                                            run.font.name = style['font_name']
                                        if style.get('font_size'):
                                            run.font.size = style['font_size']
                                        if style.get('bold') is not None:
                                            run.font.bold = style['bold']
                                        if style.get('italic') is not None:
                                            run.font.italic = style['italic']
                                        if style.get('underline') is not None:
                                            run.font.underline = style['underline']
                                            
                                        self.logger.debug(f"셀({r},{c}) 폴백 스타일 적용 완료")
                                    except Exception as se:
                                        self.logger.warning(f"셀({r},{c}) 폴백 스타일 적용 실패: {se}")
                        else:
                            cell.text = val
                            
                    except Exception as e:
                        self.logger.warning(f"셀 ({r},{c}) 설정 실패: {e}")
                        # 폴백: 기본 텍스트 설정
                        cell.text = val

            # 🎨 헤더 스타일링 강화 (첫 번째 행을 헤더로 간주)
            try:
                if headers and rows_to_fill > 0:
                    for c in range(min(cols_to_fill, len(headers))):
                        cell = tbl.cell(0, c)
                        if cell.text_frame and cell.text_frame.paragraphs:
                            para = cell.text_frame.paragraphs[0]
                            if para.runs:
                                run = para.runs[0]
                                
                                # 원본 헤더 스타일이 있으면 그것을 유지하되, 굵게 처리는 추가
                                if (0, c) in original_styles:
                                    # 원본 헤더 스타일 유지
                                    pass  # 이미 위에서 적용됨
                                
                                # 헤더는 항상 굵게 (원본이 굵지 않았어도)
                                if run.font.bold is not True:
                                    run.font.bold = True
                                
                                self.logger.debug(f"헤더 셀({0},{c}) 굵게 처리 완료")
            except Exception as e:
                self.logger.warning(f"헤더 스타일링 실패: {e}")

            self.logger.info("✅ 테이블 내용 교체 완료 (스타일 보존)")
            self.logger.info(f"📊 적용된 스타일 정보: {len(original_styles)}개 셀")
            
        except Exception as e:
            self.logger.error(f"테이블 교체 실패: {e}")

    def _recreate_table_with_data(self, shape, data_matrix: List[List[str]], headers: List[str], name: Optional[str], left, top, width, height, orig_col_widths: List[Any], orig_row_heights: List[Any]):
        """열/행 추가가 불가능한 환경에서 테이블을 재생성하여 데이터를 채움"""
        try:
            slide = shape.part.slide

            # 원본 스타일 일부 백업 (폰트)
            base_style = None
            try:
                if hasattr(shape, 'table') and shape.table and shape.table.cell(0, 0).text_frame.paragraphs and shape.table.cell(0, 0).text_frame.paragraphs[0].runs:
                    run = shape.table.cell(0, 0).text_frame.paragraphs[0].runs[0]
                    base_style = {
                        'font_color': getattr(run.font.color, 'rgb', None),
                        'font_name': run.font.name,
                        'font_size': run.font.size,
                        'bold': run.font.bold,
                    }
            except Exception:
                pass

            # 새 테이블 추가 (그래픽프레임)
            rows = len(data_matrix)
            cols = max((len(r) for r in data_matrix), default=0)
            if rows == 0 or cols == 0:
                self.logger.warning("재생성용 데이터가 비어있어 중단")
                return

            new_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            new_tbl = new_shape.table

            # 열/행 크기: 원본 총 width/height 안으로 자동 맞춤
            try:
                # 1) 우선 원본 비율을 참고하되, 합계가 원본 width/height가 되도록 균등 분배
                if width and cols > 0:
                    # 균등 분배하여 레이아웃 영역 내로 강제 맞춤
                    even_col = int(width // cols)
                    for c in range(cols):
                        new_tbl.columns[c].width = even_col
                if height and rows > 0:
                    even_row = int(height // rows)
                    for r in range(rows):
                        new_tbl.rows[r].height = even_row
            except Exception:
                pass

            # 데이터 채우기 + 기본 스타일 적용
            for r in range(rows):
                row_vals = data_matrix[r]
                for c in range(cols):
                    val = row_vals[c] if c < len(row_vals) else ''
                    cell = new_tbl.cell(r, c)
                    try:
                        if cell.text_frame and cell.text_frame.paragraphs:
                            para = cell.text_frame.paragraphs[0]
                            para.clear()
                            para.text = val
                            if para.runs:
                                run = para.runs[0]
                                if base_style:
                                    if base_style.get('font_color'):
                                        run.font.color.rgb = base_style['font_color']
                                    if base_style.get('font_name'):
                                        run.font.name = base_style['font_name']
                                    if base_style.get('font_size'):
                                        run.font.size = base_style['font_size']
                        else:
                            cell.text = val
                    except Exception as e:
                        self.logger.warning(f"재생성 셀 ({r},{c}) 설정 실패: {e}")

            # 헤더 볼드 처리
            try:
                if headers and rows > 0:
                    for c in range(min(cols, len(headers))):
                        cell = new_tbl.cell(0, c)
                        if cell.text_frame and cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                            cell.text_frame.paragraphs[0].runs[0].font.bold = True
            except Exception:
                pass

            # 원본 테이블 삭제 (XML 레벨)
            try:
                sp = shape._element
                sp.getparent().remove(sp)
            except Exception as e:
                self.logger.warning(f"원본 테이블 삭제 실패(무시): {e}")

            # 이름 유지
            try:
                if name:
                    new_shape.name = name
            except Exception:
                pass

            self.logger.info("🧱 폴백: 새 테이블로 재생성 완료")
        except Exception as e:
            self.logger.error(f"재생성 폴백 실패: {e}")
    
    def _apply_textbox_content(self, shape, mapping: Dict[str, Any]):
        """기존 텍스트박스 로직 적용 (개선된 버전)"""
        try:
            # 다양한 키 이름 지원 (하위 호환성)
            assigned_content = (
                mapping.get('assignedContent', '') or 
                mapping.get('newContent', '') or
                mapping.get('content', '')
            )
            use_original = mapping.get('useOriginal', False)
            content_source = mapping.get('contentSource', '')
            
            self.logger.info(f"🔄 기존 텍스트박스 로직 적용 시작:")
            self.logger.info(f"  - assignedContent='{assigned_content[:50]}...'")
            self.logger.info(f"  - useOriginal={use_original}")
            self.logger.info(f"  - contentSource='{content_source}'")
            
            # 기존 로직과 동일
            if use_original or content_source == 'keep_original':
                original_content = mapping.get('originalContent', '')
                if assigned_content == original_content:
                    self.logger.info(f"내용이 동일하여 건너뜀: '{original_content}'")
                    return
                else:
                    self.logger.info(f"내용이 변경되어 적용 진행: '{original_content}' -> '{assigned_content}'")
            
            # 텍스트 프레임 확인 및 내용 적용
            if not hasattr(shape, 'text_frame'):
                self.logger.warning(f"Shape에 text_frame이 없음")
                return
                
            if not shape.text_frame:
                self.logger.warning(f"text_frame이 None임")
                return
                
            if not assigned_content:
                self.logger.warning(f"적용할 내용이 비어있음")
                return
            
            # 텍스트 내용 적용
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.text = assigned_content
            self.logger.info(f"✅ 기존 텍스트박스 내용 성공적으로 적용: '{assigned_content[:100]}...'")
                
        except Exception as e:
            self.logger.error(f"텍스트박스 내용 적용 실패: {e}")
            import traceback
            self.logger.error(traceback.format_exc())