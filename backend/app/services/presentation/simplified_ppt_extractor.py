"""
간단하고 UI 친화적인 PPT 템플릿 메타데이터 추출기
"""
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path
from typing import Dict, List, Any, Optional

class SimplifiedPPTExtractor:
    """UI 친화적인 PPT 메타데이터 추출기"""
    
    def __init__(self):
        # 위치 기반 매핑을 위한 기준값 설정
        self.slide_width = 960.17  # 표준 슬라이드 너비 (px)
        self.slide_height = 720.0  # 표준 슬라이드 높이 (px)
    
    def extract_presentation(self, pptx_path: str) -> Dict[str, Any]:
        """PPT 파일을 분석하여 UI 친화적인 메타데이터 생성"""
        try:
            prs = Presentation(pptx_path)
            
            result = {
                "presentationTitle": self._extract_title(prs),
                "totalPages": len(prs.slides),
                "slides": []
            }
            
            # 슬라이드별 추출을 개별적으로 처리하여 오류 위치 파악
            for i, slide in enumerate(prs.slides):
                try:
                    slide_data = self._extract_slide(i, slide)
                    result["slides"].append(slide_data)
                except Exception as slide_e:
                    print(f"Error extracting slide {i}: {slide_e}")
                    # 기본 슬라이드 정보라도 추가
                    result["slides"].append({
                        "pageNumber": i + 1,
                        "layout": "Error",
                        "elements": []
                    })
            
            return result
            
        except Exception as e:
            print(f"Error extracting presentation: {e}")
            import traceback
            traceback.print_exc()
            return {
                "presentationTitle": "Error",
                "totalPages": 0,
                "slides": []
            }
    
    def _extract_title(self, prs) -> str:
        """프레젠테이션 제목 추출 (첫 번째 슬라이드에서)"""
        if len(prs.slides) > 0:
            first_slide = prs.slides[0]
            for shape in first_slide.shapes:
                if (hasattr(shape, 'has_text_frame') and 
                    shape.has_text_frame and 
                    shape.text.strip()):
                    return shape.text.strip()
        return "Untitled Presentation"
    
    def _extract_slide(self, index: int, slide) -> Dict[str, Any]:
        """슬라이드별 요소 추출 - 의미있는 구조로 변환"""
        try:
            # 슬라이드 레이아웃 이름을 기반으로 의미있는 레이아웃 결정
            layout_name = self._determine_meaningful_layout(slide, index)
            
            # 텍스트 요소들을 의미있는 순서로 정렬하여 추출
            elements = self._extract_meaningful_elements(slide, layout_name, index)
            
            return {
                "pageNumber": index + 1,
                "layout": layout_name,
                "elements": elements
            }
        except Exception as e:
            print(f"Error extracting slide {index}: {e}")
            return {
                "pageNumber": index + 1,
                "layout": "Basic Layout",
                "elements": []
            }
    
    def _determine_meaningful_layout(self, slide, index: int) -> str:
        """슬라이드의 의미있는 레이아웃 타입 결정"""
        if index == 0:
            return "Title Slide"
        
        # 텍스트 내용 기반 레이아웃 추측
        text_contents = [shape.text.strip().lower() for shape in slide.shapes 
                        if hasattr(shape, 'has_text_frame') and shape.has_text_frame and shape.text.strip()]
        
        combined_text = ' '.join(text_contents)
        
        if '목차' in combined_text or 'contents' in combined_text:
            return "Table of Contents"
        elif '기능' in combined_text or 'features' in combined_text:
            return "Key Features"
        elif '사양' in combined_text or 'specifications' in combined_text:
            return "Technical Specifications"  
        elif '연결' in combined_text or 'connectivity' in combined_text:
            return "Connectivity"
        elif '안전' in combined_text or 'safety' in combined_text:
            return "Safety and Risk Management"
        elif '규제' in combined_text or 'regulatory' in combined_text:
            return "Regulatory Compliance"
        elif 'thank' in combined_text or '감사' in combined_text:
            return "Closing Slide"
        else:
            return "Title and Content"
    
    def _extract_meaningful_elements(self, slide, layout_name: str, slide_index: int) -> List[Dict[str, Any]]:
        """의미있는 요소들을 추출하여 구조화"""
        elements = []
        text_shapes = []
        image_shapes = []
        other_shapes = []  # 도형, 표, 차트 등
        
        # 도형들을 타입별로 분류
        for shape in slide.shapes:
            if self._is_meaningful_shape(shape):
                shape_type = self._determine_element_type(shape)
                
                if shape_type == 'textbox':
                    text_shapes.append(shape)
                elif shape_type == 'image':
                    image_shapes.append(shape)
                else:  # shape, table, chart 등
                    other_shapes.append(shape)
        
        # 모든 요소들을 위치 기반으로 정렬
        all_shapes = text_shapes + image_shapes + other_shapes
        all_shapes.sort(key=lambda s: (s.top, s.left))
        
        # 레이아웃별 특화 처리
        if layout_name == "Title Slide":
            elements.extend(self._process_title_slide(text_shapes, image_shapes))
        elif layout_name == "Table of Contents":
            elements.extend(self._process_toc_slide(text_shapes, image_shapes))
        elif layout_name == "Key Features":
            elements.extend(self._process_features_slide(text_shapes, image_shapes))
        elif layout_name == "Technical Specifications":
            elements.extend(self._process_specs_slide(text_shapes, image_shapes))
        else:
            elements.extend(self._process_content_slide(text_shapes, image_shapes))
        
        # 모든 다른 타입의 도형들도 추가 처리
        for shape in other_shapes:
            element = self._extract_element(shape)
            if element:
                elements.append(element)
        
        return elements
    
    def _extract_element(self, shape) -> Dict[str, Any]:
        """도형을 UI 친화적인 요소로 변환"""
        element_type = self._determine_element_type(shape)
        
        # 기본 고유 ID 생성 (shape name 기반)
        element_id = getattr(shape, 'name', f'element_{id(shape)}')
        
        base_element = {
            "id": element_id,
            "type": element_type,
            "position": self._determine_position(shape),
            "name": getattr(shape, 'name', f'Unknown_{element_type}')
        }
        
        # 타입별 세부 정보 추가
        if element_type == 'textbox' and hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            base_element.update(self._extract_text_info(shape))
        elif element_type == 'image' and hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            base_element.update(self._extract_image_info(shape))
        elif element_type == 'table' and hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            base_element.update(self._extract_table_info(shape))
        elif element_type == 'shape':
            # AUTO_SHAPE, LINE 등의 기본 정보 추가
            has_text = hasattr(shape, 'has_text_frame') and shape.has_text_frame
            text_content = ""
            
            if has_text:
                text_content = getattr(shape, 'text', '').strip()
            
            base_element.update({
                "content": text_content or getattr(shape, 'name', 'Shape'),  # 텍스트 우선, 없으면 이름
                "hasText": has_text and bool(text_content)
            })
            
            # 텍스트가 있는 도형인 경우 텍스트 정보도 추가
            if has_text and text_content:
                base_element.update(self._extract_text_info(shape))
            
        return base_element
    
    def _determine_element_type(self, shape) -> str:
        """도형 타입을 UI 친화적인 타입으로 변환"""
        # python-pptx의 실제 shape_type을 기반으로 매핑
        if hasattr(shape, 'shape_type'):
            shape_type = shape.shape_type
            shape_type_name = getattr(shape_type, 'name', str(shape_type))
            
            # 실제 PPT shape type을 프론트엔드 타입으로 매핑
            type_mapping = {
                'TEXT_BOX': 'textbox',
                'AUTO_SHAPE': 'shape',
                'LINE': 'shape',
                'PICTURE': 'image', 
                'TABLE': 'table',
                'CHART': 'chart',
                'GROUP': 'shape'
            }
            
            mapped_type = type_mapping.get(shape_type_name, 'shape')
            
            # 텍스트가 있는 경우 추가 분류
            if mapped_type == 'textbox' and hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                text = shape.text.strip() if shape.text else ""
                lines = text.split('\n')
                # 리스트 형태인지 확인
                if len(lines) > 2 and any('•' in line or line.strip().startswith(('1.', '2.', '-')) for line in lines):
                    return "list"
            
            return mapped_type
        
        # Fallback: 기존 로직
        if hasattr(shape, 'has_text_frame') and shape.has_text_frame and getattr(shape, 'text', '').strip():
            # 텍스트가 리스트 형태인지 확인
            text = shape.text.strip()
            lines = text.split('\n')
            if len(lines) > 2 and any('•' in line or line.strip().startswith(('1.', '2.', '-')) for line in lines):
                return "list"
            return "textbox"
        elif hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return "image"
        elif hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return "table"
        elif hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.CHART:
            return "chart"
        else:
            return "shape"
    
    def _determine_position(self, shape) -> str:
        """도형 위치를 의미있는 위치명으로 변환"""
        try:
            # shape 속성이 None이거나 잘못된 타입인 경우 체크
            if not hasattr(shape, 'left') or not hasattr(shape, 'top') or \
               not hasattr(shape, 'width') or not hasattr(shape, 'height'):
                return "center"
            
            # EMU를 픽셀로 변환
            left_px = (shape.left / 914400) * 96
            top_px = (shape.top / 914400) * 96
            width_px = (shape.width / 914400) * 96
            height_px = (shape.height / 914400) * 96
            
            # 중심점 계산
            center_x = left_px + width_px / 2
            center_y = top_px + height_px / 2
            
            # 슬라이드를 9구역으로 나누어 위치 결정
            h_position = "left" if center_x < self.slide_width / 3 else \
                        "center" if center_x < 2 * self.slide_width / 3 else "right"
            v_position = "top" if center_y < self.slide_height / 3 else \
                        "middle" if center_y < 2 * self.slide_height / 3 else "bottom"
            
            # 특별한 경우들
            if v_position == "top" and h_position == "center":
                if height_px > 100:  # 큰 제목
                    return "top-center-header"
                else:
                    return "top-center-small"
            
            return f"{v_position}-{h_position}"
            
        except Exception as e:
            print(f"Error determining position: {e}")
            return "center"
    
    def _extract_text_info(self, shape) -> Dict[str, Any]:
        """텍스트 정보를 추출하여 반환"""
        print(f"_extract_text_info 호출됨: {shape.text[:20] if hasattr(shape, 'text') else 'No text'}")
        
        try:
            if not hasattr(shape, 'text_frame') or not shape.text_frame:
                print("  text_frame이 없음")
                return {}
            
            tf = shape.text_frame
            if not tf.paragraphs or not tf.paragraphs[0].runs:
                print("  paragraphs나 runs가 없음")
                return {
                    "content": tf.text,
                    "fontSize": None,
                    "fontFamily": None,
                    "color": None,
                    "fontWeight": None
                }
            
            # 첫 번째 run의 스타일 정보 사용
            first_run = tf.paragraphs[0].runs[0]
            font = first_run.font
            print(f"  폰트 객체: {font}")
            
            # 폰트 크기 (EMU에서 pt로 변환)
            font_size_pt = None
            if font.size:
                font_size_pt = round(font.size.pt, 1)
                print(f"  폰트 크기: {font_size_pt}pt")
            
            # 폰트 이름
            font_name = font.name
            print(f"  폰트 이름: {font_name}")
            
            # 색상 정보 추출
            color_value = None
            if font.color:
                try:
                    # RGB 색상인 경우
                    if hasattr(font.color, 'rgb') and font.color.rgb:
                        color_value = str(font.color.rgb)
                        print(f"  RGB 색상: {color_value}")
                    # 테마 색상인 경우 - 기본값으로 흰색 사용 (제목의 경우)
                    elif hasattr(font.color, 'type') and font.color.type == 2:  # SCHEME
                        # 제목 텍스트는 보통 흰색으로 설정
                        if "제목" in tf.text or font_size_pt and font_size_pt > 30:
                            color_value = "FFFFFF"  # 흰색
                        else:
                            color_value = "000000"  # 검정색
                        print(f"  테마 색상 -> {color_value}")
                    else:
                        color_value = "000000"  # 기본 검정색
                        print(f"  기본 색상: {color_value}")
                except Exception as color_error:
                    print(f"  색상 추출 오류: {color_error}")
                    color_value = "000000"
            
            # 폰트 굵기
            font_weight = "bold" if font.bold else "normal"
            print(f"  폰트 굵기: {font_weight}")
            
            result = {
                "content": tf.text,
                "fontSize": font_size_pt,
                "fontFamily": font_name,
                "color": color_value,
                "fontWeight": font_weight
            }
            
            print(f"  결과: {result}")
            return result
            
        except Exception as e:
            print(f"  텍스트 정보 추출 오류: {e}")
            import traceback
            traceback.print_exc()
            return {
                "content": getattr(shape, 'text', ''),
                "fontSize": None,
                "fontFamily": None,
                "color": None,
                "fontWeight": None
            }
        """텍스트 정보 추출"""
        text_content = shape.text.strip()
        
        # 기본 정보
        info = {"content": text_content}
        
        # 스타일 정보 추출 (첫 번째 run 기준)
        try:
            first_paragraph = shape.text_frame.paragraphs[0]
            if first_paragraph.runs:
                first_run = first_paragraph.runs[0]
                font = first_run.font
                
                style = {}
                
                # 폰트 크기
                if font.size:
                    style["fontSize"] = f"{font.size.pt}pt"
                    
                # 폰트 이름
                if font.name:
                    style["fontFamily"] = font.name
                    
                # 폰트 굵기
                if font.bold:
                    style["fontWeight"] = "bold"
                    
                # 폰트 색상 (개선된 추출 로직)
                if font.color:
                    try:
                        # RGB 색상 타입 확인
                        if hasattr(font.color, 'type') and font.color.type == 1:  # RGB type
                            if hasattr(font.color, 'rgb') and font.color.rgb:
                                rgb = font.color.rgb
                                style["color"] = f"rgb({rgb[0]}, {rgb[1]}, {rgb[2]})"
                        # 테마 색상 타입 확인  
                        elif hasattr(font.color, 'type') and font.color.type == 2:  # SCHEME type
                            if hasattr(font.color, 'theme_color') and font.color.theme_color is not None:
                                theme_color = str(font.color.theme_color)
                                if 'BACKGROUND_1' in theme_color:
                                    style["color"] = "rgb(255, 255, 255)"  # 흰색
                                elif 'TEXT_1' in theme_color:
                                    style["color"] = "rgb(0, 0, 0)"      # 검정색
                                elif 'BACKGROUND_2' in theme_color:
                                    style["color"] = "rgb(238, 236, 225)" # 연한 회색
                                elif 'TEXT_2' in theme_color:
                                    style["color"] = "rgb(68, 84, 106)"   # 어두운 회색
                                else:
                                    style["color"] = "rgb(255, 255, 255)"  # 기본 흰색
                    except Exception as color_error:
                        # 색상 추출 실패 시에도 기본값 제공
                        style["color"] = "rgb(255, 255, 255)"
                    
                # 텍스트 정렬
                if first_paragraph.alignment:
                    alignment_map = {
                        1: "center", 2: "right", 3: "justify", 0: "left"
                    }
                    style["alignment"] = alignment_map.get(first_paragraph.alignment, "left")
                
                if style:
                    info["style"] = style
                    
        except (IndexError, AttributeError):
            pass
        
        # 리스트인 경우 아이템 분리
        if self._determine_element_type(shape) == "list":
            lines = text_content.split('\n')
            info["items"] = [
                {"text": line.strip()} 
                for line in lines 
                if line.strip()
            ]
        
        return info
    
    def _extract_image_info(self, shape) -> Dict[str, Any]:
        """이미지 정보 추출"""
        return {
            "content": "이미지",
            "style": {
                "width": f"{(shape.width / 914400) * 96:.0f}px",
                "height": f"{(shape.height / 914400) * 96:.0f}px"
            }
        }
    
    def _extract_table_info(self, shape) -> Dict[str, Any]:
        """테이블 정보 추출"""
        try:
            if hasattr(shape, 'table') and shape.table:
                table = shape.table
                
                # 기본 테이블 정보
                table_info = {
                    "content": f"Table ({len(table.rows)} x {len(table.columns)})",
                    "rows": len(table.rows),
                    "columns": len(table.columns),
                    "data": []
                }
                
                # 테이블 데이터 추출 (안전하게)
                for row_idx in range(len(table.rows)):
                    row = table.rows[row_idx]
                    row_data = []
                    for col_idx in range(len(row.cells)):
                        cell = row.cells[col_idx]
                        cell_text = cell.text.strip() if hasattr(cell, 'text') else ""
                        row_data.append(cell_text)
                    table_info["data"].append(row_data)
                
                return table_info
                
            return {"content": "Table", "rows": 0, "columns": 0, "data": []}
            
        except Exception as e:
            print(f"테이블 정보 추출 오류: {e}")
            return {"content": "Table (Error)", "rows": 0, "columns": 0, "data": [], "error": str(e)}
    
    def _is_meaningful_shape(self, shape) -> bool:
        """의미있는 도형인지 판단"""
        # 텍스트가 있는 경우 (빈 텍스트라도 일단 포함)
        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            return True
            
        # 이미지, 테이블, 차트인 경우
        if shape.shape_type in [
            MSO_SHAPE_TYPE.PICTURE,
            MSO_SHAPE_TYPE.TABLE,
            MSO_SHAPE_TYPE.CHART
        ]:
            return True
            
        # AUTO_SHAPE (도형)도 포함
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            return True
            
        # LINE (선)도 포함
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            return True
            
        return False
    
    def _process_title_slide(self, text_shapes, image_shapes) -> List[Dict[str, Any]]:
        """제목 슬라이드 처리 - 실제 폰트 정보 추출"""
        elements = []
        
        for i, shape in enumerate(text_shapes):
            # 실제 텍스트 정보 추출
            text_info = self._extract_text_info(shape)
            
            # 기본 요소 정보
            base_element = {
                "type": "textbox",
                "content": shape.text.strip(),
                "position": self._determine_position(shape),
                "id": f"textbox-0-{i}"  # 슬라이드0의 i번째 텍스트박스
            }
            
            # 추출된 스타일 정보 추가
            if text_info:
                base_element.update(text_info)
            else:
                # 폴백 스타일
                if i == 0:  # 첫 번째는 메인 제목
                    base_element.update({
                        "fontSize": 44.0,
                        "fontWeight": "bold",
                        "fontFamily": "Noto Sans KR Black",
                        "color": "FFFFFF"
                    })
                elif i == 1:  # 두 번째는 부제목
                    base_element.update({
                        "fontSize": 24.0,
                        "fontWeight": "normal", 
                        "fontFamily": "Noto Sans KR",
                        "color": "FFFFFF"
                    })
                else:  # 나머지는 작은 정보
                    base_element.update({
                        "fontSize": 12.0,
                        "fontWeight": "normal",
                        "fontFamily": "나눔스퀘어",
                        "color": "6B7286"
                    })
            
            elements.append(base_element)
        
        # 이미지 추가
        for j, shape in enumerate(image_shapes):
            elements.append({
                "type": "image",
                "content": "Logo",
                "position": "bottom-right",
                "id": f"image-0-{j}",
                "style": {"width": "auto", "height": "40px"}
            })
        
        return elements
        
        return elements
    
    def _process_toc_slide(self, text_shapes, image_shapes) -> List[Dict[str, Any]]:
        """목차 슬라이드 처리"""
        elements = []
        
        for i, shape in enumerate(text_shapes):
            text = shape.text.strip()
            if i == 0 and ('목차' in text or 'contents' in text.lower()):
                elements.append({
                    "type": "textbox",
                    "content": text,
                    "position": "top-center",
                    "style": {"fontSize": "36pt", "fontWeight": "bold"}
                })
            elif i == 1:  # 부제목
                elements.append({
                    "type": "textbox", 
                    "content": text,
                    "position": "top-center-subtitle",
                    "style": {"fontSize": "18pt", "color": "grey"}
                })
            else:  # 목차 항목들
                # 리스트 형태로 변환
                lines = text.split('\n')
                items = []
                for idx, line in enumerate(lines):
                    if line.strip():
                        items.append({"index": f"{idx+1:02d}", "text": line.strip()})
                
                if items:
                    elements.append({
                        "type": "list",
                        "position": "center",
                        "items": items,
                        "style": {"fontSize": "16pt", "layout": "two-column"}
                    })
        
        return elements
    
    def _process_features_slide(self, text_shapes, image_shapes) -> List[Dict[str, Any]]:
        """주요 기능 슬라이드 처리"""
        elements = []
        
        for i, shape in enumerate(text_shapes):
            text = shape.text.strip()
            if i == 0:  # 제목
                elements.append({
                    "type": "textbox",
                    "content": text,
                    "position": "top-left-header",
                    "style": {"fontSize": "28pt", "fontWeight": "bold"}
                })
            else:
                # 기능 리스트로 처리
                lines = [line.strip() for line in text.split('\n') if line.strip()]
                if len(lines) > 2:  # 여러 기능이 있는 경우
                    items = []
                    for idx, line in enumerate(lines):
                        items.append({
                            "index": f"{idx+1:02d}",
                            "title": line,
                            "description": "상세 설명"
                        })
                    
                    elements.append({
                        "type": "list",
                        "position": "center", 
                        "items": items,
                        "style": {"layout": "grid"}
                    })
                else:
                    elements.append({
                        "type": "textbox",
                        "content": text,
                        "position": "middle-left-main",
                        "style": {"fontSize": "16pt"}
                    })
        
        return elements
    
    def _process_specs_slide(self, text_shapes, image_shapes) -> List[Dict[str, Any]]:
        """기술 사양 슬라이드 처리"""
        elements = []
        
        for i, shape in enumerate(text_shapes):
            text = shape.text.strip()
            if i == 0:  # 제목
                elements.append({
                    "type": "textbox",
                    "content": text,
                    "position": "top-left-header", 
                    "style": {"fontSize": "28pt", "fontWeight": "bold"}
                })
            elif i == 1:  # 부제목
                elements.append({
                    "type": "textbox",
                    "content": text,
                    "position": "top-left-subtitle",
                    "style": {"fontSize": "14pt"}
                })
            else:
                # 사양 정보를 표 형태로 변환
                if ':' in text:  # 키:값 형태가 있는 경우
                    lines = [line.strip() for line in text.split('\n') if line.strip() and ':' in line]
                    rows = []
                    for line in lines:
                        if ':' in line:
                            key, value = line.split(':', 1)
                            rows.append({"항목": key.strip(), "사양": value.strip()})
                    
                    if rows:
                        elements.append({
                            "type": "table",
                            "position": "center",
                            "headers": ["항목", "사양"],
                            "rows": rows
                        })
                else:
                    elements.append({
                        "type": "textbox",
                        "content": text,
                        "position": "center",
                        "style": {"fontSize": "14pt"}
                    })
        
        return elements
    
    def _process_content_slide(self, text_shapes, image_shapes) -> List[Dict[str, Any]]:
        """일반 콘텐츠 슬라이드 처리"""
        elements = []
        
        for i, shape in enumerate(text_shapes):
            text = shape.text.strip()
            
            # 텍스트 길이와 위치에 따라 역할 결정
            if i == 0:  # 첫 번째는 제목
                elements.append({
                    "type": "textbox",
                    "content": text,
                    "position": "top-left-header",
                    "style": {"fontSize": "28pt", "fontWeight": "bold"}
                })
            elif len(text) > 200:  # 긴 텍스트는 본문
                elements.append({
                    "type": "textbox",
                    "content": text,
                    "position": "middle-left-sub",
                    "style": {"fontSize": "12pt"}
                })
            elif len(text) < 100:  # 짧은 텍스트는 부제목
                elements.append({
                    "type": "textbox", 
                    "content": text,
                    "position": "middle-left-main",
                    "style": {"fontSize": "16pt", "fontWeight": "bold"}
                })
            else:  # 중간 길이는 일반 설명
                elements.append({
                    "type": "textbox",
                    "content": text,
                    "position": "center",
                    "style": {"fontSize": "14pt"}
                })
        
        # 이미지 추가
        for shape in image_shapes:
            elements.append({
                "type": "image",
                "content": "이미지",
                "position": "right-half",
                "style": {"layout": "auto"}
            })
        
        return elements

# 사용 예시
if __name__ == "__main__":
    extractor = SimplifiedPPTExtractor()
    result = extractor.extract_presentation("제품소개서 샘플.pptx")
    
    if result:
        with open("simplified_metadata.json", "w", encoding="utf-8") as f:
            json.dump(result, f, indent=2, ensure_ascii=False)
        print("✅ 간소화된 메타데이터가 생성되었습니다!")
