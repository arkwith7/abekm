"""
템플릿 디버깅 도구
PPTX 템플릿의 구조를 상세히 분석하고 로그를 출력하는 디버깅 도구
"""
import logging
from typing import Dict, Any, List
from pathlib import Path

from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)

class TemplateDebugger:
    """템플릿 디버깅 도구"""
    
    def debug_template(self, template_path: str) -> Dict[str, Any]:
        """템플릿을 상세히 분석하고 디버그 정보 반환"""
        try:
            logger.info(f"🔍 템플릿 디버깅 시작: {template_path}")
            
            if not Path(template_path).exists():
                logger.error(f"❌ 템플릿 파일이 존재하지 않음: {template_path}")
                return {"error": "파일 없음"}
            
            # PPTX 로드
            presentation = Presentation(template_path)
            
            debug_info = {
                "template_path": template_path,
                "total_slides": len(presentation.slides),
                "slides": [],
                "summary": {
                    "total_shapes": 0,
                    "text_shapes": 0,
                    "placeholder_shapes": 0,
                    "cleanable_shapes": 0
                }
            }
            
            # 각 슬라이드 분석
            for slide_idx, slide in enumerate(presentation.slides):
                slide_info = self._debug_slide(slide, slide_idx + 1)
                debug_info["slides"].append(slide_info)
                
                # 요약 정보 업데이트
                debug_info["summary"]["total_shapes"] += slide_info["total_shapes"]
                debug_info["summary"]["text_shapes"] += slide_info["text_shapes"]
                debug_info["summary"]["placeholder_shapes"] += slide_info["placeholder_shapes"]
                debug_info["summary"]["cleanable_shapes"] += slide_info["cleanable_shapes"]
            
            # 결과 로깅
            self._log_debug_summary(debug_info)
            
            return debug_info
            
        except Exception as e:
            logger.error(f"❌ 템플릿 디버깅 실패: {e}")
            return {"error": str(e)}
    
    def _debug_slide(self, slide: Slide, slide_number: int) -> Dict[str, Any]:
        """개별 슬라이드 디버깅"""
        logger.info(f"🔎 슬라이드 {slide_number} 분석 시작")
        
        slide_info = {
            "slide_number": slide_number,
            "layout_name": getattr(slide.slide_layout, 'name', 'Unknown') if hasattr(slide, 'slide_layout') else 'Unknown',
            "total_shapes": len(slide.shapes),
            "text_shapes": 0,
            "placeholder_shapes": 0,
            "cleanable_shapes": 0,
            "shapes": []
        }
        
        # 각 도형 분석
        for shape_idx, shape in enumerate(slide.shapes):
            shape_info = self._debug_shape(shape, shape_idx)
            slide_info["shapes"].append(shape_info)
            
            # 카운터 업데이트
            if shape_info["has_text"]:
                slide_info["text_shapes"] += 1
            if shape_info["is_placeholder"]:
                slide_info["placeholder_shapes"] += 1
            if shape_info["is_cleanable"]:
                slide_info["cleanable_shapes"] += 1
        
        logger.info(f"   슬라이드 {slide_number}: {slide_info['total_shapes']}개 도형 "
                   f"(텍스트: {slide_info['text_shapes']}, "
                   f"플레이스홀더: {slide_info['placeholder_shapes']}, "
                   f"정리가능: {slide_info['cleanable_shapes']})")
        
        return slide_info
    
    def _debug_shape(self, shape: BaseShape, shape_idx: int) -> Dict[str, Any]:
        """개별 도형 디버깅"""
        shape_info = {
            "index": shape_idx,
            "name": getattr(shape, 'name', f'Shape_{shape_idx}'),
            "shape_type": str(getattr(shape, 'shape_type', 'Unknown')),
            "has_text": False,
            "text_content": "",
            "text_length": 0,
            "is_placeholder": False,
            "placeholder_type": None,
            "is_cleanable": False,
            "access_methods": []
        }
        
        # 텍스트 접근 방법 테스트
        
        # 방법 1: text_frame
        try:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                text = shape.text_frame.text
                if text.strip():
                    shape_info["has_text"] = True
                    shape_info["text_content"] = text[:100]
                    shape_info["text_length"] = len(text)
                    shape_info["access_methods"].append("text_frame")
        except Exception as e:
            shape_info["access_methods"].append(f"text_frame_error: {e}")
        
        # 방법 2: 직접 text 속성
        try:
            if hasattr(shape, 'text'):
                text = shape.text
                if text and text.strip():
                    if not shape_info["has_text"]:  # 중복 방지
                        shape_info["has_text"] = True
                        shape_info["text_content"] = text[:100]
                        shape_info["text_length"] = len(text)
                    shape_info["access_methods"].append("direct_text")
        except Exception as e:
            shape_info["access_methods"].append(f"direct_text_error: {e}")
        
        # 방법 3: 테이블
        try:
            if hasattr(shape, 'table') and shape.table:
                table_texts = []
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame and cell.text_frame.text.strip():
                            table_texts.append(cell.text_frame.text)
                
                if table_texts:
                    if not shape_info["has_text"]:
                        shape_info["has_text"] = True
                        shape_info["text_content"] = " | ".join(table_texts)[:100]
                        shape_info["text_length"] = sum(len(t) for t in table_texts)
                    shape_info["access_methods"].append("table")
        except Exception as e:
            shape_info["access_methods"].append(f"table_error: {e}")
        
        # 플레이스홀더 확인
        try:
            if hasattr(shape, 'placeholder_format'):
                try:
                    placeholder_format = shape.placeholder_format
                    if placeholder_format:
                        shape_info["is_placeholder"] = True
                        shape_info["placeholder_type"] = str(placeholder_format.type)
                except Exception:
                    shape_info["is_placeholder"] = True
                    shape_info["placeholder_type"] = "access_error"
        except Exception:
            pass
        
        # 정리 가능 여부 판단
        shape_info["is_cleanable"] = shape_info["has_text"] and self._is_shape_cleanable(shape)
        
        # 상세 로그
        if shape_info["has_text"]:
            logger.debug(f"    도형 {shape_idx}: {shape_info['name']} "
                        f"(타입: {shape_info['shape_type']}, "
                        f"텍스트길이: {shape_info['text_length']}, "
                        f"플레이스홀더: {shape_info['is_placeholder']}, "
                        f"정리가능: {shape_info['is_cleanable']})")
            logger.debug(f"      텍스트: '{shape_info['text_content'][:50]}...'")
            logger.debug(f"      접근방법: {shape_info['access_methods']}")
        
        return shape_info
    
    def _is_shape_cleanable(self, shape: BaseShape) -> bool:
        """도형이 정리 가능한지 판단 (template_content_cleaner와 동일한 로직)"""
        try:
            # 보존 대상 도형 체크
            preserved_shapes = {"background", "logo", "decoration", "border"}
            if hasattr(shape, 'name') and shape.name:
                shape_name = shape.name.lower()
                for preserved in preserved_shapes:
                    if preserved in shape_name:
                        return False
            
            # 텍스트가 있는 모든 도형은 정리 가능
            return True
            
        except Exception:
            return False
    
    def _log_debug_summary(self, debug_info: Dict[str, Any]) -> None:
        """디버그 결과 요약 로깅"""
        summary = debug_info["summary"]
        
        logger.info("🎯 템플릿 디버깅 결과 요약:")
        logger.info(f"   📄 총 슬라이드: {debug_info['total_slides']}개")
        logger.info(f"   🔷 총 도형: {summary['total_shapes']}개")
        logger.info(f"   📝 텍스트 도형: {summary['text_shapes']}개")
        logger.info(f"   🎭 플레이스홀더: {summary['placeholder_shapes']}개")
        logger.info(f"   🧹 정리 가능: {summary['cleanable_shapes']}개")
        
        if summary['cleanable_shapes'] == 0:
            logger.warning("⚠️  정리 가능한 도형이 0개입니다! 콘텐츠 정리가 작동하지 않을 수 있습니다.")
        
        # 슬라이드별 요약
        for slide in debug_info["slides"]:
            if slide["cleanable_shapes"] == 0 and slide["text_shapes"] > 0:
                logger.warning(f"⚠️  슬라이드 {slide['slide_number']}: "
                              f"텍스트 도형 {slide['text_shapes']}개가 있지만 정리 가능한 도형은 0개")


# 전역 인스턴스
template_debugger = TemplateDebugger()
