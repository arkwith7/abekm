#!/usr/bin/env python3
"""
PPT 슬라이드별 상세 내용 분석
각 슬라이드의 제목, 키메시지, 내용 요소들을 분석합니다.
"""

from pptx import Presentation
import sys

def analyze_slide_content(ppt_path: str):
    """각 슬라이드의 상세 내용을 분석"""
    try:
        prs = Presentation(ppt_path)
        print(f"📊 PPT 파일: {ppt_path}")
        print(f"📄 총 슬라이드 수: {len(prs.slides)}")
        print("=" * 80)
        
        for i, slide in enumerate(prs.slides):
            print(f"\n📋 슬라이드 #{i+1}")
            print(f"레이아웃: {slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else 'Unknown'}")
            
            text_elements = []
            
            # 모든 텍스트 추출
            for j, shape in enumerate(slide.shapes):
                try:
                    text = ""
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        text = shape.text_frame.text.strip()
                    elif hasattr(shape, "text"):
                        text = shape.text.strip()
                    
                    if text:
                        text_elements.append({
                            'shape_id': j + 1,
                            'text': text,
                            'length': len(text),
                            'lines': len(text.split('\n'))
                        })
                except:
                    pass
            
            print(f"총 텍스트 요소 수: {len(text_elements)}")
            
            for k, element in enumerate(text_elements):
                print(f"\n  텍스트 요소 {k+1} (Shape #{element['shape_id']}):")
                print(f"    길이: {element['length']}자, 라인 수: {element['lines']}")
                print(f"    내용: '{element['text'][:100]}{'...' if len(element['text']) > 100 else ''}'")
                
                # 만약 여러 라인이면 각 라인 표시
                if element['lines'] > 1:
                    lines = element['text'].split('\n')
                    for l, line in enumerate(lines[:5]):  # 처음 5줄만
                        if line.strip():
                            print(f"      라인 {l+1}: '{line.strip()}'")
                    if len(lines) > 5:
                        print(f"      ... (총 {len(lines)}줄)")
            
            print("-" * 60)
        
    except Exception as e:
        print(f"❌ 분석 실패: {e}")

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("사용법: python analyze_slide_content.py <ppt_file_path>")
        sys.exit(1)
    
    analyze_slide_content(sys.argv[1])
