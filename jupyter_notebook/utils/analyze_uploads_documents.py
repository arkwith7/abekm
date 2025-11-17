#!/usr/bin/env python3
"""
업로드된 문서들을 분석하여 그라운드 트루스 데이터를 생성하는 스크립트
"""

import os
import sys
import pandas as pd
import json
from datetime import datetime
from typing import List, Dict, Tuple
import re

# 문서 처리 라이브러리
try:
    from docx import Document
    import PyPDF2
    from pptx import Presentation
except ImportError as e:
    print(f"필요한 라이브러리를 가상환경에서 설치해주세요: {e}")
    sys.exit(1)


class DocumentAnalyzer:
    """업로드된 문서들을 분석하여 내용을 추출하는 클래스"""
    
    def __init__(self, uploads_dir: str):
        self.uploads_dir = uploads_dir
        self.documents_info = []
    
    def extract_docx_content(self, file_path: str) -> Dict:
        """DOCX 파일에서 내용 추출"""
        try:
            doc = Document(file_path)
            
            # 텍스트 추출
            full_text = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    full_text.append(paragraph.text.strip())
            
            # 테이블 내용 추출
            tables_content = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_data.append(cell.text.strip())
                    if row_data:
                        table_data.append(" | ".join(row_data))
                if table_data:
                    tables_content.append("\n".join(table_data))
            
            # 키워드 추출
            content = " ".join(full_text)
            keywords = self.extract_keywords(content)
            
            return {
                "type": "docx",
                "text_content": full_text,
                "tables_content": tables_content,
                "keywords": keywords,
                "word_count": len(content.split()),
                "summary": self.generate_summary(content)
            }
            
        except Exception as e:
            return {"error": f"DOCX 처리 오류: {str(e)}"}
    
    def extract_pdf_content(self, file_path: str) -> Dict:
        """PDF 파일에서 내용 추출"""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                full_text = []
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
                    if text.strip():
                        full_text.append(text.strip())
                
                content = " ".join(full_text)
                keywords = self.extract_keywords(content)
                
                return {
                    "type": "pdf",
                    "text_content": full_text,
                    "page_count": len(pdf_reader.pages),
                    "keywords": keywords,
                    "word_count": len(content.split()),
                    "summary": self.generate_summary(content)
                }
                
        except Exception as e:
            return {"error": f"PDF 처리 오류: {str(e)}"}
    
    def extract_pptx_content(self, file_path: str) -> Dict:
        """PPTX 파일에서 내용 추출"""
        try:
            prs = Presentation(file_path)
            
            slides_content = []
            all_text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                # 텍스트 박스에서 텍스트 추출
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    slide_content = {
                        "slide_number": slide_num,
                        "content": slide_text
                    }
                    slides_content.append(slide_content)
                    all_text.extend(slide_text)
            
            content = " ".join(all_text)
            keywords = self.extract_keywords(content)
            
            return {
                "type": "pptx",
                "slides_content": slides_content,
                "slide_count": len(prs.slides),
                "keywords": keywords,
                "word_count": len(content.split()),
                "summary": self.generate_summary(content)
            }
            
        except Exception as e:
            return {"error": f"PPTX 처리 오류: {str(e)}"}
    
    def extract_keywords(self, text: str) -> List[str]:
        """텍스트에서 주요 키워드 추출"""
        # 한글, 영어 키워드 추출
        korean_words = re.findall(r'[가-힣]{2,}', text)
        english_words = re.findall(r'[a-zA-Z]{3,}', text.lower())
        
        # 불용어 제거
        stopwords = {
            '있는', '하는', '되는', '같은', '이런', '그런', '어떤', '수도', '때문',
            '그리고', '또한', '하지만', '그러나', '따라서', '그래서', '이것', '그것',
            'the', 'and', 'are', 'for', 'with', 'can', 'you', 'have', 'what',
            'this', 'that', 'will', 'from', 'they', 'been', 'said', 'each'
        }
        
        # 빈도수 계산 및 상위 키워드 선택
        all_words = korean_words + english_words
        word_freq = {}
        for word in all_words:
            if word not in stopwords and len(word) >= 2:
                word_freq[word] = word_freq.get(word, 0) + 1
        
        # 빈도수 기준 상위 10개 키워드
        top_keywords = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)[:10]
        return [word for word, freq in top_keywords]
    
    def generate_summary(self, content: str, max_length: int = 200) -> str:
        """텍스트 요약 생성 (간단한 방식)"""
        if not content:
            return ""
        
        # 문장 분리
        sentences = re.split(r'[.!?]\s+', content)
        
        if not sentences:
            return content[:max_length]
        
        # 첫 번째 의미있는 문장을 요약으로 사용
        for sentence in sentences:
            if len(sentence.strip()) > 10:
                summary = sentence.strip()
                if len(summary) > max_length:
                    summary = summary[:max_length] + "..."
                return summary
        
        return content[:max_length] + "..." if len(content) > max_length else content
    
    def analyze_all_documents(self) -> List[Dict]:
        """모든 문서 분석"""
        results = []
        
        for filename in os.listdir(self.uploads_dir):
            file_path = os.path.join(self.uploads_dir, filename)
            
            # 디렉토리는 건너뛰기
            if os.path.isdir(file_path):
                continue
            
            print(f"분석 중: {filename}")
            
            # 파일 확장자에 따라 처리
            if filename.lower().endswith('.docx'):
                content_info = self.extract_docx_content(file_path)
            elif filename.lower().endswith('.pdf'):
                content_info = self.extract_pdf_content(file_path)
            elif filename.lower().endswith('.pptx'):
                content_info = self.extract_pptx_content(file_path)
            else:
                continue
            
            # 파일 정보 추가
            file_info = {
                "filename": filename,
                "file_path": file_path,
                "file_size": os.path.getsize(file_path),
                "modified_date": datetime.fromtimestamp(os.path.getmtime(file_path)),
                **content_info
            }
            
            results.append(file_info)
        
        return results


def create_ground_truth_from_documents(documents_info: List[Dict]) -> pd.DataFrame:
    """문서 정보를 기반으로 그라운드 트루스 데이터 생성"""
    
    ground_truth_data = []
    
    for doc in documents_info:
        if "error" in doc:
            continue
        
        filename = doc["filename"]
        doc_type = doc["type"]
        keywords = doc.get("keywords", [])
        summary = doc.get("summary", "")
        
        # 1. 문서 존재 확인 질문들
        existence_questions = [
            f"{filename} 파일이 있나요?",
            f"{doc_type.upper()} 파일 중에 {keywords[0] if keywords else '관련'} 내용이 있나요?",
            f"업로드된 문서 중에 {keywords[0] if keywords else '특정'} 관련 자료가 있나요?"
        ]
        
        for question in existence_questions:
            ground_truth_data.append({
                "question": question,
                "category": "document_existence",
                "api_type": "general",
                "expected_has_reference": True,
                "expected_reference_file": filename,
                "expected_answer_type": "확인",
                "keywords": ", ".join(keywords[:3]),
                "difficulty": "easy",
                "test_purpose": "문서 존재 확인"
            })
        
        # 2. 내용 기반 질문들
        if keywords:
            content_questions = [
                f"{keywords[0]}에 대해 알려주세요",
                f"{keywords[0]} 관련 정보를 찾아주세요",
                f"{keywords[0]}의 특징이나 내용을 설명해주세요"
            ]
            
            for question in content_questions:
                ground_truth_data.append({
                    "question": question,
                    "category": "content_inquiry",
                    "api_type": "general",
                    "expected_has_reference": True,
                    "expected_reference_file": filename,
                    "expected_answer_type": "설명",
                    "keywords": ", ".join(keywords[:3]),
                    "difficulty": "medium",
                    "test_purpose": "내용 검색 및 설명"
                })
        
        # 3. PPT 생성 요청 (PPT 파일인 경우)
        if doc_type == "pptx":
            ppt_questions = [
                f"{keywords[0]} PPT 만들어주세요",
                f"{keywords[0]} 발표자료 생성해주세요",
                f"{keywords[0]} 프레젠테이션 만들어주세요"
            ]
            
            for question in ppt_questions:
                ground_truth_data.append({
                    "question": question,
                    "category": "ppt_generation",
                    "api_type": "ppt",
                    "expected_has_reference": True,
                    "expected_reference_file": filename,
                    "expected_answer_type": "PPT 생성",
                    "keywords": ", ".join(keywords[:3]),
                    "difficulty": "hard",
                    "test_purpose": "PPT 생성 기능"
                })
    
    # 4. 존재하지 않는 내용에 대한 질문들 (네거티브 케이스)
    negative_questions = [
        "양자컴퓨터 기술에 대해 알려주세요",
        "블록체인 암호화폐 투자 전략을 설명해주세요", 
        "우주항공 기술의 최신 동향은 어떤가요?",
        "심해 생물의 진화 과정을 설명해주세요"
    ]
    
    for question in negative_questions:
        ground_truth_data.append({
            "question": question,
            "category": "non_existent_content",
            "api_type": "general",
            "expected_has_reference": False,
            "expected_reference_file": "없음",
            "expected_answer_type": "자료 없음 안내",
            "keywords": "존재하지 않는 내용",
            "difficulty": "medium",
            "test_purpose": "부정확한 응답 방지"
        })
    
    return pd.DataFrame(ground_truth_data)


def main():
    """메인 실행 함수"""
    uploads_dir = "/home/admin/wkms-aws/backend/uploads"
    
    if not os.path.exists(uploads_dir):
        print(f"업로드 디렉토리를 찾을 수 없습니다: {uploads_dir}")
        return
    
    print("📁 업로드된 문서들을 분석하고 있습니다...")
    
    # 문서 분석
    analyzer = DocumentAnalyzer(uploads_dir)
    documents_info = analyzer.analyze_all_documents()
    
    print(f"✅ 총 {len(documents_info)}개 문서 분석 완료")
    
    # 문서 정보 저장
    documents_df = pd.DataFrame(documents_info)
    documents_df.to_csv("/home/admin/wkms-aws/documents_analysis.csv", index=False, encoding='utf-8-sig')
    print("📊 문서 분석 결과가 documents_analysis.csv에 저장되었습니다")
    
    # 그라운드 트루스 생성
    print("🎯 그라운드 트루스 데이터를 생성하고 있습니다...")
    ground_truth_df = create_ground_truth_from_documents(documents_info)
    
    # CSV 파일로 저장
    ground_truth_df.to_csv("/home/admin/wkms-aws/ground_truth_criteria.csv", index=False, encoding='utf-8-sig')
    print(f"✅ 그라운드 트루스 데이터가 생성되었습니다: {len(ground_truth_df)}개 테스트 케이스")
    
    # 통계 정보 출력
    print("\n📈 그라운드 트루스 통계:")
    print(f"- 총 테스트 케이스: {len(ground_truth_df)}개")
    print(f"- 카테고리별 분포:")
    category_counts = ground_truth_df['category'].value_counts()
    for category, count in category_counts.items():
        print(f"  • {category}: {count}개")
    
    print(f"- API 타입별 분포:")
    api_counts = ground_truth_df['api_type'].value_counts()
    for api_type, count in api_counts.items():
        print(f"  • {api_type}: {count}개")
    
    # 상세 분석 결과 JSON으로도 저장
    analysis_result = {
        "analysis_date": datetime.now().isoformat(),
        "total_documents": len(documents_info),
        "document_types": {},
        "total_test_cases": len(ground_truth_df),
        "documents_detail": documents_info
    }
    
    # 문서 타입별 통계
    for doc in documents_info:
        if "error" not in doc:
            doc_type = doc.get("type", "unknown")
            analysis_result["document_types"][doc_type] = analysis_result["document_types"].get(doc_type, 0) + 1
    
    with open("/home/admin/wkms-aws/documents_analysis_detail.json", "w", encoding='utf-8') as f:
        json.dump(analysis_result, f, ensure_ascii=False, indent=2, default=str)
    
    print("📄 상세 분석 결과가 documents_analysis_detail.json에 저장되었습니다")
    print("\n🎉 모든 분석이 완료되었습니다!")


if __name__ == "__main__":
    main()