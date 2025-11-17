#!/usr/bin/env python3
"""
PPT 파일 구조 분석 스크립트
목차와 실제 슬라이드 제목의 매칭 여부를 확인합니다.
"""

import sys
from pathlib import Path
from pptx import Presentation
import re

def analyze_ppt_structure(ppt_path: str):
    """PPT 파일의 구조를 분석하고 목차-슬라이드 매칭을 확인"""
    try:
        prs = Presentation(ppt_path)
        print(f"📊 PPT 파일 분석: {Path(ppt_path).name}")
        print(f"📄 총 슬라이드 수: {len(prs.slides)}")
        print("=" * 60)
        
        slides_info = []
        toc_items = []
        
        for i, slide in enumerate(prs.slides):
            slide_info = {
                'index': i + 1,
                'title': '',
                'content': [],
                'layout': slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else 'Unknown'
            }
            
            # 슬라이드의 모든 텍스트 추출
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            if not slide_info['title'] and len(text) < 100:  # 첫 번째 짧은 텍스트를 제목으로 추정
                                slide_info['title'] = text
                            slide_info['content'].append(text)
                    elif hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            if not slide_info['title'] and len(text) < 100:  # 첫 번째 짧은 텍스트를 제목으로 추정
                                slide_info['title'] = text
                            slide_info['content'].append(text)
                except:
                    pass  # 텍스트가 없는 shape는 무시
            
            slides_info.append(slide_info)
            
            # 목차 슬라이드인지 확인 (제목에 '목차'가 포함되어 있으면)
            if '목차' in slide_info['title']:
                for content in slide_info['content']:
                    # 번호가 있는 항목들을 목차 아이템으로 추출
                    lines = content.split('\n')
                    for line in lines:
                        line = line.strip()
                        if re.match(r'^\d+\.', line):  # 1., 2., 3. 등으로 시작하는 라인
                            toc_items.append(line)
        
        # 결과 출력
        print("📋 슬라이드 목록:")
        for slide in slides_info:
            print(f"  {slide['index']}. 제목: '{slide['title']}'")
            print(f"     레이아웃: {slide['layout']}")
            if len(slide['content']) > 1:
                print(f"     내용 요소 수: {len(slide['content'])}")
            print()
        
        print("📑 목차 항목들:")
        for item in toc_items:
            print(f"  {item}")
        print()
        
        # 매칭 분석
        print("🔍 목차-슬라이드 매칭 분석:")
        content_slides = [s for s in slides_info if '목차' not in s['title'] and '감사' not in s['title'] and s['index'] > 1]
        
        if toc_items and content_slides:
            print(f"  목차 항목 수: {len(toc_items)}")
            print(f"  내용 슬라이드 수: {len(content_slides)}")
            
            for i, toc_item in enumerate(toc_items):
                toc_title = re.sub(r'^\d+\.\s*', '', toc_item).strip()
                
                if i < len(content_slides):
                    slide_title = content_slides[i]['title']
                    match_status = "✅" if toc_title in slide_title or slide_title in toc_title else "❌"
                    print(f"  {match_status} 목차: '{toc_title}' ↔ 슬라이드 {content_slides[i]['index']}: '{slide_title}'")
                else:
                    print(f"  ❌ 목차: '{toc_title}' ↔ 해당 슬라이드 없음")
        
        return slides_info, toc_items
        
    except Exception as e:
        print(f"❌ PPT 분석 실패: {e}")
        return None, None

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("사용법: python analyze_ppt_structure.py <ppt_file_path>")
        sys.exit(1)
    
    ppt_path = sys.argv[1]
    if not Path(ppt_path).exists():
        print(f"❌ 파일을 찾을 수 없습니다: {ppt_path}")
        sys.exit(1)
    
    analyze_ppt_structure(ppt_path)
