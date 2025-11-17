#!/usr/bin/env python3
"""
기존 PPT 템플릿들의 shape 이름을 업데이트하는 스크립트
"""
import os
import sys
from pathlib import Path
from pptx import Presentation

# 프로젝트 루트를 Python path에 추가
project_root = Path(__file__).parent
sys.path.insert(0, str(project_root))

def update_template_shape_names(prs, template_path):
    """템플릿 파일의 텍스트박스 shape들에 ID 이름을 설정하고 저장"""
    try:
        modified = False
        for slide_idx, slide in enumerate(prs.slides):
            textbox_index = 0
            for shape in slide.shapes:
                if (hasattr(shape, 'shape_type') and 
                    shape.shape_type.name == "TEXT_BOX" and 
                    hasattr(shape, "text") and getattr(shape, "text", "").strip()):
                    
                    element_id = f"textbox-{slide_idx}-{textbox_index}"
                    if shape.name != element_id:
                        shape.name = element_id
                        modified = True
                        print(f"  Shape 이름 설정: {element_id}")
                    textbox_index += 1
        
        # 수정된 경우 템플릿 파일 저장
        if modified:
            prs.save(template_path)
            print(f"  템플릿 파일 업데이트됨: {Path(template_path).name}")
        else:
            print(f"  변경사항 없음: {Path(template_path).name}")
    except Exception as e:
        print(f"  템플릿 shape 이름 업데이트 실패: {e}")

def update_all_templates():
    """backend/uploads/templates/ 디렉토리의 모든 PPT 템플릿 업데이트"""
    # 실제 템플릿 파일이 저장된 경로
    templates_dir = project_root / "backend" / "uploads" / "templates"
    metadata_dir = templates_dir / "metadata"
    
    if not templates_dir.exists():
        print(f"템플릿 디렉토리가 존재하지 않음: {templates_dir}")
        return
    
    print(f"템플릿 디렉토리: {templates_dir}")
    print(f"메타데이터 디렉토리: {metadata_dir}")
    
    updated_count = 0
    for ppt_file in templates_dir.glob("*.pptx"):
        try:
            print(f"\n템플릿 업데이트 중: {ppt_file.name}")
            prs = Presentation(str(ppt_file))
            update_template_shape_names(prs, str(ppt_file))
            
            # 해당하는 메타데이터 파일이 있는지 확인
            # 파일명에서 공백을 언더스코어로 변환하여 메타데이터 파일명 생성
            metadata_filename = ppt_file.stem.replace(" ", "_") + "_metadata.json"
            metadata_file = metadata_dir / metadata_filename
            if metadata_file.exists():
                print(f"  연결된 메타데이터: {metadata_file.name}")
            else:
                print(f"  ⚠️ 메타데이터 파일이 없음: {metadata_filename}")
                # 다른 가능한 파일명들도 확인
                for existing_file in metadata_dir.glob("*_metadata.json"):
                    print(f"    기존 메타데이터 파일들: {existing_file.name}")
            
            updated_count += 1
            print(f"✅ 업데이트 완료: {ppt_file.name}")
        except Exception as e:
            print(f"❌ 업데이트 실패: {ppt_file.name} - {e}")
    
    print(f"\n총 {updated_count}개 템플릿 업데이트 완료")

if __name__ == "__main__":
    update_all_templates()
