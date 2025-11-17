#!/usr/bin/env python3
"""
PPT 목차 슬라이드 상세 분석
"""

from pptx import Presentation
import sys

def analyze_toc_slide(ppt_path: str):
    """목차 슬라이드의 상세 내용을 분석"""
    try:
        prs = Presentation(ppt_path)
        
        for i, slide in enumerate(prs.slides):
            slide_title = ""
            all_content = []
            
            # 모든 텍스트 추출
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            if not slide_title and len(text) < 100:
                                slide_title = text
                            all_content.append(text)
                            print(f"📝 Shape 텍스트: '{text}'")
                    elif hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            if not slide_title and len(text) < 100:
                                slide_title = text
                            all_content.append(text)
                            print(f"📝 Shape 텍스트: '{text}'")
                except:
                    pass
            
            if '목차' in slide_title:
                print(f"\n📋 목차 슬라이드 (#{i+1}) 상세 분석:")
                print(f"제목: '{slide_title}'")
                print(f"총 텍스트 요소 수: {len(all_content)}")
                
                for j, content in enumerate(all_content):
                    print(f"\n텍스트 요소 {j+1}:")
                    print(f"'{content}'")
                    
                    # 라인별 분석
                    lines = content.split('\n')
                    for k, line in enumerate(lines):
                        line = line.strip()
                        if line:
                            print(f"  라인 {k+1}: '{line}'")
                
                break
        
    except Exception as e:
        print(f"❌ 분석 실패: {e}")

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("사용법: python analyze_toc_detail.py <ppt_file_path>")
        sys.exit(1)
    
    analyze_toc_slide(sys.argv[1])
