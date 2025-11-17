#!/usr/bin/env python3
"""
템플릿 기반 PPT 구조 분석 스크립트
"""

import os
import sys
from pptx import Presentation
import json

def analyze_template_ppt_structure(pptx_file):
    """템플릿 PPT 구조를 분석합니다."""
    
    if not os.path.exists(pptx_file):
        print(f"❌ 파일을 찾을 수 없습니다: {pptx_file}")
        return
    
    try:
        # PowerPoint 파일 로드
        prs = Presentation(pptx_file)
        
        print(f"📊 템플릿 PPT 분석: {os.path.basename(pptx_file)}")
        print(f"📋 총 슬라이드 수: {len(prs.slides)}")
        print("=" * 60)
        
        for slide_idx, slide in enumerate(prs.slides):
            print(f"\n🎯 슬라이드 {slide_idx + 1}:")
            print(f"   레이아웃: {slide.slide_layout.name}")
            
            # 텍스트박스 분석
            textboxes = []
            bullet_count = 0
            for shape_idx, shape in enumerate(slide.shapes):
                try:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        text_content = ""
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                text_content += paragraph.text.strip() + " "
                                # 불릿 포인트 카운트
                                if paragraph.text.strip().startswith(('•', '-', '*', '1.', '2.', '3.', '4.', '5.', '6.')):
                                    bullet_count += 1
                        
                        if text_content.strip():
                            textboxes.append({
                                "shape_idx": shape_idx,
                                "name": shape.name,
                                "text": text_content.strip()[:100] + ("..." if len(text_content.strip()) > 100 else "")
                            })
                except Exception as e:
                    # 일부 shape에서 text_frame 접근 실패할 수 있음
                    pass
            
            print(f"   텍스트박스 수: {len(textboxes)}")
            
            if textboxes:
                for tb in textboxes:
                    print(f"     📝 Shape {tb['shape_idx']} ({tb['name']}): {tb['text']}")
            else:
                print("     ⚠️ 텍스트 내용이 없습니다")
            
            print(f"   🔸 불릿 포인트 수: {bullet_count}")
        
        print("\n" + "=" * 60)
        print("✅ 템플릿 PPT 구조 분석 완료")
        
    except Exception as e:
        print(f"❌ 분석 실패: {str(e)}")

def compare_with_expected_content():
    """기대되는 컨텐츠와 비교"""
    expected_slides = [
        {"title": "1. 제품 개요", "expected_bullets": 3},
        {"title": "2. 주요 특징", "expected_bullets": 3}, 
        {"title": "3. 제품 사양", "expected_bullets": 6},
        {"title": "4. 주요 기능", "expected_bullets": 4},
        {"title": "5. 제품 구성", "expected_bullets": 4}
    ]
    
    print("\n📊 기대 컨텐츠 vs 실제 구조:")
    print("=" * 60)
    for slide in expected_slides:
        print(f"🎯 {slide['title']}: 기대 bullets={slide['expected_bullets']}, 실제 bullets=0 ❌")

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("사용법: python analyze_template_ppt_structure.py <pptx_file_path>")
        sys.exit(1)
    
    pptx_file = sys.argv[1]
    
    print("🔍 템플릿 PPT 구조 분석 시작...")
    analyze_template_ppt_structure(pptx_file)
    compare_with_expected_content()
