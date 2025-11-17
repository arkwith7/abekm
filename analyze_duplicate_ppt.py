#!/usr/bin/env python3
"""
PPT 파일의 중복 텍스트 문제 분석 스크립트
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def analyze_ppt_duplicates(ppt_path):
    """PPT 파일의 중복 텍스트 문제를 분석합니다."""
    
    if not Path(ppt_path).exists():
        print(f"❌ 파일이 존재하지 않습니다: {ppt_path}")
        return
        
    try:
        prs = Presentation(ppt_path)
        print(f"📊 PPT 분석: {ppt_path}")
        print(f"총 슬라이드 수: {len(prs.slides)}\n")
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            print(f"🎯 슬라이드 {slide_idx}:")
            
            # 텍스트 내용 수집
            text_contents = []
            table_contents = []
            shape_info = []
            
            for shape_idx, shape in enumerate(slide.shapes):
                shape_info.append({
                    'idx': shape_idx,
                    'type': shape.shape_type,
                    'name': getattr(shape, 'name', 'Unknown')
                })
                
                # 텍스트 박스/플레이스홀더 텍스트
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        text_contents.append({
                            'shape_idx': shape_idx,
                            'shape_type': str(shape.shape_type),
                            'text': text
                        })
                
                # 표 내용
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table_text = []
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                row_text.append(cell_text)
                        if row_text:
                            table_text.append(' | '.join(row_text))
                    
                    if table_text:
                        table_contents.append({
                            'shape_idx': shape_idx,
                            'table_text': table_text
                        })
            
            # Shape 정보 출력
            print(f"  📦 Shape 개수: {len(shape_info)}")
            for shape in shape_info:
                print(f"    - Shape {shape['idx']}: {shape['type']} ({shape['name']})")
            
            # 텍스트 내용 출력
            print(f"  📝 텍스트 박스/플레이스홀더 ({len(text_contents)}개):")
            for text_item in text_contents:
                print(f"    - Shape {text_item['shape_idx']} ({text_item['shape_type']}): '{text_item['text'][:100]}{'...' if len(text_item['text']) > 100 else ''}'")
            
            # 표 내용 출력
            print(f"  📋 표 ({len(table_contents)}개):")
            for table_item in table_contents:
                print(f"    - Shape {table_item['shape_idx']} (표):")
                for row_idx, row in enumerate(table_item['table_text']):
                    print(f"      Row {row_idx + 1}: {row}")
            
            # 중복 텍스트 감지
            all_texts = [item['text'] for item in text_contents]
            all_texts.extend([' '.join(table['table_text']) for table in table_contents])
            
            duplicates = []
            for i, text1 in enumerate(all_texts):
                for j, text2 in enumerate(all_texts[i+1:], i+1):
                    # 텍스트가 유사하거나 포함 관계인 경우
                    if text1 in text2 or text2 in text1 or text1 == text2:
                        if len(text1) > 10:  # 짧은 텍스트는 제외
                            duplicates.append((i, j, text1, text2))
            
            if duplicates:
                print(f"  ⚠️  중복/유사 텍스트 발견:")
                for dup in duplicates:
                    print(f"    - 텍스트 {dup[0]} vs {dup[1]}: '{dup[2][:50]}...' ≈ '{dup[3][:50]}...'")
            else:
                print(f"  ✅ 중복 텍스트 없음")
            
            print()
            
    except Exception as e:
        print(f"❌ 분석 중 오류 발생: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    # 새로 생성된 테스트 PPT 파일 분석
    new_ppt_path = "/home/admin/wkms-aws/uploads/test_fixed_insulin_pump.pptx"
    
    if Path(new_ppt_path).exists():
        print(f"\n{'='*80}")
        print("🔍 새로 생성된 PPT 분석 (중복 문제 수정 후)")
        analyze_ppt_duplicates(new_ppt_path)
        print(f"{'='*80}\n")
    else:
        print(f"❌ 새 테스트 파일이 존재하지 않습니다: {new_ppt_path}")
    
    # 기존 문제 파일들과 비교
    old_ppt_files = [
        "/home/admin/wkms-aws/backend/uploads/quick_presentation_인슐린_펌프_제품소개자료.pptx",
        "/home/admin/wkms-aws/backend/uploads/quick_presentation_스마트_인슐린_펌프_제품_소개자료_backup.pptx"
    ]
    
    print("\n🔍 기존 문제 파일들과 비교:")
    for ppt_path in old_ppt_files:
        if Path(ppt_path).exists():
            print(f"\n{'='*80}")
            print(f"📄 기존 파일: {Path(ppt_path).name}")
            analyze_ppt_duplicates(ppt_path)
            print(f"{'='*80}\n")
        else:
            print(f"⏭️  파일 건너뜀 (존재하지 않음): {ppt_path}")
